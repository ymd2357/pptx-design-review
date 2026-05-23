#!/usr/bin/env python3
"""Build a 3-slide fixture deck demonstrating each overflow lint with both
OK (non-fire) and NG (fire) cases side-by-side.

slide 1: box_canvas_overflow
slide 2: text_box_overflow
slide 3: text_canvas_overflow

Each lint target text box is tinted with a light fill so the reader can
identify which shapes the lint inspects.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_PATH = REPO_ROOT / "tmp" / "review" / "overflow-fixtures" / "before.pptx"

SLIDE_W = 1440
SLIDE_H = 810

TINT_OK = RGBColor(0xDD, 0xEE, 0xDD)
TINT_NG = RGBColor(0xFF, 0xDD, 0xDD)
LABEL_GRAY = RGBColor(0x55, 0x55, 0x55)


def _new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Pt(SLIDE_W)
    prs.slide_height = Pt(SLIDE_H)
    return prs


def _add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_run(tf, text: str, size_pt: int, color: RGBColor | None = None) -> None:
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(size_pt)
    if color is not None:
        run.font.color.rgb = color


def _add_text_box(
    slide, x: float, y: float, w: float, h: float,
    text: str, font_pt: int, tint: RGBColor,
    word_wrap: bool = True,
) -> None:
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    tb.text_frame.word_wrap = word_wrap
    tb.fill.solid()
    tb.fill.fore_color.rgb = tint
    tb.line.color.rgb = RGBColor(0x88, 0x88, 0x88)
    _set_run(tb.text_frame, text, font_pt)


def _add_label(slide, x: float, y: float, text: str, size_pt: int = 16) -> None:
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(680), Pt(36))
    tb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    _set_run(tb.text_frame, text, size_pt, LABEL_GRAY)


def _add_title(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Pt(40), Pt(20), Pt(1360), Pt(50))
    tb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    _set_run(tb.text_frame, text, 28, RGBColor(0x22, 0x22, 0x22))


def _add_canvas_guide(slide) -> None:
    """Visualize canvas right/bottom edges with thin guides."""
    # (textbox can't draw lines well; skip — slide bounds visible already)
    pass


def build_slide_box_canvas(prs: Presentation) -> None:
    slide = _add_blank(prs)
    _add_title(slide, "box_canvas_overflow — box bbox が canvas を超える")

    # OK: 全 4 辺 canvas 内
    _add_label(slide, 60, 90, "OK: box は canvas 内に収まる")
    _add_text_box(slide, 60, 130, 360, 100, "OK 内に収まるテキスト", 24, TINT_OK)

    # NG: right
    _add_label(slide, 60, 270, "NG: box.right が canvas 右端 (1440) を超える")
    _add_text_box(slide, 1280, 310, 240, 100, "右はみ出し", 24, TINT_NG)

    # NG: bottom
    _add_label(slide, 60, 450, "NG: box.bottom が canvas 下端 (810) を超える")
    _add_text_box(slide, 60, 490, 360, 360, "下はみ出し", 24, TINT_NG)

    # NG: left
    _add_label(slide, 760, 270, "NG: box.left が負 (canvas 左外)")
    _add_text_box(slide, -120, 310, 360, 100, "左はみ出し", 24, TINT_NG)


def build_slide_text_box(prs: Presentation) -> None:
    slide = _add_blank(prs)
    _add_title(slide, "text_box_overflow — 折り返した text が box.height を超える")

    # OK: 短い text、box.height 余裕
    _add_label(slide, 60, 90, "OK: text 高さ ≤ box.height")
    _add_text_box(slide, 60, 130, 360, 120,
                  "短いテキストで box 高さに収まる", 24, TINT_OK)

    # NG: 長い text を低い box に詰め込む
    _add_label(slide, 60, 290, "NG: 折り返した text が box.height を超過")
    _add_text_box(
        slide, 60, 330, 360, 60,
        "これは折り返しが必要な長いテキストで、box の高さが足りないため複数行に渡ってあふれます。",
        24, TINT_NG,
    )

    # OK reference: autofit ON で suppress される (本ロジックでは fire しない)
    _add_label(slide, 760, 90, "OK: autofit ON は suppress (lint 対象外)")
    tb = slide.shapes.add_textbox(Pt(760), Pt(130), Pt(360), Pt(60))
    tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tb.fill.solid()
    tb.fill.fore_color.rgb = TINT_OK
    tb.line.color.rgb = RGBColor(0x88, 0x88, 0x88)
    _set_run(tb.text_frame, "autofit (TEXT_TO_FIT_SHAPE) が有効なので lint は suppress される", 24)


def build_slide_text_canvas(prs: Presentation) -> None:
    slide = _add_blank(prs)
    _add_title(slide, "text_canvas_overflow — text が canvas 右端を超える")

    # OK: word_wrap=False でも canvas 内
    _add_label(slide, 60, 90, "OK: word_wrap=False でも canvas 内")
    _add_text_box(slide, 60, 130, 360, 80,
                  "短文 word_wrap=False", 24, TINT_OK, word_wrap=False)

    # NG: word_wrap=False で text 全幅が canvas を超える
    _add_label(slide, 60, 260, "NG: word_wrap=False、text 全幅が canvas 右を超える")
    _add_text_box(
        slide, 900, 300, 200, 80,
        "とても長い1行テキストでword_wrap=Falseなのでcanvas右端を超えてしまう非常に長い文字列です",
        24, TINT_NG, word_wrap=False,
    )

    # NG: word_wrap=True だが 1 単語が canvas-remaining を超える
    _add_label(slide, 60, 450, "NG: word_wrap=True だが 1 単語が canvas-remaining を超える")
    _add_text_box(
        slide, 1100, 490, 200, 80,
        "Supercalifragilisticexpialidocious",
        24, TINT_NG, word_wrap=True,
    )


def main() -> int:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = _new_deck()
    build_slide_box_canvas(prs)
    build_slide_text_box(prs)
    build_slide_text_canvas(prs)
    prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
