#!/usr/bin/env python3
"""Regression tests for pptx_lint JSON evidence schema."""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))

import make_examples  # noqa: E402
import pptx_lint  # noqa: E402
from PIL import Image, ImageDraw  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_AUTO_SIZE  # noqa: E402
from pptx.util import Pt  # noqa: E402


SCHEMA_KEYS = {
    "check_id",
    "evidence",
    "fixability",
    "fixability_reason",
    "candidate_values",
    "manual_required_reason",
    "measurement_confidence",
}

FIXABILITY_ENUM = {
    "auto_fix_candidate",
    "manual_required",
    "not_applicable",
    "decorative_review",
}


def _make_schema_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    text = slide.shapes.add_textbox(Pt(40.03), Pt(40), Pt(320), Pt(90))
    text.fill.solid()
    text.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    text.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    run = text.text_frame.paragraphs[0].add_run()
    run.text = "Low contrast bad font"
    run.font.name = "Arial"
    run.font.size = Pt(13)
    run.font.color.rgb = RGBColor.from_string("999999")

    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(100.03), Pt(200), Pt(80), Pt(40))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor.from_string("EEEEEE")

    img = out.with_name("_schema_fixture.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        pic = slide.shapes.add_picture(str(img), Pt(420), Pt(40), width=Pt(40), height=Pt(40))
        make_examples._clear_alt_text(pic)
    finally:
        img.unlink(missing_ok=True)

    prs.save(str(out))


def _make_rendered_fixture(out: Path, image_dir: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40), Pt(40), Pt(300), Pt(80))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Rendered low contrast"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(12)
    prs.save(str(out))

    image_dir.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (720, 405), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((70, 70, 300, 92), fill="#BBBBBB")
    image.save(image_dir / "slide-01.png")


def _json_findings(path: Path, **kwargs) -> list[dict]:
    return [pptx_lint.finding_to_json_dict(f) for f in pptx_lint.lint_pptx(path, **kwargs)]


def _by_check(findings: list[dict], check: str) -> list[dict]:
    return [f for f in findings if f["check"] == check]


def main() -> int:
    failures: list[str] = []
    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        schema_fixture = tmp_dir / "schema-fixture.pptx"
        rendered_fixture = tmp_dir / "rendered-fixture.pptx"
        rendered_images = tmp_dir / "rendered-images"
        _make_schema_fixture(schema_fixture)
        _make_rendered_fixture(rendered_fixture, rendered_images)

        findings = _json_findings(schema_fixture)
        for finding in findings:
            detail = finding["detail"]
            missing = SCHEMA_KEYS - set(detail)
            if missing:
                failures.append(
                    f"{finding['check']} missing schema detail keys: {sorted(missing)}"
                )
            if detail.get("check_id") != finding["check"]:
                failures.append(f"{finding['check']} detail.check_id does not match finding.check")
            if not isinstance(detail.get("evidence"), dict):
                failures.append(f"{finding['check']} detail.evidence is not a dict")
            if detail.get("fixability") not in FIXABILITY_ENUM:
                failures.append(
                    f"{finding['check']} detail.fixability is not a schema enum string"
                )
            if detail.get("fixability") == "auto_fix_candidate":
                if not detail.get("candidate_values"):
                    failures.append(f"{finding['check']} auto-fix candidate_values missing")
            elif not detail.get("manual_required_reason"):
                failures.append(f"{finding['check']} manual_required_reason missing")

        autofit = _by_check(findings, "text_autofit_disabled")
        if not any(
            f["detail"]["fixability"] == "auto_fix_candidate"
            and f["detail"].get("fixability_rule") == "autofit"
            and f["detail"]["candidate_values"] == {"auto_size": "NONE"}
            for f in autofit
        ):
            failures.append("text_autofit_disabled did not expose auto-fix candidate")

        geometry = _by_check(findings, "geometry_rounding")
        if not any(
            f["detail"]["fixability"] == "auto_fix_candidate"
            and f["detail"].get("fixability_rule") == "geometry"
            and f["detail"]["candidate_values"].get("rounded_values_pt")
            for f in geometry
        ):
            failures.append("geometry_rounding did not expose safe auto-fix evidence")

        font_size = _by_check(findings, "font_size_scale")
        if not any(f["detail"]["candidate_values"].get("size_pt") for f in font_size):
            failures.append("font_size_scale did not expose candidate size")

        low_contrast = _by_check(findings, "low_contrast")
        if not any(
            f["detail"]["evidence"].get("foreground_hex")
            and f["detail"]["evidence"].get("background_hex")
            and f["detail"]["evidence"].get("original_run_color_hex")
            and f["detail"]["candidate_values"].get("candidate_token")
            and f["detail"]["candidate_values"].get("recalculated_ratio")
            for f in low_contrast
        ):
            failures.append("low_contrast did not preserve contrast evidence and candidate")

        alt_text = _by_check(findings, "alt_text_required")
        if not any(f["detail"].get("manual_required_reason") for f in alt_text):
            failures.append("alt_text_required did not expose manual-required reason")

        rendered_findings = _json_findings(
            rendered_fixture,
            rendered_image_dir=rendered_images,
        )
        rendered_low_contrast = [
            f
            for f in rendered_findings
            if f["check"] == "low_contrast"
            and f["detail"]["evidence"].get("measurement") == "rendered_image"
        ]
        if not rendered_low_contrast:
            failures.append("rendered fixture did not trigger rendered-image low_contrast")
        elif not any(
            f["detail"].get("measurement_confidence")
            and f["detail"]["evidence"].get("rendered_image_path")
            and f["detail"]["evidence"].get("foreground_hex")
            and f["detail"]["evidence"].get("background_hex")
            and "original_run_colors_hex" in f["detail"]["evidence"]
            for f in rendered_low_contrast
        ):
            failures.append("rendered low_contrast missing rendered evidence fields")

    if failures:
        print("FAIL:")
        for failure in failures:
            print(f"- {failure}")
        return 1
    print("OK: pptx_lint JSON evidence schema fields are present")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
