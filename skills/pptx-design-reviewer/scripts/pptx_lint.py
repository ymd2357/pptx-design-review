#!/usr/bin/env python3
"""PPTX lint MVP for slide-guideline-v1 compliance.

Checks
- box_canvas_overflow  (warning) text shape の box bbox が slide canvas からはみ出す
- text_box_overflow    (warning) box.width で折り返した text の必要 height が box.height を超える
- text_canvas_overflow (warning) canvas 端強制改行でも text が canvas 右端を超える
- overflow_shapes      (error)   shape element extends beyond slide canvas
- overflow_images      (error)   image element extends beyond slide canvas
- safe_text_area_text  (warning) text-bearing element outside safe text area
- text_autofit_disabled(error)   text-to-fit auto-size has an actual font shrink
- font_family          (warning) font name not in allowlist
- font_size_scale      (warning) font size not in allowed scale
- safe_margins        (warning) non-text element outside safe margins
- line_height         (warning) explicit paragraph line spacing not in allowed scale
- alignment_left_top  (warning) explicit text alignment differs from left/top
- geometry_rounding   (warning) geometry is not on integer pt coordinates
- image_upscale_ratio (warning) picture displayed larger than source pixels
- image_aspect_distortion (warning) picture aspect ratio differs from display box
- alt_text_required   (warning) meaningful image-like object has no alt text
- text_color_allowlist (warning) explicit text color not in allowlist
- background_color_palette (warning) explicit shape/table cell fill color not in palette
- contrast_ratio      (warning) explicit text/background colors fail contrast threshold
- low_contrast        (error)   explicit text/background colors are unreadably low contrast
- color_only_meaning (warning) similar unlabeled colored shapes rely on color alone
- heading_hierarchy_broken (warning) title/body hierarchy is structurally inconsistent
- key_area_cropped   (warning) picture crop metadata suggests important content may be clipped
- missing_required_element (warning) slide content has no machine-detected title/header
- reading_order      (warning) source order diverges from visual top-left order
- wrap_break_changes_meaning (warning) explicit line break splits a semantic unit
- text_overlap       (error)   text frames overlap each other
- object_overlap     (error)   non-text object bboxes overlap, excluding
                                structural containment
- object_gap_too_small (warning) adjacent object gap is below minimum spacing
- decorative_isolated_lines (warning) isolated decorative line/connector/arrow with no companion shape
- badge_alignment      (warning) short-text badge container is not centered (paragraph CENTER + vertical MIDDLE)
- inner_padding_imbalance (warning) child objects are unbalanced inside a container
- card_grid_consistency (warning) repeated card containers have inconsistent sizing or internal layout
- text_vertical_balance (warning) text fits but vertical padding/center balance is unnatural
- animation_present    (error)   slide contains <p:transition> or <p:timing>

Source of truth for thresholds is doc/slide-guideline-v1.yml.
Constants below mirror that file; keep them in sync if the guideline changes.

Usage
    python3 pptx_lint.py DECK.pptx
    python3 pptx_lint.py DECK.pptx --json
    python3 pptx_lint.py DECK.pptx --structure-json
    python3 pptx_lint.py DECK.pptx --severity error

Exit code
    0 = no errors (warnings allowed)
    1 = at least one error
    2 = invocation error (file missing, etc.)
"""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import sys
from collections import Counter, defaultdict
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Iterable, List, Optional

from PIL import Image
from pptx import Presentation
from pptx.dml.color import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN

# ---- Thresholds (mirror doc/slide-guideline-v1.yml) ------------------------

SLIDE_W_PT = 1440
SLIDE_H_PT = 810

SAFE_MARGIN_LEFT_PT = 81
SAFE_MARGIN_RIGHT_PT = 81
SAFE_MARGIN_TOP_PT = 40
SAFE_MARGIN_BOTTOM_PT = 80

SAFE_TEXT_AREA_PT = (
    SAFE_MARGIN_LEFT_PT,
    SAFE_MARGIN_TOP_PT,
    SLIDE_W_PT - SAFE_MARGIN_LEFT_PT - SAFE_MARGIN_RIGHT_PT,  # 1278
    SLIDE_H_PT - SAFE_MARGIN_TOP_PT - SAFE_MARGIN_BOTTOM_PT,  # 690
)

ALLOWED_FONT_FAMILIES = ("Noto Sans JP", "Avenir Next Arabic", "Nunito Sans")
ALLOWED_FONT_SIZES_PT = {80, 64, 56, 48, 40, 36, 32, 28, 24, 22, 20}
FONT_SIZE_TOL_PT = 1.0
ALLOWED_LINE_HEIGHTS_PT = {90, 66, 42, 36, 30, 24}
LINE_HEIGHT_TOL_PT = 2.0
OBJECT_OVERLAP_AREA_PT2_MIN = 1.0
STRUCTURAL_CONTAINMENT_OVERLAP_RATIO_MIN = 0.9
# Background shape ratio (canvas % covered) above which a shape is treated as
# a structural background container regardless of the child's overflow. Fixes
# the case where a canvas-sized background freeform appeared in
# object_overlap findings because the child text box overflowed the canvas
# and dropped the child's containment ratio below 0.9.
BACKGROUND_CONTAINER_CANVAS_RATIO_MIN = 0.9
OBJECT_GAP_DEFAULT_MIN_PT = 0.0  # Mirrors rules.slide.object_spacing.default_adjacent_gap_pt_min.
OBJECT_GAP_ELEMENT_MIN_PT = {
    "default": OBJECT_GAP_DEFAULT_MIN_PT,
    "peer_object": 8.0,  # Mirrors rules.slide.object_spacing.element_gap_pt_min.peer_object.
}
SEMANTIC_PAIR_GAP_PT = {
    "title_subtitle": 8.0,
    "label_group": 16.0,
}
SEMANTIC_TITLE_FONT_PT_MIN = 28.0
SEMANTIC_SUBTITLE_FONT_PT_MAX = 24.0
SEMANTIC_PAIR_X_ALIGN_TOL_PT = 24.0
SEMANTIC_PAIR_Y_ALIGN_TOL_PT = 8.0
SEMANTIC_LABEL_FONT_DIFF_PT_MAX = 1.5
INNER_PADDING_RATIO_MIN = 0.5
INNER_PADDING_RATIO_MAX = 2.0
INNER_PADDING_SIDE_MIN_PT = 4.0
INNER_PADDING_DEFAULT_TARGET_PT = 0.0
INNER_PADDING_ELEMENT_TARGET_PT = {
    "default": INNER_PADDING_DEFAULT_TARGET_PT,
    "card": 24.0,  # Mirrors tokens.component.card.padding_pt.
    "callout": 24.0,  # Mirrors tokens.component.callout.padding_pt.
}
INNER_PADDING_TARGET_TOL_PT = 4.0
CARD_GRID_GROUP_MIN = 2
CARD_GRID_CHILD_COUNT_MIN = 2
CARD_GRID_ROW_CENTER_TOL_PT = 36.0
CARD_GRID_TOP_TOL_PT = 4.0
CARD_GRID_SIZE_TOL_PT = 6.0
CARD_GRID_PADDING_TOL_PT = 8.0
CARD_GRID_CHILD_RELATIVE_TOL_PT = 14.0
DECORATIVE_LINE_PROXIMITY_PT_MAX = 24.0
DECORATIVE_THIN_BAR_THICKNESS_PT_MAX = 4.0
DECORATIVE_THIN_BAR_LENGTH_PT_MIN = 20.0
BADGE_SHAPE_MAX_WIDTH_PT = 240.0
BADGE_SHAPE_MAX_HEIGHT_PT = 120.0
BADGE_SHAPE_ASPECT_MIN = 0.5
BADGE_SHAPE_ASPECT_MAX = 2.5
BADGE_TEXT_LENGTH_MAX = 20
COVER_BRAND_MARK_MAX_W_PT = 420.0
COVER_BRAND_MARK_MAX_H_PT = 150.0
COVER_BRAND_MARK_MAX_X_PT = 100.0
COVER_BRAND_MARK_MAX_Y_PT = 120.0
TEXT_VERTICAL_BALANCE_DEAD_SPACE_RATIO_MAX = 0.40
TEXT_VERTICAL_BALANCE_DEAD_SPACE_PT_MAX = 60.0
TEXT_VERTICAL_BALANCE_MIDDLE_MARGIN_ASYMMETRY_PT_MAX = 12.0
TEXT_VERTICAL_BALANCE_CENTER_OFFSET_RATIO_MAX = 0.20
TEXT_VERTICAL_BALANCE_MIN_BOX_HEIGHT_PT = 30.0
DEFAULT_LINE_HEIGHT_MULTIPLIER = 1.2
TITLE_ZONE_BOTTOM_PT = 150.0
TITLE_FONT_SIZE_MIN_PT = 28.0
PROMINENT_TITLE_FONT_SIZE_MIN_PT = 40.0
SECTION_DIVIDER_TEXT_COUNT_MAX = 2
SECTION_DIVIDER_TITLE_Y_MIN_PT = 220.0
SECTION_DIVIDER_TITLE_Y_MAX_PT = 520.0
SECTION_DIVIDER_TITLE_WIDTH_MIN_PT = 600.0
SECTION_DIVIDER_TITLE_CENTER_TOL_PT = 140.0
SECTION_DIVIDER_NON_TEXT_COUNT_MAX = 2
COVER_TEXT_COUNT_MIN = 3
COVER_TITLE_Y_MIN_PT = 160.0
COVER_TITLE_Y_MAX_PT = 620.0
COVER_TITLE_WIDTH_MIN_PT = 600.0
HEADING_BODY_FONT_DELTA_PT = 4.0
READING_ORDER_TOP_BUCKET_PT = 24.0
READING_ORDER_INVERSION_MIN = 2
READING_ORDER_UNIT_VERTICAL_GAP_PT = 150.0
READING_ORDER_UNIT_CENTER_X_MAX_PT = 220.0
READING_ORDER_UNIT_CENTER_X_RATIO = 0.35
READING_ORDER_UNIT_EDGE_X_TOL_PT = 60.0
READING_ORDER_FOOTER_Y_MIN_PT = 720.0
READING_ORDER_FULL_WIDTH_TEXT_PT = 900.0
SAFE_TEXT_HEADER_Y_MAX_PT = 45.0
SAFE_TEXT_FOOTER_Y_MIN_PT = 720.0
SAFE_TEXT_FULL_WIDTH_MIN_PT = 1200.0
SAFE_TEXT_FULL_WIDTH_X_MAX_PT = 140.0
SAFE_TEXT_FULL_WIDTH_RIGHT_OVERFLOW_MAX_PT = 80.0
SAFE_TEXT_PAGE_NUMBER_MAX_W_PT = 80.0
SAFE_TEXT_PAGE_NUMBER_MAX_H_PT = 40.0
IMAGE_CROP_SIDE_RATIO_MAX = 0.18
IMAGE_CROP_AXIS_TOTAL_RATIO_MAX = 0.30
COLOR_ONLY_GROUP_GAP_PT = 18.0
COLOR_ONLY_SIZE_TOL_PT = 12.0
MAX_IMAGE_UPSCALE_RATIO = 1.0
MAX_IMAGE_ASPECT_DELTA_RATIO = 0.05
DECORATIVE_RASTER_KEYWORDS = (
    "background",
    "bg",
    "decoration",
    "decorative",
    "footer",
    "gradient",
    "header",
)
PX_PER_PT = 96 / 72

# ---- Color palette (DS-COLOR-001: YAML 一本化, 2026-05-20) ----------------
#
# doc/slide-guideline-v1.yml -> rules.color.lint_palette が lint の 6 色定数
# (allowed_text/fill_colors_hex, *_color_token_by_hex,
#  contrast_repair / fill_repair _color_families) の唯一の正本。
# 旧来ここに hardcoded していた Python 辞書/タプルは DS-COLOR-001 の完了に
# 伴い撤去 (2026-05-20)。design_system_loader.load_lint_palette() が
# import 時に YAML を読んで同じ名前空間に展開する。
from design_system_loader import load_lint_palette as _load_lint_palette

_LINT_PALETTE = _load_lint_palette()
ALLOWED_TEXT_COLORS_HEX = set(_LINT_PALETTE.allowed_text_colors_hex)
ALLOWED_FILL_COLORS_HEX = set(_LINT_PALETTE.allowed_fill_colors_hex)
TEXT_COLOR_TOKEN_BY_HEX = dict(_LINT_PALETTE.text_color_token_by_hex)
FILL_COLOR_TOKEN_BY_HEX = dict(_LINT_PALETTE.fill_color_token_by_hex)
CONTRAST_REPAIR_COLOR_FAMILIES = _LINT_PALETTE.contrast_repair_color_families
FILL_REPAIR_COLOR_FAMILIES = _LINT_PALETTE.fill_repair_color_families


# Background-side hue families (sorted light→dark within each family) for fill repair.


CONTRAST_RATIO_NORMAL_TEXT_MIN = 4.5
CONTRAST_RATIO_LARGE_TEXT_MIN = 3.0
LOW_CONTRAST_RATIO_MAX = 3.0
LARGE_TEXT_MIN_SIZE_PT = 36.0
DEFAULT_SLIDE_BACKGROUND_HEX = "#FFFFFF"
RENDERED_CONTRAST_PALETTE_COLORS = 16
RENDERED_CONTRAST_MIN_FOREGROUND_PIXELS = 8
RENDERED_CONTRAST_MIN_FOREGROUND_RATIO = 0.001
RENDERED_CONTRAST_MAX_SAMPLE_PX = 240
RENDERED_CONTRAST_MAX_RAW_SAMPLE_PIXELS = 1_000_000
RENDERED_CONTRAST_EXPECTED_COLOR_DISTANCE_MAX = 18
RENDERED_CONTRAST_BACKGROUND_UNIFORM_DISTANCE_MAX = 12
RENDERED_CONTRAST_BACKGROUND_UNIFORMITY_MIN = 0.70



# Tolerance for fp comparisons in pt. Geometry rounding policy is 0.5pt.
TOL_PT = 0.5
GEOMETRY_INTEGER_TOL_PT = 0.01
SLIDE_ASPECT_TOL = 0.01

PML_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# ---- Result types ----------------------------------------------------------


@dataclass
class Finding:
    severity: str  # "error" | "warning"
    check: str
    slide_index: int  # 1-based
    slide_id: Optional[int]
    shape_id: Optional[int]
    shape_name: Optional[str]
    message: str
    detail: dict = field(default_factory=dict)


@dataclass(frozen=True)
class LintPolicy:
    check_vertical_text_anchor: bool = False


LINT_PROFILES = {
    "default": LintPolicy(),
    "strict": LintPolicy(check_vertical_text_anchor=True),
}


@dataclass(frozen=True)
class LintContext:
    """Coordinate normalization from actual slide size to guideline base size."""

    actual_w_pt: float
    actual_h_pt: float
    scale_x: float
    scale_y: float
    proportional_to_base: bool
    policy: LintPolicy

    @property
    def font_scale(self) -> float:
        return (self.scale_x + self.scale_y) / 2


@dataclass
class ShapeRecord:
    shape: Any
    actual_bbox_pt: tuple
    bbox_pt: tuple
    kind: str
    source_order_index: int = 0


@dataclass(frozen=True)
class StructureRelation:
    relation: str
    container: ShapeRecord
    child: ShapeRecord
    overlap_area_pt2: float
    overlap_bbox_pt: tuple
    metadata: dict = field(default_factory=dict)


# ---- Helpers ---------------------------------------------------------------


def emu_to_pt(emu: int) -> float:
    return emu / 12700.0


def iter_shapes(shapes) -> Iterable:
    """Recursively iterate shapes, descending into groups."""
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(s.shapes)
        else:
            yield s


def animation_markers(slide) -> List[str]:
    markers: List[str] = []
    for tag in ("transition", "timing"):
        if slide.element.find(f"p:{tag}", PML_NS) is not None:
            markers.append(f"p:{tag}")
    return markers


def has_animation(slide) -> bool:
    return bool(animation_markers(slide))


def _font_name_allowed(name: str) -> bool:
    """Match family even when the run carries a weight suffix.

    Some source decks emit names like "Noto Sans JP Medium" or
    "Noto Sans JP Bold". Treat those as the same family as "Noto Sans JP".
    """
    for family in ALLOWED_FONT_FAMILIES:
        if name == family or name.startswith(family + " "):
            return True
    return False


def _run_explicit_typefaces(run) -> list[tuple[str, str]]:
    try:
        r_pr = run._r.rPr
    except AttributeError:
        return []
    if r_pr is None:
        return []
    typefaces: list[tuple[str, str]] = []
    for script in ("latin", "ea"):
        node = r_pr.find(f"{A_NS}{script}")
        if node is None:
            continue
        name = node.get("typeface")
        if name:
            typefaces.append((script, name))
    if not typefaces:
        name = run.font.name
        if name:
            typefaces.append(("font", name))
    return typefaces


def _nearest_allowed_font_size(size_pt: float) -> int:
    return min(ALLOWED_FONT_SIZES_PT, key=lambda allowed: abs(size_pt - allowed))


def _font_size_allowed(size_pt: float) -> bool:
    return abs(size_pt - _nearest_allowed_font_size(size_pt)) <= FONT_SIZE_TOL_PT


def _nearest_allowed_line_height(line_height_pt: float) -> int:
    return min(ALLOWED_LINE_HEIGHTS_PT, key=lambda allowed: abs(line_height_pt - allowed))


def _line_height_allowed(line_height_pt: float) -> bool:
    return abs(line_height_pt - _nearest_allowed_line_height(line_height_pt)) <= LINE_HEIGHT_TOL_PT


def _normalized_length_pt(value, scale: float) -> float:
    if value is None:
        return 0.0
    try:
        return value.pt * scale
    except AttributeError:
        return emu_to_pt(value) * scale


def _text_frame_dominant_font_size_pt(ctx: LintContext, text_frame) -> Optional[float]:
    first_size_pt: Optional[float] = None
    size_weights: dict[float, int] = {}
    for para in text_frame.paragraphs:
        for run in para.runs:
            size = run.font.size
            if size is None:
                continue
            size_pt = round(size.pt * ctx.font_scale, 4)
            if first_size_pt is None:
                first_size_pt = size_pt
            weight = len(run.text.strip()) or 1
            size_weights[size_pt] = size_weights.get(size_pt, 0) + weight
    if size_weights:
        return max(size_weights.items(), key=lambda item: item[1])[0]
    return first_size_pt


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text_frame.text or "").strip()


def _text_excerpt(shape, limit: int = 80) -> str:
    text = " ".join(_shape_text(shape).split())
    if len(text) <= limit:
        return text
    return text[: limit - 3] + "..."


def _estimate_text_width_pt(text: str, font_size_pt: float) -> float:
    """Rough horizontal-extent estimate for mixed CJK + Latin text at a given
    font size. CJK/Hiragana/Katakana glyphs count as full-width (1.0 × size);
    everything else (Latin / digits / spaces / punctuation) as ~0.55× width.
    The estimate is intentionally conservative on the wide side so that the
    widen-to-fit auto-fix proposes a target wide enough for the real glyphs
    rather than a value that ends up overflowing again.
    """
    if not text:
        return 0.0
    full = 0
    half = 0
    for ch in text:
        if (
            "぀" <= ch <= "ヿ"          # Hiragana + Katakana
            or "㐀" <= ch <= "鿿"        # CJK Unified Ideographs
            or "＀" <= ch <= "￯"        # Halfwidth/Fullwidth forms
            or "　" <= ch <= "〿"        # CJK Symbols & Punctuation
        ):
            full += 1
        else:
            half += 1
    return font_size_pt * (full + half * 0.55)


def _text_frame_has_inline_break(text_frame) -> bool:
    for para in text_frame.paragraphs:
        for run in para.runs:
            if "\n" in (run.text or "") or "\v" in (run.text or ""):
                return True
    return False


def _text_frame_inline_joined_text(text_frame) -> str:
    """Concatenate all inline (intra-run) break-separated text into a single
    line, simulating what `_replace_text_breaks` would produce post-fix."""
    parts: list[str] = []
    for para in text_frame.paragraphs:
        for run in para.runs:
            raw = run.text or ""
            collapsed = re.sub(r"[\n\v]+", " ", raw)
            parts.append(collapsed)
        parts.append(" ")
    return "".join(parts).strip()


def _text_frame_margin_pt_value(text_frame, attr_name: str) -> float:
    """text_frame.margin_{top,bottom,left,right} を pt で返す。
    属性が None ・取得不能なら 0pt fallback。
    """
    try:
        value = getattr(text_frame, attr_name)
    except (AttributeError, ValueError):
        return 0.0
    if value is None:
        return 0.0
    try:
        return float(value.pt)
    except (AttributeError, TypeError, ValueError):
        return 0.0


def _text_frame_margin_pt(text_frame) -> tuple[float, float]:
    left = text_frame.margin_left
    right = text_frame.margin_right
    left_pt = left.pt if left is not None else 0.0
    right_pt = right.pt if right is not None else 0.0
    return left_pt, right_pt


def _text_frame_lines(text_frame) -> list[str]:
    lines: list[str] = []
    for para in text_frame.paragraphs:
        para_lines: list[str] = []
        for run in para.runs:
            para_lines.extend(str(run.text or "").splitlines())
        if not para_lines:
            para_lines = [para.text or ""]
        lines.extend(para_lines)
    return lines


def _paragraph_line_height_pt(ctx: LintContext, para, font_size_pt: float) -> float:
    line_spacing = para.line_spacing
    if line_spacing is None or isinstance(line_spacing, float):
        return font_size_pt * DEFAULT_LINE_HEIGHT_MULTIPLIER
    try:
        return line_spacing.pt * ctx.font_scale
    except AttributeError:
        return font_size_pt * DEFAULT_LINE_HEIGHT_MULTIPLIER


def _vertical_anchor_name(anchor) -> str:
    if anchor is None:
        return "NONE"
    if anchor == MSO_VERTICAL_ANCHOR.TOP:
        return "TOP"
    if anchor == MSO_VERTICAL_ANCHOR.MIDDLE:
        return "MIDDLE"
    if anchor == MSO_VERTICAL_ANCHOR.BOTTOM:
        return "BOTTOM"
    return str(anchor)


def _shape_has_visible_text(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    return bool(shape.text_frame.text.strip())


def _shape_has_visible_fill_or_line(shape) -> bool:
    try:
        fill_type = shape.fill.type
    except AttributeError:
        fill_type = None
    if fill_type is not None and fill_type != MSO_FILL.BACKGROUND:
        return True

    try:
        line_fill_type = shape.line.fill.type
        line_width = shape.line.width
    except AttributeError:
        return False
    if line_fill_type is None:
        return False
    if line_width is None:
        return True
    return line_width.pt > TOL_PT


def _shape_kind(shape) -> str:
    if _shape_has_visible_text(shape):
        return "text"
    if getattr(shape, "has_table", False):
        return "table"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    return "shape"


def _shape_label(record: ShapeRecord) -> str:
    name = getattr(record.shape, "name", None) or record.kind
    shape_id = getattr(record.shape, "shape_id", None)
    if shape_id is None:
        return name
    return f"{name}#{shape_id}"


def _bbox_edges(bbox: tuple) -> tuple[float, float, float, float]:
    x, y, w, h = bbox
    return x, y, x + w, y + h


def _bbox_area(bbox: tuple) -> float:
    return max(0.0, bbox[2]) * max(0.0, bbox[3])


def _bbox_center(bbox: tuple) -> tuple[float, float]:
    x, y, w, h = bbox
    return x + w / 2, y + h / 2


def _point_in_bbox(point: tuple[float, float], bbox: tuple, tolerance: float = TOL_PT) -> bool:
    x, y = point
    x1, y1, x2, y2 = _bbox_edges(bbox)
    return x1 - tolerance <= x <= x2 + tolerance and y1 - tolerance <= y <= y2 + tolerance


def _bbox_overflow_sides(outer: tuple, inner: tuple) -> dict[str, float]:
    ox1, oy1, ox2, oy2 = _bbox_edges(outer)
    ix1, iy1, ix2, iy2 = _bbox_edges(inner)
    overflow = {
        "left": max(0.0, ox1 - ix1),
        "top": max(0.0, oy1 - iy1),
        "right": max(0.0, ix2 - ox2),
        "bottom": max(0.0, iy2 - oy2),
    }
    return {side: value for side, value in overflow.items() if value > TOL_PT}


def _intersection_bbox(a: tuple, b: tuple) -> Optional[tuple]:
    ax1, ay1, ax2, ay2 = _bbox_edges(a)
    bx1, by1, bx2, by2 = _bbox_edges(b)
    x1 = max(ax1, bx1)
    y1 = max(ay1, by1)
    x2 = min(ax2, bx2)
    y2 = min(ay2, by2)
    if x2 <= x1 or y2 <= y1:
        return None
    return (x1, y1, x2 - x1, y2 - y1)


def _axis_gap(a_min: float, a_max: float, b_min: float, b_max: float) -> float:
    if a_max < b_min:
        return b_min - a_max
    if b_max < a_min:
        return a_min - b_max
    return 0.0


def _design_element(record: ShapeRecord) -> str:
    title, descr = _alt_text_values(record.shape)
    haystack = " ".join(
        str(value or "")
        for value in (getattr(record.shape, "name", ""), title, descr)
    )
    match = re.search(r"(?:^|[\s;])ds:element=([A-Za-z0-9_.-]+)", haystack)
    return match.group(1) if match else "default"


def _object_gap_threshold_pt(first: ShapeRecord, second: ShapeRecord) -> tuple[float, str, str]:
    first_element = _design_element(first)
    second_element = _design_element(second)
    first_threshold = OBJECT_GAP_ELEMENT_MIN_PT.get(first_element, OBJECT_GAP_DEFAULT_MIN_PT)
    second_threshold = OBJECT_GAP_ELEMENT_MIN_PT.get(second_element, OBJECT_GAP_DEFAULT_MIN_PT)
    return max(first_threshold, second_threshold), first_element, second_element


def _shape_dominant_font_size_pt(ctx: LintContext, shape) -> Optional[float]:
    if not getattr(shape, "has_text_frame", False):
        return None
    return _text_frame_dominant_font_size_pt(ctx, shape.text_frame)


def _semantic_pair_kind(
    ctx: LintContext,
    first: ShapeRecord,
    second: ShapeRecord,
    axis: str,
) -> Optional[str]:
    """Classify (first, second) into a known semantic gap pair when applicable.

    Title-subtitle: a large-font text directly above a smaller-font text with
    similar left-x; both must have visible text.
    Label group: two text shapes with the same font size sitting in the same
    horizontal row (vertically adjacent on the cross-axis), adjacent
    horizontally. Used to enforce inter-label spacing.
    """
    if first.kind != "text" or second.kind != "text":
        return None
    if not _shape_text(first.shape) or not _shape_text(second.shape):
        return None
    size_a = _shape_dominant_font_size_pt(ctx, first.shape)
    size_b = _shape_dominant_font_size_pt(ctx, second.shape)
    if size_a is None or size_b is None:
        return None

    if axis == "vertical":
        if first.bbox_pt[1] <= second.bbox_pt[1]:
            upper_rec, upper_size, lower_rec, lower_size = first, size_a, second, size_b
        else:
            upper_rec, upper_size, lower_rec, lower_size = second, size_b, first, size_a
        if (
            upper_size >= SEMANTIC_TITLE_FONT_PT_MIN
            and lower_size <= SEMANTIC_SUBTITLE_FONT_PT_MAX
            and lower_size < upper_size
            and abs(upper_rec.bbox_pt[0] - lower_rec.bbox_pt[0]) <= SEMANTIC_PAIR_X_ALIGN_TOL_PT
        ):
            return "title_subtitle"
        return None

    # axis == "horizontal" — same-row label peers.
    if abs(size_a - size_b) > SEMANTIC_LABEL_FONT_DIFF_PT_MAX:
        return None
    cy_a = first.bbox_pt[1] + first.bbox_pt[3] / 2
    cy_b = second.bbox_pt[1] + second.bbox_pt[3] / 2
    if abs(cy_a - cy_b) > SEMANTIC_PAIR_Y_ALIGN_TOL_PT:
        return None
    return "label_group"


def _bbox_contains(outer: tuple, inner: tuple, tolerance: float = TOL_PT) -> bool:
    ox1, oy1, ox2, oy2 = _bbox_edges(outer)
    ix1, iy1, ix2, iy2 = _bbox_edges(inner)
    return (
        ix1 >= ox1 - tolerance
        and iy1 >= oy1 - tolerance
        and ix2 <= ox2 + tolerance
        and iy2 <= oy2 + tolerance
    )


def _shape_record_detail(record: ShapeRecord) -> dict:
    return {
        "shape_id": getattr(record.shape, "shape_id", None),
        "shape_name": getattr(record.shape, "name", None),
        "design_element": _design_element(record),
        "kind": record.kind,
        "source_order_index": record.source_order_index,
        "bbox_pt": [round(v, 2) for v in record.bbox_pt],
    }


def _shape_record_detail_for_card_role(record: ShapeRecord) -> dict:
    """Card-grid-flavored detail: includes fill / anchor / text length so the
    fix-side role classifier (icon / header / body / badge / decorative) has
    enough information without re-reading the PPTX.
    """
    base = _shape_record_detail(record)
    shape = record.shape
    if record.kind in {"shape", "text"}:
        try:
            base["solid_fill_hex"] = _solid_shape_fill_rgb_hex(shape)
        except (AttributeError, TypeError, ValueError):
            base["solid_fill_hex"] = None
    if getattr(shape, "has_text_frame", False):
        text = _shape_text(shape)
        base["has_text"] = bool(text)
        base["text_length"] = len(text)
        try:
            base["vertical_anchor"] = _vertical_anchor_name(shape.text_frame.vertical_anchor)
        except (AttributeError, TypeError, ValueError):
            base["vertical_anchor"] = None
    else:
        base["has_text"] = False
        base["text_length"] = 0
    return base


def _structure_relation_detail(relation: StructureRelation) -> dict:
    detail = {
        "relation": relation.relation,
        "container": _shape_record_detail(relation.container),
        "child": _shape_record_detail(relation.child),
        "overlap_area_pt2": round(relation.overlap_area_pt2, 2),
        "overlap_bbox_pt": [round(v, 2) for v in relation.overlap_bbox_pt],
    }
    detail.update(relation.metadata)
    return detail


def _containment_relation(
    first: ShapeRecord,
    second: ShapeRecord,
    overlap_bbox: tuple,
) -> Optional[StructureRelation]:
    if _bbox_contains(first.bbox_pt, second.bbox_pt):
        return StructureRelation(
            relation="contains",
            container=first,
            child=second,
            overlap_area_pt2=overlap_bbox[2] * overlap_bbox[3],
            overlap_bbox_pt=overlap_bbox,
        )
    if _bbox_contains(second.bbox_pt, first.bbox_pt):
        return StructureRelation(
            relation="contains",
            container=second,
            child=first,
            overlap_area_pt2=overlap_bbox[2] * overlap_bbox[3],
            overlap_bbox_pt=overlap_bbox,
        )
    overlap_area = overlap_bbox[2] * overlap_bbox[3]
    canvas_area = float(SLIDE_W_PT * SLIDE_H_PT)
    canvas_bbox = (0.0, 0.0, float(SLIDE_W_PT), float(SLIDE_H_PT))
    for container, child in ((first, second), (second, first)):
        # (1) Background-shaped container: if it covers >= 90% of the canvas
        # and the child's visible portion sits inside it, treat the relation
        # as structural containment regardless of how much the child overflows
        # the canvas. This is the fix for the slide-background freeform
        # appearing in object_overlap findings.
        container_canvas_ratio = _bbox_area(container.bbox_pt) / canvas_area if canvas_area else 0
        if (
            container_canvas_ratio >= BACKGROUND_CONTAINER_CANVAS_RATIO_MIN
            and _point_in_bbox(_bbox_center(child.bbox_pt), container.bbox_pt)
        ):
            return StructureRelation(
                relation="contains_background",
                container=container,
                child=child,
                overlap_area_pt2=overlap_area,
                overlap_bbox_pt=overlap_bbox,
                metadata={
                    "container_canvas_ratio": round(container_canvas_ratio, 4),
                    "threshold_canvas_ratio": BACKGROUND_CONTAINER_CANVAS_RATIO_MIN,
                },
            )

        # (2) Standard containment by ratio, but measure against the child's
        # **canvas-visible** area (= child clipped to the slide canvas). A
        # child that overflows the canvas should not lose its containment
        # status — what matters is whether the part of the child that is
        # actually rendered sits inside the container.
        clipped = _intersection_bbox(child.bbox_pt, canvas_bbox) or child.bbox_pt
        visible_area = _bbox_area(clipped)
        if visible_area <= 0:
            continue
        overlap_ratio = overlap_area / visible_area
        if (
            overlap_ratio >= STRUCTURAL_CONTAINMENT_OVERLAP_RATIO_MIN
            and _point_in_bbox(_bbox_center(child.bbox_pt), container.bbox_pt)
        ):
            return StructureRelation(
                relation="contains_with_child_overflow",
                container=container,
                child=child,
                overlap_area_pt2=overlap_area,
                overlap_bbox_pt=overlap_bbox,
                metadata={
                    "child_overlap_ratio": round(overlap_ratio, 4),
                    "child_overflow_pt": {
                        side: round(value, 2)
                        for side, value in _bbox_overflow_sides(
                            container.bbox_pt,
                            child.bbox_pt,
                        ).items()
                    },
                    "threshold_overlap_ratio": STRUCTURAL_CONTAINMENT_OVERLAP_RATIO_MIN,
                },
            )
    return None


def _rgb_hex(color_format) -> Optional[str]:
    """Return #RRGGBB for explicit RGB colors; inherited/theme colors are skipped."""
    try:
        if color_format.type != MSO_COLOR_TYPE.RGB:
            return None
        rgb = color_format.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    if rgb is None:
        return None
    return f"#{str(rgb).upper()}"


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    color = hex_color.lstrip("#")
    return int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)


def _linearized_srgb_channel(value: int) -> float:
    normalized = value / 255
    if normalized <= 0.03928:
        return normalized / 12.92
    return ((normalized + 0.055) / 1.055) ** 2.4


def _relative_luminance(hex_color: str) -> float:
    r, g, b = _hex_to_rgb(hex_color)
    return (
        0.2126 * _linearized_srgb_channel(r)
        + 0.7152 * _linearized_srgb_channel(g)
        + 0.0722 * _linearized_srgb_channel(b)
    )


def _contrast_ratio(foreground_hex: str, background_hex: str) -> float:
    fg_lum = _relative_luminance(foreground_hex)
    bg_lum = _relative_luminance(background_hex)
    lighter = max(fg_lum, bg_lum)
    darker = min(fg_lum, bg_lum)
    return (lighter + 0.05) / (darker + 0.05)


def _nearest_color_hex(source_hex: str, candidates: Iterable[str]) -> str:
    source = _hex_to_rgb(source_hex)
    return min(candidates, key=lambda candidate: _rgb_distance(source, _hex_to_rgb(candidate)))


def _contrast_repair_family(foreground_hex: str) -> tuple[str, tuple[str, ...]]:
    source = _hex_to_rgb(foreground_hex)
    best: tuple[float, str, tuple[str, ...]] | None = None
    for family_name, colors in CONTRAST_REPAIR_COLOR_FAMILIES:
        for color in colors:
            distance = _rgb_distance(source, _hex_to_rgb(color))
            if best is None or distance < best[0]:
                best = (distance, family_name, colors)
    if best is None:
        return ("neutral_black", CONTRAST_REPAIR_COLOR_FAMILIES[0][1])
    return best[1], best[2]


def _fill_repair_family(background_hex: str) -> tuple[str, tuple[str, ...]]:
    source = _hex_to_rgb(background_hex)
    best: tuple[float, str, tuple[str, ...]] | None = None
    for family_name, colors in FILL_REPAIR_COLOR_FAMILIES:
        for color in colors:
            distance = _rgb_distance(source, _hex_to_rgb(color))
            if best is None or distance < best[0]:
                best = (distance, family_name, colors)
    if best is None:
        return ("neutral_black", FILL_REPAIR_COLOR_FAMILIES[0][1])
    return best[1], best[2]


def _luminance_delta(a_hex: str, b_hex: str) -> float:
    return abs(_relative_luminance(a_hex) - _relative_luminance(b_hex))


def _contrast_foreground_option(
    foreground_hex: str,
    background_hex: str,
    required_ratio: float,
) -> Optional[dict]:
    family_name, family_colors = _contrast_repair_family(foreground_hex)
    passing = [
        color
        for color in family_colors
        if color in ALLOWED_TEXT_COLORS_HEX
        and _contrast_ratio(color, background_hex) >= required_ratio
    ]
    candidate_group = f"hue_family:{family_name}"
    if not passing:
        neutral_family = CONTRAST_REPAIR_COLOR_FAMILIES[0][1]
        passing = [
            color
            for color in neutral_family
            if color in ALLOWED_TEXT_COLORS_HEX
            and _contrast_ratio(color, background_hex) >= required_ratio
        ]
        if not passing:
            return None
        candidate_group = "neutral_fallback"
    candidate = _nearest_color_hex(foreground_hex, passing)
    return {
        "from_hex": foreground_hex,
        "to_hex": candidate,
        "to_token": TEXT_COLOR_TOKEN_BY_HEX.get(candidate),
        "recalculated_ratio": round(_contrast_ratio(candidate, background_hex), 2),
        "luminance_delta": round(_luminance_delta(foreground_hex, candidate), 4),
        "repair_candidate_group": candidate_group,
    }


def _contrast_background_option(
    foreground_hex: str,
    background_hex: str,
    required_ratio: float,
) -> Optional[dict]:
    family_name, family_colors = _fill_repair_family(background_hex)
    passing = [
        color
        for color in family_colors
        if color in ALLOWED_FILL_COLORS_HEX
        and _contrast_ratio(foreground_hex, color) >= required_ratio
    ]
    candidate_group = f"hue_family:{family_name}"
    if not passing:
        neutral_family = FILL_REPAIR_COLOR_FAMILIES[0][1]
        passing = [
            color
            for color in neutral_family
            if color in ALLOWED_FILL_COLORS_HEX
            and _contrast_ratio(foreground_hex, color) >= required_ratio
        ]
        if not passing:
            return None
        candidate_group = "neutral_fallback"
    candidate = _nearest_color_hex(background_hex, passing)
    return {
        "from_hex": background_hex,
        "to_hex": candidate,
        "to_token": FILL_COLOR_TOKEN_BY_HEX.get(candidate),
        "recalculated_ratio": round(_contrast_ratio(foreground_hex, candidate), 2),
        "luminance_delta": round(_luminance_delta(background_hex, candidate), 4),
        "repair_candidate_group": candidate_group,
    }


def _contrast_candidate(
    foreground_hex: str,
    background_hex: str,
    required_ratio: float,
) -> Optional[dict]:
    fg_option = _contrast_foreground_option(foreground_hex, background_hex, required_ratio)
    bg_option = _contrast_background_option(foreground_hex, background_hex, required_ratio)
    if not fg_option and not bg_option:
        return None

    if fg_option and bg_option:
        preferred = (
            "foreground"
            if fg_option["luminance_delta"] <= bg_option["luminance_delta"]
            else "background"
        )
    elif fg_option:
        preferred = "foreground"
    else:
        preferred = "background"

    if preferred == "foreground":
        top_fg = fg_option["to_hex"]
        top_bg = background_hex
        top_ratio = fg_option["recalculated_ratio"]
        top_group = fg_option["repair_candidate_group"]
    else:
        top_fg = foreground_hex
        top_bg = bg_option["to_hex"]
        top_ratio = bg_option["recalculated_ratio"]
        top_group = bg_option["repair_candidate_group"]

    return {
        "foreground_hex": top_fg,
        "foreground_token": TEXT_COLOR_TOKEN_BY_HEX.get(top_fg),
        "background_hex": top_bg,
        "background_token": FILL_COLOR_TOKEN_BY_HEX.get(top_bg),
        "recalculated_ratio": top_ratio,
        "required_ratio": required_ratio,
        "repair_candidate_group": top_group,
        "preferred_strategy": preferred,
        "foreground_option": fg_option,
        "background_option": bg_option,
        "selection_policy": "rules.color.repair_candidates.hue_preserving_min_luminance_delta",
    }


def _allowed_text_color_candidate(color_hex: str) -> dict:
    candidate = _nearest_color_hex(color_hex, ALLOWED_TEXT_COLORS_HEX)
    return {
        "color_hex": candidate,
        "color_token": TEXT_COLOR_TOKEN_BY_HEX.get(candidate),
        "allowed_colors_hex": sorted(ALLOWED_TEXT_COLORS_HEX),
    }


def _allowed_fill_color_candidate(color_hex: str) -> dict:
    candidate = _nearest_color_hex(color_hex, ALLOWED_FILL_COLORS_HEX)
    return {
        "color_hex": candidate,
        "color_token": FILL_COLOR_TOKEN_BY_HEX.get(candidate),
        "allowed_colors_hex": sorted(ALLOWED_FILL_COLORS_HEX),
    }


def _rgb_tuple_to_hex(rgb: tuple[int, int, int]) -> str:
    return "#{:02X}{:02X}{:02X}".format(*rgb)


def _rgb_distance(first: tuple[int, int, int], second: tuple[int, int, int]) -> float:
    return sum((a - b) ** 2 for a, b in zip(first, second)) ** 0.5


def _rendered_slide_image_path(image_dir: Path, slide_idx: int) -> Optional[Path]:
    candidates = [
        image_dir / f"slide-{slide_idx:02d}.png",
        image_dir / f"slide-{slide_idx:03d}.png",
        image_dir / f"slide_{slide_idx:02d}.png",
        image_dir / f"slide_{slide_idx:03d}.png",
        image_dir / f"{slide_idx}.png",
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def _rendered_crop_for_bbox(ctx: LintContext, image: Image.Image, bbox_pt: tuple) -> Optional[Image.Image]:
    x, y, w, h = bbox_pt
    left = int(max(0, round(x / ctx.actual_w_pt * image.width)))
    top = int(max(0, round(y / ctx.actual_h_pt * image.height)))
    right = int(min(image.width, round((x + w) / ctx.actual_w_pt * image.width)))
    bottom = int(min(image.height, round((y + h) / ctx.actual_h_pt * image.height)))
    if right <= left or bottom <= top:
        return None
    return image.crop((left, top, right, bottom)).convert("RGB")


def _sample_rendered_pixels(crop: Image.Image) -> tuple[Counter, int]:
    if crop.width <= 0 or crop.height <= 0:
        return Counter(), 0

    total_pixels = crop.width * crop.height
    if total_pixels <= RENDERED_CONTRAST_MAX_RAW_SAMPLE_PIXELS:
        pixels = crop.get_flattened_data()
        return Counter(pixels), total_pixels

    stride = int((total_pixels / RENDERED_CONTRAST_MAX_RAW_SAMPLE_PIXELS) ** 0.5) + 1
    pixels = [
        crop.getpixel((x, y))
        for y in range(0, crop.height, stride)
        for x in range(0, crop.width, stride)
    ]
    return Counter(pixels), len(pixels)


def _background_uniformity(
    counts: Counter,
    background_rgb: tuple[int, int, int],
    total: int,
) -> float:
    if total <= 0:
        return 0.0
    matching = sum(
        count
        for rgb, count in counts.items()
        if _rgb_distance(rgb, background_rgb) <= RENDERED_CONTRAST_BACKGROUND_UNIFORM_DISTANCE_MAX
    )
    return matching / total


def _background_model_detail(
    counts: Counter,
    background_rgb: tuple[int, int, int],
    total: int,
) -> dict:
    uniformity = _background_uniformity(counts, background_rgb, total)
    return {
        "background_model": "dominant_raw_color",
        "background_uniformity": round(uniformity, 4),
        "background_complexity": (
            "uniform"
            if uniformity >= RENDERED_CONTRAST_BACKGROUND_UNIFORMITY_MIN
            else "complex"
        ),
        "background_uniformity_min": RENDERED_CONTRAST_BACKGROUND_UNIFORMITY_MIN,
    }


def _dominant_rendered_contrast(
    crop: Image.Image,
    expected_foreground_hexes: Optional[Iterable[str]] = None,
) -> Optional[dict]:
    counts, total = _sample_rendered_pixels(crop)
    if total <= 0 or len(counts) < 2:
        return None

    background_rgb, background_count = counts.most_common(1)[0]
    min_pixels = max(
        RENDERED_CONTRAST_MIN_FOREGROUND_PIXELS,
        int(total * RENDERED_CONTRAST_MIN_FOREGROUND_RATIO),
    )

    expected: list[tuple[str, tuple[int, int, int]]] = []
    for hex_color in expected_foreground_hexes or []:
        try:
            expected.append((hex_color.upper(), _hex_to_rgb(hex_color)))
        except Exception:
            continue

    expected_candidates: list[tuple[str, int, float]] = []
    for hex_color, expected_rgb in expected:
        matched = sum(
            count
            for rgb, count in counts.items()
            if _rgb_distance(rgb, expected_rgb) <= RENDERED_CONTRAST_EXPECTED_COLOR_DISTANCE_MAX
        )
        if matched >= min_pixels and _rgb_distance(expected_rgb, background_rgb) >= 4:
            expected_candidates.append((hex_color, matched, _rgb_distance(expected_rgb, background_rgb)))

    if expected_candidates:
        foreground_hex, foreground_count, _ = max(
            expected_candidates,
            key=lambda item: item[2],
        )
        foreground_rgb = _hex_to_rgb(foreground_hex)
        background_counts = Counter(
            {
                rgb: count
                for rgb, count in counts.items()
                if _rgb_distance(rgb, foreground_rgb) > RENDERED_CONTRAST_EXPECTED_COLOR_DISTANCE_MAX
            }
        )
        if background_counts:
            background_rgb, background_count = background_counts.most_common(1)[0]
            background_total = sum(background_counts.values())
        else:
            background_total = total
        background_hex = _rgb_tuple_to_hex(background_rgb)
        return {
            "text_hex": foreground_hex,
            "background_hex": background_hex,
            "contrast_ratio": round(_contrast_ratio(foreground_hex, background_hex), 2),
            "foreground_pixels": foreground_count,
            "background_pixels": background_count,
            "sample_pixels": total,
            "foreground_detection": "expected_run_color",
            **_background_model_detail(background_counts or counts, background_rgb, background_total),
        }

    candidates = [
        (rgb, count)
        for rgb, count in counts.items()
        if rgb != background_rgb
        and count >= min_pixels
        and _rgb_distance(rgb, background_rgb) >= 4
    ]
    if not candidates:
        return None

    foreground_rgb, foreground_count = max(
        candidates,
        key=lambda item: _rgb_distance(item[0], background_rgb),
    )
    foreground_hex = _rgb_tuple_to_hex(foreground_rgb)
    background_hex = _rgb_tuple_to_hex(background_rgb)
    return {
        "text_hex": foreground_hex,
        "background_hex": background_hex,
        "contrast_ratio": round(_contrast_ratio(foreground_hex, background_hex), 2),
        "foreground_pixels": foreground_count,
        "background_pixels": background_count,
        "sample_pixels": total,
        "foreground_detection": "raw_pixel_distance",
        **_background_model_detail(counts, background_rgb, total),
    }


def _shape_fill_rgb_hex(shape) -> Optional[str]:
    try:
        fore_color = shape.fill.fore_color
    except (AttributeError, TypeError, ValueError):
        return None
    return _rgb_hex(fore_color)


def _solid_shape_fill_rgb_hex(shape) -> Optional[str]:
    try:
        if shape.fill.type != MSO_FILL.SOLID:
            return None
    except (AttributeError, TypeError, ValueError):
        return None
    return _shape_fill_rgb_hex(shape)


def _text_shape_background_rgb_hex(shape) -> tuple[Optional[str], str]:
    try:
        fill_type = shape.fill.type
    except (AttributeError, TypeError, ValueError):
        return DEFAULT_SLIDE_BACKGROUND_HEX, "slide_background_assumed"
    if fill_type == MSO_FILL.SOLID:
        fill_hex = _shape_fill_rgb_hex(shape)
        if fill_hex:
            return fill_hex, "shape_solid_fill"
        return None, "unknown_solid_fill"
    if fill_type in (None, MSO_FILL.BACKGROUND):
        return DEFAULT_SLIDE_BACKGROUND_HEX, "slide_background_assumed"
    return None, f"unsupported_fill_type:{fill_type}"


def _text_record_background_rgb_hex(record: ShapeRecord, records: list[ShapeRecord]) -> tuple[Optional[str], str]:
    own_background, own_source = _text_shape_background_rgb_hex(record.shape)
    if own_source == "shape_solid_fill" or own_background is None:
        return own_background, own_source

    text_center = _bbox_center(record.bbox_pt)
    candidates: list[tuple[float, ShapeRecord, str]] = []
    unknown_covering_background = False
    for candidate in records:
        if candidate is record:
            continue
        if candidate.kind == "text":
            continue
        if not _point_in_bbox(text_center, candidate.bbox_pt):
            continue
        fill_hex = _solid_shape_fill_rgb_hex(candidate.shape)
        if fill_hex:
            candidates.append((_bbox_area(candidate.bbox_pt), candidate, fill_hex))
            continue
        if candidate.kind in {"image", "shape"}:
            unknown_covering_background = True

    if candidates:
        _, candidate, fill_hex = min(candidates, key=lambda item: item[0])
        return fill_hex, f"behind_solid_fill:{_shape_label(candidate)}"
    if unknown_covering_background:
        return None, "unknown_covering_background"
    return own_background, own_source


def _table_cell_fill_rgb_hexes(shape) -> dict[str, list[dict[str, int]]]:
    if not getattr(shape, "has_table", False):
        return {}

    colors: dict[str, list[dict[str, int]]] = defaultdict(list)
    for row_idx, row in enumerate(shape.table.rows, start=1):
        for col_idx, cell in enumerate(row.cells, start=1):
            try:
                fore_color = cell.fill.fore_color
            except (AttributeError, TypeError, ValueError):
                continue
            hex_color = _rgb_hex(fore_color)
            if hex_color:
                colors[hex_color].append({"row": row_idx, "col": col_idx})
    return colors


def _iter_text_runs(shape):
    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                yield run, None

    if getattr(shape, "has_table", False):
        for row_idx, row in enumerate(shape.table.rows, start=1):
            for col_idx, cell in enumerate(row.cells, start=1):
                location = {"scope": "table_cell", "row": row_idx, "col": col_idx}
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        yield run, location


def _text_frame_explicit_run_colors(text_frame) -> list[str]:
    colors: set[str] = set()
    for para in text_frame.paragraphs:
        for run in para.runs:
            if not run.text:
                continue
            hex_color = _rgb_hex(run.font.color)
            if hex_color:
                colors.add(hex_color)
    return sorted(colors)


def _c_nv_pr(shape):
    try:
        matches = shape._element.xpath(".//p:cNvPr")
    except (AttributeError, TypeError, ValueError):
        return None
    return matches[0] if matches else None


def _alt_text_values(shape) -> tuple[str, str]:
    c_nv_pr = _c_nv_pr(shape)
    if c_nv_pr is None:
        return "", ""
    title = (c_nv_pr.get("title") or "").strip()
    descr = (c_nv_pr.get("descr") or "").strip()
    return title, descr


def _decorative_template_raster_reason(shape) -> Optional[str]:
    title, descr = _alt_text_values(shape)
    haystack = " ".join([getattr(shape, "name", ""), title, descr]).lower()
    for keyword in DECORATIVE_RASTER_KEYWORDS:
        if keyword in haystack:
            return f"keyword:{keyword}"
    return None


def _picture_crop_ratios(shape) -> Optional[dict[str, float]]:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    try:
        matches = shape._element.xpath(".//a:srcRect")
    except (AttributeError, TypeError, ValueError):
        matches = []
    if not matches:
        return None
    src_rect = matches[0]
    ratios: dict[str, float] = {}
    for side in ("l", "t", "r", "b"):
        raw = src_rect.get(side)
        if raw is None:
            ratios[side] = 0.0
            continue
        try:
            ratios[side] = max(0.0, float(raw) / 100000.0)
        except ValueError:
            ratios[side] = 0.0
    return ratios


def shape_bbox_pt(shape) -> Optional[tuple]:
    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
        return None
    return (
        emu_to_pt(shape.left),
        emu_to_pt(shape.top),
        emu_to_pt(shape.width),
        emu_to_pt(shape.height),
    )


def make_context(
    actual_w: float,
    actual_h: float,
    *,
    policy: LintPolicy,
) -> LintContext:
    scale_x = SLIDE_W_PT / actual_w if actual_w else 1.0
    scale_y = SLIDE_H_PT / actual_h if actual_h else 1.0
    base_aspect = SLIDE_W_PT / SLIDE_H_PT
    actual_aspect = actual_w / actual_h if actual_h else 0
    proportional = abs(actual_aspect - base_aspect) <= SLIDE_ASPECT_TOL
    return LintContext(
        actual_w_pt=actual_w,
        actual_h_pt=actual_h,
        scale_x=scale_x,
        scale_y=scale_y,
        proportional_to_base=proportional,
        policy=policy,
    )


def normalize_bbox(ctx: LintContext, bbox: tuple) -> tuple:
    x, y, w, h = bbox
    return (
        x * ctx.scale_x,
        y * ctx.scale_y,
        w * ctx.scale_x,
        h * ctx.scale_y,
    )


def make_finding(severity, check, slide_idx, slide_id, shape, message, detail=None) -> Finding:
    return Finding(
        severity=severity,
        check=check,
        slide_index=slide_idx,
        slide_id=slide_id,
        shape_id=getattr(shape, "shape_id", None) if shape is not None else None,
        shape_name=getattr(shape, "name", None) if shape is not None else None,
        message=message,
        detail=detail or {},
    )


# ---- Checks ----------------------------------------------------------------


def _estimate_text_render_lines(
    ctx: LintContext,
    shape,
    wrap_width_pt: float,
    font_size_pt: float,
) -> int:
    """Estimate the number of rendered lines when `shape.text_frame` is wrapped
    at `wrap_width_pt`. Each paragraph is wrapped independently (PPTX semantics);
    an empty paragraph counts as 1 visual line. CJK / Latin width mixing is
    delegated to `_estimate_text_width_pt`.
    """
    if wrap_width_pt <= 0:
        return 1
    total = 0
    for para in shape.text_frame.paragraphs:
        text = (para.text or "").strip()
        if not text:
            total += 1
            continue
        para_w = _estimate_text_width_pt(text, font_size_pt)
        total += max(1, int(math.ceil(para_w / wrap_width_pt)))
    return max(total, 1)


def _canvas_overflow_sides(x: float, y: float, w: float, h: float) -> dict[str, float]:
    out: dict[str, float] = {}
    right = x + w
    bottom = y + h
    if x < -TOL_PT:
        out["left"] = round(-x, 2)
    if y < -TOL_PT:
        out["top"] = round(-y, 2)
    if right > SLIDE_W_PT + TOL_PT:
        out["right"] = round(right - SLIDE_W_PT, 2)
    if bottom > SLIDE_H_PT + TOL_PT:
        out["bottom"] = round(bottom - SLIDE_H_PT, 2)
    return out


def check_overflow(ctx, slide_idx, slide_id, shape, bbox, findings):
    """3-way overflow lint family (DS-OVERFLOW-001 段階1+).

    1. box_canvas_overflow  : box bbox が slide canvas 外。text/non-text 共通。
                              fixer = box_canvas_clip (右下方向のみ width/height
                              切り、左/上方向は manual で残す)。
    2. text_box_overflow    : (段階2 で実装) text 描画が自分の box.height を超える。
    3. text_canvas_overflow : (段階3 で実装) canvas 端で wrap した時、text 自体
                              の描画範囲が canvas 右端を超える。

    Non-text picture / shape の box overflow は従来通り `overflow_images` /
    `overflow_shapes` として fire させる (= box_canvas_overflow とは別 lint で、
    text 系の box geometry overflow とは責務が違う)。
    """
    decorative_reason = _decorative_template_raster_reason(shape)
    normalized = normalize_bbox(ctx, bbox)
    x, y, w, h = normalized
    overflow_sides = _canvas_overflow_sides(x, y, w, h)

    if shape.has_text_frame and shape.text_frame.text.strip():
        # 1) box_canvas_overflow: box bbox が canvas 外。
        if overflow_sides:
            findings.append(
                make_finding(
                    "warning",
                    "box_canvas_overflow",
                    slide_idx,
                    slide_id,
                    shape,
                    "text box outside slide canvas: "
                    + ", ".join(f"{side}+{amount:.1f}pt" for side, amount in overflow_sides.items()),
                    {
                        "bbox_pt": [round(v, 2) for v in normalized],
                        "actual_bbox_pt": [round(v, 2) for v in bbox],
                        "overflow_sides_pt": overflow_sides,
                        "canvas_pt": [SLIDE_W_PT, SLIDE_H_PT],
                        "text_excerpt": _text_excerpt(shape),
                    },
                )
            )

        # Codex セカンドオピニオン反映 (2026-05-23):
        # - font_size_pt が None (= run.font.size が explicit でない、
        #   placeholder 継承未解決) なら text_box / text_canvas は推定不能。
        # - text_frame.auto_size が TEXT_TO_FIT_SHAPE / SHAPE_TO_FIT_TEXT
        #   なら PowerPoint が text を自動縮小 / box を自動拡張するので
        #   両 lint とも suppress。
        # - 段落ごとに lines × line_height を合算 (旧 max ×全行で過剰膨張
        #   していたのを是正)。margin_top/bottom/left/right を box から
        #   inner box にして判定する。
        font_size_pt = _text_frame_dominant_font_size_pt(ctx, shape.text_frame)
        try:
            auto_size = shape.text_frame.auto_size
        except (AttributeError, ValueError):
            auto_size = None
        autofit_suppress = auto_size in {
            MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
            MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        }
        text_render_can_estimate = font_size_pt is not None and not autofit_suppress

        if text_render_can_estimate:
            margin_top = _text_frame_margin_pt_value(shape.text_frame, "margin_top")
            margin_bottom = _text_frame_margin_pt_value(shape.text_frame, "margin_bottom")
            margin_left = _text_frame_margin_pt_value(shape.text_frame, "margin_left")
            margin_right = _text_frame_margin_pt_value(shape.text_frame, "margin_right")
            inner_h = max(0.0, h - margin_top - margin_bottom)
            inner_w = max(0.0, w - margin_left - margin_right)

            # 段落ごとに paragraph_lines × paragraph_line_height を合算。
            total_required_h = 0.0
            line_count_by_paragraph: list[int] = []
            for para in shape.text_frame.paragraphs:
                ptext = (para.text or "")
                para_line_height = _paragraph_line_height_pt(ctx, para, font_size_pt)
                if not ptext.strip():
                    para_lines = 1
                else:
                    if inner_w <= 0:
                        para_lines = 9999
                    else:
                        para_w = _estimate_text_width_pt(ptext, font_size_pt)
                        para_lines = max(1, int(math.ceil(para_w / inner_w)))
                total_required_h += para_lines * para_line_height
                line_count_by_paragraph.append(para_lines)

            # 2) text_box_overflow
            if total_required_h > inner_h + TOL_PT and inner_w > 0:
                overflow_by_pt = round(total_required_h - inner_h, 2)
                findings.append(
                    make_finding(
                        "warning",
                        "text_box_overflow",
                        slide_idx,
                        slide_id,
                        shape,
                        (
                            f"rendered text height {total_required_h:.1f}pt > inner_h "
                            f"{inner_h:.1f}pt (overflow +{overflow_by_pt:.1f}pt; "
                            f"paragraphs={line_count_by_paragraph}, "
                            f"font {font_size_pt:g}pt)"
                        ),
                        {
                            "bbox_pt": [round(v, 2) for v in normalized],
                            "actual_bbox_pt": [round(v, 2) for v in bbox],
                            "canvas_pt": [SLIDE_W_PT, SLIDE_H_PT],
                            "text_render": {
                                "font_size_pt": round(font_size_pt, 2),
                                "effective_font_source": "run_explicit",
                                "autofit": str(auto_size) if auto_size is not None else None,
                                "wrap_width_pt": round(inner_w, 2),
                                "inner_height_pt": round(inner_h, 2),
                                "box_height_pt": round(h, 2),
                                "margins_pt": {
                                    "top": round(margin_top, 2),
                                    "bottom": round(margin_bottom, 2),
                                    "left": round(margin_left, 2),
                                    "right": round(margin_right, 2),
                                },
                                "line_count_by_paragraph": line_count_by_paragraph,
                                "required_height_pt": round(total_required_h, 2),
                            },
                            "overflow_by_pt": overflow_by_pt,
                            "text_excerpt": _text_excerpt(shape),
                        },
                    )
                )

            # 3) text_canvas_overflow
            # 判定軸は box.width ではなく canvas 右端までの実残幅 (text 開始位置基準)。
            text_start_x = x + margin_left
            canvas_remaining = max(0.0, SLIDE_W_PT - text_start_x)
            text_raw = (shape.text_frame.text or "").strip()
            text_canvas_fires = False
            overflow_pt_canvas = 0.0
            longest_unit_pt = 0.0
            word_wrap = getattr(shape.text_frame, "word_wrap", True)
            word_wrap_enabled = word_wrap is None or word_wrap
            if canvas_remaining > 0 and text_raw:
                if not word_wrap_enabled:
                    total_w = _estimate_text_width_pt(text_raw, font_size_pt)
                    if total_w > canvas_remaining + TOL_PT:
                        text_canvas_fires = True
                        overflow_pt_canvas = total_w - canvas_remaining
                        longest_unit_pt = total_w
                else:
                    max_unit_w = 0.0
                    for word in re.split(r"\s+", text_raw):
                        if not word:
                            continue
                        current_run = ""
                        for ch in word:
                            if (
                                "぀" <= ch <= "ヿ"
                                or "㐀" <= ch <= "鿿"
                                or "＀" <= ch <= "￯"
                                or "　" <= ch <= "〿"
                            ):
                                if current_run:
                                    max_unit_w = max(
                                        max_unit_w,
                                        _estimate_text_width_pt(current_run, font_size_pt),
                                    )
                                    current_run = ""
                                max_unit_w = max(
                                    max_unit_w,
                                    _estimate_text_width_pt(ch, font_size_pt),
                                )
                            else:
                                current_run += ch
                        if current_run:
                            max_unit_w = max(
                                max_unit_w,
                                _estimate_text_width_pt(current_run, font_size_pt),
                            )
                    if max_unit_w > canvas_remaining + TOL_PT:
                        text_canvas_fires = True
                        overflow_pt_canvas = max_unit_w - canvas_remaining
                        longest_unit_pt = max_unit_w
            if text_canvas_fires:
                findings.append(
                    make_finding(
                        "warning",
                        "text_canvas_overflow",
                        slide_idx,
                        slide_id,
                        shape,
                        (
                            f"longest text unit {longest_unit_pt:.1f}pt exceeds "
                            f"canvas-remaining {canvas_remaining:.1f}pt "
                            f"from text_start_x={text_start_x:.1f}pt "
                            f"(word_wrap={'False' if word_wrap is False else 'True'}; "
                            f"overflow +{overflow_pt_canvas:.1f}pt)"
                        ),
                        {
                            "bbox_pt": [round(v, 2) for v in normalized],
                            "actual_bbox_pt": [round(v, 2) for v in bbox],
                            "canvas_pt": [SLIDE_W_PT, SLIDE_H_PT],
                            "text_render": {
                                "font_size_pt": round(font_size_pt, 2),
                                "effective_font_source": "run_explicit",
                                "autofit": str(auto_size) if auto_size is not None else None,
                                "text_start_x_pt": round(text_start_x, 2),
                                "canvas_remaining_pt": round(canvas_remaining, 2),
                                "longest_unit_pt": round(longest_unit_pt, 2),
                                "word_wrap": bool(word_wrap_enabled),
                                "margins_pt": {
                                    "left": round(margin_left, 2),
                                    "right": round(margin_right, 2),
                                },
                            },
                            "overflow_by_pt": round(overflow_pt_canvas, 2),
                            "text_excerpt": _text_excerpt(shape),
                        },
                    )
                )
        return

    if not overflow_sides:
        return
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        if decorative_reason:
            return
        check_id = "overflow_images"
    else:
        check_id = "overflow_shapes"
    msg = "outside slide canvas: " + ", ".join(
        f"{s}+{a:.1f}pt" for s, a in overflow_sides.items()
    )
    detail = {
        "bbox_pt": [round(v, 2) for v in normalized],
        "actual_bbox_pt": [round(v, 2) for v in bbox],
        "overflow_sides_pt": overflow_sides,
        "canvas_pt": [SLIDE_W_PT, SLIDE_H_PT],
    }
    findings.append(
        make_finding(
            "error", check_id, slide_idx, slide_id, shape, msg, detail,
        )
    )


def _is_page_number_text(shape, bbox_pt: tuple) -> bool:
    text = _shape_text(shape)
    if not text or not re.fullmatch(r"\d{1,3}", text):
        return False
    x, y, w, h = bbox_pt
    return (
        y >= SAFE_TEXT_FOOTER_Y_MIN_PT
        and w <= SAFE_TEXT_PAGE_NUMBER_MAX_W_PT
        and h <= SAFE_TEXT_PAGE_NUMBER_MAX_H_PT
    )


def _safe_text_area_exemption(ctx: LintContext, record: ShapeRecord, slide_type: dict | None) -> Optional[str]:
    shape = record.shape
    if slide_type and slide_type.get("slide_type") in {"cover", "section_divider"}:
        return f"slide_type:{slide_type.get('slide_type')}"
    x, y, w, h = record.bbox_pt
    if y <= SAFE_TEXT_HEADER_Y_MAX_PT:
        return "template_header_title"
    if y >= SAFE_TEXT_FOOTER_Y_MIN_PT:
        return "template_footer"
    if _is_page_number_text(shape, record.bbox_pt):
        return "page_number"
    sx, _, sw, _ = SAFE_TEXT_AREA_PT
    right_overflow = max(0.0, (x + w) - (sx + sw))
    if (
        x <= SAFE_TEXT_FULL_WIDTH_X_MAX_PT
        and w >= SAFE_TEXT_FULL_WIDTH_MIN_PT
        and 0 < right_overflow <= SAFE_TEXT_FULL_WIDTH_RIGHT_OVERFLOW_MAX_PT
    ):
        return "template_full_width_text"
    return None


def check_safe_text_area(ctx, slide_idx, slide_id, record: ShapeRecord, findings, slide_type: dict | None = None):
    shape = record.shape
    if not shape.has_text_frame:
        return
    if not shape.text_frame.text.strip():
        return
    exemption = _safe_text_area_exemption(ctx, record, slide_type)
    if exemption:
        return
    bbox = record.actual_bbox_pt
    normalized = normalize_bbox(ctx, bbox)
    x, y, w, h = normalized
    sx, sy, sw, sh = SAFE_TEXT_AREA_PT
    out: List[tuple] = []
    if x < sx - TOL_PT:
        out.append(("left", sx - x))
    if y < sy - TOL_PT:
        out.append(("top", sy - y))
    if x + w > sx + sw + TOL_PT:
        out.append(("right", (x + w) - (sx + sw)))
    if y + h > sy + sh + TOL_PT:
        out.append(("bottom", (y + h) - (sy + sh)))
    if not out:
        return
    msg = "text outside safe text area: " + ", ".join(f"{s}+{a:.1f}pt" for s, a in out)
    findings.append(
        make_finding(
            "warning", "safe_text_area_text", slide_idx, slide_id, shape, msg,
            {
                "bbox_pt": [round(v, 2) for v in normalized],
                "actual_bbox_pt": [round(v, 2) for v in bbox],
                "overflow_sides_pt": {side: round(amount, 2) for side, amount in out},
                "safe_text_area_pt": [round(v, 2) for v in SAFE_TEXT_AREA_PT],
            },
        )
    )


def _is_cover_brand_mark(ctx: LintContext, shape, bbox: tuple, slide_type: dict | None) -> bool:
    if not slide_type or slide_type.get("slide_type") != "cover":
        return False
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    x, y, w, h = normalize_bbox(ctx, bbox)
    return (
        x <= COVER_BRAND_MARK_MAX_X_PT
        and y <= COVER_BRAND_MARK_MAX_Y_PT
        and w <= COVER_BRAND_MARK_MAX_W_PT
        and h <= COVER_BRAND_MARK_MAX_H_PT
    )


def check_safe_margins(ctx, slide_idx, slide_id, shape, bbox, findings, slide_type: dict | None = None):
    if shape.has_text_frame and shape.text_frame.text.strip():
        return
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and _decorative_template_raster_reason(shape):
        return
    if _is_cover_brand_mark(ctx, shape, bbox, slide_type):
        return
    normalized = normalize_bbox(ctx, bbox)
    x, y, w, h = normalized
    out: List[tuple] = []
    if x < SAFE_MARGIN_LEFT_PT - TOL_PT:
        out.append(("left", SAFE_MARGIN_LEFT_PT - x))
    if y < SAFE_MARGIN_TOP_PT - TOL_PT:
        out.append(("top", SAFE_MARGIN_TOP_PT - y))
    if x + w > SLIDE_W_PT - SAFE_MARGIN_RIGHT_PT + TOL_PT:
        out.append(("right", (x + w) - (SLIDE_W_PT - SAFE_MARGIN_RIGHT_PT)))
    if y + h > SLIDE_H_PT - SAFE_MARGIN_BOTTOM_PT + TOL_PT:
        out.append(("bottom", (y + h) - (SLIDE_H_PT - SAFE_MARGIN_BOTTOM_PT)))
    if not out:
        return
    msg = "non-text outside safe margins: " + ", ".join(f"{s}+{a:.1f}pt" for s, a in out)
    findings.append(
        make_finding(
            "warning", "safe_margins", slide_idx, slide_id, shape, msg,
            {
                "bbox_pt": [round(v, 2) for v in normalized],
                "actual_bbox_pt": [round(v, 2) for v in bbox],
                "overflow_sides_pt": {side: round(amount, 2) for side, amount in out},
                "safe_margin_pt": {
                    "left": SAFE_MARGIN_LEFT_PT,
                    "right": SAFE_MARGIN_RIGHT_PT,
                    "top": SAFE_MARGIN_TOP_PT,
                    "bottom": SAFE_MARGIN_BOTTOM_PT,
                },
            },
        )
    )


def check_geometry_rounding(ctx, slide_idx, slide_id, shape, bbox, findings):
    normalized = normalize_bbox(ctx, bbox)
    names = ("x", "y", "w", "h")
    drifted = {
        name: round(value, 3)
        for name, value in zip(names, normalized)
        if abs(value - round(value)) > GEOMETRY_INTEGER_TOL_PT
    }
    if not drifted:
        return
    msg = "geometry is not rounded to integer pt: " + ", ".join(
        f"{name}={value:g}pt" for name, value in drifted.items()
    )
    findings.append(
        make_finding(
            "warning", "geometry_rounding", slide_idx, slide_id, shape, msg,
            {
                "bbox_pt": [round(v, 3) for v in normalized],
                "actual_bbox_pt": [round(v, 3) for v in bbox],
                "drifted": drifted,
            },
        )
    )


def check_image_upscale(ctx, slide_idx, slide_id, shape, bbox, findings):
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return
    try:
        source_w_px, source_h_px = shape.image.size
    except (AttributeError, ValueError):
        return
    if source_w_px <= 0 or source_h_px <= 0:
        return
    decorative_reason = _decorative_template_raster_reason(shape)
    if decorative_reason:
        return
    normalized = normalize_bbox(ctx, bbox)
    _, _, display_w_pt, display_h_pt = normalized
    if display_w_pt <= 0 or display_h_pt <= 0:
        return
    source_aspect = source_w_px / source_h_px
    display_aspect = display_w_pt / display_h_pt
    aspect_delta_ratio = abs(display_aspect / source_aspect - 1)
    if aspect_delta_ratio > MAX_IMAGE_ASPECT_DELTA_RATIO:
        findings.append(
            make_finding(
                "warning", "image_aspect_distortion", slide_idx, slide_id, shape,
                (
                    f"image aspect ratio changed from {source_aspect:.2f} to "
                    f"{display_aspect:.2f} ({aspect_delta_ratio:.0%} delta)"
                ),
                {
                    "source_px": [source_w_px, source_h_px],
                    "source_aspect": round(source_aspect, 4),
                    "display_aspect": round(display_aspect, 4),
                    "aspect_delta_ratio": round(aspect_delta_ratio, 4),
                    "actual_display_pt": [round(bbox[2], 2), round(bbox[3], 2)],
                    "normalized_display_pt": [round(display_w_pt, 2), round(display_h_pt, 2)],
                },
            )
        )
    display_w_px = display_w_pt * PX_PER_PT
    display_h_px = display_h_pt * PX_PER_PT
    ratio = max(display_w_px / source_w_px, display_h_px / source_h_px)
    if ratio <= MAX_IMAGE_UPSCALE_RATIO + 0.01:
        return
    findings.append(
        make_finding(
            "warning", "image_upscale_ratio", slide_idx, slide_id, shape,
            f"image displayed at {ratio:.2f}x source pixel size; maximum is {MAX_IMAGE_UPSCALE_RATIO:.2f}x",
            {
                "source_px": [source_w_px, source_h_px],
                "display_px": [round(display_w_px, 1), round(display_h_px, 1)],
                "actual_display_pt": [round(bbox[2], 2), round(bbox[3], 2)],
                "normalized_display_pt": [round(display_w_pt, 2), round(display_h_pt, 2)],
                "upscale_ratio": round(ratio, 3),
                "source_aspect": round(source_aspect, 4),
                "display_aspect": round(display_aspect, 4),
                "aspect_delta_ratio": round(aspect_delta_ratio, 4),
            },
        )
    )


def check_key_area_cropped(ctx, slide_idx, slide_id, shape, bbox, findings):
    crop = _picture_crop_ratios(shape)
    if crop is None:
        return
    decorative_reason = _decorative_template_raster_reason(shape)
    if decorative_reason:
        return
    horizontal_total = crop["l"] + crop["r"]
    vertical_total = crop["t"] + crop["b"]
    max_side = max(crop.values())
    triggered_rules: list[str] = []
    if max_side > IMAGE_CROP_SIDE_RATIO_MAX:
        triggered_rules.append("single_side_crop")
    if horizontal_total > IMAGE_CROP_AXIS_TOTAL_RATIO_MAX:
        triggered_rules.append("horizontal_total_crop")
    if vertical_total > IMAGE_CROP_AXIS_TOTAL_RATIO_MAX:
        triggered_rules.append("vertical_total_crop")
    if not triggered_rules:
        return
    normalized = normalize_bbox(ctx, bbox)
    findings.append(
        make_finding(
            "warning", "key_area_cropped", slide_idx, slide_id, shape,
            (
                "picture crop may remove important content: "
                + ", ".join(triggered_rules)
            ),
            {
                "evidence_source": "pptx_xml",
                "evidence_confidence": "medium",
                "measurement": "image_crop_metadata",
                "crop_ratio": {side: round(value, 4) for side, value in crop.items()},
                "horizontal_total_crop_ratio": round(horizontal_total, 4),
                "vertical_total_crop_ratio": round(vertical_total, 4),
                "max_side_crop_ratio": round(max_side, 4),
                "triggered_rules": triggered_rules,
                "thresholds": {
                    "side_ratio_max": IMAGE_CROP_SIDE_RATIO_MAX,
                    "axis_total_ratio_max": IMAGE_CROP_AXIS_TOTAL_RATIO_MAX,
                },
                "bbox_pt": [round(v, 2) for v in normalized],
                "actual_bbox_pt": [round(v, 2) for v in bbox],
            },
        )
    )


def check_alt_text(slide_idx, slide_id, shape, findings):
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return
    title, descr = _alt_text_values(shape)
    if title or descr:
        return
    findings.append(
        make_finding(
            "warning", "alt_text_required", slide_idx, slide_id, shape,
            "picture has no alt text title or description",
            {"alt_title": title, "alt_description": descr},
        )
    )


def check_autofit(slide_idx, slide_id, shape, findings):
    if not shape.has_text_frame:
        return
    af = shape.text_frame.auto_size
    if af is None or af != MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
        return
    body_pr_matches = shape._element.xpath(".//a:bodyPr")
    if not body_pr_matches:
        return
    norm_autofit = body_pr_matches[0].find(f"{A_NS}normAutofit")
    if norm_autofit is None:
        return
    font_scale_raw = norm_autofit.get("fontScale")
    line_space_reduction_raw = norm_autofit.get("lnSpcReduction")
    try:
        font_scale = int(font_scale_raw) if font_scale_raw is not None else 100000
    except ValueError:
        font_scale = 100000
    try:
        line_space_reduction = (
            int(line_space_reduction_raw) if line_space_reduction_raw is not None else 0
        )
    except ValueError:
        line_space_reduction = 0
    if font_scale >= 100000 and line_space_reduction <= 0:
        return

    run_sizes = [
        run.font.size.pt
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.font.size is not None
    ]
    base_font_size = max(run_sizes) if run_sizes else None
    font_scale_percent = round(font_scale / 1000, 2)
    effective_font_size = (
        round(base_font_size * font_scale / 100000, 2)
        if base_font_size is not None
        else None
    )
    line_space_reduction_percent = round(line_space_reduction / 1000, 2)
    findings.append(
        make_finding(
            "error", "text_autofit_disabled", slide_idx, slide_id, shape,
            (
                f"text auto-fit shrinks font to {font_scale_percent:g}%"
                + (
                    f" ({base_font_size:g}pt -> {effective_font_size:g}pt)"
                    if base_font_size is not None and effective_font_size is not None
                    else ""
                )
                + " (must be NONE)"
            ),
            {
                "auto_size": str(af),
                "autofit_mode": "TEXT_TO_FIT_SHAPE",
                "expected_auto_size": "NONE",
                "font_scale_percent": font_scale_percent,
                "font_shrink_percent": round(100 - font_scale_percent, 2),
                "line_space_reduction_percent": line_space_reduction_percent,
                "base_font_size_pt": base_font_size,
                "effective_font_size_pt": effective_font_size,
            },
        )
    )


def check_font(ctx, slide_idx, slide_id, shape, findings):
    if not shape.has_text_frame and not getattr(shape, "has_table", False):
        return
    bad_fonts: dict = {}
    bad_sizes: dict = {}
    bad_font_cells: dict = defaultdict(list)
    bad_size_cells: dict = defaultdict(list)
    for run, location in _iter_text_runs(shape):
        if not run.text:
            continue
        for script, name in _run_explicit_typefaces(run):
            if name and not _font_name_allowed(name):
                key = (script, name)
                bad_fonts[key] = bad_fonts.get(key, 0) + 1
                if location:
                    bad_font_cells[key].append(location)
        size = run.font.size
        if size is not None:
            actual_pt = round(size.pt, 1)
            normalized_pt = round(size.pt * ctx.font_scale, 1)
            if not _font_size_allowed(normalized_pt):
                key = (actual_pt, normalized_pt)
                bad_sizes[key] = bad_sizes.get(key, 0) + 1
                if location:
                    bad_size_cells[key].append(location)
    for (script, name), count in bad_fonts.items():
        findings.append(
            make_finding(
                "warning", "font_family", slide_idx, slide_id, shape,
                f"{script} font '{name}' not in allowlist {list(ALLOWED_FONT_FAMILIES)} ({count} run(s))",
                {
                    "script": script,
                    "font": name,
                    "runs": count,
                    "cells": bad_font_cells.get((script, name), []),
                },
            )
        )
    for (actual_pt, normalized_pt), count in bad_sizes.items():
        findings.append(
            make_finding(
                "warning", "font_size_scale", slide_idx, slide_id, shape,
                (
                    f"font size {normalized_pt}pt normalized from {actual_pt}pt "
                    f"not within {FONT_SIZE_TOL_PT:g}pt of scale "
                    f"{sorted(ALLOWED_FONT_SIZES_PT)} ({count} run(s))"
                ),
                {
                    "size_pt": normalized_pt,
                    "nearest_allowed_size_pt": _nearest_allowed_font_size(normalized_pt),
                    "actual_size_pt": actual_pt,
                    "normalization_scale": round(ctx.font_scale, 4),
                    "tolerance_pt": FONT_SIZE_TOL_PT,
                    "runs": count,
                    "cells": bad_size_cells.get((actual_pt, normalized_pt), []),
                },
            )
        )


def check_line_height(ctx, slide_idx, slide_id, shape, findings):
    if not shape.has_text_frame:
        return
    bad: dict = {}
    for para in shape.text_frame.paragraphs:
        if not para.text.strip():
            continue
        line_spacing = para.line_spacing
        if line_spacing is None:
            continue
        if isinstance(line_spacing, float):
            key = f"{line_spacing:g}x"
        else:
            actual_pt = round(line_spacing.pt, 1)
            normalized_pt = round(line_spacing.pt * ctx.font_scale, 1)
            if _line_height_allowed(normalized_pt):
                continue
            key = f"{normalized_pt:g}pt normalized from {actual_pt:g}pt"
        bad[key] = bad.get(key, 0) + 1
    for value, count in bad.items():
        findings.append(
            make_finding(
                "warning", "line_height", slide_idx, slide_id, shape,
                (
                    f"line spacing {value} not within {LINE_HEIGHT_TOL_PT:g}pt of "
                    f"allowed fixed pt scale {sorted(ALLOWED_LINE_HEIGHTS_PT)} "
                    f"({count} paragraph(s))"
                ),
                {
                    "line_spacing": value,
                    "allowed_line_heights_pt": sorted(ALLOWED_LINE_HEIGHTS_PT),
                    "tolerance_pt": LINE_HEIGHT_TOL_PT,
                    "paragraphs": count,
                },
            )
        )


def check_alignment(ctx, slide_idx, slide_id, shape, findings):
    if not shape.has_text_frame:
        return
    if not shape.text_frame.text.strip():
        return
    shape_bbox = shape_bbox_pt(shape)
    if shape_bbox is not None and _is_badge_like_shape(shape, normalize_bbox(ctx, shape_bbox)):
        return
    vertical = shape.text_frame.vertical_anchor
    if (
        ctx.policy.check_vertical_text_anchor
        and
        vertical is not None
        and vertical != MSO_VERTICAL_ANCHOR.TOP
    ):
        findings.append(
            make_finding(
                "warning", "alignment_left_top", slide_idx, slide_id, shape,
                f"text vertical alignment is {vertical}; expected TOP",
                {
                    "vertical_anchor": str(vertical),
                    "expected_vertical_anchor": "TOP",
                },
            )
        )
    bad: dict = {}
    for para in shape.text_frame.paragraphs:
        if not para.text.strip():
            continue
        if para.alignment is not None and para.alignment != PP_ALIGN.LEFT:
            bad[str(para.alignment)] = bad.get(str(para.alignment), 0) + 1
    for alignment, count in bad.items():
        findings.append(
            make_finding(
                "warning", "alignment_left_top", slide_idx, slide_id, shape,
                f"paragraph alignment is {alignment}; expected LEFT ({count} paragraph(s))",
                {
                    "alignment": alignment,
                    "expected_alignment": "LEFT",
                    "vertical_anchor": str(vertical) if vertical is not None else None,
                    "expected_vertical_anchor": "TOP",
                    "paragraphs": count,
                },
            )
        )


def _is_badge_like_shape(shape, bbox) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    text = _shape_text(shape)
    if not text or len(text) > BADGE_TEXT_LENGTH_MAX:
        return False
    if "\n" in text:
        return False
    try:
        if shape.fill.type != MSO_FILL.SOLID:
            return False
    except (AttributeError, TypeError, ValueError):
        return False
    x, y, w, h = bbox
    if w <= 0 or h <= 0:
        return False
    if w > BADGE_SHAPE_MAX_WIDTH_PT or h > BADGE_SHAPE_MAX_HEIGHT_PT:
        return False
    ratio = w / h
    if not (BADGE_SHAPE_ASPECT_MIN <= ratio <= BADGE_SHAPE_ASPECT_MAX):
        return False
    return True


def check_badge_alignment(ctx, slide_idx, slide_id, shape, bbox, findings):
    if not _is_badge_like_shape(shape, bbox):
        return
    vertical = shape.text_frame.vertical_anchor
    misaligned: list[dict[str, str]] = []
    if vertical is None:
        # PowerPoint default vertical anchor is TOP, which is not MIDDLE.
        misaligned.append({"axis": "vertical_anchor", "actual": "TOP_default", "expected": "MIDDLE"})
    elif vertical != MSO_VERTICAL_ANCHOR.MIDDLE:
        misaligned.append({"axis": "vertical_anchor", "actual": str(vertical), "expected": "MIDDLE"})

    horizontal_actual: Optional[str] = None
    for para in shape.text_frame.paragraphs:
        if not para.text.strip():
            continue
        if para.alignment is None:
            horizontal_actual = "LEFT_default"
        elif para.alignment != PP_ALIGN.CENTER:
            horizontal_actual = str(para.alignment)
        break
    if horizontal_actual is not None:
        misaligned.append({"axis": "alignment", "actual": horizontal_actual, "expected": "CENTER"})

    if not misaligned:
        return
    findings.append(
        make_finding(
            "warning",
            "badge_alignment",
            slide_idx,
            slide_id,
            shape,
            "badge container text is not centered (expected alignment=CENTER, vertical_anchor=MIDDLE)",
            {
                "shape_kind": "badge",
                "bbox_pt": [round(v, 2) for v in bbox],
                "text_excerpt": _text_excerpt(shape),
                "text_length": len(_shape_text(shape)),
                "aspect_ratio": round(bbox[2] / bbox[3], 3) if bbox[3] else None,
                "misaligned": misaligned,
                "evidence_source": "pptx_xml",
                "evidence_confidence": "high",
                "fixability": "auto_fix_candidate",
                "fixability_reason": "badge_center_alignment",
                "candidate_values": {
                    "alignment": "CENTER",
                    "vertical_anchor": "MIDDLE",
                },
                "measurement_confidence": "high",
                "measured_value": len(misaligned),
                "threshold": 0,
                "delta": len(misaligned),
                "unit": "misaligned_axes",
            },
        )
    )


def check_wrap_break_changes_meaning(ctx, slide_idx, slide_id, shape, bbox, findings):
    if not getattr(shape, "has_text_frame", False):
        return
    if not shape.text_frame.text.strip():
        return
    lines = _text_frame_lines(shape.text_frame)
    if len(lines) < 2:
        return

    triggered: list[dict[str, Any]] = []
    for idx, (left_raw, right_raw) in enumerate(zip(lines, lines[1:]), start=1):
        left = left_raw.strip()
        right = right_raw.strip()
        if not left or not right:
            continue
        rule = None
        if re.search(r"[A-Za-z]$", left) and re.search(r"^[a-z]", right):
            rule = "latin_word_split"
        elif re.search(r"\d$", left) and re.search(r"^(%|pt|px|mm|cm|kg|g|円|年|月|日)", right):
            rule = "number_unit_split"
        elif left[-1] in "([{/" or right[0] in ")]},、。":
            rule = "punctuation_or_bracket_orphan"
        elif len(left) <= 2 or len(right) <= 2:
            rule = "short_fragment_orphan"
        if rule is None:
            continue
        triggered.append(
            {
                "line_index": idx,
                "rule": rule,
                "left": left[-24:],
                "right": right[:24],
            }
        )
    if not triggered:
        return

    normalized = normalize_bbox(ctx, bbox)
    detail: dict[str, Any] = {
        "evidence_source": "pptx_xml",
        "evidence_confidence": "medium",
        "text_excerpt": _text_excerpt(shape),
        "line_count": len(lines),
        "triggered_rules": triggered,
        "bbox_pt": [round(v, 2) for v in normalized],
        "actual_bbox_pt": [round(v, 2) for v in bbox],
        "measured_value": len(triggered),
        "threshold": 1,
        "delta": len(triggered),
        "unit": "line_break_risks",
        "fixability": "manual_required",
        "manual_required_reason": "意味を保つ折返しは内容判断が必要なので手動で確認する。",
        "candidate_values": [],
    }

    if _text_frame_has_inline_break(shape.text_frame):
        font_size_pt = _text_frame_dominant_font_size_pt(ctx, shape.text_frame)
        if font_size_pt and font_size_pt > 0:
            joined = _text_frame_inline_joined_text(shape.text_frame)
            margin_l, margin_r = _text_frame_margin_pt(shape.text_frame)
            text_width_pt = _estimate_text_width_pt(joined, font_size_pt)
            required_width_pt = text_width_pt + margin_l + margin_r + 1.0
            current_x, current_y, current_w, current_h = normalized
            safe_right_x = SLIDE_W_PT - SAFE_MARGIN_RIGHT_PT
            available_width = safe_right_x - current_x
            if (
                required_width_pt <= available_width
                and required_width_pt > current_w + 0.5
            ):
                target_width_pt = round(required_width_pt + 4.0, 2)
                target_width_pt = min(target_width_pt, round(available_width, 2))
                target_bbox_pt = [
                    round(current_x, 2),
                    round(current_y, 2),
                    target_width_pt,
                    round(current_h, 2),
                ]
                detail["candidate_values"] = [
                    {
                        "strategy": "widen_to_fit",
                        "width_pt": target_width_pt,
                        "bbox_pt": target_bbox_pt,
                        "estimated_text_width_pt": round(text_width_pt, 2),
                        "available_width_pt": round(available_width, 2),
                        "font_size_pt": round(font_size_pt, 2),
                        "joined_text_excerpt": (
                            joined if len(joined) <= 80 else joined[:77] + "..."
                        ),
                        "replace_inline_breaks": True,
                    }
                ]
                detail["fixability"] = "auto_fix_candidate"
                detail["fixability_reason"] = "widen_to_fit_within_safe_area"
                detail.pop("manual_required_reason", None)

    findings.append(
        make_finding(
            "warning", "wrap_break_changes_meaning", slide_idx, slide_id, shape,
            "explicit line break may split a semantic unit: "
            + ", ".join(item["rule"] for item in triggered),
            detail,
        )
    )


def check_color(
    ctx,
    slide_idx,
    slide_id,
    record: ShapeRecord,
    records: list[ShapeRecord],
    findings,
    *,
    measure_structural_contrast: bool = True,
):
    shape = record.shape
    if shape.has_text_frame:
        bad_text_colors: dict = {}
        low_contrast_runs: dict = {}
        contrast_ratio_runs: dict = {}
        background_hex, background_source = _text_record_background_rgb_hex(record, records)
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text:
                    continue
                hex_color = _rgb_hex(run.font.color)
                if hex_color and hex_color not in ALLOWED_TEXT_COLORS_HEX:
                    bad_text_colors[hex_color] = bad_text_colors.get(hex_color, 0) + 1
                if not measure_structural_contrast or not hex_color or not background_hex:
                    continue
                size_pt = run.font.size.pt * ctx.font_scale if run.font.size is not None else None
                threshold = (
                    CONTRAST_RATIO_LARGE_TEXT_MIN
                    if size_pt is not None and size_pt >= LARGE_TEXT_MIN_SIZE_PT
                    else CONTRAST_RATIO_NORMAL_TEXT_MIN
                )
                ratio = _contrast_ratio(hex_color, background_hex)
                key = (
                    hex_color,
                    background_hex,
                    threshold,
                    background_source,
                    "large_text" if threshold == CONTRAST_RATIO_LARGE_TEXT_MIN else "normal_text",
                    round(size_pt, 2) if size_pt is not None else None,
                    round(ratio, 2),
                )
                if ratio < LOW_CONTRAST_RATIO_MAX:
                    low_contrast_runs[key] = low_contrast_runs.get(key, 0) + 1
                elif ratio < threshold:
                    contrast_ratio_runs[key] = contrast_ratio_runs.get(key, 0) + 1
        for hex_color, count in bad_text_colors.items():
            findings.append(
                make_finding(
                    "warning", "text_color_allowlist", slide_idx, slide_id, shape,
                    (
                        f"text color {hex_color} not in allowlist "
                        f"{sorted(ALLOWED_TEXT_COLORS_HEX)} ({count} run(s))"
                    ),
                    {"color_hex": hex_color, "runs": count},
                )
            )
        for key, count in low_contrast_runs.items():
            (
                text_hex,
                background_hex,
                threshold,
                background_source,
                text_class,
                size_pt,
                ratio,
            ) = key
            findings.append(
                make_finding(
                    "error", "low_contrast", slide_idx, slide_id, shape,
                    (
                        f"text/background contrast ratio {ratio:.2f}:1 is below "
                        f"unreadable threshold {LOW_CONTRAST_RATIO_MAX:.1f}:1 "
                        f"for {text_class} ({count} run(s))"
                    ),
                    {
                        "text_hex": text_hex,
                        "foreground_hex": text_hex,
                        "original_run_color_hex": text_hex,
                        "background_hex": background_hex,
                        "background_source": background_source,
                        "contrast_ratio": ratio,
                        "required_ratio": threshold,
                        "low_contrast_threshold": LOW_CONTRAST_RATIO_MAX,
                        "text_class": text_class,
                        "font_size_pt": size_pt,
                        "runs": count,
                    },
                )
            )
        for key, count in contrast_ratio_runs.items():
            (
                text_hex,
                background_hex,
                threshold,
                background_source,
                text_class,
                size_pt,
                ratio,
            ) = key
            findings.append(
                make_finding(
                    "warning", "contrast_ratio", slide_idx, slide_id, shape,
                    (
                        f"text/background contrast ratio {ratio:.2f}:1 is below "
                        f"required {threshold:.1f}:1 for {text_class} ({count} run(s))"
                    ),
                    {
                        "text_hex": text_hex,
                        "foreground_hex": text_hex,
                        "original_run_color_hex": text_hex,
                        "background_hex": background_hex,
                        "background_source": background_source,
                        "contrast_ratio": ratio,
                        "required_ratio": threshold,
                        "text_class": text_class,
                        "font_size_pt": size_pt,
                        "runs": count,
                    },
                )
            )

    fill_hex = _shape_fill_rgb_hex(shape)
    if fill_hex and fill_hex not in ALLOWED_FILL_COLORS_HEX:
        findings.append(
            make_finding(
                "warning", "background_color_palette", slide_idx, slide_id, shape,
                f"fill color {fill_hex} not in palette {sorted(ALLOWED_FILL_COLORS_HEX)}",
                {"color_hex": fill_hex},
            )
        )

    for fill_hex, cells in _table_cell_fill_rgb_hexes(shape).items():
        if fill_hex in ALLOWED_FILL_COLORS_HEX:
            continue
        findings.append(
            make_finding(
                "warning", "background_color_palette", slide_idx, slide_id, shape,
                (
                    f"table cell fill color {fill_hex} not in palette "
                    f"{sorted(ALLOWED_FILL_COLORS_HEX)} ({len(cells)} cell(s))"
                ),
                {"color_hex": fill_hex, "scope": "table_cell", "cells": cells},
            )
            )


def check_rendered_contrast(
    ctx,
    slide_idx,
    slide_id,
    record: ShapeRecord,
    image: Image.Image,
    image_path: Path,
    findings,
):
    shape = record.shape
    if not shape.has_text_frame:
        return
    if not shape.text_frame.text.strip():
        return
    crop = _rendered_crop_for_bbox(ctx, image, record.actual_bbox_pt)
    if crop is None:
        return
    original_run_colors = _text_frame_explicit_run_colors(shape.text_frame)
    measured = _dominant_rendered_contrast(
        crop,
        expected_foreground_hexes=original_run_colors,
    )
    if measured is None:
        return

    size_pt = _text_frame_dominant_font_size_pt(ctx, shape.text_frame)
    threshold = (
        CONTRAST_RATIO_LARGE_TEXT_MIN
        if size_pt is not None and size_pt >= LARGE_TEXT_MIN_SIZE_PT
        else CONTRAST_RATIO_NORMAL_TEXT_MIN
    )
    ratio = measured["contrast_ratio"]
    if ratio >= threshold:
        return

    text_class = "large_text" if threshold == CONTRAST_RATIO_LARGE_TEXT_MIN else "normal_text"
    check_id = "low_contrast" if ratio < LOW_CONTRAST_RATIO_MAX else "contrast_ratio"
    severity = "error" if check_id == "low_contrast" else "warning"
    if check_id == "low_contrast":
        message = (
            f"rendered text/background contrast ratio {ratio:.2f}:1 is below "
            f"unreadable threshold {LOW_CONTRAST_RATIO_MAX:.1f}:1 for {text_class}"
        )
    else:
        message = (
            f"rendered text/background contrast ratio {ratio:.2f}:1 is below "
            f"required {threshold:.1f}:1 for {text_class}"
        )
    detail = {
        **measured,
        "measurement": "rendered_image",
        "rendered_image_path": str(image_path),
        "foreground_hex": measured["text_hex"],
        "original_run_colors_hex": original_run_colors,
        "required_ratio": threshold,
        "low_contrast_threshold": LOW_CONTRAST_RATIO_MAX,
        "text_class": text_class,
        "font_size_pt": round(size_pt, 2) if size_pt is not None else None,
        "bbox_pt": [round(v, 2) for v in record.bbox_pt],
        "actual_bbox_pt": [round(v, 2) for v in record.actual_bbox_pt],
    }
    findings.append(
        make_finding(
            severity,
            check_id,
            slide_idx,
            slide_id,
            shape,
            message,
            detail,
        )
    )


def extract_structure_relations(records: list[ShapeRecord]) -> list[StructureRelation]:
    relations: list[StructureRelation] = []
    for idx, first in enumerate(records):
        for second in records[idx + 1:]:
            overlap_bbox = _intersection_bbox(first.bbox_pt, second.bbox_pt)
            if overlap_bbox is None:
                continue
            relation = _containment_relation(first, second, overlap_bbox)
            if relation is not None:
                relations.append(relation)
    return relations


def _multi_step_overlap_candidates(
    ctx: Optional[LintContext],
    first: ShapeRecord,
    second: ShapeRecord,
    overlap_bbox: tuple,
) -> list[dict]:
    candidates: list[dict] = []
    ax, ay, aw, ah = first.bbox_pt
    bx, by, bw, bh = second.bbox_pt
    ox, oy, ow, oh = overlap_bbox

    # 1) Move shape_b: nudge it past shape_a along the shorter overlap axis.
    if ow <= oh:
        # horizontal nudge
        dx = aw - (ax + aw - bx) if bx >= ax else -(ax - (bx + bw))
        new_left = bx + (ow + 4.0 if bx >= ax else -(ow + 4.0))
        candidates.append(
            {
                "strategy": "move_shape_b",
                "geometry": {
                    "left": round(new_left, 2),
                    "top": round(by, 2),
                    "width": round(bw, 2),
                    "height": round(bh, 2),
                },
                "reason": "horizontal nudge to clear overlap",
            }
        )
    else:
        new_top = by + (oh + 4.0 if by >= ay else -(oh + 4.0))
        candidates.append(
            {
                "strategy": "move_shape_b",
                "geometry": {
                    "left": round(bx, 2),
                    "top": round(new_top, 2),
                    "width": round(bw, 2),
                    "height": round(bh, 2),
                },
                "reason": "vertical nudge to clear overlap",
            }
        )

    # 2) Shrink the larger shape's font size to release space.
    if ctx is not None:
        size_a = _shape_dominant_font_size_pt(ctx, first.shape)
        size_b = _shape_dominant_font_size_pt(ctx, second.shape)
        if size_a is not None and size_b is not None:
            target_rec, target_size = (
                (first, size_a) if size_a >= size_b else (second, size_b)
            )
            smaller = [s for s in sorted(ALLOWED_FONT_SIZES_PT) if s < target_size]
            if smaller:
                candidates.append(
                    {
                        "strategy": "shrink_font_size",
                        "target_shape_id": getattr(target_rec.shape, "shape_id", None),
                        "font_size_pt": smaller[-1],
                        "from_pt": round(target_size, 2),
                        "reason": "smaller font size on the dominant shape may resolve overlap",
                    }
                )

    # 3) Shrink shape_b height by overlap height if vertical overlap dominates.
    if oh < bh and oh > 4.0:
        new_h = max(20.0, round(bh - oh - 4.0, 2))
        if new_h < bh:
            candidates.append(
                {
                    "strategy": "shrink_shape_b_height",
                    "geometry": {
                        "left": round(bx, 2),
                        "top": round(by, 2),
                        "width": round(bw, 2),
                        "height": new_h,
                    },
                    "reason": "reduce shape_b height to remove overlap area",
                }
            )

    return candidates


def check_object_relationships(slide_idx, slide_id, records: list[ShapeRecord], findings, ctx: Optional[LintContext] = None):
    for idx, first in enumerate(records):
        for second in records[idx + 1:]:
            overlap_bbox = _intersection_bbox(first.bbox_pt, second.bbox_pt)
            if overlap_bbox is not None:
                if _containment_relation(first, second, overlap_bbox) is not None:
                    continue
                overlap_area_pt2 = overlap_bbox[2] * overlap_bbox[3]
                if overlap_area_pt2 > OBJECT_OVERLAP_AREA_PT2_MIN:
                    if first.kind == "text" and second.kind == "text":
                        check_id = "text_overlap"
                        severity = "error"
                        subject = "text frames"
                    else:
                        check_id = "object_overlap"
                        severity = "error"
                        subject = "objects"
                    detail = {
                        "shape_a": _shape_record_detail(first),
                        "shape_b": _shape_record_detail(second),
                        "overlap_area_pt2": round(overlap_area_pt2, 2),
                        "overlap_bbox_pt": [round(v, 2) for v in overlap_bbox],
                        "threshold_pt2": OBJECT_OVERLAP_AREA_PT2_MIN,
                    }
                    if check_id == "text_overlap":
                        multi = _multi_step_overlap_candidates(ctx, first, second, overlap_bbox)
                        if multi:
                            detail["multi_step_candidates"] = multi
                    findings.append(
                        make_finding(
                            severity, check_id, slide_idx, slide_id, first.shape,
                            (
                                f"{subject} overlap by {overlap_area_pt2:.1f}pt^2 "
                                f"(>{OBJECT_OVERLAP_AREA_PT2_MIN:g}pt^2): "
                                f"{_shape_label(first)} / {_shape_label(second)}"
                            ),
                            detail,
                        )
                    )
                continue

            ax1, ay1, ax2, ay2 = _bbox_edges(first.bbox_pt)
            bx1, by1, bx2, by2 = _bbox_edges(second.bbox_pt)
            horizontal_gap = _axis_gap(ax1, ax2, bx1, bx2)
            vertical_gap = _axis_gap(ay1, ay2, by1, by2)
            horizontal_overlap = vertical_gap == 0.0
            vertical_overlap = horizontal_gap == 0.0
            element_threshold_pt, first_element, second_element = _object_gap_threshold_pt(
                first, second
            )
            for axis, gap, overlap_other_axis in (
                ("horizontal", horizontal_gap, horizontal_overlap),
                ("vertical", vertical_gap, vertical_overlap),
            ):
                if not overlap_other_axis or gap <= 0:
                    continue
                semantic_kind: Optional[str] = None
                if ctx is not None:
                    semantic_kind = _semantic_pair_kind(ctx, first, second, axis)
                semantic_threshold = (
                    SEMANTIC_PAIR_GAP_PT.get(semantic_kind, 0.0)
                    if semantic_kind
                    else 0.0
                )
                gap_threshold_pt = max(element_threshold_pt, semantic_threshold)
                if gap_threshold_pt <= OBJECT_GAP_DEFAULT_MIN_PT:
                    continue
                if gap >= gap_threshold_pt:
                    continue
                detail = {
                    "shape_a": _shape_record_detail(first),
                    "shape_b": _shape_record_detail(second),
                    "axis": axis,
                    "gap_pt": round(gap, 2),
                    "threshold_pt": gap_threshold_pt,
                    "element_a": first_element,
                    "element_b": second_element,
                    "default_threshold_pt": OBJECT_GAP_DEFAULT_MIN_PT,
                }
                if semantic_kind:
                    detail["semantic_pair_kind"] = semantic_kind
                    detail["semantic_threshold_pt"] = SEMANTIC_PAIR_GAP_PT[semantic_kind]
                findings.append(
                    make_finding(
                        "warning", "object_gap_too_small", slide_idx, slide_id, first.shape,
                        (
                            f"{axis} gap {gap:.1f}pt is below {gap_threshold_pt:g}pt"
                            + (f" [{semantic_kind}]" if semantic_kind else "")
                            + f": {_shape_label(first)} / {_shape_label(second)}"
                        ),
                        detail,
                    )
                )


def _classify_decorative_line(record: ShapeRecord) -> str | None:
    shape = record.shape
    try:
        shape_type = shape.shape_type
    except (AttributeError, ValueError):
        return None
    if shape_type == MSO_SHAPE_TYPE.LINE:
        return "connector_line"
    if record.kind in {"text", "image", "table"}:
        return None
    try:
        auto_shape_type = shape.auto_shape_type
    except (AttributeError, ValueError):
        auto_shape_type = None
    if auto_shape_type is not None:
        name = getattr(auto_shape_type, "name", None) or str(auto_shape_type)
        upper = name.upper() if isinstance(name, str) else ""
        if "ARROW" in upper:
            return "arrow_autoshape"
        if "LINE" in upper:
            return "line_autoshape"
    _, _, w, h = record.bbox_pt
    thickness = min(w, h)
    length = max(w, h)
    if (
        thickness <= DECORATIVE_THIN_BAR_THICKNESS_PT_MAX
        and length >= DECORATIVE_THIN_BAR_LENGTH_PT_MIN
    ):
        return "thin_decorative_bar"
    return None


def _nearest_neighbor_gap_pt(target: ShapeRecord, others: Iterable[ShapeRecord]) -> float | None:
    tx1, ty1, tx2, ty2 = _bbox_edges(target.bbox_pt)
    best: float | None = None
    for other in others:
        if other is target:
            continue
        ox1, oy1, ox2, oy2 = _bbox_edges(other.bbox_pt)
        hgap = _axis_gap(tx1, tx2, ox1, ox2)
        vgap = _axis_gap(ty1, ty2, oy1, oy2)
        gap = max(hgap, vgap)
        if best is None or gap < best:
            best = gap
    return best


def check_decorative_isolated_lines(slide_idx, slide_id, records: list[ShapeRecord], findings):
    classified: list[tuple[ShapeRecord, str]] = []
    for record in records:
        label = _classify_decorative_line(record)
        if label is not None:
            classified.append((record, label))
    if not classified:
        return
    line_ids = {id(rec) for rec, _ in classified}
    others = [r for r in records if id(r) not in line_ids]
    for line, label in classified:
        gap = _nearest_neighbor_gap_pt(line, others)
        if gap is not None and gap < DECORATIVE_LINE_PROXIMITY_PT_MAX:
            continue
        gap_value: float | None = None if gap is None else round(gap, 2)
        bbox_rounded = [round(v, 2) for v in line.bbox_pt]
        detail = {
            "shape": _shape_record_detail(line),
            "line_classification": label,
            "nearest_neighbor_gap_pt": gap_value,
            "proximity_threshold_pt": DECORATIVE_LINE_PROXIMITY_PT_MAX,
            # FIX-013 / POLICY-001 段階3: apply_mode=judgement_fix で扱う。
            # lint は manual_required + candidate (= 削除候補) を出し、SPA で
            # judgement_reason=auto_fixable に promote されたものだけ
            # pptx_fix が shape を削除する (strict gate)。
            "fixability": "manual_required",
            "fixability_reason": "decorative_isolated_line",
            "manual_required_reason": (
                "孤立した装飾 line/connector/arrow は意図的か削除候補かの判断が必要。"
                "SPA で auto_fixable と判定された場合のみ pptx_fix が削除する。"
            ),
            "candidate_values": [
                {
                    "strategy": "remove_shape",
                    "shape_id": getattr(line.shape, "shape_id", None),
                    "shape_name": getattr(line.shape, "name", None),
                    "line_classification": label,
                    "bbox_pt": bbox_rounded,
                }
            ],
            "measurement_confidence": "medium",
            "evidence_source": "pptx_xml",
            "evidence_confidence": "medium",
            "evidence": {
                "shape_kind": line.kind,
                "bbox_pt": bbox_rounded,
                "line_classification": label,
                "nearest_neighbor_gap_pt": gap_value,
                "proximity_threshold_pt": DECORATIVE_LINE_PROXIMITY_PT_MAX,
            },
        }
        message = (
            f"isolated decorative {label}: {_shape_label(line)} "
            f"(nearest neighbor "
            + (f"{gap_value:.1f}pt" if gap_value is not None else "n/a")
            + f", threshold {DECORATIVE_LINE_PROXIMITY_PT_MAX:g}pt)"
        )
        findings.append(
            make_finding(
                "warning",
                "decorative_isolated_lines",
                slide_idx,
                slide_id,
                line.shape,
                message,
                detail,
            )
        )


def _is_container_candidate(record: ShapeRecord) -> bool:
    if record.kind in {"text", "image", "table"}:
        return False
    _, _, w, h = record.bbox_pt
    return w > INNER_PADDING_SIDE_MIN_PT * 2 and h > INNER_PADDING_SIDE_MIN_PT * 2


def check_inner_padding_imbalance(slide_idx, slide_id, records: list[ShapeRecord], findings):
    children_by_container: dict[int, list[ShapeRecord]] = defaultdict(list)
    for relation in extract_structure_relations(records):
        if not _is_container_candidate(relation.container):
            continue
        children_by_container[id(relation.container)].append(relation.child)

    for container in records:
        if not _is_container_candidate(container):
            continue
        container_element = _design_element(container)
        target_padding = INNER_PADDING_ELEMENT_TARGET_PT.get(
            container_element,
            INNER_PADDING_DEFAULT_TARGET_PT,
        )
        if target_padding <= INNER_PADDING_DEFAULT_TARGET_PT:
            continue
        children = children_by_container.get(id(container), [])
        if len(children) < 2:
            continue

        child_left = min(child.bbox_pt[0] for child in children)
        child_top = min(child.bbox_pt[1] for child in children)
        child_right = max(child.bbox_pt[0] + child.bbox_pt[2] for child in children)
        child_bottom = max(child.bbox_pt[1] + child.bbox_pt[3] for child in children)
        container_left, container_top, container_right, container_bottom = _bbox_edges(container.bbox_pt)
        padding = {
            "left": child_left - container_left,
            "right": container_right - child_right,
            "top": child_top - container_top,
            "bottom": container_bottom - child_bottom,
        }
        triggered_rules: list[str] = []
        for side in ("left", "top"):
            if abs(padding[side] - target_padding) > INNER_PADDING_TARGET_TOL_PT:
                triggered_rules.append(f"{side}_padding_target_delta")
        for side in ("right", "bottom"):
            if padding[side] < target_padding - INNER_PADDING_TARGET_TOL_PT:
                triggered_rules.append(f"{side}_padding_below_target")

        horizontal_ratio = (
            padding["left"] / padding["right"]
            if padding["left"] > 0 and padding["right"] > 0
            else None
        )
        vertical_ratio = (
            padding["top"] / padding["bottom"]
            if padding["top"] > 0 and padding["bottom"] > 0
            else None
        )

        if not triggered_rules:
            continue

        findings.append(
            make_finding(
                "warning", "inner_padding_imbalance", slide_idx, slide_id, container.shape,
                (
                    f"container padding is imbalanced for {len(children)} child objects: "
                    + ", ".join(triggered_rules)
                ),
                {
                    "container": _shape_record_detail(container),
                    "children": [_shape_record_detail(child) for child in children],
                    "container_element": container_element,
                    "padding_pt": {side: round(value, 2) for side, value in padding.items()},
                    "target_padding_pt": target_padding,
                    "target_tolerance_pt": INNER_PADDING_TARGET_TOL_PT,
                    "horizontal_ratio": (
                        round(horizontal_ratio, 3) if horizontal_ratio is not None else None
                    ),
                    "vertical_ratio": (
                        round(vertical_ratio, 3) if vertical_ratio is not None else None
                    ),
                    "triggered_rules": triggered_rules,
                    "thresholds": {
                        "default_padding_pt": INNER_PADDING_DEFAULT_TARGET_PT,
                        "target_padding_pt": target_padding,
                        "target_tolerance_pt": INNER_PADDING_TARGET_TOL_PT,
                        "right_bottom_policy": "minimum_target",
                    },
                },
            )
        )


def _container_children(records: list[ShapeRecord]) -> dict[int, list[ShapeRecord]]:
    children_by_container: dict[int, list[ShapeRecord]] = defaultdict(list)
    for relation in extract_structure_relations(records):
        if _is_container_candidate(relation.container):
            children_by_container[id(relation.container)].append(relation.child)
    return children_by_container


def _container_padding(container: ShapeRecord, children: list[ShapeRecord]) -> dict[str, float]:
    child_left = min(child.bbox_pt[0] for child in children)
    child_top = min(child.bbox_pt[1] for child in children)
    child_right = max(child.bbox_pt[0] + child.bbox_pt[2] for child in children)
    child_bottom = max(child.bbox_pt[1] + child.bbox_pt[3] for child in children)
    container_left, container_top, container_right, container_bottom = _bbox_edges(container.bbox_pt)
    return {
        "left": child_left - container_left,
        "right": container_right - child_right,
        "top": child_top - container_top,
        "bottom": container_bottom - child_bottom,
    }


def _median(values: list[float]) -> float:
    ordered = sorted(values)
    mid = len(ordered) // 2
    if len(ordered) % 2:
        return ordered[mid]
    return (ordered[mid - 1] + ordered[mid]) / 2


def _card_row_groups(cards: list[ShapeRecord]) -> list[list[ShapeRecord]]:
    rows: list[list[ShapeRecord]] = []
    for card in sorted(cards, key=lambda item: (_bbox_center(item.bbox_pt)[1], item.bbox_pt[0])):
        _, cy = _bbox_center(card.bbox_pt)
        target: list[ShapeRecord] | None = None
        for row in rows:
            row_center = _median([_bbox_center(existing.bbox_pt)[1] for existing in row])
            if abs(cy - row_center) <= CARD_GRID_ROW_CENTER_TOL_PT:
                target = row
                break
        if target is None:
            rows.append([card])
        else:
            target.append(card)
    return [
        sorted(row, key=lambda item: item.bbox_pt[0])
        for row in rows
        if len(row) >= CARD_GRID_GROUP_MIN
    ]


def _first_child_relative_bbox(container: ShapeRecord, children: list[ShapeRecord]) -> Optional[list[float]]:
    if not children:
        return None
    child = sorted(children, key=lambda item: (item.bbox_pt[1], item.bbox_pt[0]))[0]
    cx, cy, _, _ = container.bbox_pt
    x, y, w, h = child.bbox_pt
    return [x - cx, y - cy, w, h]


def check_card_grid_consistency(slide_idx, slide_id, records: list[ShapeRecord], findings):
    children_by_container = _container_children(records)
    cards = [
        record
        for record in records
        if _is_container_candidate(record)
        and len(children_by_container.get(id(record), [])) >= CARD_GRID_CHILD_COUNT_MIN
    ]
    for row in _card_row_groups(cards):
        row_metrics: list[dict] = []
        for card in row:
            children = children_by_container[id(card)]
            padding = _container_padding(card, children)
            first_child = _first_child_relative_bbox(card, children)
            row_metrics.append(
                {
                    "container": card,
                    "children": children,
                    "padding": padding,
                    "first_child_relative_bbox_pt": first_child,
                }
            )

        medians = {
            "top": _median([item["container"].bbox_pt[1] for item in row_metrics]),
            "width": _median([item["container"].bbox_pt[2] for item in row_metrics]),
            "height": _median([item["container"].bbox_pt[3] for item in row_metrics]),
            "padding_left": _median([item["padding"]["left"] for item in row_metrics]),
            "padding_right": _median([item["padding"]["right"] for item in row_metrics]),
            "padding_top": _median([item["padding"]["top"] for item in row_metrics]),
            "padding_bottom": _median([item["padding"]["bottom"] for item in row_metrics]),
        }
        child_relatives = [
            item["first_child_relative_bbox_pt"]
            for item in row_metrics
            if item["first_child_relative_bbox_pt"] is not None
        ]
        child_medians = None
        if len(child_relatives) == len(row_metrics):
            child_medians = [
                _median([relative[idx] for relative in child_relatives])
                for idx in range(4)
            ]

        inconsistent: list[dict] = []
        for item in row_metrics:
            card = item["container"]
            x, y, w, h = card.bbox_pt
            deltas = {
                "top": round(abs(y - medians["top"]), 2),
                "width": round(abs(w - medians["width"]), 2),
                "height": round(abs(h - medians["height"]), 2),
                "padding_left": round(abs(item["padding"]["left"] - medians["padding_left"]), 2),
                "padding_right": round(abs(item["padding"]["right"] - medians["padding_right"]), 2),
                "padding_top": round(abs(item["padding"]["top"] - medians["padding_top"]), 2),
                "padding_bottom": round(abs(item["padding"]["bottom"] - medians["padding_bottom"]), 2),
            }
            triggered = [
                key
                for key, delta in deltas.items()
                if (
                    key == "top" and delta > CARD_GRID_TOP_TOL_PT
                    or key in {"width", "height"} and delta > CARD_GRID_SIZE_TOL_PT
                    or key.startswith("padding_") and delta > CARD_GRID_PADDING_TOL_PT
                )
            ]
            child_delta = None
            if child_medians is not None and item["first_child_relative_bbox_pt"] is not None:
                child_delta = [
                    round(abs(value - child_medians[idx]), 2)
                    for idx, value in enumerate(item["first_child_relative_bbox_pt"])
                ]
                if any(delta > CARD_GRID_CHILD_RELATIVE_TOL_PT for delta in child_delta):
                    triggered.append("first_child_relative_bbox")
            if triggered:
                inconsistent.append(
                    {
                        "container": _shape_record_detail_for_card_role(card),
                        "children": [
                            _shape_record_detail_for_card_role(child)
                            for child in item["children"]
                        ],
                        "padding_pt": {side: round(value, 2) for side, value in item["padding"].items()},
                        "deltas_from_group_median_pt": deltas,
                        "first_child_relative_bbox_pt": (
                            [round(v, 2) for v in item["first_child_relative_bbox_pt"]]
                            if item["first_child_relative_bbox_pt"] is not None
                            else None
                        ),
                        "first_child_relative_delta_pt": child_delta,
                        "triggered_rules": triggered,
                    }
                )

        if not inconsistent:
            continue
        max_delta = max(
            delta
            for item in inconsistent
            for delta in item["deltas_from_group_median_pt"].values()
        )
        findings.append(
            make_finding(
                "warning", "card_grid_consistency", slide_idx, slide_id, row[0].shape,
                (
                    f"{len(row)} repeated card containers in the same row are not visually consistent"
                ),
                {
                    "evidence_source": "structure_json",
                    "evidence_confidence": "medium",
                    "row_containers": [
                        {
                            **_shape_record_detail_for_card_role(item["container"]),
                            "children": [
                                _shape_record_detail_for_card_role(child)
                                for child in item["children"]
                            ],
                            "padding_pt": {
                                side: round(value, 2)
                                for side, value in item["padding"].items()
                            },
                        }
                        for item in row_metrics
                    ],
                    "group_medians": {key: round(value, 2) for key, value in medians.items()},
                    "inconsistent_containers": inconsistent,
                    "thresholds": {
                        "group_min": CARD_GRID_GROUP_MIN,
                        "child_count_min": CARD_GRID_CHILD_COUNT_MIN,
                        "row_center_tolerance_pt": CARD_GRID_ROW_CENTER_TOL_PT,
                        "top_tolerance_pt": CARD_GRID_TOP_TOL_PT,
                        "size_tolerance_pt": CARD_GRID_SIZE_TOL_PT,
                        "padding_tolerance_pt": CARD_GRID_PADDING_TOL_PT,
                        "first_child_relative_tolerance_pt": CARD_GRID_CHILD_RELATIVE_TOL_PT,
                    },
                    "measured_value": round(max_delta, 2),
                    "threshold": CARD_GRID_PADDING_TOL_PT,
                    "delta": round(max(0.0, max_delta - CARD_GRID_PADDING_TOL_PT), 2),
                    "unit": "pt",
                },
            )
        )


def _text_records(records: list[ShapeRecord]) -> list[ShapeRecord]:
    return [
        record
        for record in records
        if record.kind == "text" and _shape_text(record.shape)
    ]


def _font_size_for_record(ctx: LintContext, record: ShapeRecord) -> Optional[float]:
    if not getattr(record.shape, "has_text_frame", False):
        return None
    return _text_frame_dominant_font_size_pt(ctx, record.shape.text_frame)


def _title_candidates(ctx: LintContext, records: list[ShapeRecord]) -> list[tuple[ShapeRecord, float]]:
    candidates: list[tuple[ShapeRecord, float]] = []
    text_records = _text_records(records)
    if not text_records:
        return candidates
    font_sizes = [
        size
        for record in text_records
        if (size := _font_size_for_record(ctx, record)) is not None
    ]
    max_font_size = max(font_sizes) if font_sizes else None
    for record in text_records:
        size = _font_size_for_record(ctx, record)
        if size is None:
            continue
        _, y, _, _ = record.bbox_pt
        if y <= TITLE_ZONE_BOTTOM_PT and (
            size >= TITLE_FONT_SIZE_MIN_PT
            or (max_font_size is not None and abs(size - max_font_size) <= TOL_PT)
        ):
            candidates.append((record, size))
    return candidates


def _prominent_title_candidates(ctx: LintContext, records: list[ShapeRecord]) -> list[tuple[ShapeRecord, float]]:
    candidates: list[tuple[ShapeRecord, float]] = []
    text_records = _text_records(records)
    if not text_records:
        return candidates
    font_sizes = [
        size
        for record in text_records
        if (size := _font_size_for_record(ctx, record)) is not None
    ]
    max_font_size = max(font_sizes) if font_sizes else None
    for record in text_records:
        size = _font_size_for_record(ctx, record)
        if size is None:
            continue
        _, y, w, _ = record.bbox_pt
        if (
            y <= COVER_TITLE_Y_MAX_PT
            and w >= SECTION_DIVIDER_TITLE_WIDTH_MIN_PT
            and size >= PROMINENT_TITLE_FONT_SIZE_MIN_PT
            and (max_font_size is None or abs(size - max_font_size) <= TOL_PT)
        ):
            candidates.append((record, size))
    return candidates


def _slide_type_detail(ctx: LintContext, records: list[ShapeRecord]) -> dict:
    text_records = _text_records(records)
    title_candidates = _title_candidates(ctx, records)
    prominent = _prominent_title_candidates(ctx, records)
    non_text_records = [
        record
        for record in records
        if record.kind != "text"
    ]
    detail = {
        "slide_type": "content",
        "reason": "default_content_slide",
        "required_element_record": None,
        "top_title_candidates": [
            {
                "shape": _shape_record_detail(record),
                "font_size_pt": size,
                "text_excerpt": _text_excerpt(record.shape),
            }
            for record, size in title_candidates
        ],
        "prominent_title_candidates": [
            {
                "shape": _shape_record_detail(record),
                "font_size_pt": size,
                "text_excerpt": _text_excerpt(record.shape),
            }
            for record, size in prominent
        ],
        "text_shape_count": len(text_records),
        "non_text_shape_count": len(non_text_records),
    }
    if title_candidates:
        record, size = max(title_candidates, key=lambda item: item[1])
        detail.update(
            {
                "required_element_record": {
                    "shape": _shape_record_detail(record),
                    "font_size_pt": size,
                    "text_excerpt": _text_excerpt(record.shape),
                },
                "reason": "top_title_candidate_found",
            }
        )
        return detail

    if prominent:
        record, size = max(prominent, key=lambda item: item[1])
        x, y, w, h = record.bbox_pt
        center_y = y + h / 2
        centered_in_slide = abs(center_y - SLIDE_H_PT / 2) <= SECTION_DIVIDER_TITLE_CENTER_TOL_PT
        if (
            len(text_records) <= SECTION_DIVIDER_TEXT_COUNT_MAX
            and len(non_text_records) <= SECTION_DIVIDER_NON_TEXT_COUNT_MAX
            and SECTION_DIVIDER_TITLE_Y_MIN_PT <= y <= SECTION_DIVIDER_TITLE_Y_MAX_PT
            and w >= SECTION_DIVIDER_TITLE_WIDTH_MIN_PT
            and centered_in_slide
        ):
            detail.update(
                {
                    "slide_type": "section_divider",
                    "reason": "single_prominent_center_title",
                    "required_element_record": {
                        "shape": _shape_record_detail(record),
                        "font_size_pt": size,
                        "text_excerpt": _text_excerpt(record.shape),
                    },
                }
            )
            return detail
        if (
            len(text_records) >= COVER_TEXT_COUNT_MIN
            and COVER_TITLE_Y_MIN_PT <= y <= COVER_TITLE_Y_MAX_PT
            and w >= COVER_TITLE_WIDTH_MIN_PT
        ):
            detail.update(
                {
                    "slide_type": "cover",
                    "reason": "cover_prominent_title_with_metadata_text",
                    "required_element_record": {
                        "shape": _shape_record_detail(record),
                        "font_size_pt": size,
                        "text_excerpt": _text_excerpt(record.shape),
                    },
                }
            )
            return detail

    return detail


def check_missing_required_element(ctx, slide_idx, slide_id, records: list[ShapeRecord], findings):
    text_records = _text_records(records)
    if not text_records:
        return
    slide_type = _slide_type_detail(ctx, records)
    if slide_type.get("required_element_record"):
        return
    top_texts = [
        {
            "shape": _shape_record_detail(record),
            "font_size_pt": _font_size_for_record(ctx, record),
            "text_excerpt": _text_excerpt(record.shape),
        }
        for record in text_records
        if record.bbox_pt[1] <= TITLE_ZONE_BOTTOM_PT
    ]
    findings.append(
        make_finding(
            "warning", "missing_required_element", slide_idx, slide_id, None,
            "slide has content but no machine-detected title/header text",
            {
                "evidence_source": "structure_json",
                "evidence_confidence": "medium",
                "missing_element": "title",
                "slide_type": slide_type["slide_type"],
                "slide_type_reason": slide_type["reason"],
                "text_shape_count": len(text_records),
                "top_text_candidates": top_texts,
                "prominent_title_candidates": slide_type["prominent_title_candidates"],
                "required_zone_pt": {
                    "top": 0,
                    "bottom": TITLE_ZONE_BOTTOM_PT,
                },
                "required_font_size_pt_min": TITLE_FONT_SIZE_MIN_PT,
                "prominent_title_font_size_pt_min": PROMINENT_TITLE_FONT_SIZE_MIN_PT,
                "measured_value": 0,
                "threshold": 1,
                "delta": 1,
                "unit": "title_candidate_count",
            },
        )
    )


def check_heading_hierarchy(ctx, slide_idx, slide_id, records: list[ShapeRecord], findings):
    text_records = _text_records(records)
    if len(text_records) < 2:
        return
    title_candidates = _title_candidates(ctx, records)
    if not title_candidates:
        return
    title_record, title_size = max(title_candidates, key=lambda item: item[1])
    body_candidates: list[tuple[ShapeRecord, float]] = []
    for record in text_records:
        if record is title_record:
            continue
        size = _font_size_for_record(ctx, record)
        if size is None:
            continue
        body_candidates.append((record, size))
    if not body_candidates:
        return

    largest_body, body_size = max(body_candidates, key=lambda item: item[1])
    triggered_rules: list[str] = []
    if body_size >= title_size + HEADING_BODY_FONT_DELTA_PT:
        triggered_rules.append("body_larger_than_title")
    if title_record.bbox_pt[1] > min(record.bbox_pt[1] for record, _ in body_candidates) + TOL_PT:
        triggered_rules.append("title_after_body")
    if len(title_candidates) > 1:
        triggered_rules.append("multiple_title_candidates")
    if not triggered_rules:
        return

    findings.append(
        make_finding(
            "warning", "heading_hierarchy_broken", slide_idx, slide_id, title_record.shape,
            "machine-detected heading hierarchy risk: " + ", ".join(triggered_rules),
            {
                "evidence_source": "structure_json",
                "evidence_confidence": "medium",
                "title_candidate": {
                    **_shape_record_detail(title_record),
                    "font_size_pt": round(title_size, 2),
                    "text_excerpt": _text_excerpt(title_record.shape),
                },
                "largest_body_candidate": {
                    **_shape_record_detail(largest_body),
                    "font_size_pt": round(body_size, 2),
                    "text_excerpt": _text_excerpt(largest_body.shape),
                },
                "title_candidate_count": len(title_candidates),
                "triggered_rules": triggered_rules,
                "thresholds": {
                    "body_title_font_delta_pt": HEADING_BODY_FONT_DELTA_PT,
                    "title_zone_bottom_pt": TITLE_ZONE_BOTTOM_PT,
                },
                "measured_value": round(body_size - title_size, 2),
                "threshold": HEADING_BODY_FONT_DELTA_PT,
                "delta": round(max(0.0, body_size - title_size - HEADING_BODY_FONT_DELTA_PT), 2),
                "unit": "pt",
            },
        )
    )


def _visual_reading_key(record: ShapeRecord) -> tuple[int, float, float]:
    x, y, _, _ = record.bbox_pt
    return (round(y / READING_ORDER_TOP_BUCKET_PT), x, y)


def _reading_unit_related(first: ShapeRecord, second: ShapeRecord) -> bool:
    ax, ay, aw, ah = first.bbox_pt
    bx, by, bw, bh = second.bbox_pt
    if aw >= READING_ORDER_FULL_WIDTH_TEXT_PT or bw >= READING_ORDER_FULL_WIDTH_TEXT_PT:
        return False
    if min(ay, by) < TITLE_ZONE_BOTTOM_PT:
        return False
    if min(ay, by) >= READING_ORDER_FOOTER_Y_MIN_PT:
        return False
    ax1, ay1, ax2, ay2 = _bbox_edges(first.bbox_pt)
    bx1, by1, bx2, by2 = _bbox_edges(second.bbox_pt)
    vertical_gap = _axis_gap(ay1, ay2, by1, by2)
    if vertical_gap > READING_ORDER_UNIT_VERTICAL_GAP_PT:
        return False
    acx, _ = _bbox_center(first.bbox_pt)
    bcx, _ = _bbox_center(second.bbox_pt)
    center_threshold = min(
        READING_ORDER_UNIT_CENTER_X_MAX_PT,
        max(aw, bw) * READING_ORDER_UNIT_CENTER_X_RATIO,
    )
    center_aligned = abs(acx - bcx) <= center_threshold
    edge_aligned = (
        abs(ax1 - bx1) <= READING_ORDER_UNIT_EDGE_X_TOL_PT
        or abs(ax2 - bx2) <= READING_ORDER_UNIT_EDGE_X_TOL_PT
    )
    return center_aligned or edge_aligned


def _reading_units(records: list[ShapeRecord]) -> list[list[ShapeRecord]]:
    units: list[list[ShapeRecord]] = []
    for record in sorted(records, key=_visual_reading_key):
        target: list[ShapeRecord] | None = None
        for unit in units:
            if any(_reading_unit_related(record, existing) for existing in unit):
                target = unit
                break
        if target is None:
            units.append([record])
        else:
            target.append(record)
    return units


def _unit_bbox(unit: list[ShapeRecord]) -> tuple[float, float, float, float]:
    x1 = min(record.bbox_pt[0] for record in unit)
    y1 = min(record.bbox_pt[1] for record in unit)
    x2 = max(record.bbox_pt[0] + record.bbox_pt[2] for record in unit)
    y2 = max(record.bbox_pt[1] + record.bbox_pt[3] for record in unit)
    return (x1, y1, x2 - x1, y2 - y1)


def _unit_reading_key(unit: list[ShapeRecord]) -> tuple[int, float, float]:
    x, y, _, _ = _unit_bbox(unit)
    return (round(y / READING_ORDER_TOP_BUCKET_PT), x, y)


def _grouped_visual_reading_order(records: list[ShapeRecord]) -> tuple[list[ShapeRecord], list[list[ShapeRecord]]]:
    units = _reading_units(records)
    ordered_units = sorted(units, key=_unit_reading_key)
    ordered_records: list[ShapeRecord] = []
    for unit in ordered_units:
        ordered_records.extend(sorted(unit, key=_visual_reading_key))
    return ordered_records, ordered_units


def _reading_order_inversions(ordered_records: list[ShapeRecord]) -> list[dict]:
    inversions: list[dict] = []
    for idx, first in enumerate(ordered_records):
        for second in ordered_records[idx + 1:]:
            if first.source_order_index > second.source_order_index:
                inversions.append(
                    {
                        "visual_first": _shape_record_detail(first),
                        "visual_second": _shape_record_detail(second),
                    }
                )
            if len(inversions) >= READING_ORDER_INVERSION_MIN:
                break
        if len(inversions) >= READING_ORDER_INVERSION_MIN:
            break
    return inversions


def check_reading_order(ctx, slide_idx, slide_id, records: list[ShapeRecord], findings):
    text_records = _text_records(records)
    if len(text_records) < 3:
        return
    slide_type = _slide_type_detail(ctx, records)
    if slide_type.get("slide_type") in {"cover", "section_divider"}:
        return
    visual_order, reading_units = _grouped_visual_reading_order(text_records)
    inversions = _reading_order_inversions(visual_order)
    if len(inversions) < READING_ORDER_INVERSION_MIN:
        return
    findings.append(
        make_finding(
            "warning", "reading_order", slide_idx, slide_id, visual_order[0].shape,
            (
                f"selection/source order disagrees with visual top-left order "
                f"({len(inversions)} inversion(s))"
            ),
            {
                "evidence_source": "structure_json",
                "evidence_confidence": "medium",
                "visual_order": [_shape_record_detail(record) for record in visual_order],
                "reading_units": [
                    {
                        "bbox_pt": [round(v, 2) for v in _unit_bbox(unit)],
                        "records": [_shape_record_detail(record) for record in sorted(unit, key=_visual_reading_key)],
                    }
                    for unit in sorted(reading_units, key=_unit_reading_key)
                ],
                "source_order": [
                    _shape_record_detail(record)
                    for record in sorted(text_records, key=lambda item: item.source_order_index)
                ],
                "inversions": inversions,
                "measured_value": len(inversions),
                "threshold": READING_ORDER_INVERSION_MIN,
                "delta": len(inversions) - READING_ORDER_INVERSION_MIN + 1,
                "unit": "inversions",
            },
        )
    )


def check_color_only_meaning(slide_idx, slide_id, records: list[ShapeRecord], findings):
    shape_records = [
        record
        for record in records
        if record.kind == "shape"
        and _solid_shape_fill_rgb_hex(record.shape)
        and not _shape_text(record.shape)
        and _bbox_area(record.bbox_pt) > 0
    ]
    for idx, first in enumerate(shape_records):
        group = [first]
        first_w = first.bbox_pt[2]
        first_h = first.bbox_pt[3]
        for second in shape_records[idx + 1:]:
            if abs(second.bbox_pt[2] - first_w) > COLOR_ONLY_SIZE_TOL_PT:
                continue
            if abs(second.bbox_pt[3] - first_h) > COLOR_ONLY_SIZE_TOL_PT:
                continue
            ax1, ay1, ax2, ay2 = _bbox_edges(first.bbox_pt)
            bx1, by1, bx2, by2 = _bbox_edges(second.bbox_pt)
            gap = min(
                _axis_gap(ax1, ax2, bx1, bx2),
                _axis_gap(ay1, ay2, by1, by2),
            )
            if gap <= COLOR_ONLY_GROUP_GAP_PT:
                group.append(second)
        if len(group) < 2:
            continue
        colors = sorted({
            color
            for record in group
            if (color := _solid_shape_fill_rgb_hex(record.shape))
        })
        if len(colors) < 2:
            continue
        findings.append(
            make_finding(
                "warning", "color_only_meaning", slide_idx, slide_id, first.shape,
                (
                    "similar unlabeled shapes differ only by color; "
                    "add text, icon, pattern, or shape cue"
                ),
                {
                    "evidence_source": "structure_json",
                    "evidence_confidence": "low",
                    "group": [
                        {
                            **_shape_record_detail(record),
                            "fill_hex": _solid_shape_fill_rgb_hex(record.shape),
                        }
                        for record in group
                    ],
                    "colors": colors,
                    "non_color_cue_present": False,
                    "triggered_rules": ["unlabeled_similar_shapes_distinct_colors"],
                    "thresholds": {
                        "group_gap_pt": COLOR_ONLY_GROUP_GAP_PT,
                        "size_tolerance_pt": COLOR_ONLY_SIZE_TOL_PT,
                    },
                    "measured_value": len(colors),
                    "threshold": 2,
                    "delta": len(colors) - 1,
                    "unit": "distinct_colors",
                },
            )
        )
        return


def check_text_vertical_balance(ctx, slide_idx, slide_id, shape, bbox, findings):
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    if not text_frame.text.strip():
        return
    if not _shape_has_visible_fill_or_line(shape):
        return

    normalized = normalize_bbox(ctx, bbox)
    _, box_y_pt, box_w_pt, box_h_pt = normalized
    if box_h_pt < TEXT_VERTICAL_BALANCE_MIN_BOX_HEIGHT_PT:
        return

    font_size_pt = _text_frame_dominant_font_size_pt(ctx, text_frame)
    if font_size_pt is None or font_size_pt <= 0:
        return

    non_empty_paragraphs = [para for para in text_frame.paragraphs if para.text.strip()]
    if not non_empty_paragraphs:
        return
    anchor = text_frame.vertical_anchor
    if (
        anchor is None
        and len(non_empty_paragraphs) == 1
        and box_w_pt >= SAFE_TEXT_AREA_PT[2] - TOL_PT
    ):
        return

    line_heights_pt = [
        _paragraph_line_height_pt(ctx, para, font_size_pt)
        for para in non_empty_paragraphs
    ]
    estimated_text_height_pt = sum(line_heights_pt)
    margin_top_pt = _normalized_length_pt(text_frame.margin_top, ctx.scale_y)
    margin_bottom_pt = _normalized_length_pt(text_frame.margin_bottom, ctx.scale_y)
    inner_height_pt = box_h_pt - margin_top_pt - margin_bottom_pt
    if estimated_text_height_pt <= 0 or inner_height_pt <= 0:
        return
    if estimated_text_height_pt > inner_height_pt + TOL_PT:
        return

    free_space_pt = max(0.0, inner_height_pt - estimated_text_height_pt)
    anchor_name = _vertical_anchor_name(anchor)

    if anchor == MSO_VERTICAL_ANCHOR.MIDDLE:
        anchor_top_offset_pt = free_space_pt / 2
        anchor_bottom_offset_pt = free_space_pt / 2
    elif anchor == MSO_VERTICAL_ANCHOR.BOTTOM:
        anchor_top_offset_pt = free_space_pt
        anchor_bottom_offset_pt = 0.0
    else:
        anchor_top_offset_pt = 0.0
        anchor_bottom_offset_pt = free_space_pt

    actual_top_pad_pt = margin_top_pt + anchor_top_offset_pt
    actual_bottom_pad_pt = margin_bottom_pt + anchor_bottom_offset_pt
    visual_center_y_pt = box_y_pt + margin_top_pt + anchor_top_offset_pt + estimated_text_height_pt / 2
    box_geometric_center_y_pt = box_y_pt + box_h_pt / 2
    visual_center_offset_pt = visual_center_y_pt - box_geometric_center_y_pt
    free_space_ratio = free_space_pt / inner_height_pt
    visual_center_offset_ratio = abs(visual_center_offset_pt) / inner_height_pt

    triggered_rules: list[str] = []
    if (
        free_space_ratio > TEXT_VERTICAL_BALANCE_DEAD_SPACE_RATIO_MAX
        and anchor != MSO_VERTICAL_ANCHOR.MIDDLE
    ):
        triggered_rules.append("dead_space_ratio_non_middle")
    if anchor in (None, MSO_VERTICAL_ANCHOR.TOP) and actual_bottom_pad_pt > TEXT_VERTICAL_BALANCE_DEAD_SPACE_PT_MAX:
        triggered_rules.append("top_anchor_bottom_dead_space")
    if anchor == MSO_VERTICAL_ANCHOR.BOTTOM and actual_top_pad_pt > TEXT_VERTICAL_BALANCE_DEAD_SPACE_PT_MAX:
        triggered_rules.append("bottom_anchor_top_dead_space")
    if (
        anchor == MSO_VERTICAL_ANCHOR.MIDDLE
        and abs(margin_top_pt - margin_bottom_pt)
        > TEXT_VERTICAL_BALANCE_MIDDLE_MARGIN_ASYMMETRY_PT_MAX
    ):
        triggered_rules.append("middle_anchor_margin_asymmetry")
    if visual_center_offset_ratio > TEXT_VERTICAL_BALANCE_CENTER_OFFSET_RATIO_MAX:
        triggered_rules.append("visual_center_offset")

    if not triggered_rules:
        return

    findings.append(
        make_finding(
            "warning", "text_vertical_balance", slide_idx, slide_id, shape,
            (
                "text fits but vertical padding/visual center is imbalanced: "
                f"{', '.join(triggered_rules)}"
            ),
            {
                "anchor": anchor_name,
                "font_size_pt": round(font_size_pt, 2),
                "line_height_pt": round(estimated_text_height_pt / len(non_empty_paragraphs), 2),
                "paragraphs": len(non_empty_paragraphs),
                "estimated_text_height_pt": round(estimated_text_height_pt, 2),
                "inner_height_pt": round(inner_height_pt, 2),
                "free_space_pt": round(free_space_pt, 2),
                "margin_top_pt": round(margin_top_pt, 2),
                "margin_bottom_pt": round(margin_bottom_pt, 2),
                "visual_center_offset_pt": round(visual_center_offset_pt, 2),
                "triggered_rules": triggered_rules,
                "bbox_pt": [round(v, 2) for v in normalized],
                "box_width_pt": round(box_w_pt, 2),
                "actual_padding_pt": {
                    "top": round(actual_top_pad_pt, 2),
                    "bottom": round(actual_bottom_pad_pt, 2),
                },
                "thresholds": {
                    "dead_space_ratio_max": TEXT_VERTICAL_BALANCE_DEAD_SPACE_RATIO_MAX,
                    "top_or_bottom_dead_space_pt_max": TEXT_VERTICAL_BALANCE_DEAD_SPACE_PT_MAX,
                    "middle_anchor_padding_asymmetry_pt_max": (
                        TEXT_VERTICAL_BALANCE_MIDDLE_MARGIN_ASYMMETRY_PT_MAX
                    ),
                    "visual_center_offset_ratio_max": TEXT_VERTICAL_BALANCE_CENTER_OFFSET_RATIO_MAX,
                    "min_box_height_pt": TEXT_VERTICAL_BALANCE_MIN_BOX_HEIGHT_PT,
                },
            },
        )
    )


# ---- Driver ----------------------------------------------------------------


def lint_pptx(
    path: Path,
    *,
    profile: str = "default",
    policy: Optional[LintPolicy] = None,
    rendered_image_dir: Optional[Path] = None,
) -> List[Finding]:
    if policy is None:
        try:
            policy = LINT_PROFILES[profile]
        except KeyError as exc:
            raise ValueError(f"unknown lint profile: {profile}") from exc
    prs = Presentation(str(path))
    findings: List[Finding] = []

    actual_w = emu_to_pt(prs.slide_width)
    actual_h = emu_to_pt(prs.slide_height)
    ctx = make_context(actual_w, actual_h, policy=policy)
    rendered_images: dict[int, tuple[Path, Image.Image]] = {}
    if rendered_image_dir is not None:
        for idx in range(1, len(prs.slides) + 1):
            image_path = _rendered_slide_image_path(rendered_image_dir, idx)
            if image_path is not None:
                rendered_images[idx] = (image_path, Image.open(image_path).convert("RGB"))
    if not ctx.proportional_to_base:
        findings.append(
            Finding(
                severity="warning",
                check="slide_size",
                slide_index=0,
                slide_id=None,
                shape_id=None,
                shape_name=None,
                message=(
                    f"slide size {actual_w:.1f}x{actual_h:.1f}pt is not proportional to "
                    f"guideline base {SLIDE_W_PT}x{SLIDE_H_PT}pt"
                ),
                detail={
                    "actual_pt": [actual_w, actual_h],
                    "base_pt": [SLIDE_W_PT, SLIDE_H_PT],
                    "scale": [round(ctx.scale_x, 4), round(ctx.scale_y, 4)],
                },
            )
        )

    for idx, slide in enumerate(prs.slides, start=1):
        slide_id = getattr(slide, "slide_id", None)
        markers = animation_markers(slide)
        if markers:
            findings.append(
                make_finding(
                    "error", "animation_present", idx, slide_id, None,
                    "slide contains animation/transition XML: " + ", ".join(markers),
                    {"markers": markers},
                )
            )
        records: list[ShapeRecord] = []
        for source_order_index, shape in enumerate(iter_shapes(slide.shapes), start=1):
            bbox = shape_bbox_pt(shape)
            if bbox is None:
                continue
            records.append(
                ShapeRecord(
                    shape=shape,
                    actual_bbox_pt=bbox,
                    bbox_pt=normalize_bbox(ctx, bbox),
                    kind=_shape_kind(shape),
                    source_order_index=source_order_index,
                )
            )
        slide_type = _slide_type_detail(ctx, records)
        check_object_relationships(idx, slide_id, records, findings, ctx=ctx)
        check_decorative_isolated_lines(idx, slide_id, records, findings)
        check_inner_padding_imbalance(idx, slide_id, records, findings)
        check_card_grid_consistency(idx, slide_id, records, findings)
        check_missing_required_element(ctx, idx, slide_id, records, findings)
        check_heading_hierarchy(ctx, idx, slide_id, records, findings)
        check_reading_order(ctx, idx, slide_id, records, findings)
        check_color_only_meaning(idx, slide_id, records, findings)
        for record in records:
            shape = record.shape
            bbox = record.actual_bbox_pt
            before_overflow_count = len(findings)
            check_overflow(ctx, idx, slide_id, shape, bbox, findings)
            overflowed = any(f.check in ("box_canvas_overflow", "overflow_shapes", "overflow_images") for f in findings[before_overflow_count:])
            check_safe_text_area(ctx, idx, slide_id, record, findings, slide_type)
            check_safe_margins(ctx, idx, slide_id, shape, bbox, findings, slide_type)
            check_geometry_rounding(ctx, idx, slide_id, shape, bbox, findings)
            check_image_upscale(ctx, idx, slide_id, shape, bbox, findings)
            check_key_area_cropped(ctx, idx, slide_id, shape, bbox, findings)
            check_alt_text(idx, slide_id, shape, findings)
            check_autofit(idx, slide_id, shape, findings)
            check_font(ctx, idx, slide_id, shape, findings)
            check_line_height(ctx, idx, slide_id, shape, findings)
            check_alignment(ctx, idx, slide_id, shape, findings)
            check_badge_alignment(ctx, idx, slide_id, shape, record.bbox_pt, findings)
            check_wrap_break_changes_meaning(ctx, idx, slide_id, shape, bbox, findings)
            check_color(
                ctx,
                idx,
                slide_id,
                record,
                records,
                findings,
                measure_structural_contrast=rendered_image_dir is None,
            )
            if idx in rendered_images:
                image_path, rendered_image = rendered_images[idx]
                check_rendered_contrast(
                    ctx,
                    idx,
                    slide_id,
                    record,
                    rendered_image,
                    image_path,
                    findings,
                )
            if not overflowed:
                check_text_vertical_balance(ctx, idx, slide_id, shape, bbox, findings)

    return findings


def extract_pptx_structure(
    path: Path,
    *,
    profile: str = "default",
    policy: Optional[LintPolicy] = None,
) -> list[dict]:
    if policy is None:
        try:
            policy = LINT_PROFILES[profile]
        except KeyError as exc:
            raise ValueError(f"unknown lint profile: {profile}") from exc
    prs = Presentation(str(path))
    actual_w = emu_to_pt(prs.slide_width)
    actual_h = emu_to_pt(prs.slide_height)
    ctx = make_context(actual_w, actual_h, policy=policy)

    structure: list[dict] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_id = getattr(slide, "slide_id", None)
        records: list[ShapeRecord] = []
        for source_order_index, shape in enumerate(iter_shapes(slide.shapes), start=1):
            bbox = shape_bbox_pt(shape)
            if bbox is None:
                continue
            records.append(
                ShapeRecord(
                    shape=shape,
                    actual_bbox_pt=bbox,
                    bbox_pt=normalize_bbox(ctx, bbox),
                    kind=_shape_kind(shape),
                    source_order_index=source_order_index,
                )
            )
        for relation in extract_structure_relations(records):
            detail = _structure_relation_detail(relation)
            detail["slide_index"] = idx
            detail["slide_id"] = slide_id
            structure.append(detail)
    return structure


# ---- Consolidation (recurring template-level issues) ----------------------


def _slide_range(slides: List[int]) -> str:
    """Format a sorted list of ints as compact ranges, e.g. [2,3,4,7,9,10] -> '2-4, 7, 9-10'."""
    if not slides:
        return ""
    runs: List[tuple] = []
    start = prev = slides[0]
    for s in slides[1:]:
        if s == prev + 1:
            prev = s
        else:
            runs.append((start, prev))
            start = prev = s
    runs.append((start, prev))
    return ", ".join(str(a) if a == b else f"{a}-{b}" for a, b in runs)


def _group_key(f: Finding) -> tuple:
    """Return a grouping key that survives shape_name drift across slides.

    PowerPoint assigns shape_name like "Freeform 7" based on insertion order, so
    the same logical template element can be "Freeform 6" on one slide and
    "Freeform 7" on the next. When a finding carries a bbox in its detail, use
    that as the stable identifier; otherwise fall back to shape_name.
    """
    bbox = f.detail.get("bbox_pt") if f.detail else None
    if bbox:
        bbox_sig = ",".join(f"{round(v):d}" for v in bbox)
        return (f.check, f"bbox:{bbox_sig}", f.message)
    return (f.check, f.shape_name or "", f.message)


def consolidate_recurring(findings: List[Finding], min_slides: int = 3) -> List[Finding]:
    """Collapse identical findings that appear on many slides into one deck-level finding.

    When the same key (see _group_key) appears on >= min_slides distinct slides,
    replace the per-slide findings with a single deck-level entry that lists the
    affected slide range.
    """
    deck_level = [f for f in findings if f.slide_index == 0]
    per_slide = [f for f in findings if f.slide_index != 0]

    groups: dict = defaultdict(list)
    for f in per_slide:
        groups[_group_key(f)].append(f)

    consolidated: List[Finding] = []
    individual: List[Finding] = []
    for key, group in groups.items():
        slides = sorted({f.slide_index for f in group})
        if len(slides) >= min_slides:
            example = group[0]
            consolidated.append(
                Finding(
                    severity=example.severity,
                    check=example.check,
                    slide_index=0,
                    slide_id=None,
                    shape_id=None,
                    shape_name=example.shape_name,
                    message=(
                        f"{example.message} "
                        f"(recurring on {len(slides)} slides: {_slide_range(slides)})"
                    ),
                    detail={
                        **(example.detail or {}),
                        "affected_slides": slides,
                        "occurrences": len(group),
                        "example_slide_index": example.slide_index,
                        "example_shape_id": example.shape_id,
                    },
                )
            )
        else:
            individual.extend(group)

    consolidated.sort(key=lambda f: (f.check, f.shape_name or ""))
    return deck_level + consolidated + individual


# ---- Output ----------------------------------------------------------------


def format_text(findings: List[Finding]) -> str:
    if not findings:
        return "OK: no issues found.\n"
    err_n = sum(1 for f in findings if f.severity == "error")
    warn_n = sum(1 for f in findings if f.severity == "warning")
    by_slide: dict = {}
    for f in findings:
        by_slide.setdefault(f.slide_index, []).append(f)
    lines = [f"Found {len(findings)} issues ({err_n} errors, {warn_n} warnings)", ""]
    for idx in sorted(by_slide):
        header = "Deck-level" if idx == 0 else f"Slide {idx}"
        lines.append(f"--- {header} ---")
        for f in by_slide[idx]:
            tag = "[ERR ]" if f.severity == "error" else "[WARN]"
            if f.shape_name and f.shape_id is not None:
                loc = f"{f.shape_name} (id={f.shape_id})"
            elif f.shape_name:
                loc = f.shape_name
            else:
                loc = "-"
            lines.append(f"{tag} {f.check}: {f.message}  [{loc}]")
        lines.append("")
    return "\n".join(lines)


def filter_by_severity(findings: List[Finding], min_sev: str) -> List[Finding]:
    if min_sev == "all":
        return findings
    if min_sev == "error":
        return [f for f in findings if f.severity == "error"]
    if min_sev == "warning":
        return findings  # warning is the lower bar; include errors too
    return findings


def _load_fix_policy() -> dict[str, dict]:
    """Read rules.lint.fix_policy from doc/slide-guideline-v1.yml (POLICY-001).

    Single source of truth: pptx_fix.py also reads the same table. PyYAML is
    expected to be present in the project venv; fall back to an empty dict
    if it's missing so that imports do not break in unconfigured envs
    (callers will then take the manual_required default path).
    """
    guideline = Path(__file__).resolve().parents[3] / "doc" / "slide-guideline-v1.yml"
    try:
        import yaml  # type: ignore

        data = yaml.safe_load(guideline.read_text(encoding="utf-8"))
    except Exception:
        return {}
    return (data.get("rules") or {}).get("lint", {}).get("fix_policy") or {}


FIX_POLICY = _load_fix_policy()


MANUAL_REQUIRED_REASONS = {
    "slide_size": "slide canvas change can reflow all content",
    "animation_present": "static delivery impact requires author decision",
    "box_canvas_overflow": "text box bbox extends past the slide canvas; clip width/height on the right/bottom side only, never relocate the box",
    "text_box_overflow": "rendered text overflows its own box.height; resolve by font_size shrink / line_height compress / box.height expand (multi_step)",
    "text_canvas_overflow": "text content cannot be wrapped to fit within canvas right edge; resolve by enabling word_wrap / shrinking box.width to canvas / shrinking font_size (multi_step)",
    "overflow_shapes": "requires visual layout decision",
    "overflow_images": "requires crop or placement decision",
    "safe_text_area_text": "requires layout hierarchy decision",
    "safe_margins": "requires placement decision",
    "image_aspect_distortion": "requires source image or crop decision",
    "image_upscale_ratio": "requires higher-resolution image or size decision",
    "alt_text_required": "requires semantic description from author",
    "font_family": "font replacement can affect text metrics and brand intent",
    "font_size_scale": "font size snapping requires layout fit review",
    "line_height": "line spacing change requires text fit review",
    "alignment_left_top": "alignment can be intentional and needs visual review",
    "text_color_allowlist": "color substitution requires semantic and brand review",
    "background_color_palette": "fill color substitution requires brand review",
    "contrast_ratio": "contrast repair requires foreground/background design decision",
    "low_contrast": "unreadable contrast requires visual design decision",
    "color_only_meaning": "non-color cue choice requires semantic design decision",
    "heading_hierarchy_broken": "hierarchy repair requires template/content intent review",
    "key_area_cropped": "crop repair requires important-area decision",
    "missing_required_element": "missing element repair requires content/template decision",
    "reading_order": "order repair requires logical content decision",
    "wrap_break_changes_meaning": "line break repair requires copy/layout decision",
    "text_overlap": "overlap repair requires layout decision",
    "object_overlap": "overlap repair requires layout decision",
    "object_gap_too_small": "spacing repair requires layout decision",
    "inner_padding_imbalance": "container padding repair requires composition review",
    "card_grid_consistency": "repeated card repair requires template grouping intent review",
    "text_vertical_balance": "vertical balance repair requires visual review",
}


def _geometry_auto_fixable(evidence: dict) -> bool:
    if evidence.get("affected_slides"):
        return False
    drifted = evidence.get("drifted")
    if not isinstance(drifted, dict) or not drifted:
        return False
    return all(abs(value - round(value)) < 0.1 for value in drifted.values())


def _card_grid_auto_fixable(check: str, evidence: dict) -> bool:
    if check != "card_grid_consistency":
        return False
    medians = evidence.get("group_medians")
    inconsistent = evidence.get("inconsistent_containers")
    if not isinstance(medians, dict) or not isinstance(inconsistent, list) or not inconsistent:
        return False
    for key in ("top", "width", "height"):
        if key not in medians:
            return False
    return True


def _contrast_auto_fixable(check: str, evidence: dict) -> bool:
    if check not in {"contrast_ratio", "low_contrast"}:
        return False
    if evidence.get("affected_slides"):
        return False
    if evidence.get("background_complexity") == "complex":
        return False
    foreground = evidence.get("foreground_hex") or evidence.get("text_hex")
    background = evidence.get("background_hex")
    required = evidence.get("required_ratio")
    if not foreground or not background or not required:
        return False
    return _contrast_candidate(foreground, background, required) is not None


_AUTO_FIX_REASONS = {
    "text_autofit_disabled": "mechanical text_frame.auto_size NONE change",
    "font_size_scale": "nearest allowed font size can be applied when pptx_fix fit checks pass",
    "line_height": "nearest allowed fixed line height can be applied when pptx_fix fit checks pass",
    "alignment_left_top": "mechanical paragraph LEFT and text-frame TOP alignment change",
    "badge_alignment": "mechanical paragraph CENTER and text-frame MIDDLE alignment change for badge container",
}

_JUDGEMENT_AUTO_FIX_REASONS = {
    "geometry_rounding": "all drifted geometry values are within 0.1pt of integer coordinates",
    "wrap_break_changes_meaning": "widen_to_fit_within_safe_area",
    "low_contrast": "nearest allowed foreground color passes the required contrast ratio",
    "contrast_ratio": "nearest allowed foreground color passes the required contrast ratio",
    "card_grid_consistency": (
        "row containers can be aligned to the group median geometry (top/width/height + padding)"
    ),
}


def _judgement_auto_fixable(check: str, evidence: dict) -> bool:
    """Per-check predicate that decides whether a `judgement`-policy finding
    qualifies for immediate auto_fix_candidate promotion (instead of waiting
    for SPA judgement). Mirrors the historical branching: presence of a
    concrete fix candidate / drift within tolerance / etc.
    """
    if check == "geometry_rounding":
        return _geometry_auto_fixable(evidence)
    if check == "wrap_break_changes_meaning":
        return _widen_to_fit_candidate(evidence) is not None
    if check in {"low_contrast", "contrast_ratio"}:
        return _contrast_auto_fixable(check, evidence)
    if check == "card_grid_consistency":
        return _card_grid_auto_fixable(check, evidence)
    return False


def _judgement_manual_reason(check: str, evidence: dict) -> str:
    """Reason string when a judgement-policy check stays at manual_required.
    Captures the small set of context-aware overrides the previous
    if/elif had for contrast 系.
    """
    if check in {"contrast_ratio", "low_contrast"}:
        if evidence.get("affected_slides"):
            return "recurring contrast finding must be expanded with --no-consolidate before auto-fix"
        if evidence.get("background_complexity") == "complex":
            return "complex background requires local-contrast review before auto-fix"
        return MANUAL_REQUIRED_REASONS.get(check, "contrast repair requires visual review")
    return MANUAL_REQUIRED_REASONS.get(check, "requires manual review")


def _fixability_for_json(check: str, evidence: dict) -> dict:
    """Resolve the fixability/rule/reason triple for `check` using the
    single source of truth `rules.lint.fix_policy` in
    doc/slide-guideline-v1.yml (POLICY-001 段階 2).

    apply_mode dispatch:
      - no_fix         : never auto-fixable (the lint check exists but no
                         machine-repair path is provided).
                         decorative_isolated_lines surfaces as
                         `decorative_review`; everything else stays
                         `manual_required`.
      - auto_fix       : always emit `auto_fix_candidate` with the policy
                         fix_rule and a per-check reason from
                         _AUTO_FIX_REASONS.
      - judgement_fix  : emit `auto_fix_candidate` only when the per-check
                         predicate `_judgement_auto_fixable` is satisfied;
                         otherwise stay at `manual_required` (the SPA-driven
                         promotion path takes over).
    """
    entry = FIX_POLICY.get(check, {})
    apply_mode = entry.get("apply_mode")
    fix_rule = entry.get("fix_rule")

    if apply_mode == "no_fix" or apply_mode is None:
        reason = MANUAL_REQUIRED_REASONS.get(check, "requires manual review")
        return {
            "fixability": "manual_required",
            "fixability_rule": None,
            "fixability_reason": reason,
            "manual_required_reason": reason,
        }

    if apply_mode == "auto_fix":
        return {
            "fixability": "auto_fix_candidate",
            "fixability_rule": fix_rule,
            "fixability_reason": _AUTO_FIX_REASONS.get(check, f"mechanical {fix_rule} change"),
        }

    # apply_mode == "judgement_fix"
    if _judgement_auto_fixable(check, evidence):
        return {
            "fixability": "auto_fix_candidate",
            "fixability_rule": fix_rule,
            "fixability_reason": _JUDGEMENT_AUTO_FIX_REASONS.get(check, f"auto-fix candidate for {fix_rule}"),
        }
    reason = _judgement_manual_reason(check, evidence)
    return {
        "fixability": "manual_required",
        "fixability_rule": fix_rule,
        "fixability_reason": reason,
        "manual_required_reason": reason,
    }


def _widen_to_fit_candidate(evidence: dict) -> Optional[dict]:
    candidates = evidence.get("candidate_values")
    if not isinstance(candidates, list):
        return None
    for entry in candidates:
        if not isinstance(entry, dict):
            continue
        if entry.get("strategy") != "widen_to_fit":
            continue
        bbox = entry.get("bbox_pt")
        if not isinstance(bbox, (list, tuple)) or len(bbox) != 4:
            continue
        return entry
    return None


def _candidate_values_for_json(check: str, evidence: dict) -> Optional[dict]:
    if check == "text_autofit_disabled":
        return {"auto_size": "NONE"}
    if check == "geometry_rounding":
        candidate: dict = {}
        if isinstance(evidence.get("drifted"), dict):
            candidate["rounded_values_pt"] = {
                key: round(value) for key, value in evidence["drifted"].items()
            }
        if isinstance(evidence.get("bbox_pt"), list) and len(evidence["bbox_pt"]) == 4:
            candidate["rounded_bbox_pt"] = [round(value) for value in evidence["bbox_pt"]]
        return candidate or None
    if check == "font_family":
        return {
            "allowed_font_families": list(ALLOWED_FONT_FAMILIES),
            "candidate_font_family": ALLOWED_FONT_FAMILIES[0],
        }
    if check == "font_size_scale":
        return {
            "size_pt": evidence.get("nearest_allowed_size_pt"),
            "allowed_size_scale_pt": sorted(ALLOWED_FONT_SIZES_PT),
        }
    if check == "line_height":
        return {"allowed_line_heights_pt": sorted(ALLOWED_LINE_HEIGHTS_PT)}
    if check == "alignment_left_top":
        candidate = {}
        if "alignment" in evidence:
            candidate["alignment"] = "LEFT"
        if "vertical_anchor" in evidence:
            candidate["vertical_anchor"] = "TOP"
        return candidate or None
    if check == "text_color_allowlist" and evidence.get("color_hex"):
        return _allowed_text_color_candidate(evidence["color_hex"])
    if check == "background_color_palette" and evidence.get("color_hex"):
        return _allowed_fill_color_candidate(evidence["color_hex"])
    if check in {"contrast_ratio", "low_contrast"}:
        foreground = evidence.get("foreground_hex") or evidence.get("text_hex")
        background = evidence.get("background_hex")
        required = evidence.get("required_ratio")
        if foreground and background and required:
            candidate = _contrast_candidate(foreground, background, required)
            if candidate is not None:
                candidate["candidate_token"] = candidate.get("foreground_token")
            return candidate
    if check == "card_grid_consistency":
        medians = evidence.get("group_medians")
        inconsistent = evidence.get("inconsistent_containers")
        if isinstance(medians, dict) and isinstance(inconsistent, list) and inconsistent:
            return {
                "group_medians": medians,
                "container_count_to_align": len(inconsistent),
                "selection_policy": "row_group_median",
            }
    if check == "wrap_break_changes_meaning":
        widen = _widen_to_fit_candidate(evidence)
        if widen is not None:
            return widen
    if check == "badge_alignment":
        return {"alignment": "CENTER", "vertical_anchor": "MIDDLE"}
    if check == "box_canvas_overflow":
        bbox = evidence.get("bbox_pt") or []
        sides = evidence.get("overflow_sides_pt") or {}
        if len(bbox) == 4 and isinstance(sides, dict):
            x, y, w, h = bbox
            cand: dict = {"overflow_sides_pt": sides}
            if sides.get("right"):
                cand["target_width_pt"] = max(20.0, round(w - float(sides["right"]), 2))
            if sides.get("bottom"):
                cand["target_height_pt"] = max(20.0, round(h - float(sides["bottom"]), 2))
            return cand
        return None
    if check == "text_canvas_overflow":
        tr = evidence.get("text_render") or {}
        bbox = evidence.get("bbox_pt") or []
        if not isinstance(tr, dict) or len(bbox) != 4:
            return None
        x, y, w, h = bbox
        candidates: list[dict] = []
        # 設計原則 ([[feedback-overflow-fix-priority]]):
        # word_wrap=False (改行なし): 1.右に伸ばす 2.下に伸ばす 3.enable_word_wrap 4.font shrink
        # word_wrap=True (= 1単語が wrap 不可で canvas 超え): box 拡張しても
        #   text_start_x 起点の canvas-remaining は変わらないので font shrink 一択。
        max_w_to_canvas = max(0.0, SLIDE_W_PT - float(x))
        max_h_to_canvas = max(0.0, SLIDE_H_PT - float(y))
        is_no_wrap = tr.get("word_wrap") is False
        longest_unit_v = tr.get("longest_unit_pt")
        text_start_x = tr.get("text_start_x_pt")
        # 「box 拡張だけで text 全幅が canvas に収まるか」を予測 (= 効く strategy だけ
        # 提示する: 空振りする default は出さない)。
        expand_width_will_fit = (
            isinstance(longest_unit_v, (int, float))
            and isinstance(text_start_x, (int, float))
            and float(text_start_x) + float(longest_unit_v) <= SLIDE_W_PT + 0.5
        )
        if is_no_wrap:
            # Strategy 1: 右に伸ばす (= 拡張で実際に収まる場合のみ default 候補)
            if max_w_to_canvas > float(w) + 0.5 and expand_width_will_fit:
                candidates.append(
                    {
                        "strategy": "expand_box_width_to_canvas",
                        "target_width_pt": round(max_w_to_canvas, 2),
                        "from_pt": round(w, 2),
                        "reason": "expand box.width to canvas right edge (まず右に伸ばす)",
                    }
                )
            # Strategy 2: enable_word_wrap (text を折り返して text_box に流す)
            # → text 量が多くて 1 行で canvas に収まらない場合の現実的な default。
            candidates.append(
                {
                    "strategy": "enable_word_wrap",
                    "from_word_wrap": False,
                    "to_word_wrap": True,
                    "reason": "turn on word_wrap so the text reflows inside the box (1 行で canvas に収まらない場合)",
                }
            )
            # Strategy 3: 下に伸ばす (word_wrap が後で True になる前提なら box.height 余裕)
            if max_h_to_canvas > float(h) + 0.5:
                candidates.append(
                    {
                        "strategy": "expand_box_height",
                        "target_height_pt": round(max_h_to_canvas, 2),
                        "from_pt": round(h, 2),
                        "fits_within_canvas": True,
                        "reason": "expand box.height to canvas bottom",
                    }
                )
        # Strategy 4 (legacy): box.right > canvas のときは shrink_box_width_to_canvas
        if w > max_w_to_canvas + 0.5 and max_w_to_canvas > 20.0:
            candidates.append(
                {
                    "strategy": "shrink_box_width_to_canvas",
                    "target_width_pt": round(max_w_to_canvas, 2),
                    "from_pt": round(w, 2),
                    "reason": "clip box.width so its right edge sits at canvas right edge",
                }
            )
        # Strategy C: font_size 縮小 → 実際に収まる **最大** の allowed font。
        # 1 段階下では足りないケースがあるため (longest_unit × target/cur
        # ≤ canvas_remaining を満たす最大値を選ぶ)。
        cur_font = tr.get("font_size_pt")
        longest_unit = tr.get("longest_unit_pt")
        canvas_remaining = tr.get("canvas_remaining_pt")
        if (
            isinstance(cur_font, (int, float))
            and isinstance(longest_unit, (int, float))
            and isinstance(canvas_remaining, (int, float))
            and float(cur_font) > 0
        ):
            cur_f = float(cur_font)
            target_pt: Optional[float] = None
            for size in sorted([s for s in ALLOWED_FONT_SIZES_PT if s < cur_f], reverse=True):
                if float(longest_unit) * (size / cur_f) <= float(canvas_remaining) + 0.5:
                    target_pt = float(size)
                    break
            if target_pt is None:
                smaller = [s for s in sorted(ALLOWED_FONT_SIZES_PT) if s < cur_f]
                if smaller:
                    target_pt = float(smaller[0])
            if target_pt is not None:
                candidates.append(
                    {
                        "strategy": "shrink_font_size",
                        "font_size_pt": target_pt,
                        "from_pt": round(cur_f, 2),
                        "reason": "shrink to largest allowed font size that fits longest_unit within canvas_remaining",
                    }
                )
        return candidates or None
    if check == "text_box_overflow":
        tr = evidence.get("text_render") or {}
        bbox = evidence.get("bbox_pt") or []
        if not isinstance(tr, dict) or len(bbox) != 4:
            return None
        cur_font = tr.get("font_size_pt")
        cur_lh = tr.get("line_height_pt")
        lines = tr.get("lines") or 1
        required_h = float(tr.get("required_height_pt") or 0)
        candidates: list[dict] = []
        # 優先順位 ([[feedback-overflow-fix-priority]]): 見た目を変えない方向から:
        # 1) box.height 拡張 (下に伸ばす) ← default
        # 2) box.width 拡張 (右に伸ばす)
        # 3) line_height 圧縮 / font shrink (最終手段、見た目変わる)
        x, y, w, h = bbox
        # Strategy 1: box.height 拡張 → required_h + margin(top+bottom) + 1 行
        # buffer まで box を下に伸ばす (required_h は inner_h ベースなので
        # margin 足し戻し + PowerPoint 実描画でちょうど 1 行溢れる現象に対し
        # line_height 分の余裕を加算 [[feedback-overflow-fix-priority]])。
        margins_pt = tr.get("margins_pt") or {}
        margin_top = float(margins_pt.get("top", 7.2) or 7.2)
        margin_bot = float(margins_pt.get("bottom", 7.2) or 7.2)
        line_h_buffer = float(cur_lh) if isinstance(cur_lh, (int, float)) else (
            float(cur_font) * DEFAULT_LINE_HEIGHT_MULTIPLIER if isinstance(cur_font, (int, float)) else 28.8
        )
        max_h = max(0.0, SLIDE_H_PT - float(y))
        target_h_with_margin = round(required_h + margin_top + margin_bot + line_h_buffer, 2)
        target_h = min(target_h_with_margin, round(max_h, 2)) if max_h > 0 else target_h_with_margin
        if target_h > h + 0.5:
            candidates.append(
                {
                    "strategy": "expand_box_height",
                    "target_height_pt": target_h,
                    "from_pt": round(h, 2),
                    "fits_within_canvas": target_h <= max_h,
                    "reason": "expand box.height (下に伸ばす) — visual impact 最小",
                }
            )
        # Strategy 2: box.width を canvas 右端まで拡張する。
        max_w = max(0.0, SLIDE_W_PT - float(x))
        if max_w > w + 0.5:
            candidates.append(
                {
                    "strategy": "expand_box_width_to_canvas",
                    "target_width_pt": round(max_w, 2),
                    "from_pt": round(w, 2),
                    "reason": "expand box.width to canvas right edge so wrapped text needs fewer lines",
                }
            )
        # Strategy 3: line_height 圧縮 → 1 段階下の allowed line height。
        if isinstance(cur_lh, (int, float)):
            smaller_lh = [lh for lh in sorted(ALLOWED_LINE_HEIGHTS_PT) if lh < float(cur_lh)]
            if smaller_lh:
                candidates.append(
                    {
                        "strategy": "compress_line_height",
                        "line_height_pt": smaller_lh[-1],
                        "from_pt": round(float(cur_lh), 2),
                        "reason": "step down to nearest smaller allowed line height",
                    }
                )
        # Strategy 4: font_size shrink → 1 段階下の allowed font size に (最終手段)。
        if isinstance(cur_font, (int, float)):
            smaller = [s for s in sorted(ALLOWED_FONT_SIZES_PT) if s < float(cur_font)]
            if smaller:
                target_font = smaller[-1]
                candidates.append(
                    {
                        "strategy": "shrink_font_size",
                        "font_size_pt": target_font,
                        "from_pt": round(float(cur_font), 2),
                        "reason": "step down to nearest smaller allowed font size (見た目変わる最終手段)",
                    }
                )
        return candidates or None
    if check == "decorative_isolated_lines":
        candidates = evidence.get("candidate_values")
        if isinstance(candidates, list) and candidates:
            return candidates
    return None


def _measurement_confidence_for_json(evidence: dict) -> Optional[dict]:
    if evidence.get("measurement") != "rendered_image":
        return None
    sample_pixels = evidence.get("sample_pixels") or 0
    foreground_pixels = evidence.get("foreground_pixels") or 0
    foreground_ratio = foreground_pixels / sample_pixels if sample_pixels else 0
    background_complexity = evidence.get("background_complexity")
    if background_complexity == "complex":
        level = "low"
    elif foreground_pixels >= 50 and foreground_ratio >= 0.005:
        level = "medium"
    else:
        level = "low"
    return {
        "level": level,
        "method": evidence.get("foreground_detection", "raw_pixel_distance"),
        "background_model": evidence.get("background_model"),
        "background_uniformity": evidence.get("background_uniformity"),
        "background_complexity": background_complexity,
        "foreground_pixels": foreground_pixels,
        "background_pixels": evidence.get("background_pixels"),
        "sample_pixels": sample_pixels,
        "foreground_ratio": round(foreground_ratio, 4),
        "limitations": (
            "foreground is measured from raw rendered pixels; background is a dominant "
            "raw color model and complex backgrounds require local-contrast follow-up"
        ),
    }


def _evidence_for_json(check: str, detail: dict) -> dict:
    evidence = dict(detail or {})
    if "text_hex" in evidence and "foreground_hex" not in evidence:
        evidence["foreground_hex"] = evidence["text_hex"]
    if evidence.get("foreground_hex"):
        evidence["foreground_token"] = TEXT_COLOR_TOKEN_BY_HEX.get(evidence["foreground_hex"])
    if evidence.get("background_hex"):
        evidence["background_token"] = (
            FILL_COLOR_TOKEN_BY_HEX.get(evidence["background_hex"])
            or TEXT_COLOR_TOKEN_BY_HEX.get(evidence["background_hex"])
        )
    if check in {"contrast_ratio", "low_contrast"} and "original_run_color_hex" not in evidence:
        evidence["original_run_color_hex"] = evidence.get("foreground_hex")
    return evidence


def finding_to_json_dict(finding: Finding) -> dict:
    raw = asdict(finding)
    evidence = _evidence_for_json(finding.check, finding.detail)
    fixability = _fixability_for_json(finding.check, evidence)
    candidate_values = _candidate_values_for_json(finding.check, evidence)
    measurement_confidence = _measurement_confidence_for_json(evidence)
    schema_detail = {
        **evidence,
        "check_id": finding.check,
        "evidence": evidence,
        "fixability": fixability["fixability"],
        "fixability_rule": fixability.get("fixability_rule"),
        "fixability_reason": fixability.get("fixability_reason"),
        "candidate_values": candidate_values,
        "manual_required_reason": fixability.get("manual_required_reason"),
        "measurement_confidence": measurement_confidence,
    }
    raw["detail"] = schema_detail
    return raw


# ---- CLI -------------------------------------------------------------------


def main(argv: List[str]) -> int:
    ap = argparse.ArgumentParser(description="PPTX lint (v1 guideline)")
    ap.add_argument("pptx", type=Path, help="path to .pptx")
    ap.add_argument("--json", action="store_true", help="emit findings as JSON")
    ap.add_argument(
        "--structure-json",
        action="store_true",
        help="emit structural containment metadata as JSON instead of lint findings",
    )
    ap.add_argument(
        "--severity",
        choices=["all", "warning", "error"],
        default="all",
        help="minimum severity to report (default: all)",
    )
    ap.add_argument(
        "--no-consolidate",
        action="store_true",
        help="disable deck-level consolidation of recurring identical findings",
    )
    ap.add_argument(
        "--min-recurring-slides",
        type=int,
        default=3,
        help="threshold for consolidating recurring findings (default: 3)",
    )
    ap.add_argument(
        "--profile",
        choices=sorted(LINT_PROFILES),
        default="default",
        help="lint policy profile (default: default)",
    )
    ap.add_argument(
        "--rendered-image-dir",
        type=Path,
        help=(
            "directory containing PowerPoint-rendered slide PNGs "
            "(slide-01.png, slide-001.png, or 1.png) for rendered contrast lint"
        ),
    )
    args = ap.parse_args(argv)

    if not args.pptx.exists():
        print(f"file not found: {args.pptx}", file=sys.stderr)
        return 2

    if args.structure_json:
        payload = extract_pptx_structure(args.pptx, profile=args.profile)
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 0

    findings = lint_pptx(
        args.pptx,
        profile=args.profile,
        rendered_image_dir=args.rendered_image_dir,
    )
    if not args.no_consolidate:
        findings = consolidate_recurring(findings, min_slides=args.min_recurring_slides)
    selected = filter_by_severity(findings, args.severity)

    if args.json:
        payload: List[Any] = [finding_to_json_dict(f) for f in selected]
        print(json.dumps(payload, ensure_ascii=False, indent=2))
    else:
        print(format_text(selected), end="")

    return 1 if any(f.severity == "error" for f in findings) else 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
