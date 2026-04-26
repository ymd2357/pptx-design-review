#!/usr/bin/env python3
"""Generate good.pptx and bad.pptx fixtures for pptx_lint regression testing.

The .pptx files are excluded by .gitignore, so this script is the source of
truth: re-run it whenever you need fresh fixtures.

good.pptx
- 1 slide, 16:9 (1440x810pt), no animation
- title and body within safe text area, allowed font + size, autofit NONE
- expected lint output: 0 findings

bad.pptx
- 1 slide, 16:9
- shape A: overflow_text (right edge), font_family ("Arial"), font_size_scale (28pt),
           text_autofit_disabled (SHAPE_TO_FIT_TEXT), text_color_allowlist
- shape B: safe_text_area_text violation (positioned at x=10, left of x=81 boundary),
           background_color_palette
- shape C: safe_margins violation for non-text content
- shape D: line_height and alignment_left_top violations
- shape E: geometry_rounding violation
- picture F: image_upscale_ratio and alt_text_required violations
- slide: transition XML violation
- expected lint output: each of {overflow_text, safe_text_area_text, text_autofit_disabled,
                                 font_family, font_size_scale, safe_margins, line_height,
                                 alignment_left_top, geometry_rounding, image_upscale_ratio,
                                 alt_text_required, text_color_allowlist, background_color_palette,
                                 animation_present}
  fires at least once
"""

from __future__ import annotations

import argparse
import base64
import zipfile
from xml.etree import ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
ET.register_namespace("p", PML_NS)

TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAFgwJ/l9x1YQAAAABJRU5ErkJggg=="
)


def _new_169_deck():
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    return prs


def _add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _set_run(tf, text: str, font_name: str, size_pt: int):
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    return run


def _inject_slide_xml_child(path: Path, child_name: str) -> None:
    """Add a direct p:sld child such as p:transition to slide 1."""
    slide_part = "ppt/slides/slide1.xml"
    tmp = path.with_suffix(path.suffix + ".tmp")
    order = {
        f"{{{PML_NS}}}cSld": 0,
        f"{{{PML_NS}}}clrMapOvr": 1,
        f"{{{PML_NS}}}transition": 2,
        f"{{{PML_NS}}}timing": 3,
        f"{{{PML_NS}}}extLst": 4,
    }
    child_tag = f"{{{PML_NS}}}{child_name}"

    with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == slide_part:
                root = ET.fromstring(data)
                if root.find(f"p:{child_name}", {"p": PML_NS}) is None:
                    child_order = order[child_tag]
                    insert_at = len(root)
                    for idx, existing in enumerate(root):
                        if order.get(existing.tag, 99) > child_order:
                            insert_at = idx
                            break
                    root.insert(insert_at, ET.Element(child_tag))
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            zout.writestr(item, data)

    tmp.replace(path)


def _clear_alt_text(shape) -> None:
    c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
    c_nv_pr.set("descr", "")
    c_nv_pr.attrib.pop("title", None)


def make_good(out: Path) -> None:
    prs = _new_169_deck()
    slide = _add_blank_slide(prs)

    title = slide.shapes.add_textbox(Pt(81), Pt(40), Pt(1278), Pt(120))
    title.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    _set_run(title.text_frame, "Compliant title", "Noto Sans JP", 56)

    body = slide.shapes.add_textbox(Pt(81), Pt(200), Pt(1278), Pt(400))
    body.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    _set_run(body.text_frame, "Compliant body content.", "Noto Sans JP", 24)

    prs.save(str(out))


def make_bad(out: Path) -> None:
    prs = _new_169_deck()
    slide = _add_blank_slide(prs)

    a = slide.shapes.add_textbox(Pt(1300), Pt(40), Pt(200), Pt(120))
    a.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    a_run = _set_run(a.text_frame, "Bad shape A", "Arial", 28)
    a_run.font.color.rgb = RGBColor(255, 0, 0)

    b = slide.shapes.add_textbox(Pt(10), Pt(300), Pt(300), Pt(50))
    b.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    b.fill.solid()
    b.fill.fore_color.rgb = RGBColor(255, 204, 0)
    _set_run(b.text_frame, "Outside safe area", "Noto Sans JP", 24)

    c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(10), Pt(120), Pt(48), Pt(48))
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor(238, 238, 238)

    d = slide.shapes.add_textbox(Pt(200), Pt(420), Pt(300), Pt(80))
    d.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    d.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
    d_p = d.text_frame.paragraphs[0]
    d_p.alignment = PP_ALIGN.RIGHT
    d_p.line_spacing = Pt(25)
    d_run = d_p.add_run()
    d_run.text = "Bad line spacing and alignment"
    d_run.font.name = "Noto Sans JP"
    d_run.font.size = Pt(24)

    e = slide.shapes.add_textbox(Emu(round(81.5 * 12700)), Pt(540), Pt(180), Pt(60))
    e.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    _set_run(e.text_frame, "Half-point geometry", "Noto Sans JP", 24)

    tiny = out.with_name("_lint_tiny.png")
    tiny.write_bytes(TINY_PNG)
    try:
        f = slide.shapes.add_picture(str(tiny), Pt(600), Pt(420), width=Pt(96), height=Pt(96))
        _clear_alt_text(f)
    finally:
        tiny.unlink(missing_ok=True)

    prs.save(str(out))
    _inject_slide_xml_child(out, "transition")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument(
        "--outdir",
        type=Path,
        default=Path(__file__).parent.parent / "examples",
        help="output directory (default: skills/pptx-design-reviewer/examples)",
    )
    args = ap.parse_args()

    args.outdir.mkdir(parents=True, exist_ok=True)
    good = args.outdir / "good.pptx"
    bad = args.outdir / "bad.pptx"
    make_good(good)
    make_bad(bad)
    print(f"Wrote {good}")
    print(f"Wrote {bad}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
