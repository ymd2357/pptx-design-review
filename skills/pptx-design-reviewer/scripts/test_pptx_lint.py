#!/usr/bin/env python3
"""Smoke test for pptx_lint against the generated examples.

Asserts:
- good.pptx produces zero findings (errors and warnings)
- bad.pptx triggers each expected implemented check
- every check emitted by pptx_lint.py is defined in rules.lint.checks

Exit code: 0 on success, 1 on any failed assertion.
"""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))

import make_examples  # noqa: E402
import pptx_lint  # noqa: E402
from PIL import Image, ImageDraw  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_AUTO_SIZE  # noqa: E402
from pptx.enum.text import MSO_VERTICAL_ANCHOR  # noqa: E402
from pptx.util import Pt  # noqa: E402


A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


EXPECTED_BAD_CHECKS = {
    "overflow_text",
    "safe_text_area_text",
    "text_autofit_disabled",
    "font_family",
    "font_size_scale",
    "safe_margins",
    "line_height",
    "alignment_left_top",
    "geometry_rounding",
    "image_upscale_ratio",
    "alt_text_required",
    "text_color_allowlist",
    "background_color_palette",
    "contrast_ratio",
    "low_contrast",
    "animation_present",
}

LINT005_CHECKS = {
    "text_overlap",
    "object_overlap",
    "object_gap_too_small",
    "inner_padding_imbalance",
}

KNOWN_EMITTED_CHECKS = EXPECTED_BAD_CHECKS | {
    "card_grid_consistency",
    "slide_size",
    "overflow_shapes",
    "overflow_images",
    "image_aspect_distortion",
    "text_vertical_balance",
    "color_only_meaning",
    "heading_hierarchy_broken",
    "key_area_cropped",
    "missing_required_element",
    "reading_order",
    "wrap_break_changes_meaning",
} | LINT005_CHECKS

LINT004_POLICY = {
    "image_upscale_ratio": "automated",
    "contrast_ratio": "automated",
    "color_only_meaning": "automated",
    "alt_text_required": "automated",
    "reading_order": "automated",
}


def _guideline_lint_check_keys() -> set[str]:
    guideline = HERE.parents[2] / "doc" / "slide-guideline-v1.yml"
    lines = guideline.read_text(encoding="utf-8").splitlines()
    for idx, line in enumerate(lines):
        if line == "  lint:":
            keys: set[str] = set()
            in_checks = False
            for child in lines[idx + 1:]:
                if child and not child.startswith("    "):
                    break
                if child == "    checks:":
                    in_checks = True
                    continue
                if in_checks:
                    if child.startswith("      ") and child.endswith(":"):
                        keys.add(child.strip()[:-1])
                    elif child.startswith("    ") and child.strip() and not child.startswith("      "):
                        break
            return keys
    raise AssertionError("rules.lint.checks was not found in doc/slide-guideline-v1.yml")


def _guideline_lint_check_automation() -> dict[str, str]:
    guideline = HERE.parents[2] / "doc" / "slide-guideline-v1.yml"
    lines = guideline.read_text(encoding="utf-8").splitlines()
    automation: dict[str, str] = {}
    current_key: str | None = None
    in_checks = False
    for line in lines:
        if line == "    checks:":
            in_checks = True
            continue
        if in_checks and line.startswith("      ") and line.endswith(":"):
            current_key = line.strip()[:-1]
            continue
        if in_checks and current_key and line.startswith("        automation:"):
            automation[current_key] = line.split(":", 1)[1].strip().strip('"')
            continue
        if in_checks and line and not line.startswith("      ") and not line.startswith("        "):
            break
    return automation


def _make_scaled_good(out: Path) -> None:
    """Create a 720x405 deck that is compliant after 1440x810 normalization."""
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Pt(40.5), Pt(20), Pt(639), Pt(60))
    title.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Scaled compliant title"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(28)

    body = slide.shapes.add_textbox(Pt(40.5), Pt(100), Pt(639), Pt(200))
    body.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = body.text_frame.paragraphs[0].add_run()
    run.text = "Scaled compliant body."
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(12)

    prs.save(str(out))


def _make_centered_single_line_good(out: Path) -> None:
    """Create a deck with an intentionally vertically centered one-line label."""
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    label = slide.shapes.add_textbox(Pt(40.5), Pt(20), Pt(300), Pt(30))
    label.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    label.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    run = label.text_frame.paragraphs[0].add_run()
    run.text = "Centered single-line label"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(12)
    prs.save(str(out))


def _make_bad_table_cell_fill(out: Path) -> None:
    """Create a deck where only a table cell uses a non-palette fill color."""
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 2, Pt(40.5), Pt(40), Pt(400), Pt(120))
    table = table_shape.table
    for row in table.rows:
        for cell in row.cells:
            cell.text = "Cell"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    table.cell(0, 0).fill.fore_color.rgb = RGBColor.from_string("123456")
    prs.save(str(out))


def _make_rendered_low_contrast_case(out: Path, image_dir: Path) -> None:
    """Create a deck whose rendered PNG has low contrast not encoded in PPTX colors."""
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40.5), Pt(40), Pt(300), Pt(80))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Rendered low contrast"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(12)
    prs.save(str(out))

    image_dir.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (720, 405), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((70, 70, 300, 92), fill="#BBBBBB")
    image.save(image_dir / "slide-01.png")


def _make_allowed_latin_fonts_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Pt(40.5), Pt(20), Pt(639), Pt(60))
    title.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "SHIFT AI"
    run.font.name = "Avenir Next Arabic"
    run.font.size = Pt(28)

    body = slide.shapes.add_textbox(Pt(40.5), Pt(100), Pt(639), Pt(80))
    body.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = body.text_frame.paragraphs[0].add_run()
    run.text = "Fallback Latin"
    run.font.name = "Nunito Sans"
    run.font.size = Pt(12)
    prs.save(str(out))


def _make_scaled_near_font_sizes_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for i, size in enumerate((11.25, 13.5, 15.75, 31.5)):
        box = slide.shapes.add_textbox(Pt(40.5), Pt(20 + i * 60), Pt(500), Pt(50))
        box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"Near size {size}"
        run.font.name = "Noto Sans JP"
        run.font.size = Pt(size)
    prs.save(str(out))


def _make_scaled_near_line_heights_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for i, line_height in enumerate((20.25, 33.75, 44.1)):
        box = slide.shapes.add_textbox(Pt(40.5), Pt(20 + i * 70), Pt(500), Pt(60))
        box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        para = box.text_frame.paragraphs[0]
        para.line_spacing = Pt(line_height)
        run = para.add_run()
        run.text = f"Near line height {line_height}"
        run.font.name = "Noto Sans JP"
        run.font.size = Pt(12)
    prs.save(str(out))


def _make_bad_table_cell_font(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(1, 1, Pt(40.5), Pt(40), Pt(300), Pt(60))
    cell = table_shape.table.cell(0, 0)
    cell.text = ""
    run = cell.text_frame.paragraphs[0].add_run()
    run.text = "Bad table font"
    run.font.name = "Arial"
    run.font.size = Pt(12)
    prs.save(str(out))


def _make_bad_east_asian_font(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40.5), Pt(40), Pt(500), Pt(60))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "日本語フォント"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(12)
    r_pr = run._r.get_or_add_rPr()
    ea = r_pr.find(f"{A_NS}ea")
    if ea is None:
        ea = r_pr.makeelement(f"{A_NS}ea")
        r_pr.append(ea)
    ea.set("typeface", "Meiryo")
    prs.save(str(out))


def _make_aspect_distorted_image(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = out.with_name("_aspect_fixture.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        pic = slide.shapes.add_picture(str(img), Pt(40), Pt(40), width=Pt(120), height=Pt(40))
        make_examples._clear_alt_text(pic)
    finally:
        img.unlink(missing_ok=True)
    prs.save(str(out))


def _make_decorative_distorted_image(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = out.with_name("_decorative_aspect_fixture.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        pic = slide.shapes.add_picture(str(img), Pt(0), Pt(0), width=Pt(720), height=Pt(405))
        c_nv_pr = pic._element.xpath(".//p:cNvPr")[0]
        c_nv_pr.set("descr", "header-gradient.png")
    finally:
        img.unlink(missing_ok=True)
    prs.save(str(out))


def _make_decorative_overflow_image(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = out.with_name("_decorative_overflow_fixture.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        pic = slide.shapes.add_picture(str(img), Pt(500), Pt(-120), width=Pt(320), height=Pt(280))
        c_nv_pr = pic._element.xpath(".//p:cNvPr")[0]
        c_nv_pr.set("descr", "decoration_top_right.png")
    finally:
        img.unlink(missing_ok=True)
    prs.save(str(out))


def _make_content_overflow_image(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = out.with_name("_content_overflow_fixture.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        pic = slide.shapes.add_picture(str(img), Pt(500), Pt(-120), width=Pt(320), height=Pt(280))
        make_examples._clear_alt_text(pic)
    finally:
        img.unlink(missing_ok=True)
    prs.save(str(out))


def _make_safe_text_area_template_exempt_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_semantic_text(slide, x=45, y=13.5, width=360, text="Template header title", size=32)
    _add_semantic_text(slide, x=90, y=153, width=1323, text="Template full width subtitle", size=24)
    _add_semantic_text(slide, x=1400, y=762, width=24, text="12", size=20)
    prs.save(str(out))


def _make_safe_text_area_body_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_semantic_text(slide, x=20, y=240, width=300, text="Body outside safe text area", size=24)
    prs.save(str(out))


def _add_lint005_text_box(slide, *, x: int, y: int, width: int, text: str):
    box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(width), Pt(42))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    box.text_frame.margin_top = Pt(6)
    box.text_frame.margin_bottom = Pt(6)
    para = box.text_frame.paragraphs[0]
    para.line_spacing = Pt(30)
    run = para.add_run()
    run.text = text
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    return box


def _add_lint005_rect(slide, *, x: int, y: int, width: int, height: int):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(width), Pt(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string("EEEEEE")
    return shape


def _make_object_relationships_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_lint005_text_box(slide, x=120, y=100, width=300, text="Aligned text one")
    _add_lint005_text_box(slide, x=120, y=180, width=300, text="Aligned text two")
    _add_lint005_rect(slide, x=520, y=100, width=120, height=80)
    _add_lint005_rect(slide, x=680, y=100, width=120, height=80)
    prs.save(str(out))


def _make_structural_containment_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_lint005_rect(slide, x=500, y=300, width=360, height=180)
    _add_lint005_rect(slide, x=530, y=330, width=60, height=40)
    _add_lint005_rect(slide, x=770, y=410, width=60, height=40)
    prs.save(str(out))


def _make_structural_containment_overflow(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_lint005_rect(slide, x=500, y=300, width=360, height=180)
    _add_lint005_rect(slide, x=530, y=330, width=340, height=40)
    prs.save(str(out))


def _make_object_relationships_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_lint005_text_box(slide, x=120, y=100, width=260, text="Text overlap A")
    _add_lint005_text_box(slide, x=140, y=120, width=260, text="Text overlap B")

    _add_lint005_rect(slide, x=500, y=100, width=120, height=80)
    _add_lint005_rect(slide, x=560, y=130, width=120, height=80)

    _add_lint005_rect(slide, x=900, y=100, width=80, height=80)
    _add_lint005_rect(slide, x=984, y=100, width=80, height=80)

    _add_lint005_rect(slide, x=120, y=300, width=100, height=60)
    _add_lint005_rect(slide, x=125, y=390, width=100, height=60)
    _add_lint005_text_box(slide, x=300, y=520, width=260, text="Peer text one")
    _add_lint005_text_box(slide, x=305, y=610, width=260, text="Peer text two")

    _add_lint005_rect(slide, x=500, y=300, width=360, height=180)
    _add_lint005_rect(slide, x=502, y=330, width=60, height=40)
    _add_lint005_rect(slide, x=580, y=360, width=60, height=40)
    prs.save(str(out))


def _add_card_with_text(slide, *, x: int, y: int, width: int, height: int, child_y: int = 26):
    card = _add_lint005_rect(slide, x=x, y=y, width=width, height=height)
    title = slide.shapes.add_textbox(Pt(x + 24), Pt(y + child_y), Pt(width - 48), Pt(42))
    title.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    title_run = title.text_frame.paragraphs[0].add_run()
    title_run.text = "Card title"
    title_run.font.name = "Noto Sans JP"
    title_run.font.size = Pt(24)
    body = slide.shapes.add_textbox(Pt(x + 24), Pt(y + child_y + 54), Pt(width - 48), Pt(56))
    body.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    body_run = body.text_frame.paragraphs[0].add_run()
    body_run.text = "Card body text"
    body_run.font.name = "Noto Sans JP"
    body_run.font.size = Pt(20)
    return card


def _make_card_grid_consistency_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for x in (120, 520, 920):
        _add_card_with_text(slide, x=x, y=220, width=300, height=180)
    prs.save(str(out))


def _make_card_grid_consistency_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_card_with_text(slide, x=120, y=220, width=300, height=180)
    _add_card_with_text(slide, x=520, y=220, width=300, height=180)
    _add_card_with_text(slide, x=920, y=220, width=330, height=180, child_y=52)
    prs.save(str(out))


def _add_lint006_text_box(slide, *, y: int, height: int, anchor, margin_top: int, margin_bottom: int):
    box = slide.shapes.add_textbox(Pt(120), Pt(y), Pt(520), Pt(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string("F7F7F7")
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    box.text_frame.margin_top = Pt(margin_top)
    box.text_frame.margin_bottom = Pt(margin_bottom)
    if anchor is not None:
        box.text_frame.vertical_anchor = anchor
    para = box.text_frame.paragraphs[0]
    para.line_spacing = Pt(30)
    run = para.add_run()
    run.text = "Text vertical balance sample"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    return box


def _make_invisible_text_box_vertical_balance_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(120), Pt(120), Pt(520), Pt(220))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    box.text_frame.margin_top = Pt(0)
    box.text_frame.margin_bottom = Pt(0)
    box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    para = box.text_frame.paragraphs[0]
    para.line_spacing = Pt(30)
    run = para.add_run()
    run.text = "Invisible text box with extra selection height"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    prs.save(str(out))


def _make_text_vertical_balance_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_lint006_text_box(
        slide,
        y=120,
        height=42,
        anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        margin_top=6,
        margin_bottom=6,
    )
    prs.save(str(out))


def _make_top_anchor_bottom_void_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_lint006_text_box(
        slide,
        y=120,
        height=250,
        anchor=MSO_VERTICAL_ANCHOR.TOP,
        margin_top=6,
        margin_bottom=6,
    )
    prs.save(str(out))


def _make_middle_anchor_asymmetric_margin_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_lint006_text_box(
        slide,
        y=120,
        height=120,
        anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        margin_top=30,
        margin_bottom=6,
    )
    prs.save(str(out))


def _make_oversized_box_top_anchor_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_lint006_text_box(
        slide,
        y=120,
        height=220,
        anchor=None,
        margin_top=6,
        margin_bottom=6,
    )
    prs.save(str(out))


def _add_semantic_text(slide, *, x: int, y: int, width: int, text: str, size: int):
    box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(width), Pt(60))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(size)
    return box


def _make_missing_title_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_semantic_text(slide, x=120, y=260, width=900, text="Body without a title", size=24)
    prs.save(str(out))


def _make_cover_slide_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = out.with_name("_cover_logo_fixture.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        slide.shapes.add_picture(str(img), Pt(26), Pt(25), width=Pt(158), height=Pt(44))
    finally:
        img.unlink(missing_ok=True)
    _add_semantic_text(slide, x=120, y=180, width=720, text="Customer Name", size=32)
    _add_semantic_text(slide, x=120, y=320, width=1000, text="Prominent proposal title", size=48)
    _add_semantic_text(slide, x=120, y=650, width=520, text="2026 | SHIFT AI", size=20)
    prs.save(str(out))


def _make_section_divider_good(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_semantic_text(slide, x=260, y=330, width=920, text="Section divider title", size=56)
    prs.save(str(out))


def _make_heading_hierarchy_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_semantic_text(slide, x=120, y=50, width=900, text="Small title", size=32)
    _add_semantic_text(slide, x=120, y=220, width=900, text="Oversized body", size=40)
    prs.save(str(out))


def _make_reading_order_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_semantic_text(slide, x=120, y=340, width=900, text="Third visually", size=24)
    _add_semantic_text(slide, x=120, y=200, width=900, text="Second visually", size=24)
    _add_semantic_text(slide, x=120, y=50, width=900, text="First visually", size=32)
    prs.save(str(out))


def _make_wrap_break_bad(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_semantic_text(slide, x=120, y=50, width=900, text="Title", size=32)
    box = _add_semantic_text(slide, x=120, y=180, width=900, text="", size=24)
    run = box.text_frame.paragraphs[0].runs[0]
    run.text = "Auto\nmation improves review speed"
    prs.save(str(out))


def _make_key_area_cropped_bad(out: Path, *, decorative: bool = False) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_semantic_text(slide, x=120, y=50, width=900, text="Title", size=32)
    img = out.with_name("_cropped_fixture.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        pic = slide.shapes.add_picture(str(img), Pt(120), Pt(180), width=Pt(320), height=Pt(180))
        pic.crop_left = 0.25
        if decorative:
            c_nv_pr = pic._element.xpath(".//p:cNvPr")[0]
            c_nv_pr.set("descr", "decorative-background.png")
    finally:
        img.unlink(missing_ok=True)
    prs.save(str(out))


def _make_color_only_bad(out: Path, *, labeled: bool = False) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_semantic_text(slide, x=120, y=50, width=900, text="Title", size=32)
    for idx, color in enumerate(("FF5757", "039578")):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Pt(120 + idx * 118),
            Pt(180),
            Pt(100),
            Pt(48),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color)
        if labeled:
            shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            run = shape.text_frame.paragraphs[0].add_run()
            run.text = "NG" if idx == 0 else "OK"
            run.font.name = "Noto Sans JP"
            run.font.size = Pt(20)
    prs.save(str(out))


def main() -> int:
    failures: list = []

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        good = tmp_dir / "good.pptx"
        bad = tmp_dir / "bad.pptx"
        scaled_good = tmp_dir / "scaled-good.pptx"
        centered_single_line_good = tmp_dir / "centered-single-line-good.pptx"
        bad_table_cell_fill = tmp_dir / "bad-table-cell-fill.pptx"
        rendered_low_contrast = tmp_dir / "rendered-low-contrast.pptx"
        rendered_low_contrast_images = tmp_dir / "rendered-low-contrast-images"
        allowed_latin_fonts_good = tmp_dir / "allowed-latin-fonts-good.pptx"
        scaled_near_font_sizes_good = tmp_dir / "scaled-near-font-sizes-good.pptx"
        bad_table_cell_font = tmp_dir / "bad-table-cell-font.pptx"
        scaled_near_line_heights_good = tmp_dir / "scaled-near-line-heights-good.pptx"
        bad_east_asian_font = tmp_dir / "bad-east-asian-font.pptx"
        aspect_distorted_image = tmp_dir / "aspect-distorted-image.pptx"
        decorative_distorted_image = tmp_dir / "decorative-distorted-image.pptx"
        decorative_overflow_image = tmp_dir / "decorative-overflow-image.pptx"
        content_overflow_image = tmp_dir / "content-overflow-image.pptx"
        safe_text_area_template_exempt_good = tmp_dir / "safe-text-area-template-exempt-good.pptx"
        safe_text_area_body_bad = tmp_dir / "safe-text-area-body-bad.pptx"
        object_relationships_good = tmp_dir / "object-relationships-good.pptx"
        structural_containment_good = tmp_dir / "structural-containment-good.pptx"
        structural_containment_overflow = tmp_dir / "structural-containment-overflow.pptx"
        object_relationships_bad = tmp_dir / "object-relationships-bad.pptx"
        card_grid_consistency_good = tmp_dir / "card-grid-consistency-good.pptx"
        card_grid_consistency_bad = tmp_dir / "card-grid-consistency-bad.pptx"
        text_vertical_balance_good = tmp_dir / "text-vertical-balance-good.pptx"
        invisible_text_box_vertical_balance_good = (
            tmp_dir / "invisible-text-box-vertical-balance-good.pptx"
        )
        top_anchor_bottom_void_bad = tmp_dir / "top-anchor-bottom-void-bad.pptx"
        middle_anchor_asymmetric_margin_bad = tmp_dir / "middle-anchor-asymmetric-margin-bad.pptx"
        oversized_box_top_anchor_bad = tmp_dir / "oversized-box-top-anchor-bad.pptx"
        missing_title_bad = tmp_dir / "missing-title-bad.pptx"
        cover_slide_good = tmp_dir / "cover-slide-good.pptx"
        section_divider_good = tmp_dir / "section-divider-good.pptx"
        heading_hierarchy_bad = tmp_dir / "heading-hierarchy-bad.pptx"
        reading_order_bad = tmp_dir / "reading-order-bad.pptx"
        wrap_break_bad = tmp_dir / "wrap-break-bad.pptx"
        key_area_cropped_bad = tmp_dir / "key-area-cropped-bad.pptx"
        key_area_cropped_decorative = tmp_dir / "key-area-cropped-decorative.pptx"
        color_only_bad = tmp_dir / "color-only-bad.pptx"
        color_only_labeled = tmp_dir / "color-only-labeled.pptx"
        make_examples.make_good(good)
        make_examples.make_bad(bad)
        _make_scaled_good(scaled_good)
        _make_centered_single_line_good(centered_single_line_good)
        _make_bad_table_cell_fill(bad_table_cell_fill)
        _make_rendered_low_contrast_case(rendered_low_contrast, rendered_low_contrast_images)
        _make_allowed_latin_fonts_good(allowed_latin_fonts_good)
        _make_scaled_near_font_sizes_good(scaled_near_font_sizes_good)
        _make_bad_table_cell_font(bad_table_cell_font)
        _make_scaled_near_line_heights_good(scaled_near_line_heights_good)
        _make_bad_east_asian_font(bad_east_asian_font)
        _make_aspect_distorted_image(aspect_distorted_image)
        _make_decorative_distorted_image(decorative_distorted_image)
        _make_decorative_overflow_image(decorative_overflow_image)
        _make_content_overflow_image(content_overflow_image)
        _make_safe_text_area_template_exempt_good(safe_text_area_template_exempt_good)
        _make_safe_text_area_body_bad(safe_text_area_body_bad)
        _make_object_relationships_good(object_relationships_good)
        _make_structural_containment_good(structural_containment_good)
        _make_structural_containment_overflow(structural_containment_overflow)
        _make_object_relationships_bad(object_relationships_bad)
        _make_card_grid_consistency_good(card_grid_consistency_good)
        _make_card_grid_consistency_bad(card_grid_consistency_bad)
        _make_text_vertical_balance_good(text_vertical_balance_good)
        _make_invisible_text_box_vertical_balance_good(invisible_text_box_vertical_balance_good)
        _make_top_anchor_bottom_void_bad(top_anchor_bottom_void_bad)
        _make_middle_anchor_asymmetric_margin_bad(middle_anchor_asymmetric_margin_bad)
        _make_oversized_box_top_anchor_bad(oversized_box_top_anchor_bad)
        _make_missing_title_bad(missing_title_bad)
        _make_cover_slide_good(cover_slide_good)
        _make_section_divider_good(section_divider_good)
        _make_heading_hierarchy_bad(heading_hierarchy_bad)
        _make_reading_order_bad(reading_order_bad)
        _make_wrap_break_bad(wrap_break_bad)
        _make_key_area_cropped_bad(key_area_cropped_bad)
        _make_key_area_cropped_bad(key_area_cropped_decorative, decorative=True)
        _make_color_only_bad(color_only_bad)
        _make_color_only_bad(color_only_labeled, labeled=True)

        good_findings = pptx_lint.lint_pptx(good)
        if good_findings:
            failures.append(
                "good.pptx had {n} findings; expected 0:\n  {lines}".format(
                    n=len(good_findings),
                    lines="\n  ".join(f"[{f.severity}] {f.check}: {f.message}" for f in good_findings),
                )
            )

        scaled_good_findings = pptx_lint.lint_pptx(scaled_good)
        if scaled_good_findings:
            failures.append(
                "scaled-good.pptx had {n} findings; expected 0 after normalization:\n  {lines}".format(
                    n=len(scaled_good_findings),
                    lines="\n  ".join(
                        f"[{f.severity}] {f.check}: {f.message}" for f in scaled_good_findings
                    ),
                )
            )

        centered_single_line_findings = pptx_lint.lint_pptx(centered_single_line_good)
        if centered_single_line_findings:
            failures.append(
                "centered-single-line-good.pptx had {n} findings; expected 0:\n  {lines}".format(
                    n=len(centered_single_line_findings),
                    lines="\n  ".join(
                        f"[{f.severity}] {f.check}: {f.message}"
                        for f in centered_single_line_findings
                    ),
                )
            )
        centered_single_line_strict_findings = pptx_lint.lint_pptx(
            centered_single_line_good,
            profile="strict",
        )
        if not any(f.check == "alignment_left_top" for f in centered_single_line_strict_findings):
            failures.append(
                "centered-single-line-good.pptx did not trigger vertical anchor alignment in strict profile"
            )

        table_cell_findings = [
            f
            for f in pptx_lint.lint_pptx(bad_table_cell_fill)
            if f.check == "background_color_palette"
        ]
        if not any(f.detail.get("scope") == "table_cell" for f in table_cell_findings):
            failures.append(
                "bad-table-cell-fill.pptx did not trigger background_color_palette for table cell fill"
            )

        rendered_without_images_findings = [
            f
            for f in pptx_lint.lint_pptx(rendered_low_contrast)
            if f.check in {"low_contrast", "contrast_ratio"}
        ]
        if rendered_without_images_findings:
            failures.append(
                "rendered-low-contrast.pptx triggered contrast without rendered images:\n  "
                + "\n  ".join(f"{f.check}: {f.message}" for f in rendered_without_images_findings)
            )
        rendered_with_images_findings = [
            f
            for f in pptx_lint.lint_pptx(
                rendered_low_contrast,
                rendered_image_dir=rendered_low_contrast_images,
            )
            if f.check == "low_contrast" and f.detail.get("measurement") == "rendered_image"
        ]
        if not rendered_with_images_findings:
            failures.append(
                "rendered-low-contrast.pptx did not trigger rendered-image low_contrast"
            )
        thin_text_crop = Image.new("RGB", (1764, 54), "white")
        draw = ImageDraw.Draw(thin_text_crop)
        for x in range(0, 1764, 14):
            draw.rectangle((x, 18, x + 2, 36), fill="#707070")
            draw.rectangle((x + 3, 18, x + 4, 36), fill="#A8A8A8")
        thin_text_measurement = pptx_lint._dominant_rendered_contrast(
            thin_text_crop,
            expected_foreground_hexes=["#707070"],
        )
        if not thin_text_measurement or thin_text_measurement.get("text_hex") != "#707070":
            failures.append(
                "rendered contrast foreground detection should prefer observed run color "
                f"over antialias pixels; got {thin_text_measurement}"
            )
        elif thin_text_measurement.get("contrast_ratio", 0) < 4.5:
            failures.append(
                "rendered contrast foreground detection used an antialias color instead of "
                f"#707070: {thin_text_measurement}"
            )
        elif thin_text_measurement.get("background_complexity") != "uniform":
            failures.append(
                "rendered contrast should mark plain white crop background as uniform: "
                f"{thin_text_measurement}"
            )

        complex_background_crop = Image.new("RGB", (240, 80), "white")
        complex_draw = ImageDraw.Draw(complex_background_crop)
        for x in range(240):
            shade = 40 + int(190 * x / 239)
            complex_draw.line((x, 0, x, 79), fill=(shade, shade, shade))
        complex_draw.rectangle((80, 30, 160, 42), fill="#707070")
        complex_measurement = pptx_lint._dominant_rendered_contrast(
            complex_background_crop,
            expected_foreground_hexes=["#707070"],
        )
        if not complex_measurement or complex_measurement.get("background_complexity") != "complex":
            failures.append(
                "rendered contrast should expose complex background evidence for gradients: "
                f"{complex_measurement}"
            )
        recurring_rendered_findings = [
            pptx_lint.Finding(
                severity="error",
                check="low_contrast",
                slide_index=idx,
                slide_id=idx,
                shape_id=idx + 100,
                shape_name=f"Text {idx}",
                message="rendered text/background contrast ratio 2.50:1 is below unreadable threshold 3.0:1 for normal_text",
                detail={
                    "measurement": "rendered_image",
                    "rendered_image_path": f"slide-{idx:02d}.png",
                    "foreground_hex": "#A3A3A3",
                    "background_hex": "#FFFFFF",
                    "contrast_ratio": 2.5,
                    "required_ratio": 4.5,
                    "low_contrast_threshold": 3.0,
                    "bbox_pt": [10, 20, 100, 30],
                },
            )
            for idx in (1, 2, 3)
        ]
        consolidated_rendered = pptx_lint.consolidate_recurring(recurring_rendered_findings)
        if not consolidated_rendered or consolidated_rendered[0].detail.get("measurement") != "rendered_image":
            failures.append(
                "recurring rendered-image contrast consolidation dropped measurement evidence"
            )
        if not consolidated_rendered or "contrast_ratio" not in consolidated_rendered[0].detail:
            failures.append(
                "recurring rendered-image contrast consolidation dropped contrast evidence"
            )

        allowed_latin_font_findings = [
            f
            for f in pptx_lint.lint_pptx(allowed_latin_fonts_good)
            if f.check == "font_family"
        ]
        if allowed_latin_font_findings:
            failures.append(
                "allowed-latin-fonts-good.pptx triggered font_family for design guideline fonts:\n  "
                + "\n  ".join(f.message for f in allowed_latin_font_findings)
            )

        scaled_near_font_size_findings = [
            f
            for f in pptx_lint.lint_pptx(scaled_near_font_sizes_good)
            if f.check == "font_size_scale"
        ]
        if scaled_near_font_size_findings:
            failures.append(
                "scaled-near-font-sizes-good.pptx triggered font_size_scale inside tolerance:\n  "
                + "\n  ".join(f.message for f in scaled_near_font_size_findings)
            )

        scaled_near_line_height_findings = [
            f
            for f in pptx_lint.lint_pptx(scaled_near_line_heights_good)
            if f.check == "line_height"
        ]
        if scaled_near_line_height_findings:
            failures.append(
                "scaled-near-line-heights-good.pptx triggered line_height inside tolerance:\n  "
                + "\n  ".join(f.message for f in scaled_near_line_height_findings)
            )

        table_cell_font_findings = [
            f
            for f in pptx_lint.lint_pptx(bad_table_cell_font)
            if f.check == "font_family"
        ]
        if not any(f.detail.get("cells") for f in table_cell_font_findings):
            failures.append(
                "bad-table-cell-font.pptx did not trigger font_family for table cell text"
            )

        east_asian_font_findings = [
            f
            for f in pptx_lint.lint_pptx(bad_east_asian_font)
            if f.check == "font_family"
        ]
        if not any(f.detail.get("script") == "ea" and f.detail.get("font") == "Meiryo" for f in east_asian_font_findings):
            failures.append(
                "bad-east-asian-font.pptx did not trigger font_family for explicit East Asian font"
            )

        aspect_findings = [
            f
            for f in pptx_lint.lint_pptx(aspect_distorted_image)
            if f.check == "image_aspect_distortion"
        ]
        if not aspect_findings:
            failures.append("aspect-distorted-image.pptx did not trigger image_aspect_distortion")

        decorative_image_findings = [
            f
            for f in pptx_lint.lint_pptx(decorative_distorted_image)
            if f.check in {"image_aspect_distortion", "image_upscale_ratio"}
        ]
        if decorative_image_findings:
            failures.append(
                "decorative-distorted-image.pptx triggered raster quality checks for a template raster:\n  "
                + "\n  ".join(f"{f.check}: {f.message}" for f in decorative_image_findings)
            )

        decorative_overflow_findings = [
            f
            for f in pptx_lint.lint_pptx(decorative_overflow_image)
            if f.check in {"overflow_images", "safe_margins"}
        ]
        if decorative_overflow_findings:
            failures.append(
                "decorative-overflow-image.pptx triggered overflow/safe margin checks:\n  "
                + "\n  ".join(f"{f.check}: {f.message}" for f in decorative_overflow_findings)
            )

        content_overflow_findings = [
            f
            for f in pptx_lint.lint_pptx(content_overflow_image)
            if f.check in {"overflow_images", "safe_margins"}
        ]
        content_overflow_checks = {f.check for f in content_overflow_findings}
        if {"overflow_images", "safe_margins"} - content_overflow_checks:
            failures.append(
                "content-overflow-image.pptx did not trigger overflow/safe margin checks; got "
                f"{sorted(content_overflow_checks)}"
            )

        safe_text_area_template_findings = [
            f
            for f in pptx_lint.lint_pptx(safe_text_area_template_exempt_good)
            if f.check == "safe_text_area_text"
        ]
        if safe_text_area_template_findings:
            failures.append(
                "safe-text-area-template-exempt-good.pptx triggered safe_text_area_text:\n  "
                + "\n  ".join(f.message for f in safe_text_area_template_findings)
            )
        safe_text_area_body_findings = [
            f
            for f in pptx_lint.lint_pptx(safe_text_area_body_bad)
            if f.check == "safe_text_area_text"
        ]
        if not safe_text_area_body_findings:
            failures.append("safe-text-area-body-bad.pptx did not trigger safe_text_area_text")

        lint005_good_findings = [
            f for f in pptx_lint.lint_pptx(object_relationships_good) if f.check in LINT005_CHECKS
        ]
        if lint005_good_findings:
            failures.append(
                "object-relationships-good.pptx triggered LINT-005 checks:\n  "
                + "\n  ".join(f"{f.check}: {f.message}" for f in lint005_good_findings)
            )
        structural_containment_findings = [
            f
            for f in pptx_lint.lint_pptx(structural_containment_good)
            if f.check in {"object_overlap", "inner_padding_imbalance"}
        ]
        if structural_containment_findings:
            failures.append(
                "structural-containment-good.pptx treated containment as a lint issue:\n  "
                + "\n  ".join(
                    f"{f.check}: {f.message}" for f in structural_containment_findings
                )
            )
        structure = pptx_lint.extract_pptx_structure(structural_containment_good)
        if not any(
            relation["relation"] == "contains"
            and relation["container"]["shape_name"]
            and relation["child"]["shape_name"]
            for relation in structure
        ):
            failures.append("object-relationships-good.pptx did not expose containment structure")
        structural_overflow_object_overlap = [
            f
            for f in pptx_lint.lint_pptx(structural_containment_overflow)
            if f.check == "object_overlap"
        ]
        if structural_overflow_object_overlap:
            failures.append(
                "structural-containment-overflow.pptx treated mostly-contained child as object_overlap:\n  "
                + "\n  ".join(f.message for f in structural_overflow_object_overlap)
            )
        overflow_structure = pptx_lint.extract_pptx_structure(structural_containment_overflow)
        if not any(
            relation["relation"] == "contains_with_child_overflow"
            and relation.get("child_overflow_pt")
            for relation in overflow_structure
        ):
            failures.append(
                "structural-containment-overflow.pptx did not expose child overflow structure"
            )

        lint005_bad_check_set = {
            f.check for f in pptx_lint.lint_pptx(object_relationships_bad) if f.check in LINT005_CHECKS
        }
        missing_lint005 = LINT005_CHECKS - lint005_bad_check_set
        if missing_lint005:
            failures.append(
                f"object-relationships-bad.pptx did not trigger {sorted(missing_lint005)}; "
                f"got {sorted(lint005_bad_check_set)}"
            )

        card_grid_good_findings = [
            f
            for f in pptx_lint.lint_pptx(card_grid_consistency_good)
            if f.check == "card_grid_consistency"
        ]
        if card_grid_good_findings:
            failures.append(
                "card-grid-consistency-good.pptx triggered card_grid_consistency:\n  "
                + "\n  ".join(f.message for f in card_grid_good_findings)
            )
        card_grid_bad_findings = [
            f
            for f in pptx_lint.lint_pptx(card_grid_consistency_bad)
            if f.check == "card_grid_consistency"
        ]
        if not card_grid_bad_findings:
            failures.append("card-grid-consistency-bad.pptx did not trigger card_grid_consistency")

        bad_findings = pptx_lint.lint_pptx(bad)
        bad_check_set = {f.check for f in bad_findings}
        missing = EXPECTED_BAD_CHECKS - bad_check_set
        if missing:
            failures.append(
                f"bad.pptx did not trigger {sorted(missing)}; got {sorted(bad_check_set)}"
            )
        guideline_checks = _guideline_lint_check_keys()
        unknown = KNOWN_EMITTED_CHECKS - guideline_checks
        if unknown:
            failures.append(
                f"pptx_lint.py emitted checks missing from rules.lint.checks: {sorted(unknown)}"
            )
        automation = _guideline_lint_check_automation()
        for check, expected in LINT004_POLICY.items():
            got = automation.get(check)
            if got != expected:
                failures.append(
                    f"rules.lint.checks.{check}.automation expected {expected!r}; got {got!r}"
                )

        vertical_balance_good_findings = [
            f
            for f in pptx_lint.lint_pptx(text_vertical_balance_good)
            if f.check == "text_vertical_balance"
        ]
        if vertical_balance_good_findings:
            failures.append(
                "text-vertical-balance-good.pptx triggered text_vertical_balance:\n  "
                + "\n  ".join(f.message for f in vertical_balance_good_findings)
            )
        invisible_text_box_vertical_balance_findings = [
            f
            for f in pptx_lint.lint_pptx(invisible_text_box_vertical_balance_good)
            if f.check == "text_vertical_balance"
        ]
        if invisible_text_box_vertical_balance_findings:
            failures.append(
                "invisible-text-box-vertical-balance-good.pptx triggered text_vertical_balance:\n  "
                + "\n  ".join(f.message for f in invisible_text_box_vertical_balance_findings)
            )

        for fixture, label in (
            (top_anchor_bottom_void_bad, "top-anchor-bottom-void-bad.pptx"),
            (middle_anchor_asymmetric_margin_bad, "middle-anchor-asymmetric-margin-bad.pptx"),
            (oversized_box_top_anchor_bad, "oversized-box-top-anchor-bad.pptx"),
        ):
            vertical_balance_findings = [
                f
                for f in pptx_lint.lint_pptx(fixture)
                if f.check == "text_vertical_balance"
            ]
            if not vertical_balance_findings:
                failures.append(f"{label} did not trigger text_vertical_balance")

        for fixture, check, label in (
            (missing_title_bad, "missing_required_element", "missing-title-bad.pptx"),
            (heading_hierarchy_bad, "heading_hierarchy_broken", "heading-hierarchy-bad.pptx"),
            (reading_order_bad, "reading_order", "reading-order-bad.pptx"),
            (wrap_break_bad, "wrap_break_changes_meaning", "wrap-break-bad.pptx"),
            (key_area_cropped_bad, "key_area_cropped", "key-area-cropped-bad.pptx"),
            (color_only_bad, "color_only_meaning", "color-only-bad.pptx"),
        ):
            findings = [f for f in pptx_lint.lint_pptx(fixture) if f.check == check]
            if not findings:
                failures.append(f"{label} did not trigger {check}")

        for fixture, label in (
            (cover_slide_good, "cover-slide-good.pptx"),
            (section_divider_good, "section-divider-good.pptx"),
        ):
            findings = [
                f
                for f in pptx_lint.lint_pptx(fixture)
                if f.check == "missing_required_element"
            ]
            if findings:
                failures.append(
                    f"{label} should satisfy missing_required_element by slide type:\n  "
                    + "\n  ".join(f.message for f in findings)
                )
        cover_safe_margin_findings = [
            f
            for f in pptx_lint.lint_pptx(cover_slide_good)
            if f.check == "safe_margins"
        ]
        if cover_safe_margin_findings:
            failures.append(
                "cover-slide-good.pptx triggered safe_margins for cover brand mark:\n  "
                + "\n  ".join(f.message for f in cover_safe_margin_findings)
            )

        decorative_crop_findings = [
            f
            for f in pptx_lint.lint_pptx(key_area_cropped_decorative)
            if f.check == "key_area_cropped"
        ]
        if decorative_crop_findings:
            failures.append(
                "key-area-cropped-decorative.pptx triggered key_area_cropped:\n  "
                + "\n  ".join(f.message for f in decorative_crop_findings)
            )

        labeled_color_findings = [
            f
            for f in pptx_lint.lint_pptx(color_only_labeled)
            if f.check == "color_only_meaning"
        ]
        if labeled_color_findings:
            failures.append(
                "color-only-labeled.pptx triggered color_only_meaning:\n  "
                + "\n  ".join(f.message for f in labeled_color_findings)
            )

    if failures:
        print("FAIL:")
        for line in failures:
            print(f"- {line}")
        return 1

    print(f"OK: good.pptx clean, bad.pptx triggered all {len(EXPECTED_BAD_CHECKS)} expected checks")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
