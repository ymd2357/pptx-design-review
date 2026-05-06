#!/usr/bin/env python3
"""Summarize PPTX review findings as P0-P3 action priorities.

This is intentionally not a finding counter. It groups mechanical lint results
into decision-level issues that a reviewer can act on.
"""

from __future__ import annotations

import argparse
import json
import sys
from collections import defaultdict
from dataclasses import asdict, dataclass, field
from pathlib import Path
from types import SimpleNamespace
from typing import Any, Iterable

from pptx import Presentation

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))

import pptx_lint  # noqa: E402


PRIORITY_ORDER = {"P0": 0, "P1": 1, "P2": 2, "P3": 3}


@dataclass
class PriorityIssue:
    priority: str
    title: str
    slides: list[int] = field(default_factory=list)
    evidence: list[str] = field(default_factory=list)
    action: str = ""
    checks: list[str] = field(default_factory=list)


def _slide_range(slides: Iterable[int]) -> str:
    ordered = sorted(set(s for s in slides if s > 0))
    if not ordered:
        return "deck"
    runs: list[tuple[int, int]] = []
    start = prev = ordered[0]
    for slide in ordered[1:]:
        if slide == prev + 1:
            prev = slide
        else:
            runs.append((start, prev))
            start = prev = slide
    runs.append((start, prev))
    return ", ".join(str(a) if a == b else f"{a}-{b}" for a, b in runs)


def _clip(text: str, limit: int = 80) -> str:
    text = " ".join(text.split())
    if len(text) <= limit:
        return text
    return text[: limit - 1] + "..."


def load_priority_catalog(guideline_path: Path | None = None) -> dict[str, SimpleNamespace]:
    if guideline_path is None:
        guideline_path = HERE.parents[2] / "doc" / "slide-guideline-v1.yml"

    required_fields = {
        "priority",
        "detection",
        "fix_policy",
        "impact_axis",
        "issue_title",
        "action",
    }

    def strip_yaml_scalar(raw: str) -> str:
        value = raw.strip()
        if not value:
            return ""
        if value[0] in {"'", '"'} and value[-1:] == value[0]:
            return value[1:-1]
        return value

    lines = guideline_path.read_text(encoding="utf-8").splitlines()
    entries: dict[str, dict[str, str]] = {}
    current_check: str | None = None
    in_catalog = False

    for line in lines:
        if line == "    priority_catalog:":
            in_catalog = True
            continue
        if not in_catalog:
            continue
        if not line.strip() or line.lstrip().startswith("#"):
            continue
        if line.startswith("    ") and not line.startswith("      "):
            break
        if line.startswith("      ") and not line.startswith("        ") and line.strip().endswith(":"):
            current_check = line.strip()[:-1]
            entries[current_check] = {}
            continue
        if current_check and line.startswith("        ") and ":" in line:
            key, raw_value = line.strip().split(":", 1)
            entries[current_check][key] = strip_yaml_scalar(raw_value)

    if not entries:
        raise RuntimeError(f"rules.lint.priority_catalog was not found in {guideline_path}")

    catalog: dict[str, SimpleNamespace] = {}
    for check, values in entries.items():
        missing = sorted(required_fields - values.keys())
        if missing:
            raise RuntimeError(f"priority_catalog.{check} is missing fields: {missing}")
        catalog[check] = SimpleNamespace(
            check=check,
            priority=values["priority"],
            detection=values["detection"],
            fix_policy=values["fix_policy"],
            impact_axis=values["impact_axis"],
            issue_title=values["issue_title"],
            action=values["action"],
            status=values.get("status", ""),
        )
    return catalog


def _issue(
    priority: str,
    title: str,
    slides: Iterable[int],
    evidence: Iterable[str],
    action: str,
    checks: Iterable[str],
) -> PriorityIssue:
    return PriorityIssue(
        priority=priority,
        title=title,
        slides=sorted(set(slides)),
        evidence=list(dict.fromkeys(evidence)),
        action=action,
        checks=sorted(set(checks)),
    )


def _text_quality_issues(path: Path) -> list[PriorityIssue]:
    prs = Presentation(str(path))
    replacement_slides: dict[int, list[str]] = defaultdict(list)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if "\ufffd" in text:
                replacement_slides[slide_idx].append(_clip(text))

    issues: list[PriorityIssue] = []
    if replacement_slides:
        evidence = [
            f"slide {slide}: {sample}"
            for slide, samples in replacement_slides.items()
            for sample in samples[:2]
        ]
        issues.append(
            _issue(
                "P0",
                "文字化けしており、内容を正しく読めない",
                replacement_slides.keys(),
                evidence,
                "該当テキストを元原稿または正しい日本語から差し替える。",
                ["text_encoding"],
            )
        )
    return issues


def _slides_for(findings: list[pptx_lint.Finding], checks: set[str]) -> list[int]:
    slides: list[int] = []
    for finding in findings:
        if finding.check not in checks:
            continue
        affected = finding.detail.get("affected_slides") if finding.detail else None
        if affected:
            slides.extend(int(s) for s in affected)
        else:
            slides.append(finding.slide_index)
    return sorted(set(s for s in slides if s > 0))


SCHEMA_EVIDENCE_KEYS = {
    "review_status",
    "shape_kind",
    "text_excerpt",
    "target_description",
    "element_description",
    "bbox_pt",
    "actual_bbox_pt",
    "rendered_bbox_px",
    "measured_value",
    "threshold",
    "delta",
    "unit",
    "evidence_source",
    "evidence_confidence",
    "fixability",
    "manual_required_reason",
    "manual_required_reasons",
    "candidate_values",
    "recommended_value",
    "group_key",
    "artifact_refs",
}

LEGACY_DETAIL_EVIDENCE_KEYS = {
    "color_hex",
    "contrast_ratio",
    "required_ratio",
    "low_contrast_threshold",
    "measurement",
    "text_class",
    "font_size_pt",
    "scope",
}


def _present(value: Any) -> bool:
    return value is not None and value != "" and value != [] and value != {}


def _compact_value(value: Any, limit: int = 80) -> str:
    if isinstance(value, float):
        return f"{value:.2f}".rstrip("0").rstrip(".")
    if isinstance(value, (list, tuple)) and all(isinstance(item, (int, float)) for item in value):
        return "[" + ", ".join(_compact_value(item) for item in value) + "]"
    if isinstance(value, (dict, list, tuple)):
        return _clip(json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":")), limit)
    return _clip(str(value), limit)


def _first_present(detail: dict[str, Any], *keys: str) -> Any:
    for key in keys:
        value = detail.get(key)
        if _present(value):
            return value
    return None


def _has_detail_evidence(detail: dict[str, Any]) -> bool:
    return any(key in detail for key in SCHEMA_EVIDENCE_KEYS | LEGACY_DETAIL_EVIDENCE_KEYS)


def _schema_evidence_for(finding: pptx_lint.Finding) -> str | None:
    detail = finding.detail or {}
    if not _has_detail_evidence(detail):
        return None

    slide = "deck" if finding.slide_index == 0 else f"slide {finding.slide_index}"
    parts = [f"{slide}: {finding.check}"]

    target = _first_present(detail, "text_excerpt", "target_description", "element_description")
    if target is not None:
        label = "text" if "text_excerpt" in detail and _present(detail.get("text_excerpt")) else "target"
        parts.append(f'{label}="{_clip(str(target), 48)}"')

    shape_kind = _first_present(detail, "shape_kind")
    if shape_kind is not None:
        shape_label = _compact_value(shape_kind)
        shape_name = _first_present(detail, "shape_name")
        if shape_name is not None:
            shape_label += f"/{_compact_value(shape_name, 40)}"
        parts.append(f"shape={shape_label}")

    unit = _first_present(detail, "unit")
    unit_text = _compact_value(unit, 16) if unit is not None else ""
    measured = _first_present(detail, "measured_value", "contrast_ratio", "color_hex")
    if measured is not None:
        parts.append(f"measured={_compact_value(measured)}{unit_text}")
    threshold = _first_present(detail, "threshold", "required_ratio", "low_contrast_threshold")
    if threshold is not None:
        parts.append(f"threshold={_compact_value(threshold)}{unit_text}")
    delta = _first_present(detail, "delta")
    if delta is not None:
        parts.append(f"delta={_compact_value(delta)}{unit_text}")

    source = _first_present(detail, "evidence_source", "measurement")
    confidence = _first_present(detail, "evidence_confidence")
    if source is not None:
        source_text = _compact_value(source, 32)
        if confidence is not None:
            source_text += f"/{_compact_value(confidence, 16)}"
        parts.append(f"source={source_text}")

    review_status = _first_present(detail, "review_status")
    if review_status is not None:
        parts.append(f"review={_compact_value(review_status, 32)}")

    fixability = _first_present(detail, "fixability")
    if fixability is not None:
        fixability_text = _compact_value(fixability, 40)
        reason = _first_present(detail, "manual_required_reason", "manual_required_reasons", "reason")
        if reason is not None and fixability == "manual_required":
            fixability_text += f" ({_compact_value(reason, 80)})"
        parts.append(f"fixability={fixability_text}")

    recommended = _first_present(detail, "recommended_value")
    if recommended is not None:
        parts.append(f"recommended={_compact_value(recommended, 80)}")
    candidates = _first_present(detail, "candidate_values")
    if candidates is not None:
        parts.append(f"candidates={_compact_value(candidates, 100)}")

    artifacts = _first_present(detail, "artifact_refs")
    if artifacts is not None:
        parts.append(f"artifacts={_compact_value(artifacts, 100)}")

    bbox = _first_present(detail, "bbox_pt", "actual_bbox_pt", "rendered_bbox_px")
    if bbox is not None and measured is None:
        parts.append(f"bbox={_compact_value(bbox)}")

    return _clip("; ".join(parts), 500)


def _evidence_for(findings: list[pptx_lint.Finding], checks: set[str], limit: int = 5) -> list[str]:
    evidence: list[str] = []
    for finding in findings:
        if finding.check not in checks:
            continue
        slide = "deck" if finding.slide_index == 0 else f"slide {finding.slide_index}"
        schema_evidence = _schema_evidence_for(finding)
        if schema_evidence is not None:
            evidence.append(schema_evidence)
        else:
            evidence.append(f"{slide}: {finding.check}: {_clip(finding.message)}")
        if len(evidence) >= limit:
            break
    return evidence


def _summarize_findings(
    raw_findings: list[pptx_lint.Finding],
    issues: list[PriorityIssue] | None = None,
) -> list[PriorityIssue]:
    findings = pptx_lint.consolidate_recurring(raw_findings)
    issues = list(issues or [])
    p0_checks = {"overflow_text", "text_autofit_disabled", "text_overlap", "low_contrast"}
    slides = _slides_for(findings, p0_checks)
    if slides:
        issues.append(
            _issue(
                "P0",
                "テキストが読めない、または自動縮小で読みにくくなるリスクがある",
                slides,
                _evidence_for(raw_findings, p0_checks),
                "該当スライドを目視確認し、文字色・背景色・テキスト枠・文字量・改行・重なりを調整する。",
                p0_checks,
            )
        )

    p1_motion = {"animation_present"}
    slides = _slides_for(findings, p1_motion)
    if slides:
        issues.append(
            _issue(
                "P1",
                "静的配布で失われる可能性があるアニメーション/遷移が含まれる",
                slides,
                _evidence_for(raw_findings, p1_motion),
                "PPTX/PDF 配布で必要な情報が静止状態でも読めるか確認し、不要な動きは削除する。",
                p1_motion,
            )
        )

    p1_contrast = {"contrast_ratio"}
    slides = _slides_for(findings, p1_contrast)
    if slides:
        issues.append(
            _issue(
                "P1",
                "コントラスト不足で可読性にリスクがある",
                slides,
                _evidence_for(raw_findings, p1_contrast),
                "文字色と背景色を調整し、通常文字 4.5:1 / 大きい文字 3.0:1 以上を満たす。",
                p1_contrast,
            )
        )

    p1_structure = {"object_overlap"}
    slides = _slides_for(findings, p1_structure)
    if slides:
        issues.append(
            _issue(
                "P1",
                "オブジェクト重なりで内容や構造を追いにくい",
                slides,
                _evidence_for(raw_findings, p1_structure),
                "重なり、近接不足、レイヤー順を確認し、意図しない被りを解消する。",
                p1_structure,
            )
        )

    p1_semantic = {
        "color_only_meaning",
        "heading_hierarchy_broken",
        "key_area_cropped",
        "missing_required_element",
        "reading_order",
        "wrap_break_changes_meaning",
    }
    slides = _slides_for(findings, p1_semantic)
    if slides:
        issues.append(
            _issue(
                "P1",
                "構造、順序、欠落、色依存、改行、画像トリミングが理解を阻害するリスクがある",
                slides,
                _evidence_for(raw_findings, p1_semantic),
                "機械 evidence の対象要素を確認し、テンプレート構造、読む順、非色手掛かり、見出し、必須要素、改行、画像トリミングを直す。",
                p1_semantic,
            )
        )

    p2_checks = {
        "overflow_images",
        "overflow_shapes",
        "safe_margins",
        "safe_text_area_text",
        "alignment_left_top",
        "card_grid_consistency",
        "font_size_scale",
        "font_family",
        "inner_padding_imbalance",
        "line_height",
        "image_aspect_distortion",
        "image_upscale_ratio",
        "object_gap_too_small",
        "text_color_allowlist",
        "background_color_palette",
        "slide_size",
        "alt_text_required",
    }
    slides = _slides_for(findings, p2_checks)
    if slides:
        issues.append(
            _issue(
                "P2",
                "テンプレート/ブランドルールから外れているが、許容判断できる項目がある",
                slides,
                _evidence_for(raw_findings, p2_checks),
                "SHIFT AI テンプレートとして意図した表現なら許容ルールへ寄せ、そうでなければマスター/スタイルを直す。",
                p2_checks,
            )
        )

    p3_checks = {"geometry_rounding"}
    slides = _slides_for(findings, p3_checks)
    if slides:
        issues.append(
            _issue(
                "P3",
                "座標の微小な丸めズレがある",
                slides,
                _evidence_for(raw_findings, p3_checks),
                "見た目に問題がなければ後回し。テンプレート整備時に一括補正する。",
                p3_checks,
            )
        )

    issues.sort(key=lambda issue: (PRIORITY_ORDER[issue.priority], issue.title))
    return issues


def summarize_priorities(
    path: Path,
    *,
    rendered_image_dir: Path | None = None,
) -> list[PriorityIssue]:
    raw_findings = pptx_lint.lint_pptx(path, rendered_image_dir=rendered_image_dir)
    return _summarize_findings(raw_findings, _text_quality_issues(path))


def format_markdown(issues: list[PriorityIssue]) -> str:
    if not issues:
        return "OK: P0-P3 review issues were not found.\n"

    lines: list[str] = []
    for priority in ("P0", "P1", "P2", "P3"):
        bucket = [issue for issue in issues if issue.priority == priority]
        if not bucket:
            continue
        lines.append(f"## {priority}")
        for issue in bucket:
            lines.append(f"- {issue.title}")
            lines.append(f"  - slides: {_slide_range(issue.slides)}")
            if issue.evidence:
                lines.append(f"  - evidence: {issue.evidence[0]}")
                for extra in issue.evidence[1:3]:
                    lines.append(f"    / {extra}")
            lines.append(f"  - action: {issue.action}")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def main(argv: list[str]) -> int:
    ap = argparse.ArgumentParser(description="Summarize PPTX review as P0-P3 priorities")
    ap.add_argument("pptx", type=Path, help="path to .pptx")
    ap.add_argument("--json", action="store_true", help="emit priority issues as JSON")
    ap.add_argument(
        "--rendered-image-dir",
        type=Path,
        help="directory containing PowerPoint-rendered slide PNGs for rendered contrast lint",
    )
    args = ap.parse_args(argv)

    if not args.pptx.exists():
        print(f"file not found: {args.pptx}", file=sys.stderr)
        return 2

    issues = summarize_priorities(args.pptx, rendered_image_dir=args.rendered_image_dir)
    if args.json:
        print(json.dumps([asdict(issue) for issue in issues], ensure_ascii=False, indent=2))
    else:
        print(format_markdown(issues), end="")
    return 1 if any(issue.priority == "P0" for issue in issues) else 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
