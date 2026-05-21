#!/usr/bin/env python3
"""Smoke test for P0-P3 PPTX review priority summarization."""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))

import make_examples  # noqa: E402
import pptx_review_priorities  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import MSO_AUTO_SIZE  # noqa: E402
from pptx.util import Pt  # noqa: E402


REQUIRED_CATALOG_CHECKS = {
    "alignment_left_top",
    "alt_text_required",
    "animation_present",
    "background_color_palette",
    "card_grid_consistency",
    "color_only_meaning",
    "contrast_ratio",
    "font_family",
    "font_size_scale",
    "geometry_rounding",
    "heading_hierarchy_broken",
    "image_aspect_distortion",
    "image_upscale_ratio",
    "inner_padding_imbalance",
    "key_area_cropped",
    "line_height",
    "low_contrast",
    "missing_required_element",
    "object_gap_too_small",
    "object_overlap",
    "overflow_images",
    "overflow_shapes",
    "box_canvas_overflow",
    "reading_order",
    "safe_margins",
    "safe_text_area_text",
    "slide_size",
    "text_autofit_disabled",
    "text_color_allowlist",
    "text_encoding",
    "text_overlap",
    "wrap_break_changes_meaning",
}

def _make_mojibake_deck(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40.5), Pt(20), Pt(639), Pt(60))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "プ��ジェクト"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(28)
    prs.save(str(out))


def _make_explicit_break_deck(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40.5), Pt(20), Pt(300), Pt(80))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    first = box.text_frame.paragraphs[0].add_run()
    first.text = "「実装方針の相談」"
    first.font.name = "Noto Sans JP"
    first.font.size = Pt(12)
    second_para = box.text_frame.add_paragraph()
    second = second_para.add_run()
    second.text = "→ 規約に沿った回答"
    second.font.name = "Noto Sans JP"
    second.font.size = Pt(12)
    prs.save(str(out))


def _make_unknown_table_cell_color_deck(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 2, Pt(40.5), Pt(40), Pt(400), Pt(120))
    table = table_shape.table
    for row in table.rows:
        for cell in row.cells:
            cell.text = "Cell"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    table.cell(0, 0).fill.fore_color.rgb = RGBColor.from_string("123456")
    prs.save(str(out))


def _make_unknown_table_cell_font_deck(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(1, 1, Pt(40.5), Pt(40), Pt(300), Pt(60))
    cell = table_shape.table.cell(0, 0)
    cell.text = ""
    run = cell.text_frame.paragraphs[0].add_run()
    run.text = "Bad table font"
    run.font.name = "Arial"
    run.font.size = Pt(12)
    prs.save(str(out))


def main() -> int:
    failures: list[str] = []
    catalog = pptx_review_priorities.load_priority_catalog()
    missing_catalog = REQUIRED_CATALOG_CHECKS - catalog.keys()
    if missing_catalog:
        failures.append(f"priority_catalog is missing checks: {sorted(missing_catalog)}")
    if catalog["overflow_images"].priority != "P1":
        failures.append("priority_catalog.overflow_images should be P1 delivery-quality risk")
    if catalog["geometry_rounding"].priority != "P3" or catalog["geometry_rounding"].fix_policy != "auto_fix":
        failures.append("priority_catalog.geometry_rounding should be P3 auto_fix cleanup")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        bad = tmp_dir / "bad.pptx"
        mojibake = tmp_dir / "mojibake.pptx"
        explicit_break = tmp_dir / "explicit-break.pptx"
        unknown_table_cell_color = tmp_dir / "unknown-table-cell-color.pptx"
        unknown_table_cell_font = tmp_dir / "unknown-table-cell-font.pptx"
        make_examples.make_bad(bad)
        _make_mojibake_deck(mojibake)
        _make_explicit_break_deck(explicit_break)
        _make_unknown_table_cell_color_deck(unknown_table_cell_color)
        _make_unknown_table_cell_font_deck(unknown_table_cell_font)

        bad_issues = pptx_review_priorities.summarize_priorities(bad)
        bad_priorities = {issue.priority for issue in bad_issues}
        for expected in ("P0", "P1", "P2", "P3"):
            if expected not in bad_priorities:
                failures.append(f"bad.pptx did not produce {expected}; got {sorted(bad_priorities)}")

        mojibake_issues = pptx_review_priorities.summarize_priorities(mojibake)
        if not any(issue.priority == "P0" and "文字化け" in issue.title for issue in mojibake_issues):
            failures.append("mojibake.pptx did not produce a P0 mojibake issue")

        explicit_break_issues = pptx_review_priorities.summarize_priorities(explicit_break)
        if any(issue.priority == "P1" for issue in explicit_break_issues):
            failures.append("explicit-break.pptx produced P1 for an authored paragraph break")

        table_color_issues = pptx_review_priorities.summarize_priorities(unknown_table_cell_color)
        if not any(
            issue.priority == "P2" and "background_color_palette" in issue.checks
            for issue in table_color_issues
        ):
            failures.append("unknown-table-cell-color.pptx did not produce P2 background_color_palette")

        table_font_issues = pptx_review_priorities.summarize_priorities(unknown_table_cell_font)
        if not any(
            issue.priority == "P2" and "font_family" in issue.checks
            for issue in table_font_issues
        ):
            failures.append("unknown-table-cell-font.pptx did not produce P2 font_family")

    if failures:
        print("FAIL:")
        for failure in failures:
            print(f"- {failure}")
        return 1

    print("OK: priority summarization produced expected P0-P3 buckets")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
