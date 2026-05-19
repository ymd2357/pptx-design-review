#!/usr/bin/env python3
"""Build a Pn-n review matrix from intentional fixtures and deck results.

This is not a replacement for visual review. It answers two separate questions:

1. Does each Pn-n have an intentionally bad fixture proving the detector can fire?
2. What does the current target deck still require after lint/fix/render gates?
"""

from __future__ import annotations

import argparse
import csv
import html
import json
import re
import shutil
import subprocess
import sys
import zipfile
from collections import Counter, defaultdict
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Callable
from xml.etree import ElementTree as ET

from PIL import Image, ImageChops, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

HERE = Path(__file__).parent
ROOT = HERE.parents[2]
VSCODE_RENDER_SCRIPT = HERE / "render_with_vscode_pptx_viewer.js"
VSCODE_CAPTURE_SCRIPT = HERE / "capture_vscode_pptx_viewer.js"
sys.path.insert(0, str(HERE))

import make_examples  # noqa: E402
import pptx_fix  # noqa: E402
import pptx_lint  # noqa: E402
import pptx_review_priorities  # noqa: E402
import test_pptx_lint as lint_fixtures  # noqa: E402


TASK_ROW_RE = re.compile(
    r"^\|\s*(P[0-3]-\d+)\s*\|\s*`([^`]+)`\s*\|\s*([^|]+?)\s*\|\s*([^|]+?)\s*\|\s*([^|]+?)\s*\|\s*([^|]+?)\s*\|"
)
PRIORITY_ORDER = {"P0": 0, "P1": 1, "P2": 2, "P3": 3}
RESULT_ONLY_IF_ZERO = {
    "alignment_left_top",
    "alt_text_required",
    "animation_present",
    "background_color_palette",
    "color_only_meaning",
    "font_family",
    "font_size_scale",
    "heading_hierarchy_broken",
    "image_aspect_distortion",
    "image_upscale_ratio",
    "key_area_cropped",
    "line_height",
    "object_gap_too_small",
    "object_overlap",
    "overflow_images",
    "overflow_shapes",
    "overflow_text",
    "slide_size",
    "text_autofit_disabled",
    "text_color_allowlist",
    "text_encoding",
    "text_overlap",
    "text_vertical_balance",
    "wrap_break_changes_meaning",
}


@dataclass
class CatalogItem:
    pn: str
    check: str
    priority: str
    declared_status: str
    detection: str
    fix_policy: str
    viewpoint: str


@dataclass
class FixtureSpec:
    name: str
    checks: tuple[str, ...]
    builder: Callable[[Path, Path], None]
    detector: str = "lint"


@dataclass
class FixtureResult:
    fixture: str
    detector: str
    path: str
    checks: list[str]
    found_counts: dict[str, int]
    status: str


@dataclass
class ReviewRow:
    pn: str
    priority: str
    check: str
    fixture_status: str
    fixture_count: int
    fixture_paths: str
    deck_decision: str
    before_count: int
    after_count: int
    fixed_actions: int
    slides_after: str
    reviewer_mode: str
    next_action: str


@dataclass
class FixEvidenceRow:
    pn: str
    priority: str
    check: str
    fixture: str
    fixture_status: str
    before_count: int
    after_count: int
    applied_actions: int
    manual_actions: int
    outcome: str
    outcome_reason: str
    render_status: str
    diff_nonempty_slides: int
    work_dir: str
    before_pptx: str
    after_pptx: str
    report_href: str
    issue_state: str = ""
    expected_state: str = ""


def _set_text(shape, text: str, *, size: int = 24, color: str | None = None) -> None:
    shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = RGBColor.from_string(color)


def _make_text_encoding_bad(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Pt(120), Pt(60), Pt(900), Pt(80))
    _set_text(title, "文字化け\ufffdサンプル", size=32)
    prs.save(str(out))


def _make_text_encoding_good(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Pt(120), Pt(60), Pt(900), Pt(80))
    _set_text(title, "正しい日本語サンプル", size=32)
    prs.save(str(out))


def _make_slide_size_bad(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1000)
    prs.slide_height = Pt(500)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Pt(80), Pt(60), Pt(600), Pt(80))
    _set_text(title, "Non proportional slide", size=24)
    prs.save(str(out))


def _make_safe_font_size_bad(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40), Pt(40), Pt(300), Pt(80))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Short label"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(14.75)
    prs.save(str(out))


def _make_safe_line_height_bad(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40), Pt(40), Pt(420), Pt(90))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    para = box.text_frame.paragraphs[0]
    para.line_spacing = Pt(16.4)
    run = para.add_run()
    run.text = "Line height fixture"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(12)
    prs.save(str(out))


def _make_safe_alignment_bad(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40), Pt(40), Pt(300), Pt(80))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    para = box.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "Alignment fixture"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(12)
    prs.save(str(out))


def _make_safe_geometry_rounding_bad(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(round(81.05 * 12700)), Pt(40), Pt(200), Pt(50))
    _set_text(box, "Geometry fixture", size=24)
    prs.save(str(out))


def _make_overflow_shape_bad(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(1380), Pt(200), Pt(120), Pt(80))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string("EEEEEE")
    prs.save(str(out))


def _make_overflow_shape_good(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(120), Pt(200), Pt(120), Pt(80))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string("EEEEEE")
    prs.save(str(out))


def _set_norm_autofit_font_scale(shape, font_scale: int) -> None:
    body_pr = shape._element.xpath(".//a:bodyPr")[0]
    norm_autofit = body_pr.find(f"{pptx_lint.A_NS}normAutofit")
    if norm_autofit is not None:
        norm_autofit.set("fontScale", str(font_scale))


def _make_text_autofit_shrink_bad(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(120), Pt(120), Pt(260), Pt(54))
    box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    _set_norm_autofit_font_scale(box, 65000)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "This text is auto-shrunk"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(30)
    prs.save(str(out))


def _make_content_overflow_image_good(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = out.with_name("_content_image_good.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        slide.shapes.add_picture(str(img), Pt(80), Pt(80), width=Pt(120), height=Pt(120))
    finally:
        img.unlink(missing_ok=True)
    prs.save(str(out))


def _make_aspect_image_good(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = out.with_name("_aspect_image_good.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        slide.shapes.add_picture(str(img), Pt(80), Pt(80), width=Pt(120), height=Pt(120))
    finally:
        img.unlink(missing_ok=True)
    prs.save(str(out))


def _make_missing_title_good(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(50), Pt(900), Pt(70)), "Title", size=36)
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(180), Pt(900), Pt(80)), "Body with required title", size=24)
    prs.save(str(out))


def _make_heading_hierarchy_good(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(50), Pt(900), Pt(70)), "Heading slot uses heading size", size=40)
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(220), Pt(900), Pt(70)), "Body slot uses body size", size=24)
    prs.save(str(out))


def _make_reading_order_good(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(50), Pt(900), Pt(70)), "First visually", size=32)
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(200), Pt(900), Pt(70)), "Second visually", size=24)
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(340), Pt(900), Pt(70)), "Third visually", size=24)
    prs.save(str(out))


def _make_wrap_break_good(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(50), Pt(900), Pt(70)), "Title", size=32)
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(180), Pt(900), Pt(70)), "Automation improves review speed", size=24)
    prs.save(str(out))


def _make_key_area_cropped_good(out: Path, _work: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_text(slide.shapes.add_textbox(Pt(120), Pt(50), Pt(900), Pt(70)), "Title", size=32)
    img = out.with_name("_key_area_good.png")
    img.write_bytes(make_examples.TINY_PNG)
    try:
        slide.shapes.add_picture(str(img), Pt(120), Pt(180), width=Pt(320), height=Pt(320))
    finally:
        img.unlink(missing_ok=True)
    prs.save(str(out))


def _make_color_only_good(out: Path, _work: Path) -> None:
    lint_fixtures._make_color_only_bad(out, labeled=True)


def _builder(fn: Callable[[Path], None]) -> Callable[[Path, Path], None]:
    return lambda out, _work: fn(out)


def _expected_builder_for_check(check: str) -> Callable[[Path, Path], None]:
    object_relationships_checks = {
        "inner_padding_imbalance",
        "object_gap_too_small",
        "object_overlap",
        "text_overlap",
    }
    generic_good_checks = {
        "animation_present",
        "alt_text_required",
        "alignment_left_top",
        "background_color_palette",
        "contrast_ratio",
        "font_family",
        "font_size_scale",
        "geometry_rounding",
        "image_upscale_ratio",
        "line_height",
        "low_contrast",
        "safe_margins",
        "safe_text_area_text",
        "text_autofit_disabled",
        "text_color_allowlist",
        "overflow_text",
        "slide_size",
    }
    if check == "text_encoding":
        return _make_text_encoding_good
    if check in object_relationships_checks:
        return _builder(lint_fixtures._make_object_relationships_good)
    if check == "overflow_images":
        return _make_content_overflow_image_good
    if check == "overflow_shapes":
        return _make_overflow_shape_good
    if check == "image_aspect_distortion":
        return _make_aspect_image_good
    if check == "missing_required_element":
        return _make_missing_title_good
    if check == "heading_hierarchy_broken":
        return _make_heading_hierarchy_good
    if check == "reading_order":
        return _make_reading_order_good
    if check == "wrap_break_changes_meaning":
        return _make_wrap_break_good
    if check == "key_area_cropped":
        return _make_key_area_cropped_good
    if check == "color_only_meaning":
        return _make_color_only_good
    if check == "text_vertical_balance":
        return _builder(lint_fixtures._make_text_vertical_balance_good)
    if check == "card_grid_consistency":
        return _builder(lint_fixtures._make_card_grid_consistency_good)
    if check in generic_good_checks:
        return _builder(make_examples.make_good)
    return _builder(make_examples.make_good)


def _rendered_low_contrast(out: Path, work: Path) -> None:
    lint_fixtures._make_rendered_low_contrast_case(out, work / "rendered-low-contrast-images")


def _set_shape_text_style(shape, *, font_size: int | None = None, font_name: str | None = None, color: str | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if font_size is not None:
                run.font.size = Pt(font_size)
            if font_name is not None:
                run.font.name = font_name
            if color is not None:
                run.font.color.rgb = RGBColor.from_string(color)


def _replace_text_preserving_runs(shape, old: str, new: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)


def _picture_shapes(slide) -> list:
    return [shape for shape in slide.shapes if getattr(shape, "shape_type", None) == 13]


def _remove_slide_timing(path: Path) -> None:
    ns = {"p": make_examples.PML_NS}
    tmp = path.with_suffix(path.suffix + ".tmp")
    with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                root = ET.fromstring(data)
                for tag in ("transition", "timing"):
                    child = root.find(f"p:{tag}", ns)
                    if child is not None:
                        root.remove(child)
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            zout.writestr(item, data)
    tmp.replace(path)


def _reorder_shapes_top_to_bottom(slide) -> None:
    shape_elements = [shape._element for shape in slide.shapes]
    ordered_shapes = sorted(slide.shapes, key=lambda shape: (shape.top, shape.left))
    tree = slide.shapes._spTree
    for element in shape_elements:
        tree.remove(element)
    for shape in ordered_shapes:
        tree.append(shape._element)


def _shape_by_id(slide, shape_id: Any):
    try:
        target_id = int(shape_id)
    except (TypeError, ValueError):
        return None
    for shape in slide.shapes:
        if getattr(shape, "shape_id", None) == target_id:
            return shape
    return None


def _detail_shape_id(detail: dict, *keys: str) -> int | None:
    current: Any = detail
    for key in keys:
        if not isinstance(current, dict):
            return None
        current = current.get(key)
    if isinstance(current, dict):
        value = current.get("shape_id")
    else:
        value = current
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _move_shape_to_bbox(shape, bbox: list[float] | tuple[float, float, float, float], *, dx: float = 0, dy: float = 0) -> None:
    x, y, w, h = bbox
    shape.left = Pt(x + dx)
    shape.top = Pt(y + dy)
    shape.width = Pt(w)
    shape.height = Pt(h)


def _apply_mechanical_expected_preview(check: str, before: Path, expected: Path, findings_json: list[dict]) -> None:
    shutil.copy2(before, expected)
    rules = _fix_rules_for_check(check)
    if rules:
        pptx_fix.fix_pptx(expected, apply=True, rules=rules, findings=findings_json)
        return

    prs = Presentation(str(expected))
    slide = prs.slides[0]
    shapes = list(slide.shapes)
    detail = _first_detail(findings_json)

    if check == "text_encoding":
        for shape in shapes:
            _replace_text_preserving_runs(shape, "\ufffd", "正")
    elif check == "text_overlap":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "shape_b"))
        bbox = (detail.get("shape_b") or {}).get("bbox_pt")
        if shape is not None and bbox:
            _move_shape_to_bbox(shape, bbox, dy=64)
    elif check == "object_overlap":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "shape_b"))
        bbox = (detail.get("shape_b") or {}).get("bbox_pt")
        if shape is not None and bbox:
            _move_shape_to_bbox(shape, bbox, dx=140, dy=-30)
    elif check == "object_gap_too_small":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "shape_b"))
        bbox = (detail.get("shape_b") or {}).get("bbox_pt")
        gap = float(detail.get("gap_pt") or 0)
        threshold = float(detail.get("threshold_pt") or 8)
        axis = detail.get("axis")
        if shape is not None and bbox:
            if axis == "horizontal":
                _move_shape_to_bbox(shape, bbox, dx=max(0, threshold - gap + 4))
            else:
                _move_shape_to_bbox(shape, bbox, dy=max(0, threshold - gap + 4))
    elif check == "inner_padding_imbalance":
        children = detail.get("children") or []
        if len(children) >= 2:
            first = _shape_by_id(slide, children[0].get("shape_id"))
            second = _shape_by_id(slide, children[1].get("shape_id"))
            if first is not None:
                first.left = Pt(530)
                first.top = Pt(330)
            if second is not None:
                second.left = Pt(770)
                second.top = Pt(410)
    elif check == "overflow_text":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "target"))
        if shape is None and shapes:
            shape = shapes[0]
        if shape is not None:
            shape.left = Pt(980)
            shape.width = Pt(360)
    elif check == "alt_text_required":
        for pic in _picture_shapes(slide):
            c_nv_pr = pic._element.xpath(".//p:cNvPr")[0]
            c_nv_pr.set("descr", "Content image")
            c_nv_pr.set("title", "Content image")
    elif check == "color_only_meaning":
        labels = ("NG", "OK")
        idx = 0
        for shape in shapes:
            if getattr(shape, "has_text_frame", False) and not shape.text.strip() and idx < len(labels):
                shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                run = shape.text_frame.paragraphs[0].add_run()
                run.text = labels[idx]
                run.font.name = "Noto Sans JP"
                run.font.size = Pt(20)
                idx += 1
    elif check == "heading_hierarchy_broken":
        title = _shape_by_id(slide, _detail_shape_id(detail, "title_candidate"))
        body = _shape_by_id(slide, _detail_shape_id(detail, "largest_body_candidate"))
        _set_shape_text_style(title, font_size=40)
        _set_shape_text_style(body, font_size=24)
    elif check == "image_aspect_distortion":
        for pic in _picture_shapes(slide):
            pic.height = pic.width
    elif check == "key_area_cropped":
        for pic in _picture_shapes(slide):
            pic.crop_left = 0
            pic.crop_right = 0
            pic.crop_top = 0
            pic.crop_bottom = 0
    elif check == "missing_required_element":
        title = slide.shapes.add_textbox(Pt(120), Pt(50), Pt(900), Pt(70))
        _set_text(title, "Title", size=36)
    elif check == "overflow_images":
        for pic in _picture_shapes(slide):
            pic.left = Pt(80)
            pic.top = Pt(80)
    elif check == "overflow_shapes":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "target"))
        if shape is not None:
            shape.left = Pt(120)
            shape.top = Pt(200)
    elif check == "reading_order":
        _reorder_shapes_top_to_bottom(slide)
    elif check == "wrap_break_changes_meaning":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "target"))
        if shape is None:
            shape = _shape_by_id(slide, detail.get("shape_id"))
        if shape is not None:
            _replace_text_preserving_runs(shape, "Auto\nmation", "Automation")
    elif check == "background_color_palette":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "target"))
        if shape is not None and getattr(shape, "fill", None) is not None:
            try:
                shape.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
            except Exception:
                pass
    elif check == "font_family":
        for shape in shapes:
            _set_shape_text_style(shape, font_name="Noto Sans JP")
    elif check == "safe_margins":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "target"))
        if shape is not None:
            shape.left = max(shape.left, Pt(81))
            shape.top = max(shape.top, Pt(40))
    elif check == "safe_text_area_text":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "target"))
        if shape is not None:
            shape.left = max(shape.left, Pt(81))
    elif check == "text_color_allowlist":
        shape = _shape_by_id(slide, _detail_shape_id(detail, "target"))
        _set_shape_text_style(shape, color="202020")
    elif check == "card_grid_consistency":
        for shape in shapes:
            if 900 <= shape.left.pt <= 980:
                shape.width = Pt(300)
            if getattr(shape, "has_text_frame", False) and shape.left.pt >= 930:
                if "Card title" in shape.text:
                    shape.top = Pt(246)
                elif "Card body" in shape.text:
                    shape.top = Pt(300)
    elif check == "text_vertical_balance" and shapes:
        shapes[0].height = Pt(42)
        if getattr(shapes[0], "has_text_frame", False):
            shapes[0].text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    prs.save(str(expected))
    if check == "animation_present":
        _remove_slide_timing(expected)


def fixture_specs() -> list[FixtureSpec]:
    return [
        FixtureSpec("text-encoding-bad", ("text_encoding",), _make_text_encoding_bad, "priority"),
        FixtureSpec("safe-font-size-bad", ("font_size_scale",), _make_safe_font_size_bad),
        FixtureSpec("safe-line-height-bad", ("line_height",), _make_safe_line_height_bad),
        FixtureSpec("safe-alignment-bad", ("alignment_left_top",), _make_safe_alignment_bad),
        FixtureSpec("safe-geometry-rounding-bad", ("geometry_rounding",), _make_safe_geometry_rounding_bad),
        FixtureSpec("text-autofit-shrink-bad", ("text_autofit_disabled",), _make_text_autofit_shrink_bad),
        FixtureSpec(
            "bad-multi-check",
            (
                "animation_present",
                "alt_text_required",
                "alignment_left_top",
                "background_color_palette",
                "contrast_ratio",
                "font_family",
                "font_size_scale",
                "geometry_rounding",
                "image_upscale_ratio",
                "line_height",
                "low_contrast",
                "safe_margins",
                "safe_text_area_text",
                "text_color_allowlist",
                "overflow_text",
            ),
            _builder(make_examples.make_bad),
        ),
        FixtureSpec(
            "object-relationships-bad",
            (
                "inner_padding_imbalance",
                "object_gap_too_small",
                "object_overlap",
                "text_overlap",
            ),
            _builder(lint_fixtures._make_object_relationships_bad),
        ),
        FixtureSpec("content-overflow-image", ("overflow_images",), _builder(lint_fixtures._make_content_overflow_image)),
        FixtureSpec("overflow-shape-bad", ("overflow_shapes",), _make_overflow_shape_bad),
        FixtureSpec("slide-size-bad", ("slide_size",), _make_slide_size_bad),
        FixtureSpec("aspect-distorted-image", ("image_aspect_distortion",), _builder(lint_fixtures._make_aspect_distorted_image)),
        FixtureSpec("missing-title-bad", ("missing_required_element",), _builder(lint_fixtures._make_missing_title_bad)),
        FixtureSpec("heading-hierarchy-bad", ("heading_hierarchy_broken",), _builder(lint_fixtures._make_heading_hierarchy_bad)),
        FixtureSpec("reading-order-bad", ("reading_order",), _builder(lint_fixtures._make_reading_order_bad)),
        FixtureSpec("wrap-break-bad", ("wrap_break_changes_meaning",), _builder(lint_fixtures._make_wrap_break_bad)),
        FixtureSpec("key-area-cropped-bad", ("key_area_cropped",), _builder(lint_fixtures._make_key_area_cropped_bad)),
        FixtureSpec("color-only-bad", ("color_only_meaning",), _builder(lint_fixtures._make_color_only_bad)),
        FixtureSpec("top-anchor-bottom-void-bad", ("text_vertical_balance",), _builder(lint_fixtures._make_top_anchor_bottom_void_bad)),
        FixtureSpec("card-grid-consistency-bad", ("card_grid_consistency",), _builder(lint_fixtures._make_card_grid_consistency_bad)),
        FixtureSpec("rendered-low-contrast", ("low_contrast",), _rendered_low_contrast),
    ]


def load_json(path: Path | None, default: Any) -> Any:
    if path is None or not path.exists():
        return default
    return json.loads(path.read_text(encoding="utf-8"))


def iter_catalog_lines(tasks_md: Path) -> list[str]:
    lines: list[str] = []
    in_catalog = False
    for line in tasks_md.read_text(encoding="utf-8").splitlines():
        if line.startswith("## チェック観点一覧"):
            in_catalog = True
            continue
        if in_catalog and line.startswith("## "):
            break
        if in_catalog:
            lines.append(line)
    if not lines:
        raise RuntimeError(f"Pn-n catalog section was not found in {tasks_md}")
    return lines


def load_catalog(tasks_md: Path) -> list[CatalogItem]:
    items: list[CatalogItem] = []
    seen: set[tuple[str, str]] = set()
    for line in iter_catalog_lines(tasks_md):
        match = TASK_ROW_RE.match(line)
        if not match:
            continue
        pn, check, status, detection, fix_policy, viewpoint = match.groups()
        key = (pn.strip(), check.strip())
        if key in seen:
            continue
        seen.add(key)
        items.append(
            CatalogItem(
                pn=pn.strip(),
                check=check.strip(),
                priority=pn.split("-", 1)[0],
                declared_status=status.strip(),
                detection=detection.strip(),
                fix_policy=fix_policy.strip(),
                viewpoint=viewpoint.strip(),
            )
        )
    if not items:
        raise RuntimeError(f"Pn-n catalog rows were not found in {tasks_md}")
    return sorted(items, key=lambda item: (PRIORITY_ORDER[item.priority], int(item.pn.split("-")[1])))


def run_fixture(spec: FixtureSpec, fixture_dir: Path) -> FixtureResult:
    work = fixture_dir / spec.name
    work.mkdir(parents=True, exist_ok=True)
    pptx_path = work / f"{spec.name}.pptx"
    spec.builder(pptx_path, work)

    if spec.detector == "priority":
        issues = pptx_review_priorities.summarize_priorities(pptx_path)
        found = Counter(check for issue in issues for check in issue.checks)
    else:
        rendered_dir = work / "rendered-low-contrast-images"
        rendered_arg = rendered_dir if rendered_dir.exists() else None
        findings = pptx_lint.lint_pptx(pptx_path, rendered_image_dir=rendered_arg)
        found = Counter(finding.check for finding in findings)

    expected = set(spec.checks)
    found_counts = {check: found.get(check, 0) for check in spec.checks}
    missing = [check for check, count in found_counts.items() if count <= 0]
    return FixtureResult(
        fixture=spec.name,
        detector=spec.detector,
        path=str(pptx_path),
        checks=sorted(expected),
        found_counts=found_counts,
        status="pass" if not missing else "fail",
    )


def findings_by_check(findings: list[dict]) -> dict[str, list[dict]]:
    by: dict[str, list[dict]] = defaultdict(list)
    for finding in findings:
        by[finding.get("check", "")].append(finding)
    return by


def slides_for(findings: list[dict]) -> list[int]:
    slides: set[int] = set()
    for finding in findings:
        detail = finding.get("detail") or {}
        affected = detail.get("affected_slides")
        if isinstance(affected, list):
            slides.update(int(s) for s in affected if int(s) > 0)
        else:
            slide = finding.get("slide_index")
            if isinstance(slide, int) and slide > 0:
                slides.add(slide)
    return sorted(slides)


def slide_range(slides: list[int]) -> str:
    if not slides:
        return ""
    runs: list[tuple[int, int]] = []
    start = prev = slides[0]
    for slide in slides[1:]:
        if slide == prev + 1:
            prev = slide
            continue
        runs.append((start, prev))
        start = prev = slide
    runs.append((start, prev))
    return ", ".join(str(a) if a == b else f"{a}-{b}" for a, b in runs)


def action_counts_by_check(actions_json: dict) -> Counter:
    counts: Counter = Counter()
    for action in actions_json.get("actions", []):
        if action.get("status") != "apply":
            continue
        for update in (action.get("after") or {}).get("updates", []):
            check = update.get("check")
            if check:
                counts[check] += 1
    return counts


def fixability_counts(findings: list[dict]) -> tuple[int, int]:
    auto = 0
    manual = 0
    for finding in findings:
        detail = finding.get("detail") or {}
        if detail.get("fixability") == "auto_fix_candidate":
            auto += 1
        elif detail.get("fixability") == "manual_required":
            manual += 1
    return auto, manual


def classify_deck(item: CatalogItem, before: list[dict], after: list[dict], fixed_actions: int) -> tuple[str, str, str]:
    if after and fixability_counts(after)[0] == len(after):
        return "auto_fix_ready", "result_then_auto_fix", "fixer を適用し、after lint=0 と diff 対象範囲を確認する。"
    if after:
        if item.priority in {"P0", "P1"}:
            return "manual_review_required", "finding_detail_and_visual", "finding detail とレンダ画像で修正/許容/対象外を判断する。"
        if item.priority == "P2":
            return "accept_or_fix_decision_required", "sampled_visual_or_acceptance", "代表スライドでテンプレート意図として許容するか修正方針を決める。"
        return "defer_or_batch_fix", "result_only_unless_touching_layout", "見た目に問題がなければ defer。レイアウト修正時だけ一括確認する。"
    if fixed_actions:
        return "fixed_verified", "result_and_diff", "fix actions、after lint=0、before/after diff の対象スライドだけ確認する。"
    mode = "result_only_with_fixture" if item.check in RESULT_ONLY_IF_ZERO else "result_only"
    return "pass_by_result", mode, "after lint=0。ただし intentionally bad fixture が pass していることを前提に結果だけ見る。"


def build_rows(
    catalog: list[CatalogItem],
    before_lint: list[dict],
    after_lint: list[dict],
    actions_json: dict,
    fixture_results: list[FixtureResult],
) -> list[ReviewRow]:
    before_by = findings_by_check(before_lint)
    after_by = findings_by_check(after_lint)
    actions_by = action_counts_by_check(actions_json)
    fixtures_by_check: dict[str, list[FixtureResult]] = defaultdict(list)
    for result in fixture_results:
        for check, count in result.found_counts.items():
            if count > 0:
                fixtures_by_check[check].append(result)

    rows: list[ReviewRow] = []
    for item in catalog:
        before = before_by.get(item.check, [])
        after = after_by.get(item.check, [])
        fixed = actions_by[item.check]
        decision, mode, next_action = classify_deck(item, before, after, fixed)
        fixtures = fixtures_by_check.get(item.check, [])
        rows.append(
            ReviewRow(
                pn=item.pn,
                priority=item.priority,
                check=item.check,
                fixture_status="pass" if fixtures else "missing",
                fixture_count=sum(result.found_counts.get(item.check, 0) for result in fixtures),
                fixture_paths="; ".join(result.path for result in fixtures),
                deck_decision=decision,
                before_count=len(before),
                after_count=len(after),
                fixed_actions=fixed,
                slides_after=slide_range(slides_for(after)),
                reviewer_mode=mode,
                next_action=next_action,
            )
        )
    return rows


def write_tsv(path: Path, rows: list[ReviewRow]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=list(asdict(rows[0]).keys()), delimiter="\t")
        writer.writeheader()
        for row in rows:
            writer.writerow(asdict(row))


def summary(rows: list[ReviewRow], fixtures: list[FixtureResult]) -> dict[str, Any]:
    return {
        "total_items": len(rows),
        "fixture_pass_items": [row.pn for row in rows if row.fixture_status == "pass"],
        "fixture_missing_items": [row.pn for row in rows if row.fixture_status == "missing"],
        "deck_decisions": dict(Counter(row.deck_decision for row in rows)),
        "result_only_pass_items": [row.pn for row in rows if row.deck_decision == "pass_by_result"],
        "manual_review_items": [
            row.pn
            for row in rows
            if row.deck_decision in {"manual_review_required", "accept_or_fix_decision_required"}
        ],
        "auto_fix_ready_items": [row.pn for row in rows if row.deck_decision == "auto_fix_ready"],
        "fixed_verified_items": [row.pn for row in rows if row.deck_decision == "fixed_verified"],
        "fixture_failures": [asdict(result) for result in fixtures if result.status != "pass"],
    }


def write_json(path: Path, rows: list[ReviewRow], fixtures: list[FixtureResult]) -> None:
    payload = {
        "summary": summary(rows, fixtures),
        "rows": [asdict(row) for row in rows],
        "fixtures": [asdict(result) for result in fixtures],
    }
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def write_markdown(path: Path, rows: list[ReviewRow], fixtures: list[FixtureResult]) -> None:
    s = summary(rows, fixtures)
    result_only = [row for row in rows if row.deck_decision == "pass_by_result"]
    active = [row for row in rows if row.deck_decision != "pass_by_result"]
    lines = [
        "# Pn-n Review Orchestrator Report",
        "",
        "## Summary",
        "",
        f"- total Pn-n items: {s['total_items']}",
        f"- intentionally bad fixture pass: {len(s['fixture_pass_items'])}",
        f"- intentionally bad fixture missing: {len(s['fixture_missing_items'])}",
        f"- result-only deck pass: {len(result_only)}",
        f"- active deck queue: {len(active)}",
        "",
        "## Test Flow",
        "",
        "1. Catalog gate: `doc/tasks.md` の `チェック観点一覧` から Pn-n/check を固定する。",
        "2. Intentional fixture gate: 各 check の bad fixture を生成し、検出器が意図通り fire するか確認する。",
        "3. Deck measurement gate: 対象 deck の before/after lint と fix actions を check 別に突き合わせる。",
        "4. Review routing gate: fixture pass かつ after=0 の項目は result-only、残件は auto-fix/manual/acceptance/defer に振り分ける。",
        "",
    ]
    if s["fixture_missing_items"]:
        lines.extend(["## Fixture Missing", ""])
        for row in rows:
            if row.fixture_status != "pass":
                lines.append(f"- {row.pn} `{row.check}`")
        lines.append("")

    lines.extend(["## Result-Only Deck Pass With Fixture", ""])
    for row in result_only:
        lines.append(
            f"- {row.pn} `{row.check}`: deck after=0, fixture_count={row.fixture_count}, fixture={row.fixture_paths}"
        )
    lines.append("")

    lines.extend(["## Active Deck Queue", ""])
    for row in active:
        lines.append(
            f"- {row.pn} `{row.check}`: {row.deck_decision}, before={row.before_count}, after={row.after_count}, "
            f"slides={row.slides_after or '-'}, fixture_count={row.fixture_count}"
        )
        lines.append(f"  - 次: {row.next_action}")
    lines.append("")

    lines.extend(["## Full Matrix", ""])
    lines.append("| Pn | check | fixture | fixture_count | deck_decision | after | slides | next |")
    lines.append("| --- | --- | --- | ---: | --- | ---: | --- | --- |")
    for row in rows:
        lines.append(
            f"| {row.pn} | `{row.check}` | {row.fixture_status} | {row.fixture_count} | "
            f"{row.deck_decision} | {row.after_count} | {row.slides_after or '-'} | {row.next_action} |"
        )
    lines.append("")
    path.write_text("\n".join(lines), encoding="utf-8")


def write_html(path: Path, markdown_path: Path) -> None:
    text = markdown_path.read_text(encoding="utf-8")
    escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    path.write_text(
        f"""<!doctype html>
<html lang="ja">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width,initial-scale=1">
  <title>Pn-n Review Orchestrator</title>
  <style>
    body{{font-family:system-ui,-apple-system,BlinkMacSystemFont,"Segoe UI",sans-serif;margin:24px;color:#222;line-height:1.55}}
    pre{{white-space:pre-wrap;background:#f7f7f7;border:1px solid #ddd;padding:16px;overflow:auto}}
  </style>
</head>
<body>
<pre>{escaped}</pre>
</body>
</html>
""",
        encoding="utf-8",
    )


def _slug(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9_.-]+", "-", value).strip("-")


def _fixture_spec_by_check(specs: list[FixtureSpec]) -> dict[str, FixtureSpec]:
    mapping: dict[str, FixtureSpec] = {}
    for spec in specs:
        for check in spec.checks:
            mapping.setdefault(check, spec)
    return mapping


def _lint_fixture(path: Path, work: Path) -> list:
    rendered_dir = work / "rendered-low-contrast-images"
    rendered_arg = rendered_dir if rendered_dir.exists() else None
    return pptx_lint.lint_pptx(path, rendered_image_dir=rendered_arg)


def _count_check(item: CatalogItem, pptx_path: Path, findings: list) -> int:
    if item.check == "text_encoding":
        return sum(
            1
            for issue in pptx_review_priorities.summarize_priorities(pptx_path)
            if item.check in issue.checks
        )
    return sum(
        1
        for finding in findings
        if (
            finding.get("check")
            if isinstance(finding, dict)
            else getattr(finding, "check", None)
        )
        == item.check
    )


def _fix_rules_for_check(check: str) -> tuple[str, ...]:
    rule = pptx_fix.CHECK_TO_RULE.get(check)
    return (rule,) if rule else ()


def _fixability_counts_for_findings(findings_json: list[dict]) -> Counter:
    counts: Counter = Counter()
    for finding in findings_json:
        detail = finding.get("detail") or {}
        fixability = detail.get("fixability")
        if fixability:
            counts[fixability] += 1
    return counts


def _fixture_fix_outcome(
    *,
    before_count: int,
    after_count: int,
    rules: tuple[str, ...],
    applied: int,
    manual: int,
    findings_json: list[dict],
    render_status: str,
) -> tuple[str, str]:
    if render_status != "ok":
        return "render_failed", "before/after render failed; visual evidence is incomplete"
    if before_count == 0:
        return "fixture_not_detected", "intentional bad fixture did not trigger the expected check"
    if applied and after_count == 0:
        return "autofixed", "fixer applied a mechanical rule and the check count dropped to zero"
    if applied and after_count > 0:
        return "fixer_incomplete_for_mechanical_candidate", "fixer applied an action but the check still remains"

    fixability = _fixability_counts_for_findings(findings_json)
    if fixability["auto_fix_candidate"] > 0:
        if manual:
            return (
                "fixer_blocked_for_mechanical_candidate",
                "lint marked this as auto-fixable, but fixer returned manual_required",
            )
        return (
            "fixer_missing_for_mechanical_candidate",
            "lint marked this as auto-fixable, but no fixer action was applied",
        )
    if fixability["manual_required"] > 0 or not rules:
        return (
            "tested_manual_decision_required",
            "fixture, lint, render, and diff were tested; source repair requires semantic or layout judgment",
        )
    if manual:
        return "tested_manual_decision_required", "fixer reported manual_required after safety checks"
    return "no_autofix_rule", "no mechanical auto-fix rule is declared for this check"


OUTCOME_JA = {
    "autofixed": "自動修正済み",
    "tested_manual_decision_required": "テスト済み・手動判断が必要",
    "fixer_missing_for_mechanical_candidate": "バグ: 機械修正候補だが fixer 未実装",
    "fixer_blocked_for_mechanical_candidate": "バグ: 機械修正候補だが fixer が手動扱い",
    "fixer_incomplete_for_mechanical_candidate": "バグ: 自動修正後も検出が残存",
    "fixture_not_detected": "バグ: fixture が検出されない",
    "render_failed": "バグ: レンダリング失敗",
    "no_autofix_rule": "自動修正ルールなし",
}

OUTCOME_REASON_JA = {
    "autofixed": "機械修正を適用し、after の検出数が 0 になった。",
    "tested_manual_decision_required": "fixture・lint・render はテスト済み。check 固有の人手判断が必要なため、機械修正しない。",
    "fixer_missing_for_mechanical_candidate": "lint は自動修正候補としているのに、fixer が何も適用していない。fixer の未実装として扱う。",
    "fixer_blocked_for_mechanical_candidate": "lint は自動修正候補としているのに、fixer が manual_required を返した。fixer 側の安全判定または実装を見直す。",
    "fixer_incomplete_for_mechanical_candidate": "fixer は修正を適用したが、after でも同じ check が残っている。修正不完全として扱う。",
    "fixture_not_detected": "意図的に壊した fixture が期待 check を発火していない。検出器または fixture の問題。",
    "render_failed": "before/after のレンダリングに失敗しており、目視 evidence が不完全。",
    "no_autofix_rule": "この check には機械修正 rule が定義されていない。",
}


PRIORITY_CATALOG = pptx_review_priorities.load_priority_catalog()


def _check_title(row: FixEvidenceRow) -> str:
    item = PRIORITY_CATALOG.get(row.check)
    return getattr(item, "issue_title", "") or row.check


def _check_action(row: FixEvidenceRow) -> str:
    item = PRIORITY_CATALOG.get(row.check)
    return getattr(item, "action", "") or "対応方針が priority_catalog に未定義。"


def _check_label_html(row: FixEvidenceRow) -> str:
    return (
        f"<code>{html.escape(row.check)}</code>"
        f"<br><span class='muted'>{html.escape(_check_title(row))}</span>"
    )


def _outcome_label(row: FixEvidenceRow) -> str:
    return f"{OUTCOME_JA.get(row.outcome, row.outcome)} ({row.outcome})"


def _outcome_reason(row: FixEvidenceRow) -> str:
    reason = OUTCOME_REASON_JA.get(row.outcome, row.outcome_reason)
    if row.outcome == "tested_manual_decision_required":
        return f"{reason} 対象: {_check_title(row)}。必要な判断: {_check_action(row)}"
    if row.outcome in {"fixer_missing_for_mechanical_candidate", "fixer_blocked_for_mechanical_candidate", "fixer_incomplete_for_mechanical_candidate"}:
        return f"{reason} 対象: {_check_title(row)}。期待する対応: {_check_action(row)}"
    return reason


def _check_explanation(row: FixEvidenceRow) -> str:
    return f"check の意味: {_check_title(row)} / 対応方針: {_check_action(row)}"


def _first_detail(findings_json: list[dict]) -> dict:
    if not findings_json:
        return {}
    return findings_json[0].get("detail") or {}


def _fmt_num(value: Any) -> str:
    if value is None:
        return "未記録"
    if isinstance(value, float):
        return f"{value:.2f}".rstrip("0").rstrip(".")
    return str(value)


def _autofit_mode_label(auto_size: Any) -> str:
    value = str(auto_size or "")
    if "SHAPE_TO_FIT_TEXT" in value:
        return "SHAPE_TO_FIT_TEXT（文字に合わせて枠が伸縮する設定）"
    if "TEXT_TO_FIT_SHAPE" in value:
        return "TEXT_TO_FIT_SHAPE（枠に合わせて文字が縮小される設定）"
    return value or "未記録"


def _issue_state_for_check(check: str, findings_json: list[dict], before_count: int) -> str:
    detail = _first_detail(findings_json)
    if check == "text_encoding":
        return "文字化け検出。元原稿と照合しないと、どの文字へ戻すべきかは決められない。"
    if check == "animation_present":
        markers = ", ".join(detail.get("markers") or [])
        return f"PPTX XML上に静的配布で問題になる可能性のある要素がある（{markers or 'animation/transition'}）。特定shapeではなくslide XMLの検出なので、画像上の枠は付けない。"
    if check == "heading_hierarchy_broken":
        title = detail.get("title_candidate") or {}
        body = detail.get("largest_body_candidate") or {}
        title_size = title.get("font_size_pt")
        body_size = body.get("font_size_pt")
        rules = ", ".join(detail.get("triggered_rules") or [])
        return (
            f"デザインシステム上の見出し/本文階層に対し、見出し候補が {_fmt_num(title_size)}pt、"
            f"本文候補が {_fmt_num(body_size)}pt。本文が見出しを視覚的に上回っている可能性がある"
            f"（検出規則: {rules or 'n/a'}）。"
        )
    if check == "key_area_cropped":
        crop = detail.get("crop_ratio") or {}
        thresholds = detail.get("thresholds") or {}
        side = max(crop, key=lambda item: crop.get(item) or 0) if crop else "n/a"
        side_label = {"l": "左", "r": "右", "t": "上", "b": "下"}.get(side, side)
        max_side = detail.get("max_side_crop_ratio")
        side_limit = thresholds.get("side_ratio_max")
        return (
            f"非装飾画像の{side_label}側 crop が {_fmt_num(max_side)}。片側 crop の閾値 "
            f"{_fmt_num(side_limit)} を超えているため、重要部が切れているリスクとして検出している。"
            "機械は被写体の意味までは確定せず、crop metadata から要確認として出している。"
        )
    if check == "text_autofit_disabled":
        mode = _autofit_mode_label(detail.get("auto_size"))
        base_size = detail.get("base_font_size_pt")
        effective_size = detail.get("effective_font_size_pt")
        font_scale = detail.get("font_scale_percent")
        shrink = detail.get("font_shrink_percent")
        size_text = (
            f"文字サイズは {_fmt_num(base_size)}pt から {_fmt_num(effective_size)}pt 相当に縮小（{_fmt_num(font_scale)}%、縮小量 {_fmt_num(shrink)}%）。"
            if base_size is not None and effective_size is not None
            else f"fontScale={_fmt_num(font_scale)}%。"
        )
        return (
            f"PowerPoint の text_frame.auto_size が NONE ではなく {mode}。"
            f"{size_text}"
            "開く環境や編集で枠・改行・文字サイズが変わる余地を消すため、機械的に NONE にする。"
        )
    if check in {"text_overlap", "object_overlap"}:
        area = detail.get("overlap_area_pt2")
        threshold = detail.get("threshold_pt2") or detail.get("threshold")
        return f"枠同士の重なり面積が {_fmt_num(area)}pt^2。閾値 {_fmt_num(threshold)} を超えている。"
    if check == "object_gap_too_small":
        axis = detail.get("axis")
        axis_label = {"horizontal": "横方向", "vertical": "縦方向"}.get(axis, axis or "隣接方向")
        gap = detail.get("gap_pt")
        threshold = detail.get("threshold_pt")
        element_a = detail.get("element_a") or "default"
        element_b = detail.get("element_b") or "default"
        default_threshold = detail.get("default_threshold_pt")
        def shape_ref(value: Any, fallback: str) -> str:
            if not isinstance(value, dict):
                return str(value or fallback)
            name = value.get("shape_name") or fallback
            shape_id = value.get("shape_id")
            return f"{name}#{shape_id}" if shape_id is not None else str(name)

        shape_a = shape_ref(detail.get("shape_a"), "要素A")
        shape_b = shape_ref(detail.get("shape_b"), "要素B")
        return (
            f"{shape_a} と {shape_b} の{axis_label}の間隔が {_fmt_num(gap)}pt。"
            f"ds:element={element_a}/{element_b} に対するデザインシステムの隣接要素最小分離 "
            f"{_fmt_num(threshold)}pt を下回っている。"
            f"未指定要素は default {_fmt_num(default_threshold)}pt なので、P2-1 では落とさない。"
        )
    if check == "overflow_images":
        sides = detail.get("overflow_sides_pt") or {}
        side_text = ", ".join(f"{side}={_fmt_num(amount)}pt" for side, amount in sides.items())
        bbox = detail.get("bbox_pt")
        return (
            f"PPTX上の画像ボックスがスライド外へ出ている（{side_text or 'はみ出しあり'}）。"
            f"bbox={_fmt_num(bbox)}。見えている画像の縁は実ボックスの枠ではなく、スライド境界でクリップされた結果。"
            "ASISははみ出した辺、測定図はスライド境界と実ボックスを示す。"
        )
    if check in {"overflow_text", "overflow_shapes", "safe_margins", "safe_text_area_text"}:
        overflow = detail.get("overflow_by_pt")
        if overflow is None:
            return "対象要素が許容領域から外れている。枠線で対象要素を示している。"
        return f"対象要素が許容領域からはみ出している。はみ出し量: {_fmt_num(overflow)}pt。"
    if check == "missing_required_element":
        zone = detail.get("required_zone_pt") or {}
        return f"スライドタイプ上、必要なタイトル要素が不足している。期待ゾーンは上端 { _fmt_num(zone.get('top')) }pt から { _fmt_num(zone.get('bottom')) }pt。存在しない要素の検出なので、既存shapeの枠は付けない。"
    if check == "reading_order":
        inversions = detail.get("inversions")
        inversion_count = len(inversions) if isinstance(inversions, list) else inversions
        return f"ソース順と視覚順が一致していない。読み順の逆転数: {_fmt_num(inversion_count)}。"
    if check == "color_only_meaning":
        return "同種要素の違いが色だけで表現されている可能性がある。ラベルや形状差が不足している。"
    if check == "text_color_allowlist":
        color = detail.get("color_hex")
        candidate = detail.get("candidate_values") or {}
        candidate_color = candidate.get("color_hex") if isinstance(candidate, dict) else None
        candidate_token = candidate.get("color_token") if isinstance(candidate, dict) else None
        return (
            f"文字色 {_fmt_num(color)} がデザインシステムの許可リスト外。"
            f"機械候補としては {_fmt_num(candidate_color)}"
            f"{f'（{candidate_token}）' if candidate_token else ''} が出ているが、"
            "design-system に従うことを優先し、この候補へ機械修正する。"
        )
    if check == "alignment_left_top":
        alignment = detail.get("alignment")
        expected_alignment = detail.get("expected_alignment") or "LEFT"
        vertical = detail.get("vertical_anchor")
        expected_vertical = detail.get("expected_vertical_anchor") or "TOP"
        return (
            f"テキスト配置が design-system の text_box 既定と合っていない。"
            f"段落配置={_fmt_num(alignment)}（期待={expected_alignment}）、"
            f"縦配置={_fmt_num(vertical)}（期待={expected_vertical}）。"
            "通常テキストとして扱う fixture なので、左揃え・上揃えへ機械修正する。"
        )
    if check == "card_grid_consistency":
        return "同じカード群として扱うべき要素のサイズ、位置、内側配置が揃っていない。"
    if check == "image_upscale_ratio":
        source_px = detail.get("source_px")
        display_px = detail.get("display_px")
        ratio = detail.get("upscale_ratio")
        return (
            f"画像の元解像度が {_fmt_num(source_px)}px なのに、表示上は {_fmt_num(display_px)}px 相当まで拡大されている。"
            f"拡大率 {_fmt_num(ratio)}倍で、デザインシステムの最大 1倍を超えている。"
            "表示サイズを1px相当に縮めるのは修正ではないため、機械修正せず高解像度画像への差し替え、または意図した表示サイズの再設計が必要。"
        )
    if check == "inner_padding_imbalance":
        element = detail.get("container_element") or (detail.get("container") or {}).get("design_element") or "default"
        padding = detail.get("padding_pt") or {}
        target = detail.get("target_padding_pt") or (detail.get("thresholds") or {}).get("target_padding_pt")
        triggered = ", ".join(detail.get("triggered_rules") or [])
        padding_text = ", ".join(
            f"{label}={_fmt_num(padding.get(side))}pt"
            for side, label in (("left", "左"), ("right", "右"), ("top", "上"), ("bottom", "下"))
        )
        return (
            f"ds:element={element} のコンテナ内余白が、デザインシステムの component padding "
            f"{_fmt_num(target)}pt と合っていない。"
            f"実測 padding: {padding_text}。検出理由: {triggered or 'n/a'}。"
            "未指定要素は default 0pt なので P2-6 では落とさない。"
        )
    if check == "slide_size":
        actual = detail.get("actual_pt")
        base = detail.get("base_pt")
        return (
            f"デックのスライドサイズが基準比率と合っていない。actual={_fmt_num(actual)}、base={_fmt_num(base)}。"
            "修正方針は、キャンバスを基準サイズへ変えて中身を再配置する方法と、現サイズを維持して基準比率に合わせる方法の少なくとも2パターンがある。"
            "どちらが正しいかは納品形式と既存レイアウト意図で決まるため、機械では決めない。"
        )
    measured = detail.get("measured_value")
    threshold = detail.get("threshold")
    unit = detail.get("unit")
    if measured is not None or threshold is not None:
        return f"測定値 {_fmt_num(measured)} が基準 {_fmt_num(threshold)} と合っていない。単位: {_fmt_num(unit)}。"
    if before_count:
        return f"{before_count}件の検出がある。詳細は枠付きの現状画像と check の意味を確認する。"
    return "この fixture では対象 check が検出されていない。"


def _expected_state_for_check(row: FixEvidenceRow, findings_json: list[dict] | None = None) -> str:
    detail = _first_detail(findings_json or [])
    check = row.check
    if check == "text_encoding":
        return "OK条件: 表示文字が元原稿と一致し、置換文字や文字化けがない。確認: 正しい原稿から文字列を復元する。"
    if check == "text_overlap":
        threshold = detail.get("threshold_pt2") or detail.get("threshold") or 1
        return f"OK条件: テキスト枠同士の重なりがなく、重なり面積が {_fmt_num(threshold)}pt^2 以下。確認: 文字が互いに被らず、上から読んで意味が通る。"
    if check == "low_contrast":
        return "OK条件: 文字と背景のコントラストが読める水準を満たす。確認: 背景が画像やグラデーションでも、最終レンダリング上の文字領域で測る。"
    if check == "overflow_text":
        return "OK条件: テキストが枠内とスライド内に収まり、末尾まで見えている。確認: 枠拡大、文字量削減、改行調整のどれで直すかを選ぶ。"
    if check == "animation_present":
        return "OK条件: 静的なPPTX/PDF配布でも全情報が読める。確認: アニメーションでしか現れない情報を固定表示にするか、不要な動きを削除する。"
    if check == "alt_text_required":
        return "OK条件: 意味を持つ画像には内容や役割を説明する代替テキストがある。確認: 装飾画像なら装飾扱い、情報画像なら説明を入れる。"
    if check == "color_only_meaning":
        return "OK条件: 色を見分けられなくても意味が分かる。確認: ラベル、番号、記号、形状差など色以外の手掛かりを足す。"
    if check == "contrast_ratio":
        return "OK条件: 文字と背景のコントラストが警告閾値以上。確認: 最終レンダリング画像で、対象文字の実背景に対して測る。"
    if check == "heading_hierarchy_broken":
        return "OK条件: スライドタイプのテンプレートスロットに従い、見出しは title/heading 系、本文は body 系としてサイズ・位置・順序が揃う。確認: 本文が見出しより強く見えない。"
    if check == "image_aspect_distortion":
        return "OK条件: 画像が元の縦横比で表示され、人物・画面・図表が不自然に伸縮していない。確認: 必要なら比率を保ったままトリミングで収める。"
    if check == "key_area_cropped":
        thresholds = detail.get("thresholds") or {}
        side_limit = thresholds.get("side_ratio_max", 0.18)
        axis_limit = thresholds.get("axis_total_ratio_max", 0.3)
        return (
            f"OK条件: 非装飾画像の重要な被写体、画面操作、図表ラベルが見切れていない。"
            f"機械基準: 片側 crop は原則 {_fmt_num(side_limit)} 以下、同一軸合計 crop は {_fmt_num(axis_limit)} 以下。"
            "確認: 閾値内でも意味上重要な部分が切れていればNG、閾値超えでも装飾なら対象外にする。"
        )
    if check == "line_height":
        return "OK条件: 行間がデザインシステムの行高スケールに乗り、本文が詰まりすぎず間延びもしない。確認: 文字サイズ、行数、枠高さをセットで見る。"
    if check == "missing_required_element":
        return "OK条件: スライドタイプに必要なタイトル、本文、注記、出典などのスロットが埋まっている。確認: 不要なスロットは例外理由を持つ。"
    if check == "object_overlap":
        threshold = detail.get("threshold_pt2") or detail.get("threshold") or 1
        return f"OK条件: 意図しないオブジェクト被りがなく、重なり面積が {_fmt_num(threshold)}pt^2 以下。確認: 背景・カード内包・装飾レイヤーは例外として扱う。"
    if check == "object_gap_too_small":
        threshold = detail.get("threshold_pt") or 8
        return f"OK条件: 隣接要素の間隔がデザインシステムの最小分離 {_fmt_num(threshold)}pt 以上で、別要素として読み分けられる。確認: 意図した密着、グループ表現、内包関係でなければ最低間隔を確保する。"
    if check == "overflow_images":
        return "OK条件: 画像の必要情報がスライド内に収まり、意図しない外にはみ出しがない。確認: フルブリード画像や装飾背景は例外として扱う。"
    if check == "overflow_shapes":
        return "OK条件: 情報を持つ図形がスライド内に収まる。確認: 装飾として外へ逃がしている図形か、情報要素かを分ける。"
    if check == "reading_order":
        return "OK条件: Selection Pane のソース順が、視覚的に読む順番と論理構造に一致する。確認: 上から下、左から右、見出しから本文へ読める順に並べる。"
    if check == "text_autofit_disabled":
        return "OK条件: TEXT_TO_FIT_SHAPE による自動縮小がなく、文字サイズが指定値のまま読める。確認: overflow がある場合は枠、改行、文字量を明示的に調整する。"
    if check == "wrap_break_changes_meaning":
        return "OK条件: 改行や折返しで単語、係り受け、数値単位の意味が分断されない。確認: 意味単位で改行する。"
    if check == "background_color_palette":
        return "OK条件: 背景・塗り色が許可パレット内、または意図した例外として記録されている。確認: ブランド表現か一時的な逸脱かを判断する。"
    if check == "font_family":
        return "OK条件: 日本語・英数字ともテンプレート指定フォントに揃っている。確認: 代替フォントが必要な場合は理由を記録する。"
    if check == "font_size_scale":
        return "OK条件: 文字サイズがデザインシステムのサイズスケールに乗っている。確認: サイズだけでなく、枠幅・改行・階層も同時に見る。"
    if check == "image_upscale_ratio":
        return "OK条件: 表示サイズに対して元画像解像度が不足していない。確認: 高解像度画像へ差し替えるか、表示サイズを下げる。"
    if check == "inner_padding_imbalance":
        target = detail.get("target_padding_pt") or (detail.get("thresholds") or {}).get("target_padding_pt") or 24
        return f"OK条件: design-system 対象コンテナでは、左上の内余白が component padding {_fmt_num(target)}pt に合い、右下も最低 {_fmt_num(target)}pt を下回らない。サイズ変更が必要な場合は手動判断。"
    if check == "card_grid_consistency":
        return "OK条件: 同一行・同種カードの幅、高さ、上端、内余白、主要子要素、同格テキスト位置が揃っている。確認: 個別要素ではなくカード群として比較する。"
    if check == "safe_margins":
        return "OK条件: 非テキスト情報要素が安全余白内にある。確認: 表紙ロゴ、背景、フルブリード装飾は例外として切り分ける。"
    if check == "safe_text_area_text":
        return "OK条件: 本文テキストが安全テキスト領域内にあり、投影・LMS表示でも読める。確認: 表紙や区切りページのテンプレート例外は別扱いにする。"
    if check == "slide_size":
        return "OK条件: 納品形式に合うスライドサイズ方針が選ばれている。確認: 1) 基準キャンバスへ変換して全要素を再配置する、2) 現キャンバスを維持しつつ比率・余白・出力先を調整する、のどちらを採るかを人が決める。"
    if check == "text_color_allowlist":
        return "OK条件: 文字色が許可リスト内で、背景との可読性も満たす。確認: 許可色でも背景上で読めなければ contrast 側で直す。"
    if check == "alignment_left_top":
        return "OK条件: テンプレート想定の左揃え・上揃えに戻っている。確認: 中央揃えや中央垂直配置が意図した表現なら例外扱いにする。"
    if check == "text_vertical_balance":
        return "OK条件: テキストボックス内の縦余白が自然で、文字が上や下に偏って見えない。確認: font size / line height / box height / margin / anchor を対等に比較する。"
    if check == "geometry_rounding":
        return "OK条件: 座標・サイズがテンプレートの丸め基準に揃っている。確認: 見た目に影響しない微小ズレはまとめて機械補正する。"
    return f"OK条件: {_check_action(row)}"


def _run(cmd: list[str]) -> tuple[bool, str]:
    try:
        completed = subprocess.run(
            cmd,
            check=False,
            cwd=ROOT,
            text=True,
            capture_output=True,
        )
    except Exception as exc:  # noqa: BLE001 - keep evidence generation non-fatal.
        return False, str(exc)
    output = "\n".join(part for part in (completed.stdout, completed.stderr) if part)
    return completed.returncode == 0, output


def _slide_part_names(pptx_path: Path) -> list[str]:
    try:
        prs = Presentation(str(pptx_path))
        return [str(slide.part.partname).lstrip("/") for slide in prs.slides]
    except Exception:
        return []


def _shape_xfrm_emu(shape_el: ET.Element) -> tuple[int | None, int | None, int | None, int | None]:
    xfrm = shape_el.find(".//p:spPr/a:xfrm", {"p": make_examples.PML_NS, "a": pptx_lint.A_NS.strip("{}")})
    if xfrm is None:
        return (None, None, None, None)
    off = xfrm.find(f"{pptx_lint.A_NS}off")
    ext = xfrm.find(f"{pptx_lint.A_NS}ext")
    if off is None or ext is None:
        return (None, None, None, None)
    try:
        return (
            int(off.get("x", "0")),
            int(off.get("y", "0")),
            int(ext.get("cx", "0")),
            int(ext.get("cy", "0")),
        )
    except ValueError:
        return (None, None, None, None)


def _norm_autofit_scales_by_shape(pptx_path: Path) -> dict[int, list[tuple[tuple[int | None, int | None, int | None, int | None], float]]]:
    by_slide: dict[int, list[tuple[tuple[int | None, int | None, int | None, int | None], float]]] = defaultdict(list)
    part_names = _slide_part_names(pptx_path)
    if not part_names:
        return by_slide
    ns = {"p": make_examples.PML_NS, "a": pptx_lint.A_NS.strip("{}")}
    with zipfile.ZipFile(pptx_path) as zf:
        for slide_index, part_name in enumerate(part_names):
            if part_name not in zf.namelist():
                continue
            root = ET.fromstring(zf.read(part_name))
            for shape_el in root.findall(".//p:sp", ns):
                norm_autofit = shape_el.find(".//a:bodyPr/a:normAutofit", ns)
                if norm_autofit is None:
                    continue
                try:
                    font_scale = int(norm_autofit.get("fontScale", "100000"))
                except ValueError:
                    font_scale = 100000
                if font_scale >= 100000:
                    continue
                by_slide[slide_index].append((_shape_xfrm_emu(shape_el), font_scale / 100000))
    return by_slide


def _write_viewer_json_and_embedded(viewer_dir: Path, data: dict) -> None:
    slides_json = viewer_dir / "slides.json"
    slides_json.write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
    index_html = viewer_dir / "index.html"
    if not index_html.exists():
        return
    embedded_data = json.dumps(data, ensure_ascii=False).replace("<", "\\u003c")
    lines = index_html.read_text(encoding="utf-8").splitlines()
    for idx, line in enumerate(lines):
        if "window.__viewerDataPromise = Promise.resolve(" in line:
            indent = line[: len(line) - len(line.lstrip())]
            lines[idx] = f"{indent}window.__viewerDataPromise = Promise.resolve({embedded_data});"
            break
    index_html.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _split_viewer_paragraph_explicit_newlines(paragraph: dict) -> list[dict]:
    output: list[dict] = [{"runs": []}]
    for run in paragraph.get("runs") or []:
        text = str(run.get("text") or "")
        parts = re.split(r"[\n\v]", text)
        for idx, part in enumerate(parts):
            if idx > 0:
                output.append({"runs": []})
            if part:
                next_run = dict(run)
                next_run["text"] = part
                output[-1]["runs"].append(next_run)
    for key, value in paragraph.items():
        if key == "runs":
            continue
        for item in output:
            item.setdefault(key, value)
    return [item for item in output if item.get("runs")]


def _apply_explicit_newlines_to_viewer_json(viewer_dir: Path) -> int:
    slides_json = viewer_dir / "slides.json"
    if not slides_json.exists():
        return 0
    data = json.loads(slides_json.read_text(encoding="utf-8"))
    patched = 0
    for slide in data.get("slides") or []:
        for element in slide.get("elements") or []:
            text_body = element.get("textBody")
            if not text_body:
                continue
            next_paragraphs: list[dict] = []
            for paragraph in text_body.get("paragraphs") or []:
                if any(("\n" in str(run.get("text") or "") or "\v" in str(run.get("text") or "")) for run in paragraph.get("runs") or []):
                    split = _split_viewer_paragraph_explicit_newlines(paragraph)
                    next_paragraphs.extend(split)
                    patched += max(1, len(split) - 1)
                else:
                    next_paragraphs.append(paragraph)
            text_body["paragraphs"] = next_paragraphs
    if patched:
        _write_viewer_json_and_embedded(viewer_dir, data)
    return patched


def _apply_norm_autofit_to_viewer_json(pptx_path: Path, viewer_dir: Path) -> int:
    scales_by_slide = _norm_autofit_scales_by_shape(pptx_path)
    if not scales_by_slide:
        return 0
    slides_json = viewer_dir / "slides.json"
    if not slides_json.exists():
        return 0
    data = json.loads(slides_json.read_text(encoding="utf-8"))
    patched = 0
    for slide_index, slide in enumerate(data.get("slides") or []):
        pending = list(scales_by_slide.get(slide_index) or [])
        if not pending:
            continue
        for element in slide.get("elements") or []:
            if element.get("type") != "shape" or not element.get("textBody"):
                continue
            match_index = None
            element_box = (
                int(element.get("x", -1)),
                int(element.get("y", -1)),
                int(element.get("width", -1)),
                int(element.get("height", -1)),
            )
            for idx, (shape_box, _scale) in enumerate(pending):
                if shape_box == element_box:
                    match_index = idx
                    break
            if match_index is None:
                match_index = 0
            _shape_box, scale = pending.pop(match_index)
            for paragraph in element["textBody"].get("paragraphs") or []:
                for run in paragraph.get("runs") or []:
                    if isinstance(run.get("fontSize"), int | float):
                        run["fontSize"] = round(float(run["fontSize"]) * scale, 4)
                        run["pptxFontScale"] = round(scale * 100, 2)
                        patched += 1
            if not pending:
                break
    if patched:
        _write_viewer_json_and_embedded(viewer_dir, data)
    return patched


def _render_pptx(pptx_path: Path, work: Path, label: str) -> tuple[bool, str, Path]:
    viewer_dir = work / f"{label}-viewer"
    render_dir = work / label
    ok, output = _run(["node", str(VSCODE_RENDER_SCRIPT), str(pptx_path), str(viewer_dir)])
    if not ok:
        return False, output, render_dir
    patched_breaks = _apply_explicit_newlines_to_viewer_json(viewer_dir)
    if patched_breaks:
        output += f"\nviewer_explicit_newline_patched_paragraphs={patched_breaks}\n"
    patched_runs = _apply_norm_autofit_to_viewer_json(pptx_path, viewer_dir)
    if patched_runs:
        output += f"\nviewer_norm_autofit_fontscale_patched_runs={patched_runs}\n"
    ok, capture_output = _run(["node", str(VSCODE_CAPTURE_SCRIPT), str(viewer_dir), str(render_dir)])
    return ok, output + "\n" + capture_output, render_dir


def _diff_renders(before_dir: Path, after_dir: Path, diff_dir: Path) -> int:
    diff_dir.mkdir(parents=True, exist_ok=True)
    nonempty = 0
    for before_path in sorted(before_dir.glob("slide-*.png")):
        after_path = after_dir / before_path.name
        if not after_path.exists():
            continue
        before = Image.open(before_path).convert("RGB")
        after = Image.open(after_path).convert("RGB")
        if before.size != after.size:
            after = after.resize(before.size)
        diff = ImageChops.difference(before, after)
        if diff.getbbox():
            nonempty += 1
        diff.point(lambda v: min(255, v * 4)).save(diff_dir / before_path.name)
    return nonempty


def _shape_bbox_detail(prs: Presentation, shape) -> dict:
    scale_x = 1440.0 / (prs.slide_width.pt or 1440.0)
    scale_y = 810.0 / (prs.slide_height.pt or 810.0)
    bbox = [
        round(shape.left.pt * scale_x, 2),
        round(shape.top.pt * scale_y, 2),
        round(shape.width.pt * scale_x, 2),
        round(shape.height.pt * scale_y, 2),
    ]
    actual = [
        round(shape.left.pt, 2),
        round(shape.top.pt, 2),
        round(shape.width.pt, 2),
        round(shape.height.pt, 2),
    ]
    return {
        "shape_id": getattr(shape, "shape_id", None),
        "shape_name": getattr(shape, "name", ""),
        "bbox_pt": bbox,
        "actual_bbox_pt": actual,
    }


def _synthetic_text_encoding_findings(path: Path) -> list[dict]:
    prs = Presentation(str(path))
    findings: list[dict] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and "\ufffd" in shape.text:
                target = _shape_bbox_detail(prs, shape)
                findings.append(
                    {
                        "severity": "error",
                        "check": "text_encoding",
                        "slide_index": slide_index,
                        "slide_id": getattr(slide, "slide_id", slide_index),
                        "shape_id": target["shape_id"],
                        "shape_name": target["shape_name"],
                        "message": "replacement character is present in text",
                        "detail": {
                            "check_id": "text_encoding",
                            "target": target,
                            "bbox_pt": target["bbox_pt"],
                            "actual_bbox_pt": target["actual_bbox_pt"],
                            "text_excerpt": shape.text[:80],
                            "fixability": "manual_required",
                            "manual_required_reason": "requires source text confirmation",
                        },
                    }
                )
    return findings


def _augment_findings_with_shape_bbox(path: Path, findings_json: list[dict]) -> list[dict]:
    if not findings_json:
        return findings_json
    prs = Presentation(str(path))
    shape_by_key: dict[tuple[int, int], Any] = {}
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            shape_by_key[(slide_index, getattr(shape, "shape_id", -1))] = shape

    augmented: list[dict] = []
    for finding in findings_json:
        detail = dict(finding.get("detail") or {})
        shape = shape_by_key.get((int(finding.get("slide_index") or 1), int(finding.get("shape_id") or -1)))
        if shape is not None:
            target = _shape_bbox_detail(prs, shape)
            detail.setdefault("target", target)
            detail.setdefault("bbox_pt", target["bbox_pt"])
            detail.setdefault("actual_bbox_pt", target["actual_bbox_pt"])
            finding = {**finding, "detail": detail}
        augmented.append(finding)
    return augmented


def _valid_bbox(value: Any) -> bool:
    return (
        isinstance(value, list)
        and len(value) == 4
        and all(isinstance(item, int | float) for item in value)
        and value[2] > 0
        and value[3] > 0
    )


def _finding_bboxes(value: Any) -> list[tuple[str, tuple[float, float, float, float]]]:
    boxes: list[tuple[str, tuple[float, float, float, float]]] = []
    if isinstance(value, dict):
        for key, child in value.items():
            if key in {"bbox_pt", "overlap_bbox_pt", "target_bbox_pt"} and _valid_bbox(child):
                boxes.append((key, tuple(float(item) for item in child)))
            else:
                boxes.extend(_finding_bboxes(child))
    elif isinstance(value, list) and not _valid_bbox(value):
        for child in value:
            boxes.extend(_finding_bboxes(child))
    return boxes


def _annotate_before_render(
    before_dir: Path,
    findings_json: list[dict],
    annotated_dir: Path,
    *,
    fallback_full_slide: bool = False,
) -> None:
    annotated_dir.mkdir(parents=True, exist_ok=True)
    findings_by_slide: dict[int, list[dict]] = defaultdict(list)
    for finding in findings_json:
        findings_by_slide[int(finding.get("slide_index") or 1)].append(finding)

    for before_path in sorted(before_dir.glob("slide-*.png")):
        match = re.search(r"slide-(\d+)\.png$", before_path.name)
        slide_index = int(match.group(1)) if match else 1
        image = Image.open(before_path).convert("RGBA")
        overlay = Image.new("RGBA", image.size, (255, 255, 255, 0))
        draw = ImageDraw.Draw(overlay)
        scale_x = image.width / 1440.0
        scale_y = image.height / 810.0
        line_width = max(4, round(min(image.size) / 220))
        seen: set[tuple[str, tuple[float, float, float, float]]] = set()

        slide_has_box = False
        for finding in findings_by_slide.get(slide_index, []):
            detail = finding.get("detail") or {}
            overflow_sides = detail.get("overflow_sides_pt") if isinstance(detail, dict) else None
            for key, bbox in _finding_bboxes(finding.get("detail") or {}):
                rounded = tuple(round(value, 2) for value in bbox)
                dedupe_key = ("overlap_bbox_pt" if key == "overlap_bbox_pt" else "bbox", rounded)
                if dedupe_key in seen:
                    continue
                seen.add(dedupe_key)
                x, y, w, h = bbox
                pad = 4 if key == "overlap_bbox_pt" else 2
                left = max(0, round((x - pad) * scale_x))
                top = max(0, round((y - pad) * scale_y))
                right = min(image.width - 1, round((x + w + pad) * scale_x))
                bottom = min(image.height - 1, round((y + h + pad) * scale_y))
                color = (220, 0, 45, 255) if key == "overlap_bbox_pt" else (255, 140, 0, 255)
                if isinstance(overflow_sides, dict) and overflow_sides:
                    marker_width = max(line_width * 2, 10)
                    marker_color = (220, 0, 45, 255)
                    unclipped_left = round((x - pad) * scale_x)
                    unclipped_top = round((y - pad) * scale_y)
                    unclipped_right = round((x + w + pad) * scale_x)
                    unclipped_bottom = round((y + h + pad) * scale_y)
                    if "top" in overflow_sides:
                        draw.rectangle(
                            [left, 0, right, min(image.height - 1, marker_width)],
                            fill=marker_color,
                        )
                    if "bottom" in overflow_sides:
                        draw.rectangle(
                            [left, max(0, image.height - 1 - marker_width), right, image.height - 1],
                            fill=marker_color,
                        )
                    if "left" in overflow_sides:
                        draw.rectangle(
                            [0, top, min(image.width - 1, marker_width), bottom],
                            fill=marker_color,
                        )
                    if "right" in overflow_sides:
                        draw.rectangle(
                            [max(0, image.width - 1 - marker_width), top, image.width - 1, bottom],
                            fill=marker_color,
                        )
                    if (
                        unclipped_left >= image.width
                        or unclipped_right <= 0
                        or unclipped_top >= image.height
                        or unclipped_bottom <= 0
                    ):
                        draw.rectangle([0, 0, image.width - 1, image.height - 1], outline=marker_color, width=line_width)
                    slide_has_box = True
                    continue
                fill = (220, 0, 45, 42) if key == "overlap_bbox_pt" else (255, 140, 0, 30)
                draw.rectangle([left, top, right, bottom], outline=color, width=line_width, fill=fill)
                slide_has_box = True

        Image.alpha_composite(image, overlay).convert("RGB").save(annotated_dir / before_path.name)


def _draw_overflow_measurement_render(
    before_dir: Path,
    findings_json: list[dict],
    measured_dir: Path,
) -> None:
    findings_by_slide: dict[int, list[tuple[dict[str, float], tuple[float, float, float, float], tuple[float, float]]]] = defaultdict(list)
    for finding in findings_json:
        detail = finding.get("detail") or {}
        if not isinstance(detail, dict):
            continue
        overflow_sides = detail.get("overflow_sides_pt")
        bbox = detail.get("bbox_pt")
        canvas = detail.get("canvas_pt") or [1440, 810]
        if isinstance(overflow_sides, dict) and overflow_sides and _valid_bbox(bbox):
            findings_by_slide[int(finding.get("slide_index") or 1)].append(
                (
                    {str(side): float(amount) for side, amount in overflow_sides.items()},
                    tuple(float(item) for item in bbox),
                    (float(canvas[0]), float(canvas[1])),
                )
            )

    if not findings_by_slide:
        return

    measured_dir.mkdir(parents=True, exist_ok=True)
    for before_path in sorted(before_dir.glob("slide-*.png")):
        match = re.search(r"slide-(\d+)\.png$", before_path.name)
        slide_index = int(match.group(1)) if match else 1
        entries = findings_by_slide.get(slide_index)
        if not entries:
            continue

        source = Image.open(before_path)
        canvas_w, canvas_h = entries[0][2]
        min_x = min([0.0] + [bbox[0] for _, bbox, _ in entries])
        min_y = min([0.0] + [bbox[1] for _, bbox, _ in entries])
        max_x = max([canvas_w] + [bbox[0] + bbox[2] for _, bbox, _ in entries])
        max_y = max([canvas_h] + [bbox[1] + bbox[3] for _, bbox, _ in entries])
        pad = 72
        scale = min((source.width - pad * 2) / (max_x - min_x), (source.height - pad * 2) / (max_y - min_y))
        diagram = Image.new("RGBA", source.size, (246, 246, 246, 255))
        draw = ImageDraw.Draw(diagram)

        def pt(value_x: float, value_y: float) -> tuple[int, int]:
            return (round((value_x - min_x) * scale + pad), round((value_y - min_y) * scale + pad))

        slide_left, slide_top = pt(0, 0)
        slide_right, slide_bottom = pt(canvas_w, canvas_h)
        draw.rectangle([slide_left, slide_top, slide_right, slide_bottom], fill=(255, 255, 255, 255), outline=(34, 34, 34, 255), width=4)

        for overflow_sides, bbox, _canvas in entries:
            x, y, w, h = bbox
            obj_left, obj_top = pt(x, y)
            obj_right, obj_bottom = pt(x + w, y + h)
            draw.rectangle([obj_left, obj_top, obj_right, obj_bottom], fill=(220, 0, 45, 38), outline=(220, 0, 45, 255), width=5)

            visible_left = max(x, 0)
            visible_top = max(y, 0)
            visible_right = min(x + w, canvas_w)
            visible_bottom = min(y + h, canvas_h)
            if visible_left < visible_right and visible_top < visible_bottom:
                v_left, v_top = pt(visible_left, visible_top)
                v_right, v_bottom = pt(visible_right, visible_bottom)
                draw.rectangle([v_left, v_top, v_right, v_bottom], fill=(255, 140, 0, 48), outline=(255, 140, 0, 255), width=4)

            marker_width = 12
            if "top" in overflow_sides:
                a_left, a_top = pt(max(x, 0), 0)
                a_right, _ = pt(min(x + w, canvas_w), 0)
                draw.rectangle([a_left, a_top - marker_width, a_right, a_top + marker_width], fill=(220, 0, 45, 255))
            if "bottom" in overflow_sides:
                a_left, a_bottom = pt(max(x, 0), canvas_h)
                a_right, _ = pt(min(x + w, canvas_w), canvas_h)
                draw.rectangle([a_left, a_bottom - marker_width, a_right, a_bottom + marker_width], fill=(220, 0, 45, 255))
            if "left" in overflow_sides:
                a_left, a_top = pt(0, max(y, 0))
                _, a_bottom = pt(0, min(y + h, canvas_h))
                draw.rectangle([a_left - marker_width, a_top, a_left + marker_width, a_bottom], fill=(220, 0, 45, 255))
            if "right" in overflow_sides:
                a_right, a_top = pt(canvas_w, max(y, 0))
                _, a_bottom = pt(canvas_w, min(y + h, canvas_h))
                draw.rectangle([a_right - marker_width, a_top, a_right + marker_width, a_bottom], fill=(220, 0, 45, 255))

        diagram.convert("RGB").save(measured_dir / before_path.name)


def _write_fixture_evidence_html(path: Path, rows: list[FixEvidenceRow], outdir: Path) -> None:
    html_dir = path.parent

    lines = [
        "<!doctype html>",
        '<html lang="ja">',
        "<head>",
        '<meta charset="utf-8">',
        '<meta name="viewport" content="width=device-width,initial-scale=1">',
        "<title>Pn-n Fixture Fix Evidence</title>",
        "<style>",
        "body{font-family:system-ui,-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;margin:24px;color:#222;background:#fafafa}",
        "h1{font-size:24px;margin:0 0 16px} h2{font-size:18px;margin:28px 0 8px}",
        "table{border-collapse:collapse;width:100%;background:white}td,th{border:1px solid #ddd;padding:6px 8px;text-align:left;font-size:13px}",
        ".compare{display:grid;grid-template-columns:repeat(3,minmax(0,1fr));gap:12px;margin:8px 0 28px}.panel{background:white;border:1px solid #ddd;padding:8px}",
        ".panel h3{font-size:14px;margin:0 0 8px}.state{background:#fff;border-left:4px solid #222;padding:10px;margin:8px 0;font-size:13px}.state b{display:block;margin-bottom:4px}",
        "img{width:100%;height:auto;display:block}.muted{color:#666}.fail{color:#b00020}.ok{color:#0b6b2b}.warn{color:#8a4b00}",
        "</style>",
        "</head><body>",
        "<h1>Pn-n Fixture Fix Evidence</h1>",
        "<p class='muted'>各Pn-nの intentionally bad fixture を生成し、検出箇所を枠付きで表示し、「現状」と「あるべき」を切り替えて確認できるようにしています。LibreOfficeは未使用です。</p>",
        "<p class='muted'>outcome は「自動修正済み」「テスト済み・手動判断が必要」「バグ: 機械修正候補だが fixer 未実装」などに分けています。manual 系も fixture/lint/render はテスト済みです。</p>",
        "<table><thead><tr><th>Pn</th><th>check / 意味</th><th>fixture</th><th>現状検出</th><th>auto後検出</th><th>applied</th><th>manual</th><th>outcome</th><th>reason</th><th>render</th></tr></thead><tbody>",
    ]
    for row in rows:
        cls = "ok" if row.render_status == "ok" else "fail"
        outcome_cls = "fail" if "mechanical_candidate" in row.outcome else "ok" if row.outcome == "autofixed" else "warn"
        lines.append(
            "<tr>"
            f"<td>{html.escape(row.pn)}</td>"
            f"<td>{_check_label_html(row)}</td>"
            f"<td>{html.escape(row.fixture)}</td>"
            f"<td>{row.before_count}</td><td>{row.after_count}</td>"
            f"<td>{row.applied_actions}</td><td>{row.manual_actions}</td>"
            f"<td class='{outcome_cls}'>{html.escape(_outcome_label(row))}</td>"
            f"<td>{html.escape(_outcome_reason(row))}</td>"
            f"<td class='{cls}'>{html.escape(row.render_status)}</td>"
            "</tr>"
        )
    lines.append("</tbody></table>")

    for row in rows:
        work = Path(row.work_dir)
        lines.append(f"<h2>{html.escape(row.pn)} <code>{html.escape(row.check)}</code> {html.escape(_check_title(row))}</h2>")
        lines.append(
            f"<p class='muted'>fixture={html.escape(row.fixture)} / "
            f"現状検出={row.before_count}, auto後検出={row.after_count}, "
            f"applied={row.applied_actions}, manual={row.manual_actions}, "
            f"outcome={html.escape(_outcome_label(row))}, render={html.escape(row.render_status)}</p>"
        )
        lines.append(f"<p>{html.escape(_check_explanation(row))}</p>")
        lines.append(f"<p class='muted'>分類: {html.escape(_outcome_reason(row))}</p>")
        marked_img = work / "before-marked" / "slide-01.png"
        measured_img = work / "before-measured" / "slide-01.png"
        expected_img = work / "expected" / "slide-01.png"
        diff_img = work / "expected-diff" / "slide-01.png"
        if marked_img.exists() and expected_img.exists() and diff_img.exists():
            marked_rel = marked_img.relative_to(html_dir).as_posix()
            expected_rel = expected_img.relative_to(html_dir).as_posix()
            diff_rel = diff_img.relative_to(html_dir).as_posix()
            measured_rel = measured_img.relative_to(html_dir).as_posix() if measured_img.exists() else ""
            lines.append("<div class='compare'>")
            lines.append("<div class='panel'>")
            lines.append("<h3>ASIS</h3>")
            lines.append(f"<img src='{html.escape(marked_rel)}'>")
            lines.append(f"<div class='state'><b>どうなっているか</b>{html.escape(row.issue_state)}</div>")
            lines.append("</div>")
            lines.append("<div class='panel'>")
            lines.append("<h3>TOBE</h3>")
            lines.append(f"<img src='{html.escape(expected_rel)}'>")
            lines.append(f"<div class='state'><b>機械修正するとどうなるか</b>{html.escape(row.expected_state)}</div>")
            lines.append("</div>")
            lines.append("<div class='panel'>")
            lines.append("<h3>DIFF / 説明</h3>")
            if measured_rel:
                lines.append(f"<img src='{html.escape(measured_rel)}'>")
                lines.append("<div class='state'><b>測定図</b>白=スライド境界、赤=PowerPoint上の実ボックス、橙=スライド内に実際に見えている部分。通常のスライド画像では、スライド外へ出た部分は描画されません。</div>")
            lines.append(f"<img src='{html.escape(diff_rel)}'>")
            lines.append(f"<div class='state'><b>差分の見方</b>ASIS の検出箇所を、TOBE の状態へ機械的に直した差分です。対象 check: {html.escape(row.check)}</div>")
            lines.append("</div>")
            lines.append("</div>")
        else:
            lines.append("<p class='fail'>render image missing</p>")
    lines.append("</body></html>")
    path.write_text("\n".join(lines), encoding="utf-8")


def write_fixture_fix_evidence(outdir: Path, catalog: list[CatalogItem]) -> list[FixEvidenceRow]:
    evidence_dir = outdir / "fixture-fix-evidence"
    evidence_dir.mkdir(parents=True, exist_ok=True)
    spec_by_check = _fixture_spec_by_check(fixture_specs())
    rows: list[FixEvidenceRow] = []
    for item in catalog:
        spec = spec_by_check.get(item.check)
        work = evidence_dir / f"{item.pn}-{_slug(item.check)}"
        work.mkdir(parents=True, exist_ok=True)
        before_pptx = work / "before.pptx"
        after_pptx = work / "after.pptx"
        expected_pptx = work / "expected.pptx"
        if spec is None:
            rows.append(
                FixEvidenceRow(
                    pn=item.pn,
                    priority=item.priority,
                    check=item.check,
                    fixture="",
                    fixture_status="missing",
                    before_count=0,
                    after_count=0,
                    applied_actions=0,
                    manual_actions=0,
                    outcome="fixture_not_detected",
                    outcome_reason="no intentional bad fixture is registered for this check",
                    render_status="missing_fixture",
                    diff_nonempty_slides=0,
                    work_dir=str(work),
                    before_pptx=str(before_pptx),
                    after_pptx=str(after_pptx),
                    report_href="",
                    issue_state="fixture がないため現状画像を作れない。",
                    expected_state=f"本来: {_check_action(FixEvidenceRow(item.pn, item.priority, item.check, '', 'missing', 0, 0, 0, 0, '', '', '', 0, '', '', '', ''))}",
                )
            )
            continue

        spec.builder(before_pptx, work)
        shutil.copy2(before_pptx, after_pptx)
        before_findings = _lint_fixture(before_pptx, work)
        before_count = _count_check(item, before_pptx, before_findings)
        before_findings_json = [
            pptx_lint.finding_to_json_dict(finding)
            for finding in before_findings
            if finding.check == item.check
        ]
        if item.check == "text_encoding":
            before_findings_json = _synthetic_text_encoding_findings(before_pptx)
        else:
            before_findings_json = _augment_findings_with_shape_bbox(before_pptx, before_findings_json)
        _apply_mechanical_expected_preview(item.check, before_pptx, expected_pptx, before_findings_json)
        rules = _fix_rules_for_check(item.check)
        actions = pptx_fix.fix_pptx(
            after_pptx,
            apply=True,
            rules=rules,
            findings=before_findings_json,
        )
        after_findings = _lint_fixture(after_pptx, work)
        after_count = _count_check(item, after_pptx, after_findings)
        applied = sum(1 for action in actions if action.status == "apply")
        manual = sum(1 for action in actions if action.status == "manual_required")
        if before_count > 0 and not rules:
            manual = before_count

        before_ok, before_render_output, before_dir = _render_pptx(before_pptx, work, "before")
        after_ok, after_render_output, after_dir = _render_pptx(after_pptx, work, "after")
        expected_ok, expected_render_output, _expected_dir = _render_pptx(expected_pptx, work, "expected")
        (work / "render.log").write_text(
            before_render_output
            + "\n--- after ---\n"
            + after_render_output
            + "\n--- expected ---\n"
            + expected_render_output,
            encoding="utf-8",
        )
        if before_ok and after_ok and expected_ok:
            diff_nonempty = _diff_renders(before_dir, after_dir, work / "diff")
            _diff_renders(before_dir, _expected_dir, work / "expected-diff")
            _annotate_before_render(
                before_dir,
                before_findings_json,
                work / "before-marked",
                fallback_full_slide=False,
            )
            _draw_overflow_measurement_render(before_dir, before_findings_json, work / "before-measured")
            render_status = "ok"
        else:
            diff_nonempty = 0
            render_status = "render_failed"
        outcome, outcome_reason = _fixture_fix_outcome(
            before_count=before_count,
            after_count=after_count,
            rules=rules,
            applied=applied,
            manual=manual,
            findings_json=before_findings_json,
            render_status=render_status,
        )

        rows.append(
            FixEvidenceRow(
                pn=item.pn,
                priority=item.priority,
                check=item.check,
                fixture=spec.name,
                fixture_status="pass" if before_count > 0 else "missing",
                before_count=before_count,
                after_count=after_count,
                applied_actions=applied,
                manual_actions=manual,
                outcome=outcome,
                outcome_reason=outcome_reason,
                render_status=render_status,
                diff_nonempty_slides=diff_nonempty,
                work_dir=str(work),
                before_pptx=str(before_pptx),
                after_pptx=str(after_pptx),
                report_href="fixture-fix-evidence.html",
                issue_state=_issue_state_for_check(item.check, before_findings_json, before_count),
                expected_state=_expected_state_for_check(
                    FixEvidenceRow(
                        pn=item.pn,
                        priority=item.priority,
                        check=item.check,
                        fixture=spec.name,
                        fixture_status="pass" if before_count > 0 else "missing",
                        before_count=before_count,
                        after_count=after_count,
                        applied_actions=applied,
                        manual_actions=manual,
                        outcome=outcome,
                        outcome_reason=outcome_reason,
                        render_status=render_status,
                        diff_nonempty_slides=diff_nonempty,
                        work_dir=str(work),
                        before_pptx=str(before_pptx),
                        after_pptx=str(after_pptx),
                        report_href="fixture-fix-evidence.html",
                    ),
                    before_findings_json,
                ),
            )
        )

    (outdir / "fixture-fix-evidence.json").write_text(
        json.dumps([asdict(row) for row in rows], ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    _write_fixture_evidence_html(outdir / "fixture-fix-evidence.html", rows, evidence_dir)
    return rows


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--tasks", type=Path, default=ROOT / "doc/tasks.md")
    ap.add_argument("--before-lint", type=Path)
    ap.add_argument("--after-lint", type=Path)
    ap.add_argument("--actions-json", type=Path)
    ap.add_argument("--outdir", type=Path, required=True)
    ap.add_argument(
        "--fixture-fix-evidence",
        action="store_true",
        help="generate per-Pn fixture auto-fix before/after/diff evidence with vscode-pptx-viewer",
    )
    args = ap.parse_args()

    catalog = load_catalog(args.tasks)
    fixture_dir = args.outdir / "fixtures"
    fixture_results = [run_fixture(spec, fixture_dir) for spec in fixture_specs()]
    rows = build_rows(
        catalog,
        load_json(args.before_lint, []),
        load_json(args.after_lint, []),
        load_json(args.actions_json, {"actions": []}),
        fixture_results,
    )

    args.outdir.mkdir(parents=True, exist_ok=True)
    write_tsv(args.outdir / "pn-review-orchestrator-matrix.tsv", rows)
    write_json(args.outdir / "pn-review-orchestrator-report.json", rows, fixture_results)
    markdown_path = args.outdir / "pn-review-orchestrator-report.md"
    write_markdown(markdown_path, rows, fixture_results)
    write_html(args.outdir / "index.html", markdown_path)
    if args.fixture_fix_evidence:
        evidence_rows = write_fixture_fix_evidence(args.outdir, catalog)
        print(args.outdir / "fixture-fix-evidence.html")
        print(args.outdir / "fixture-fix-evidence.json")
        print(f"fixture_fix_evidence_rows={len(evidence_rows)}")
    print(markdown_path)
    print(args.outdir / "index.html")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
