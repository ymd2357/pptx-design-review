#!/usr/bin/env python3
"""PPTX auto-fixer for safe mechanical violations of slide-guideline-v1.

Fixes (mechanical, no semantic decisions):
- autofit   text frame auto-size != NONE -> set to MSO_AUTO_SIZE.NONE
- geometry  shape left/top/width/height in EMU rounded to nearest 1pt,
            but only when the current value is within 0.1pt of an integer
            (catches float drift; preserves intentional sub-pt placements)
- contrast  text run color replaced with the lint-provided candidate color
            only when --rules contrast and --findings-json are supplied
- font_size text run size snapped to the nearest allowed scale size only when
            explicitly requested and all safety checks pass; otherwise reported
            as manual_required
- line_height fixed paragraph line height snapped to the nearest allowed scale
            only when explicitly requested and fit checks pass
- alignment  paragraph alignment set to LEFT and text-frame vertical anchor set
            to TOP when explicitly requested

Out of fixer scope (require human judgment):
- font_family, overflow, safe_text_area, animation_present,
  alt_text_required, reading_order, semantic color cues, card-grid composition

After --apply, the script re-reads the saved file and re-detects pending
actions; any residual means the change was not durable on disk. Known
causes include corrupted source PPTX (duplicate zip entries -- watch for
"Duplicate name: ppt/slides/slideN.xml" from zipfile) and inherited bodyPr
that python-pptx cannot override at slide level. Residuals are printed as
a warning and exit code becomes 2.

Usage
    python3 pptx_fix.py DECK.pptx                     # dry-run, prints plan
    python3 pptx_fix.py DECK.pptx --apply             # write in-place
    python3 pptx_fix.py DECK.pptx --auto --apply      # apply every lint-declared auto-fix candidate
    python3 pptx_fix.py DECK.pptx --apply --backup    # write DECK.pptx.bak if absent
    python3 pptx_fix.py DECK.pptx --apply --rules autofit
    python3 pptx_fix.py DECK.pptx --apply --rules contrast --findings-json lint.json
    python3 pptx_fix.py DECK.pptx --rules font_size
    python3 pptx_fix.py DECK.pptx --rules line_height,alignment
    python3 pptx_fix.py DECK.pptx --json

Exit code
    0 = success (or nothing to fix)
    1 = invocation error (missing file, unknown rule)
    2 = applied but self-check found residual actions on the saved file
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Callable, Iterable, List, Optional, Sequence

from pptx import Presentation
from pptx.dml.color import MSO_COLOR_TYPE, RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt


EMU_PER_PT = 12700
GEOMETRY_ROUND_TOL_PT = 0.1  # only fix when drift is well below half-pt
SLIDE_W_PT = 1440
SLIDE_H_PT = 810
ALLOWED_FONT_SIZES_PT = {80, 64, 56, 48, 40, 36, 32, 28, 24, 22, 20}
ALLOWED_LINE_HEIGHTS_PT = {90, 66, 42, 36, 30, 24}
FONT_SIZE_FIX_DELTA_MAX_PT = 1.0
LINE_HEIGHT_FIX_DELTA_MAX_PT = 4.0
FONT_SIZE_FIXER_ENABLED = True

DEFAULT_RULES = ("autofit", "geometry")
ALL_RULES = (
    "autofit",
    "geometry",
    "font_size",
    "line_height",
    "alignment",
    "contrast",
    "font_family",
    "text_color",
    "fill_color",
    "bbox_fit",
    "overlap",
    "spacing",
    "image_crop",
    "image_aspect",
    "reading_order",
    "animation",
    "text_wrap",
    "heading_hierarchy",
    "inner_padding",
    "card_grid",
    "text_vertical_balance",
    "badge_alignment",
    "decorative_remove",
    "box_canvas_clip",
    "text_box_resize",
    "text_canvas_reflow",
)


def _load_fix_policy() -> dict[str, dict]:
    """Load rules.lint.fix_policy from doc/slide-guideline-v1.yml.

    Single source of truth (POLICY-001). Returns the full {check: {apply_mode,
    fix_rule, ...}} mapping, or an empty dict when PyYAML is unavailable.
    Callers must tolerate the empty case and fall back to defaults.
    """
    guideline = Path(__file__).resolve().parents[3] / "doc" / "slide-guideline-v1.yml"
    try:
        import yaml  # type: ignore

        data = yaml.safe_load(guideline.read_text(encoding="utf-8"))
        return (data.get("rules") or {}).get("lint", {}).get("fix_policy") or {}
    except Exception:
        return {}


def _load_check_to_rule() -> dict[str, str]:
    """Derive CHECK_TO_RULE from rules.lint.fix_policy.

    Only checks with apply_mode in {auto_fix, judgement_fix} expose a fix_rule;
    apply_mode=no_fix checks intentionally do not appear in CHECK_TO_RULE.
    Falls back to a minimal text parser so pptx_fix.py can still import in
    environments without PyYAML.
    """
    fix_policy = _load_fix_policy()
    if fix_policy:
        mapping: dict[str, str] = {}
        for check, entry in fix_policy.items():
            if not isinstance(entry, dict):
                continue
            if entry.get("apply_mode") in {"auto_fix", "judgement_fix"}:
                fix_rule = entry.get("fix_rule")
                if isinstance(fix_rule, str):
                    mapping[check] = fix_rule
        return mapping
    guideline = Path(__file__).resolve().parents[3] / "doc" / "slide-guideline-v1.yml"
    return _parse_check_to_rule_textually(guideline)


def _parse_check_to_rule_textually(guideline: Path) -> dict[str, str]:
    """Minimal text-mode reader for rules.lint.fix_policy used when PyYAML
    is not installed. Only extracts {check: fix_rule} pairs.
    """
    mapping: dict[str, str] = {}
    in_section = False
    current_check: str | None = None
    current_mode: str | None = None
    current_rule: str | None = None

    def flush() -> None:
        if current_check and current_mode in {"auto_fix", "judgement_fix"} and current_rule:
            mapping[current_check] = current_rule

    for raw in guideline.read_text(encoding="utf-8").splitlines():
        if raw.lstrip().startswith("#"):
            continue
        if raw == "    fix_policy:":
            in_section = True
            continue
        if not in_section:
            continue
        if raw.startswith("    ") and not raw.startswith("      "):
            # Left the fix_policy block.
            flush()
            current_check = current_mode = current_rule = None
            break
        if raw.startswith("      ") and not raw.startswith("        ") and raw.rstrip().endswith(":"):
            flush()
            current_check = raw.strip().rstrip(":")
            current_mode = current_rule = None
            continue
        if current_check and raw.startswith("        "):
            stripped = raw.strip()
            if stripped.startswith("apply_mode:"):
                current_mode = stripped.split(":", 1)[1].strip().strip('"').strip("'")
            elif stripped.startswith("fix_rule:"):
                value = stripped.split(":", 1)[1].strip().strip('"').strip("'")
                current_rule = None if value in {"null", "~", ""} else value
    flush()
    return mapping


FIX_POLICY = _load_fix_policy()
CHECK_TO_RULE = _load_check_to_rule()


def _apply_mode_for_check(check: str | None) -> str | None:
    """Lookup apply_mode (auto_fix / judgement_fix / no_fix) for a check, or
    None when the check is not declared in rules.lint.fix_policy.
    """
    if not check:
        return None
    entry = FIX_POLICY.get(check)
    if not isinstance(entry, dict):
        return None
    return entry.get("apply_mode")


@dataclass
class FixAction:
    rule: str
    slide_index: int  # 1-based
    slide_id: Optional[int]
    shape_id: Optional[int]
    shape_name: Optional[str]
    status: str = "apply"  # "apply" | "manual_required"
    reasons: list[str] = field(default_factory=list)
    before: dict = field(default_factory=dict)
    after: dict = field(default_factory=dict)


# ---- Shape traversal -------------------------------------------------------


def _iter_shapes(shapes) -> Iterable:
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(s.shapes)
        else:
            yield s


def _walk(prs):
    for idx, slide in enumerate(prs.slides, start=1):
        slide_id = getattr(slide, "slide_id", None)
        for shape in _iter_shapes(slide.shapes):
            yield shape, idx, slide_id, slide


# ---- Per-rule detectors (pure: no mutation) --------------------------------


def _detect_autofit(shape, slide_idx, slide_id, slide=None) -> Optional[FixAction]:
    if not shape.has_text_frame:
        return None
    af = shape.text_frame.auto_size
    if af is None or af == MSO_AUTO_SIZE.NONE:
        return None
    return FixAction(
        rule="autofit",
        slide_index=slide_idx,
        slide_id=slide_id,
        shape_id=getattr(shape, "shape_id", None),
        shape_name=getattr(shape, "name", None),
        before={"auto_size": str(af)},
        after={"auto_size": "NONE"},
    )


def _round_emu_to_pt(value_emu: int) -> Optional[int]:
    """Return EMU rounded to nearest 1pt if drift < 0.1pt, else None."""
    pt = value_emu / EMU_PER_PT
    nearest = round(pt)
    if abs(pt - nearest) >= GEOMETRY_ROUND_TOL_PT:
        return None
    new_emu = nearest * EMU_PER_PT
    if new_emu == value_emu:
        return None
    return new_emu


def _detect_geometry(shape, slide_idx, slide_id, slide=None) -> Optional[FixAction]:
    coords = ("left", "top", "width", "height")
    raw = {c: getattr(shape, c) for c in coords}
    if any(v is None for v in raw.values()):
        return None
    new_vals = {}
    for c, v in raw.items():
        nv = _round_emu_to_pt(v)
        if nv is not None:
            new_vals[c] = nv
    if not new_vals:
        return None
    before = {c: round(raw[c] / EMU_PER_PT, 4) for c in new_vals}
    after = {c: new_vals[c] / EMU_PER_PT for c in new_vals}
    return FixAction(
        rule="geometry",
        slide_index=slide_idx,
        slide_id=slide_id,
        shape_id=getattr(shape, "shape_id", None),
        shape_name=getattr(shape, "name", None),
        before=before,
        after=after,
    )


def _nearest_allowed_font_size(size_pt: float) -> int:
    return min(ALLOWED_FONT_SIZES_PT, key=lambda allowed: abs(size_pt - allowed))


def _nearest_allowed_line_height(size_pt: float) -> int:
    return min(ALLOWED_LINE_HEIGHTS_PT, key=lambda allowed: abs(size_pt - allowed))


def _slide_scale(slide) -> float:
    prs = slide.part.package.presentation_part.presentation
    w_pt = prs.slide_width / EMU_PER_PT
    h_pt = prs.slide_height / EMU_PER_PT
    scale_x = SLIDE_W_PT / w_pt if w_pt else 1.0
    scale_y = SLIDE_H_PT / h_pt if h_pt else 1.0
    return (scale_x + scale_y) / 2


def _bbox_pt(shape) -> tuple[float, float, float, float]:
    return (
        shape.left / EMU_PER_PT,
        shape.top / EMU_PER_PT,
        shape.width / EMU_PER_PT,
        shape.height / EMU_PER_PT,
    )


def _overlaps(a: tuple[float, float, float, float], b: tuple[float, float, float, float]) -> bool:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0, min(ay + ah, by + bh) - max(ay, by))
    return ix * iy > 1.0


def _slide_by_index(prs, slide_index: Any):
    try:
        idx = int(slide_index)
    except (TypeError, ValueError):
        return None
    if idx < 1 or idx > len(prs.slides):
        return None
    return prs.slides[idx - 1]


def _shape_by_id(slide, shape_id: Any):
    try:
        target_id = int(shape_id)
    except (TypeError, ValueError):
        return None
    for shape in _iter_shapes(slide.shapes):
        if getattr(shape, "shape_id", None) == target_id:
            return shape
    return None


def _detail_shape_id(value: Any) -> Optional[int]:
    if isinstance(value, dict):
        value = value.get("shape_id")
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _slide_scale_xy(prs) -> tuple[float, float]:
    actual_w = prs.slide_width / EMU_PER_PT
    actual_h = prs.slide_height / EMU_PER_PT
    return actual_w / SLIDE_W_PT, actual_h / SLIDE_H_PT


def _norm_to_actual_bbox(prs, bbox: Sequence[float]) -> tuple[float, float, float, float]:
    sx, sy = _slide_scale_xy(prs)
    return (bbox[0] * sx, bbox[1] * sy, bbox[2] * sx, bbox[3] * sy)


def _valid_bbox(value: Any) -> bool:
    return (
        isinstance(value, Sequence)
        and not isinstance(value, (str, bytes))
        and len(value) == 4
        and all(isinstance(item, (int, float)) for item in value)
        and value[2] > 0
        and value[3] > 0
    )


def _card_grid_group_bbox(children: Sequence[Any]) -> Optional[tuple[float, float, float, float]]:
    boxes = [
        child.get("bbox_pt")
        for child in children
        if isinstance(child, dict) and _valid_bbox(child.get("bbox_pt"))
    ]
    if not boxes:
        return None
    left = min(float(box[0]) for box in boxes)
    top = min(float(box[1]) for box in boxes)
    right = max(float(box[0]) + float(box[2]) for box in boxes)
    bottom = max(float(box[1]) + float(box[3]) for box in boxes)
    return (left, top, right - left, bottom - top)


def _card_classify_role(child: dict, container: dict) -> str:
    """Classify a card child into one of the FIX-012 role categories.

    Returns one of: `leading_icon`, `header`, `body`, `badge_bg`, `badge_text`,
    `illustration`, `decorative`, `unknown`. The classifier intentionally
    refuses to guess when evidence is mixed (= `unknown`), so the alignment
    solver can skip rather than touch.

    Heuristics use relative position within the container (= card) plus the
    extended attributes the lint side now emits (`solid_fill_hex`,
    `vertical_anchor`, `has_text`, `text_length`).
    """
    kind = (child.get("kind") or "").lower()
    bbox = child.get("bbox_pt") or []
    container_bbox = container.get("bbox_pt") or []
    if len(bbox) != 4 or len(container_bbox) != 4:
        return "unknown"
    x, y, w, h = [float(v) for v in bbox]
    cx, cy, cw, ch = [float(v) for v in container_bbox]
    if cw <= 0 or ch <= 0:
        return "unknown"
    rel_top = (y - cy) / ch
    rel_h = h / ch
    rel_w = w / cw
    has_fill = bool(child.get("solid_fill_hex"))
    has_text = bool(child.get("has_text"))
    anchor = (child.get("vertical_anchor") or "").upper()

    if kind == "image":
        if rel_h < 0.25 and rel_w < 0.3 and rel_top < 0.55:
            return "leading_icon"
        return "illustration"

    if kind == "shape":
        if has_fill and rel_top >= 0.45 and rel_h < 0.3:
            return "badge_bg"
        return "decorative"

    if kind == "text" and has_text:
        # Badge text: anchored MIDDLE in lower portion of card.
        if anchor == "MIDDLE" and rel_top >= 0.45 and rel_h < 0.3:
            return "badge_text"
        # Zone-based classification by the child's TOP (not height): the row
        # may have headers of varying line counts whose bbox.h differs, but
        # they all start at the same rel_top within the card.
        if rel_top < 0.25:
            return "header"
        if rel_top < 0.6:
            return "body"
        return "decorative"

    return "unknown"


def _card_pair_badge_shapes(roles_by_index: dict[int, str], children: list[dict]) -> dict[int, str]:
    """If `badge_bg` and `badge_text` share an identical bbox, treat both as a
    coherent `badge` group. This avoids splitting the pair across the solver.
    Returns a new role map with paired entries reset to role=`badge`.
    """
    out = dict(roles_by_index)
    for i, ci in enumerate(children):
        if out.get(i) != "badge_bg":
            continue
        for j, cj in enumerate(children):
            if i == j or out.get(j) != "badge_text":
                continue
            if ci.get("bbox_pt") and ci["bbox_pt"] == cj.get("bbox_pt"):
                out[i] = "badge"
                out[j] = "badge"
    return out


def _card_role_map(container: dict) -> dict[int, str]:
    """Build child-index → role map for a single card container record."""
    children = container.get("children") or []
    raw = {i: _card_classify_role(child, container) for i, child in enumerate(children)}
    return _card_pair_badge_shapes(raw, children)


def _card_grid_alignment_moves(row_containers: list[dict]) -> dict:
    """Solve cross-card alignment targets and emit per-card-per-child moves.

    Returns a dict keyed by `(card_index, child_index)` whose value is a
    `(target_left, target_top)` tuple in guideline pt. Only children that
    actually need to move are included.

    Rules (anchored to observation, see doc/tasks.md FIX-012):
      - `leading_icon`: icon.center_y = header.first_line.center_y of the
        SAME card (Option B from the FIX-012 design discussion). When the
        card has no header, no move is emitted for icon.
      - `header`: not moved.
      - `body`: body.top = max(card.top + padding + header.h) across the row
        (= 上揃え driven by tallest header).
      - `badge` / `badge_bg` / `badge_text`: badge.bottom = card.bottom
        − min(observed badge padding_bottom across the row). Tightest
        observed padding becomes the row target so under-spaced badges
        (e.g. slide 13 card 1) move down without crowding overflowing
        bodies in other cards.
      - `illustration` / `decorative` / `unknown`: not moved.
    """
    moves: dict[tuple[int, int], tuple[float, float]] = {}
    if not row_containers:
        return moves

    role_maps: list[dict[int, str]] = []
    for container in row_containers:
        if not isinstance(container, dict):
            role_maps.append({})
        else:
            role_maps.append(_card_role_map(container))

    # Header height per card (used by body 上揃え).
    header_heights: list[float] = []
    for container, role_map in zip(row_containers, role_maps):
        children = container.get("children") or [] if isinstance(container, dict) else []
        header_h = 0.0
        for idx, role in role_map.items():
            if role != "header":
                continue
            bbox = children[idx].get("bbox_pt") or []
            if len(bbox) == 4:
                header_h = max(header_h, float(bbox[3]))
        header_heights.append(header_h)

    max_header_h = max(header_heights) if header_heights else 0.0

    # Header first-line center per card (= card.top + padding_top + line_h / 2,
    # where line_h ≈ header bbox height / line count; for our deck headers are
    # top-aligned so line_h equals the smallest header height in the row
    # whenever the row has a 1-line header. Fallback: header.h / 2).
    one_line_header_h = (
        min(h for h in header_heights if h > 0) if any(h > 0 for h in header_heights) else 0.0
    )

    def _header_first_line_center_y(container, role_map):
        children = container.get("children") or []
        for idx, role in role_map.items():
            if role != "header":
                continue
            bbox = children[idx].get("bbox_pt") or []
            if len(bbox) != 4:
                continue
            top = float(bbox[1])
            line_h = one_line_header_h or float(bbox[3])
            return top + line_h / 2.0
        return None

    # body.top target (上揃え): card_with_max_header's natural body.top.
    body_top_target: float | None = None
    for container, role_map in zip(row_containers, role_maps):
        children = container.get("children") or [] if isinstance(container, dict) else []
        for idx, role in role_map.items():
            if role != "body":
                continue
            bbox = children[idx].get("bbox_pt") or []
            if len(bbox) != 4:
                continue
            header_h_here = 0.0
            for jdx, r2 in role_map.items():
                if r2 == "header":
                    hb = children[jdx].get("bbox_pt") or []
                    if len(hb) == 4:
                        header_h_here = max(header_h_here, float(hb[3]))
            if abs(header_h_here - max_header_h) < 0.5:
                # This card's body is at the natural lowest position.
                body_top_target = float(bbox[1])
                break
        if body_top_target is not None:
            break

    # badge target: card.bottom − min(observed badge padding_bottom).
    badge_paddings: list[float] = []
    for container, role_map in zip(row_containers, role_maps):
        if not isinstance(container, dict):
            continue
        cbbox = container.get("bbox_pt") or []
        if len(cbbox) != 4:
            continue
        card_bottom = float(cbbox[1]) + float(cbbox[3])
        children = container.get("children") or []
        badge_bottoms: list[float] = []
        for idx, role in role_map.items():
            if role not in {"badge", "badge_bg", "badge_text"}:
                continue
            bbox = children[idx].get("bbox_pt") or []
            if len(bbox) == 4:
                badge_bottoms.append(float(bbox[1]) + float(bbox[3]))
        if badge_bottoms:
            badge_paddings.append(card_bottom - max(badge_bottoms))
    badge_padding_target = min(badge_paddings) if badge_paddings else None

    # Apply moves per card.
    for card_idx, (container, role_map) in enumerate(zip(row_containers, role_maps)):
        if not isinstance(container, dict):
            continue
        cbbox = container.get("bbox_pt") or []
        if len(cbbox) != 4:
            continue
        card_top = float(cbbox[1])
        card_bottom = card_top + float(cbbox[3])
        children = container.get("children") or []
        header_center = _header_first_line_center_y(container, role_map)

        for idx, role in role_map.items():
            bbox = children[idx].get("bbox_pt") or []
            if len(bbox) != 4:
                continue
            cur_left = float(bbox[0])
            cur_top = float(bbox[1])
            cur_h = float(bbox[3])

            target_top = cur_top

            if role == "leading_icon" and header_center is not None:
                target_top = header_center - cur_h / 2.0
            elif role == "body" and body_top_target is not None:
                target_top = body_top_target
            elif role in {"badge", "badge_bg", "badge_text"} and badge_padding_target is not None:
                target_top = card_bottom - badge_padding_target - cur_h
            else:
                continue

            if abs(target_top - cur_top) > 0.5:
                moves[(card_idx, idx)] = (cur_left, target_top)

    return moves


def _card_grid_item_manual_reasons(item: dict, medians: dict) -> list[str]:
    container = item.get("container") or {}
    container_bbox = container.get("bbox_pt") or []
    if not _valid_bbox(container_bbox):
        return ["invalid_card_grid_container_bbox"]
    child_group = _card_grid_group_bbox(item.get("children") or [])
    if child_group is None:
        return ["missing_card_grid_child_bbox"]

    target_width = float(medians.get("width", container_bbox[2]))
    target_height = float(medians.get("height", container_bbox[3]))
    padding = item.get("padding_pt") if isinstance(item.get("padding_pt"), dict) else {}
    padding_left = float(medians.get("padding_left", padding.get("left", 0)))
    padding_right = float(medians.get("padding_right", padding.get("right", 0)))
    padding_top = float(medians.get("padding_top", padding.get("top", 0)))
    padding_bottom = float(medians.get("padding_bottom", padding.get("bottom", 0)))
    target_inner_width = target_width - padding_left - padding_right
    target_inner_height = target_height - padding_top - padding_bottom
    if child_group[2] > target_inner_width + 1.0 or child_group[3] > target_inner_height + 1.0:
        container_id = container.get("shape_id")
        return [
            (
                "card_grid_child_group_exceeds_inner_box"
                f": container_id={container_id}"
                f", child_group={child_group[2]:.2f}x{child_group[3]:.2f}pt"
                f", inner={target_inner_width:.2f}x{target_inner_height:.2f}pt"
            )
        ]
    return []


def _shape_bbox_from_detail(prs, detail: dict) -> Optional[tuple[float, float, float, float]]:
    actual = detail.get("actual_bbox_pt")
    if _valid_bbox(actual):
        return tuple(float(item) for item in actual)
    bbox = detail.get("bbox_pt")
    if _valid_bbox(bbox):
        return _norm_to_actual_bbox(prs, [float(item) for item in bbox])
    return None


def _shape_from_finding(prs, finding: Any, nested_key: str | None = None):
    slide = _slide_by_index(prs, _finding_field(finding, "slide_index"))
    if slide is None:
        return None
    detail = _finding_detail(finding)
    if nested_key:
        nested = detail.get(nested_key)
        shape = _shape_by_id(slide, _detail_shape_id(nested))
        if shape is not None:
            return shape
    target = detail.get("target")
    shape = _shape_by_id(slide, _detail_shape_id(target))
    if shape is not None:
        return shape
    shape = _shape_by_id(slide, _finding_field(finding, "shape_id"))
    if shape is not None:
        return shape
    nested = detail.get("shape") if isinstance(detail.get("shape"), dict) else None
    return _shape_by_id(slide, _detail_shape_id(nested))


def _shape_geometry_pt(shape) -> dict:
    """shape の **絶対座標 (slide canvas 上の pt)** を返す。group 内 shape も
    親 group の transform を適用した絶対値で返すので、action.after.geometry
    の入出力単位が一貫する (= _apply_geometry が逆変換で raw に戻す)。"""
    # _absolute_pt_to_raw の逆 = group transform 適用
    EMU_PER_PT_LOCAL = 12700
    left = shape.left / EMU_PER_PT_LOCAL
    top = shape.top / EMU_PER_PT_LOCAL
    width = shape.width / EMU_PER_PT_LOCAL
    height = shape.height / EMU_PER_PT_LOCAL
    parent = shape._element.getparent()
    while parent is not None:
        tag = parent.tag.split("}", 1)[-1] if "}" in parent.tag else parent.tag
        if tag != "grpSp":
            break
        gxfrm = parent.find(f".//{_DML_NS_FIX}xfrm")
        if gxfrm is not None:
            goff = gxfrm.find(f"{_DML_NS_FIX}off")
            gchoff = gxfrm.find(f"{_DML_NS_FIX}chOff")
            gext = gxfrm.find(f"{_DML_NS_FIX}ext")
            gchext = gxfrm.find(f"{_DML_NS_FIX}chExt")
            if goff is not None and gchoff is not None:
                gox = float(goff.get("x", "0")) / EMU_PER_PT_LOCAL
                goy = float(goff.get("y", "0")) / EMU_PER_PT_LOCAL
                gcx = float(gchoff.get("x", "0")) / EMU_PER_PT_LOCAL
                gcy = float(gchoff.get("y", "0")) / EMU_PER_PT_LOCAL
                sx = sy = 1.0
                if gext is not None and gchext is not None:
                    ex = float(gext.get("cx", "0")) / EMU_PER_PT_LOCAL
                    ey = float(gext.get("cy", "0")) / EMU_PER_PT_LOCAL
                    cex = float(gchext.get("cx", "1") or "1") / EMU_PER_PT_LOCAL
                    cey = float(gchext.get("cy", "1") or "1") / EMU_PER_PT_LOCAL
                    if cex:
                        sx = ex / cex
                    if cey:
                        sy = ey / cey
                left = gox + (left - gcx) * sx
                top = goy + (top - gcy) * sy
                width = width * sx
                height = height * sy
        parent = parent.getparent()
    return {
        "left": round(left, 4),
        "top": round(top, 4),
        "width": round(width, 4),
        "height": round(height, 4),
    }


def _geometry_action(rule: str, finding: Any, shape, geometry: dict, reasons: list[str] | None = None) -> FixAction:
    return FixAction(
        rule=rule,
        slide_index=int(_finding_field(finding, "slide_index") or 1),
        slide_id=_finding_field(finding, "slide_id"),
        shape_id=getattr(shape, "shape_id", None),
        shape_name=getattr(shape, "name", None),
        status="apply",
        reasons=reasons or [],
        before=_shape_geometry_pt(shape),
        after={"geometry": geometry},
    )


_DML_NS_FIX = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _absolute_pt_to_raw(shape, abs_left, abs_top, abs_width, abs_height):
    """絶対座標 (slide canvas 上の pt) を shape の raw 座標 (group local) に
    逆変換する。group 内 shape の場合、親 group の xfrm 逆変換を適用する。"""
    EMU_PER_PT_LOCAL = 12700
    left = float(abs_left) if abs_left is not None else shape.left / EMU_PER_PT_LOCAL
    top = float(abs_top) if abs_top is not None else shape.top / EMU_PER_PT_LOCAL
    width = float(abs_width) if abs_width is not None else shape.width / EMU_PER_PT_LOCAL
    height = float(abs_height) if abs_height is not None else shape.height / EMU_PER_PT_LOCAL
    parent = shape._element.getparent()
    while parent is not None:
        tag = parent.tag.split("}", 1)[-1] if "}" in parent.tag else parent.tag
        if tag != "grpSp":
            break
        gxfrm = parent.find(f".//{_DML_NS_FIX}xfrm")
        if gxfrm is not None:
            goff = gxfrm.find(f"{_DML_NS_FIX}off")
            gchoff = gxfrm.find(f"{_DML_NS_FIX}chOff")
            gext = gxfrm.find(f"{_DML_NS_FIX}ext")
            gchext = gxfrm.find(f"{_DML_NS_FIX}chExt")
            if goff is not None and gchoff is not None:
                gox = float(goff.get("x", "0")) / EMU_PER_PT_LOCAL
                goy = float(goff.get("y", "0")) / EMU_PER_PT_LOCAL
                gcx = float(gchoff.get("x", "0")) / EMU_PER_PT_LOCAL
                gcy = float(gchoff.get("y", "0")) / EMU_PER_PT_LOCAL
                sx = sy = 1.0
                if gext is not None and gchext is not None:
                    ex = float(gext.get("cx", "0")) / EMU_PER_PT_LOCAL
                    ey = float(gext.get("cy", "0")) / EMU_PER_PT_LOCAL
                    cex = float(gchext.get("cx", "1") or "1") / EMU_PER_PT_LOCAL
                    cey = float(gchext.get("cy", "1") or "1") / EMU_PER_PT_LOCAL
                    if cex:
                        sx = ex / cex
                    if cey:
                        sy = ey / cey
                # 逆変換: child_raw = (abs - group.off) / scale + group.chOff
                if sx:
                    left = (left - gox) / sx + gcx
                    width = width / sx
                if sy:
                    top = (top - goy) / sy + gcy
                    height = height / sy
        parent = parent.getparent()
    return left, top, width, height


def _apply_geometry(shape, geometry: dict) -> None:
    # action.after.geometry は **絶対座標 (slide canvas 上の pt)** を想定。
    # group 内 shape では raw 値に逆変換して set する。
    raw_left, raw_top, raw_w, raw_h = _absolute_pt_to_raw(
        shape,
        geometry.get("left"),
        geometry.get("top"),
        geometry.get("width"),
        geometry.get("height"),
    )
    if "left" in geometry:
        shape.left = Pt(raw_left)
    if "top" in geometry:
        shape.top = Pt(raw_top)
    if "width" in geometry:
        shape.width = Pt(raw_w)
    if "height" in geometry:
        shape.height = Pt(raw_h)


def _shape_text_frames(shape):
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame
    if getattr(shape, "has_table", False):
        for row_idx in range(len(shape.table.rows)):
            for col_idx in range(len(shape.table.columns)):
                yield shape.table.cell(row_idx, col_idx).text_frame


def _set_text_shape_font_family(shape, font_name: str) -> int:
    count = 0
    for tf in _shape_text_frames(shape):
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text:
                    run.font.name = font_name
                    count += 1
    return count


def _set_text_shape_color(shape, before_hex: Optional[str], after_hex: str) -> int:
    count = 0
    normalized_before = _normalize_hex(before_hex) if before_hex else None
    for tf in _shape_text_frames(shape):
        for para in tf.paragraphs:
            for run in para.runs:
                if not run.text:
                    continue
                current = _run_rgb_hex(run)
                if normalized_before and current != normalized_before:
                    continue
                run.font.color.rgb = RGBColor.from_string(after_hex.lstrip("#"))
                count += 1
    return count


def _candidate_hex(candidate: Any, *keys: str) -> Optional[str]:
    if not isinstance(candidate, dict):
        return None
    for key in keys:
        value = _normalize_hex(candidate.get(key))
        if value:
            return value
    return None


def _fit_bbox_into_area(
    current: dict,
    area: tuple[float, float, float, float],
) -> dict:
    left, top, width, height = current["left"], current["top"], current["width"], current["height"]
    area_left, area_top, area_width, area_height = area
    width = min(width, area_width)
    height = min(height, area_height)
    if left < area_left:
        left = area_left
    if top < area_top:
        top = area_top
    if left + width > area_left + area_width:
        left = area_left + area_width - width
    if top + height > area_top + area_height:
        top = area_top + area_height - height
    return {
        "left": round(left, 4),
        "top": round(top, 4),
        "width": round(width, 4),
        "height": round(height, 4),
    }


def _overlap_delta(
    a: Sequence[float],
    b: Sequence[float],
    *,
    pad: float = 4.0,
) -> tuple[float, float]:
    ax, ay, aw, ah = [float(v) for v in a]
    bx, by, bw, bh = [float(v) for v in b]
    right = (ax + aw) - bx + pad
    left = -((bx + bw) - ax + pad)
    down = (ay + ah) - by + pad
    up = -((by + bh) - ay + pad)
    candidates = [(right, 0.0), (left, 0.0), (0.0, down), (0.0, up)]
    return min(candidates, key=lambda item: abs(item[0]) + abs(item[1]))


FINDING_DRIVEN_RULES = {
    "font_family",
    "text_color",
    "fill_color",
    "bbox_fit",
    "overlap",
    "spacing",
    "image_crop",
    "image_aspect",
    "reading_order",
    "animation",
    "text_wrap",
    "heading_hierarchy",
    "inner_padding",
    "card_grid",
    "badge_alignment",
    "text_vertical_balance",
    "decorative_remove",
    "box_canvas_clip",
    "text_box_resize",
    "text_canvas_reflow",
}


_OVERFLOW_ALLOWED_FONT_SIZES_PT = (10, 11, 12, 14, 16, 18, 20, 22, 24, 28, 32, 36, 40, 44, 48, 54, 60, 72)


def _compound_resolve_text_box_overflow(shape, slide_height_pt: float = 810.0) -> None:
    """expand_box_width_to_canvas で box.width を拡張した直後に呼ぶ。優先順位:
    1) box.height を canvas 下端まで拡張 ([[feedback-overflow-fix-priority]])
    2) それでも収まらなければ font を 1 段階ずつ shrink (最終手段)
    font が explicit に取れない / 元から canvas に収まっている場合は no-op。"""
    import math
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    def _margin(attr: str) -> float:
        v = getattr(tf, attr, None)
        if v is None:
            return 7.2
        try:
            return float(v.pt)
        except (AttributeError, TypeError, ValueError):
            return 7.2
    EMU_PER_PT_LOCAL = 12700
    margin_l = _margin("margin_left")
    margin_r = _margin("margin_right")
    margin_t = _margin("margin_top")
    margin_b = _margin("margin_bottom")
    inner_w = shape.width / EMU_PER_PT_LOCAL - margin_l - margin_r
    if inner_w <= 0:
        return
    cur_font_pt: Optional[float] = None
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                cur_font_pt = float(run.font.size.pt)
                break
        if cur_font_pt is not None:
            break
    if cur_font_pt is None:
        return

    def required_h(font_pt: float) -> float:
        line_h = font_pt * 1.2
        total = 0.0
        for para in tf.paragraphs:
            ptext = (para.text or "")
            if not ptext.strip():
                total += line_h
                continue
            w = _estimate_text_width_pt(ptext, font_pt)
            lines = max(1, math.ceil(w / inner_w))
            total += lines * line_h
        return total

    def inner_height_now() -> float:
        return shape.height / EMU_PER_PT_LOCAL - margin_t - margin_b

    # Step A: current font で必要な inner_h を満たすか。OK なら no-op。
    needed_h = required_h(cur_font_pt)
    if needed_h <= inner_height_now() + 0.5:
        return

    # Step B: box.height 拡張で対応できるなら canvas 下端まで伸ばす (見た目変えない方向)。
    top_pt = shape.top / EMU_PER_PT_LOCAL
    max_h_pt = max(0.0, slide_height_pt - top_pt)
    target_outer_h = needed_h + margin_t + margin_b
    new_h_pt = min(target_outer_h, max_h_pt)
    if new_h_pt > shape.height / EMU_PER_PT_LOCAL + 0.5:
        shape.height = Pt(new_h_pt)
    if needed_h <= inner_height_now() + 0.5:
        return

    # Step C: それでも収まらなければ font を 1 段階ずつ shrink (最終手段)。
    target_font: Optional[float] = None
    for size in sorted([s for s in _OVERFLOW_ALLOWED_FONT_SIZES_PT if s < cur_font_pt], reverse=True):
        if required_h(size) <= inner_height_now() + 0.5:
            target_font = float(size)
            break
    if target_font is None:
        target_font = float(_OVERFLOW_ALLOWED_FONT_SIZES_PT[0])
    _set_shape_text_size(shape, target_font)


def _estimate_text_width_pt(text: str, font_size_pt: float) -> float:
    width = 0.0
    for ch in text:
        code = ord(ch)
        if ch.isspace():
            width += 0.33
        elif 0x3040 <= code <= 0x30FF or 0x3400 <= code <= 0x9FFF:
            width += 1.0
        elif ch.isupper():
            width += 0.62
        else:
            width += 0.54
    return width * font_size_pt


def _text_containers(shape):
    if getattr(shape, "has_text_frame", False):
        yield {
            "kind": "shape",
            "text_frame": shape.text_frame,
            "width_pt": shape.width / EMU_PER_PT,
            "height_pt": shape.height / EMU_PER_PT,
            "offset": None,
        }

    if getattr(shape, "has_table", False):
        table = shape.table
        top = shape.top / EMU_PER_PT
        for row_idx, row in enumerate(table.rows, start=1):
            left = shape.left / EMU_PER_PT
            row_height = row.height / EMU_PER_PT
            for col_idx, col in enumerate(table.columns, start=1):
                cell = table.cell(row_idx - 1, col_idx - 1)
                col_width = col.width / EMU_PER_PT
                yield {
                    "kind": "table_cell",
                    "row": row_idx,
                    "col": col_idx,
                    "text_frame": cell.text_frame,
                    "width_pt": col_width,
                    "height_pt": row_height,
                    "offset": (left, top, col_width, row_height),
                }
                left += col_width
            top += row_height


def _non_empty_paragraphs(text_frame) -> list:
    return [
        para
        for para in text_frame.paragraphs
        if any(run.text for run in para.runs)
    ]


def _normalize_hex(value: Any) -> Optional[str]:
    if not isinstance(value, str):
        return None
    raw = value.strip().upper()
    if not raw:
        return None
    if not raw.startswith("#"):
        raw = "#" + raw
    if len(raw) != 7:
        return None
    try:
        int(raw[1:], 16)
    except ValueError:
        return None
    return raw


def _run_rgb_hex(run) -> Optional[str]:
    try:
        color = run.font.color
        if color.type != MSO_COLOR_TYPE.RGB:
            return None
        rgb = color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    if rgb is None:
        return None
    return f"#{str(rgb).upper()}"


def _detect_font_size(shape, slide_idx, slide_id, slide=None) -> Optional[FixAction]:
    if slide is None or (not getattr(shape, "has_text_frame", False) and not getattr(shape, "has_table", False)):
        return None

    scale = _slide_scale(slide)
    updates: list[dict] = []
    reasons: list[str] = []
    before_sizes: list[dict] = []
    after_sizes: list[dict] = []
    affected_containers: list[dict] = []

    for container in _text_containers(shape):
        paragraphs = _non_empty_paragraphs(container["text_frame"])
        container_updates: list[dict] = []
        container_text = " ".join(
            run.text
            for para in paragraphs
            for run in para.runs
            if run.text
        )
        for p_idx, para in enumerate(container["text_frame"].paragraphs, start=1):
            for r_idx, run in enumerate(para.runs, start=1):
                if not run.text or run.font.size is None:
                    continue
                actual = run.font.size.pt
                normalized = actual * scale
                nearest = _nearest_allowed_font_size(normalized)
                target_actual = nearest / scale
                if abs(normalized - nearest) <= 1.0:
                    continue
                update = {
                    "kind": container["kind"],
                    "paragraph": p_idx,
                    "run": r_idx,
                    "before_size_pt": round(actual, 4),
                    "after_size_pt": round(target_actual, 4),
                    "normalized_size_pt": round(normalized, 4),
                    "target_normalized_size_pt": nearest,
                    "text": run.text[:80],
                }
                if container["kind"] == "table_cell":
                    update["row"] = container["row"]
                    update["col"] = container["col"]
                container_updates.append(update)

        if not container_updates:
            continue

        affected_containers.append({
            k: container[k]
            for k in ("kind", "row", "col")
            if k in container
        })
        updates.extend(container_updates)
        before_sizes.extend(
            {
                "size_pt": u["before_size_pt"],
                "normalized_size_pt": u["normalized_size_pt"],
                "text": u["text"],
            }
            for u in container_updates
        )
        after_sizes.extend(
            {
                "size_pt": u["after_size_pt"],
                "normalized_size_pt": u["target_normalized_size_pt"],
                "text": u["text"],
            }
            for u in container_updates
        )

        if len(paragraphs) != 1:
            reasons.append("not_single_line")
        elif "\n" in container_text or "\v" in container_text:
            reasons.append("not_single_line")

        max_delta = max(abs(u["after_size_pt"] - u["before_size_pt"]) for u in container_updates)
        if max_delta > FONT_SIZE_FIX_DELTA_MAX_PT:
            reasons.append("font_size_delta_too_large")

        target_size = max(u["after_size_pt"] for u in container_updates)
        required_height = target_size * 1.2
        if required_height > container["height_pt"] * 0.85:
            reasons.append("insufficient_box_height")
        if _estimate_text_width_pt(container_text, target_size) > container["width_pt"] * 0.95:
            reasons.append("insufficient_box_width")

        box = container["offset"] or _bbox_pt(shape)
        prs = slide.part.package.presentation_part.presentation
        slide_w = prs.slide_width / EMU_PER_PT
        slide_h = prs.slide_height / EMU_PER_PT
        x, y, w, h = box
        if x < 0 or y < 0 or x + w > slide_w or y + h > slide_h:
            reasons.append("would_overflow_slide")

    if not updates:
        return None

    shape_box = _bbox_pt(shape)
    shape_id = getattr(shape, "shape_id", None)
    for other in _iter_shapes(slide.shapes):
        if other is shape or getattr(other, "shape_id", None) == shape_id:
            continue
        try:
            if _overlaps(shape_box, _bbox_pt(other)):
                reasons.append("shape_bbox_overlaps_other_shape")
                break
        except (AttributeError, TypeError, ValueError):
            continue

    status = "manual_required" if reasons else "apply"
    return FixAction(
        rule="font_size",
        slide_index=slide_idx,
        slide_id=slide_id,
        shape_id=getattr(shape, "shape_id", None),
        shape_name=getattr(shape, "name", None),
        status=status,
        reasons=sorted(set(reasons)),
        before={"runs": before_sizes, "scale": round(scale, 4)},
        after={"runs": after_sizes, "updates": updates},
    )


def _detect_line_height(shape, slide_idx, slide_id, slide=None) -> Optional[FixAction]:
    if slide is None or not getattr(shape, "has_text_frame", False):
        return None

    scale = _slide_scale(slide)
    updates: list[dict] = []
    reasons: list[str] = []
    paragraph_count = 0
    text = shape.text_frame.text.strip()

    for p_idx, para in enumerate(shape.text_frame.paragraphs, start=1):
        if not para.text.strip():
            continue
        paragraph_count += 1
        line_spacing = para.line_spacing
        if line_spacing is None:
            continue
        if isinstance(line_spacing, float):
            reasons.append("relative_line_spacing_requires_review")
            continue
        actual = line_spacing.pt
        normalized = actual * scale
        target_normalized = _nearest_allowed_line_height(normalized)
        if abs(normalized - target_normalized) <= 2.0:
            continue
        target_actual = target_normalized / scale
        updates.append(
            {
                "paragraph": p_idx,
                "before_line_height_pt": round(actual, 4),
                "after_line_height_pt": round(target_actual, 4),
                "normalized_line_height_pt": round(normalized, 4),
                "target_normalized_line_height_pt": target_normalized,
                "text": para.text[:80],
            }
        )

    if not updates:
        return None

    max_delta = max(abs(u["after_line_height_pt"] - u["before_line_height_pt"]) for u in updates)
    if max_delta > LINE_HEIGHT_FIX_DELTA_MAX_PT:
        reasons.append("line_height_delta_too_large")

    target_lines = paragraph_count or len(updates)
    target_line_height = max(u["after_line_height_pt"] for u in updates)
    required_height = target_lines * target_line_height
    if required_height > (shape.height / EMU_PER_PT) * 0.9:
        reasons.append("insufficient_box_height")
    if "\n" in text or "\v" in text:
        reasons.append("multi_line_text_requires_review")

    return FixAction(
        rule="line_height",
        slide_index=slide_idx,
        slide_id=slide_id,
        shape_id=getattr(shape, "shape_id", None),
        shape_name=getattr(shape, "name", None),
        status="manual_required" if reasons else "apply",
        reasons=sorted(set(reasons)),
        before={
            "paragraphs": [
                {
                    "line_height_pt": u["before_line_height_pt"],
                    "normalized_line_height_pt": u["normalized_line_height_pt"],
                    "text": u["text"],
                }
                for u in updates
            ],
            "scale": round(scale, 4),
        },
        after={
            "paragraphs": [
                {
                    "line_height_pt": u["after_line_height_pt"],
                    "normalized_line_height_pt": u["target_normalized_line_height_pt"],
                    "text": u["text"],
                }
                for u in updates
            ],
            "updates": updates,
        },
    )


def _detect_alignment(shape, slide_idx, slide_id, slide=None) -> Optional[FixAction]:
    if not getattr(shape, "has_text_frame", False):
        return None
    if not shape.text_frame.text.strip():
        return None

    updates: list[dict] = []
    if shape.text_frame.vertical_anchor not in (None, MSO_VERTICAL_ANCHOR.TOP):
        updates.append(
            {
                "target": "text_frame",
                "property": "vertical_anchor",
                "before": str(shape.text_frame.vertical_anchor),
                "after": "TOP",
            }
        )

    for p_idx, para in enumerate(shape.text_frame.paragraphs, start=1):
        if not para.text.strip():
            continue
        if para.alignment is not None and para.alignment != PP_ALIGN.LEFT:
            updates.append(
                {
                    "target": "paragraph",
                    "paragraph": p_idx,
                    "property": "alignment",
                    "before": str(para.alignment),
                    "after": "LEFT",
                    "text": para.text[:80],
                }
            )

    if not updates:
        return None

    return FixAction(
        rule="alignment",
        slide_index=slide_idx,
        slide_id=slide_id,
        shape_id=getattr(shape, "shape_id", None),
        shape_name=getattr(shape, "name", None),
        before={"updates": [{k: v for k, v in update.items() if k != "after"} for update in updates]},
        after={"updates": updates},
    )


def _contrast_candidate_hex(finding: Any) -> Optional[str]:
    detail = _finding_detail(finding)
    candidates = detail.get("candidate_values")
    if not isinstance(candidates, dict):
        return None
    return (
        _normalize_hex(candidates.get("foreground_hex"))
        or _normalize_hex(candidates.get("color_hex"))
    )


def _contrast_repair_spec(finding: Any) -> Optional[dict]:
    """Resolve preferred strategy + target colors from candidate_values.

    Returns dict with keys:
      mode: "foreground" | "background"
      to_hex: target hex
      from_hex: original hex (foreground for fg-mode, background for bg-mode)
      background_source: source label from finding evidence (informational)
    """
    detail = _finding_detail(finding)
    candidates = detail.get("candidate_values")
    if not isinstance(candidates, dict):
        return None
    preferred = candidates.get("preferred_strategy")
    bg_source = detail.get("background_source")
    if preferred == "background":
        option = candidates.get("background_option") or {}
        to_hex = _normalize_hex(option.get("to_hex"))
        from_hex = _normalize_hex(option.get("from_hex") or detail.get("background_hex"))
        if not to_hex or not from_hex:
            return None
        return {
            "mode": "background",
            "to_hex": to_hex,
            "from_hex": from_hex,
            "background_source": bg_source,
        }
    # default = foreground (backward compatible: missing preferred_strategy)
    option = candidates.get("foreground_option") or {}
    to_hex = (
        _normalize_hex(option.get("to_hex"))
        or _normalize_hex(candidates.get("foreground_hex"))
        or _normalize_hex(candidates.get("color_hex"))
    )
    if not to_hex:
        return None
    return {
        "mode": "foreground",
        "to_hex": to_hex,
        "from_hex": None,
        "background_source": bg_source,
    }


def _shape_solid_fill_hex(shape) -> Optional[str]:
    fill = getattr(shape, "fill", None)
    if fill is None:
        return None
    try:
        fore = fill.fore_color
    except (AttributeError, TypeError, ValueError):
        return None
    try:
        if fore.type != MSO_COLOR_TYPE.RGB:
            return None
    except (AttributeError, TypeError, ValueError):
        return None
    try:
        rgb = fore.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    if rgb is None:
        return None
    return f"#{str(rgb).upper()}"


def _find_bg_repair_target(text_shape, slide, from_hex_bg: str, background_source: Optional[str]):
    """Locate the shape whose solid fill should be repaired for bg-mode contrast fix.

    Currently supports `shape_solid_fill` (text shape's own fill) and
    `behind_solid_fill:<label>` (smallest enclosing shape with matching fill).
    Returns the shape, or None if not found.
    """
    own_hex = _shape_solid_fill_hex(text_shape)
    if own_hex and own_hex.upper() == from_hex_bg.upper():
        return text_shape
    if slide is None:
        return None
    if not (background_source or "").startswith("behind_solid_fill"):
        return None
    text_left = float(getattr(text_shape, "left", 0) or 0)
    text_top = float(getattr(text_shape, "top", 0) or 0)
    text_width = float(getattr(text_shape, "width", 0) or 0)
    text_height = float(getattr(text_shape, "height", 0) or 0)
    cx = text_left + text_width / 2.0
    cy = text_top + text_height / 2.0
    best = None
    for candidate in slide.shapes:
        if candidate is text_shape:
            continue
        cand_hex = _shape_solid_fill_hex(candidate)
        if not cand_hex or cand_hex.upper() != from_hex_bg.upper():
            continue
        c_left = float(getattr(candidate, "left", 0) or 0)
        c_top = float(getattr(candidate, "top", 0) or 0)
        c_w = float(getattr(candidate, "width", 0) or 0)
        c_h = float(getattr(candidate, "height", 0) or 0)
        if not (c_left <= cx <= c_left + c_w and c_top <= cy <= c_top + c_h):
            continue
        area = c_w * c_h
        if best is None or area < best[0]:
            best = (area, candidate)
    return best[1] if best else None


def _contrast_original_hexes(finding: Any) -> set[str]:
    detail = _finding_detail(finding)
    originals: set[str] = set()
    for key in ("foreground_hex", "text_hex", "original_run_color_hex"):
        normalized = _normalize_hex(detail.get(key))
        if normalized:
            originals.add(normalized)

    raw_colors = detail.get("original_run_colors_hex")
    if isinstance(raw_colors, str):
        raw_colors = [raw_colors]
    if isinstance(raw_colors, Sequence) and not isinstance(raw_colors, (str, bytes)):
        for value in raw_colors:
            normalized = _normalize_hex(value)
            if normalized:
                originals.add(normalized)
    return originals


def _contrast_finding_rejection(finding: Any) -> Optional[str]:
    detail = _finding_detail(finding)
    fixability = detail.get("fixability")
    if fixability == "auto_fix_candidate":
        if _contrast_candidate_hex(finding):
            return None
        return "missing_contrast_candidate_values"
    if fixability is None:
        if _contrast_candidate_hex(finding):
            return None
        return "missing_contrast_candidate_values"
    return detail.get("manual_required_reason") or detail.get("fixability_reason") or str(fixability)


def _matching_contrast_findings(
    shape,
    slide_idx: int,
    findings: Optional[Sequence[Any]],
) -> list[Any]:
    if not findings:
        return []
    shape_id = getattr(shape, "shape_id", None)
    shape_name = getattr(shape, "name", None)
    matched: list[Any] = []
    for finding in findings:
        if _finding_rule(finding) != "contrast":
            continue
        finding_slide = _finding_field(finding, "slide_index")
        if finding_slide is not None and finding_slide != slide_idx:
            continue
        finding_shape_id = _finding_field(finding, "shape_id")
        if finding_shape_id is not None and shape_id is not None:
            if finding_shape_id != shape_id:
                continue
        else:
            finding_shape_name = _finding_field(finding, "shape_name")
            if finding_shape_name is not None and shape_name is not None:
                if finding_shape_name != shape_name:
                    continue
        matched.append(finding)
    return matched


def _detect_contrast(
    shape,
    slide_idx,
    slide_id,
    slide=None,
    *,
    findings: Optional[Sequence[Any]] = None,
) -> Optional[FixAction]:
    matched_findings = _matching_contrast_findings(shape, slide_idx, findings)
    if not matched_findings:
        return None

    reasons: list[str] = []
    updates: list[dict] = []
    seen: set[tuple] = set()

    for finding in matched_findings:
        rejection = _contrast_finding_rejection(finding)
        if rejection:
            reasons.append(rejection)
            continue

        spec = _contrast_repair_spec(finding)
        if spec is None:
            reasons.append("missing_contrast_original_or_candidate")
            continue

        if spec["mode"] == "background":
            target_shape = _find_bg_repair_target(
                shape, slide, spec["from_hex"], spec.get("background_source")
            )
            if target_shape is None:
                reasons.append("background_repair_target_not_found")
                continue
            key = (
                "background_fill",
                getattr(target_shape, "shape_id", None),
                spec["from_hex"],
                spec["to_hex"],
            )
            if key in seen:
                continue
            seen.add(key)
            updates.append(
                {
                    "mode": "background_fill",
                    "kind": "shape_fill",
                    "shape_id": getattr(target_shape, "shape_id", None),
                    "shape_name": getattr(target_shape, "name", None),
                    "before_hex": spec["from_hex"],
                    "after_hex": spec["to_hex"],
                    "check": _finding_field(finding, "check"),
                }
            )
            continue

        # foreground mode (default)
        target_hex = spec["to_hex"]
        originals = _contrast_original_hexes(finding)
        if not originals:
            reasons.append("missing_contrast_original_or_candidate")
            continue

        for container in _text_containers(shape):
            for p_idx, para in enumerate(container["text_frame"].paragraphs, start=1):
                for r_idx, run in enumerate(para.runs, start=1):
                    before_hex = _run_rgb_hex(run)
                    if before_hex is None or before_hex not in originals:
                        continue
                    if before_hex == target_hex:
                        continue
                    key = (
                        container["kind"],
                        container.get("row"),
                        container.get("col"),
                        p_idx,
                        r_idx,
                        before_hex,
                        target_hex,
                    )
                    if key in seen:
                        continue
                    seen.add(key)
                    update = {
                        "mode": "foreground_run",
                        "kind": container["kind"],
                        "paragraph": p_idx,
                        "run": r_idx,
                        "before_hex": before_hex,
                        "after_hex": target_hex,
                        "text": run.text[:80],
                        "check": _finding_field(finding, "check"),
                    }
                    if container["kind"] == "table_cell":
                        update["row"] = container["row"]
                        update["col"] = container["col"]
                    updates.append(update)

    if not updates and not reasons:
        return None

    return FixAction(
        rule="contrast",
        slide_index=slide_idx,
        slide_id=slide_id,
        shape_id=getattr(shape, "shape_id", None),
        shape_name=getattr(shape, "name", None),
        status="manual_required" if reasons and not updates else "apply",
        reasons=sorted(set(reasons)),
        before={"runs": [
            {"color_hex": u["before_hex"], "text": u.get("text", "")}
            for u in updates if u.get("mode") == "foreground_run"
        ]},
        after={
            "runs": [
                {"color_hex": u["after_hex"], "text": u.get("text", "")}
                for u in updates if u.get("mode") == "foreground_run"
            ],
            "updates": updates,
        },
    )


DETECTORS: dict = {
    "autofit": _detect_autofit,
    "geometry": _detect_geometry,
    "font_size": _detect_font_size,
    "line_height": _detect_line_height,
    "alignment": _detect_alignment,
}

RULE_ENABLED: dict[str, bool] = {
    "autofit": True,
    "geometry": True,
    "font_size": FONT_SIZE_FIXER_ENABLED,
    "contrast": True,
}


def _enabled_rules(rules: Sequence[str]) -> tuple[str, ...]:
    return tuple(rule for rule in rules if RULE_ENABLED.get(rule, True))


def _finding_field(finding: Any, key: str, default: Any = None) -> Any:
    if isinstance(finding, dict):
        return finding.get(key, default)
    return getattr(finding, key, default)


def _finding_detail(finding: Any) -> dict:
    detail = _finding_field(finding, "detail", {})
    return detail if isinstance(detail, dict) else {}


def _finding_detail_value(finding: Any, key: str, default: Any = None) -> Any:
    detail = _finding_detail(finding)
    if key in detail:
        return detail[key]
    return _finding_field(finding, key, default)


def _finding_rule(finding: Any) -> Optional[str]:
    check = _finding_field(finding, "check")
    if check is None:
        return None
    return CHECK_TO_RULE.get(check, check)


def _finding_matches_action(finding: Any, action: FixAction) -> bool:
    if _finding_rule(finding) != action.rule:
        return False

    finding_slide = _finding_field(finding, "slide_index")
    if finding_slide is not None and finding_slide != action.slide_index:
        return False

    finding_shape_id = _finding_field(finding, "shape_id")
    if finding_shape_id is not None and action.shape_id is not None:
        return finding_shape_id == action.shape_id

    finding_shape_name = _finding_field(finding, "shape_name")
    if finding_shape_name is not None and action.shape_name is not None:
        return finding_shape_name == action.shape_name

    return True


def _candidate_values_present(candidate_values: Any) -> bool:
    if candidate_values is None:
        return False
    if isinstance(candidate_values, (list, tuple, set, dict)):
        return len(candidate_values) > 0
    return True


def _finding_reasons(finding: Any) -> list[str]:
    reasons: list[str] = []
    for key in (
        "manual_required_reason",
        "manual_reason",
        "fixability_reason",
        "reason",
        "reasons",
    ):
        value = _finding_detail_value(finding, key)
        if value is None:
            continue
        if isinstance(value, str):
            reasons.append(value)
        elif isinstance(value, (list, tuple, set)):
            reasons.extend(str(item) for item in value if item is not None)
        else:
            reasons.append(str(value))
    return reasons


def _fixability_decision_from_finding(
    finding: Any,
    *,
    judgement_gate: bool = True,
) -> Optional[tuple[str, list[str]]]:
    """Return an action status override from evidence-schema findings.

    Legacy findings that do not declare fixability fall back to the existing
    detector-only behavior.

    When ``judgement_gate=True`` (default, POLICY-001 段階 3), apply_mode=
    judgement_fix findings whose lint output is still `manual_required` are
    kept at manual_required regardless of candidate_values: the SPA
    judgement (auto_fixable) is required before pptx_fix will touch the
    deck. Set ``judgement_gate=False`` to recover the pre-段階 3 behavior
    that applied any manual_required finding with concrete candidate_values.
    apply_mode=auto_fix and apply_mode=no_fix paths are unaffected by the
    gate.
    """
    fixability = _finding_detail_value(finding, "fixability")
    if fixability is None:
        return None

    if fixability == "auto_fix_candidate":
        candidate_values = _finding_detail_value(finding, "candidate_values")
        if _candidate_values_present(candidate_values):
            return "apply", []
        return "manual_required", ["missing_candidate_values"]

    reasons = _finding_reasons(finding)
    if fixability == "manual_required":
        check = _finding_field(finding, "check")
        if check in {"image_upscale_ratio", "slide_size"}:
            return "manual_required", reasons or ["finding_marked_manual_required"]
        apply_mode = _apply_mode_for_check(check)
        if judgement_gate and apply_mode == "judgement_fix":
            gate_reasons = list(reasons) if reasons else []
            gate_reasons.append("judgement_fix_gate_requires_spa_judgement")
            return "manual_required", gate_reasons
        candidate_values = _finding_detail_value(finding, "candidate_values")
        if _candidate_values_present(candidate_values):
            return "apply", reasons or ["manual_required_overridden_by_design_system_candidate"]
        return "apply", reasons or ["finding_marked_manual_required"]

    return "manual_required", reasons or [f"fixability_{fixability}"]


def _apply_finding_fixability(
    action: FixAction,
    finding: Any,
    *,
    judgement_gate: bool = True,
) -> FixAction:
    decision = _fixability_decision_from_finding(finding, judgement_gate=judgement_gate)
    if decision is None:
        return action

    status, reasons = decision
    if reasons:
        action.reasons = sorted(set(action.reasons + reasons))
    if status == "manual_required":
        action.status = "manual_required"
    elif action.status != "manual_required":
        action.status = status
    return action


def _apply_matching_finding_fixability(
    action: FixAction,
    findings: Optional[Sequence[Any]],
    *,
    judgement_gate: bool = True,
) -> FixAction:
    if not findings:
        return action

    for finding in findings:
        if _finding_matches_action(finding, action):
            return _apply_finding_fixability(
                action, finding, judgement_gate=judgement_gate
            )
    return action


def _detect_finding_action(prs, finding: Any) -> Optional[FixAction | list[FixAction]]:
    check = _finding_field(finding, "check")
    rule = _finding_rule(finding)
    detail = _finding_detail(finding)
    slide = _slide_by_index(prs, _finding_field(finding, "slide_index"))

    if slide is None:
        return None

    if rule == "animation":
        markers = detail.get("markers") or detail.get("animation_markers")
        if not markers:
            markers = ["p:transition", "p:timing"]
        return FixAction(
            rule="animation",
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=None,
            shape_name=None,
            before={"markers": markers},
            after={"remove_markers": markers},
        )

    if rule == "reading_order":
        return FixAction(
            rule="reading_order",
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=None,
            shape_name=None,
            before={"source_order": detail.get("source_order")},
            after={"visual_order": detail.get("visual_order")},
        )

    shape = _shape_from_finding(prs, finding)

    if rule == "overlap" and check in {"text_overlap", "object_overlap"}:
        shape = _shape_from_finding(prs, finding, "shape_b")
        shape_a = detail.get("shape_a") or {}
        shape_b = detail.get("shape_b") or {}
        if shape is None or not _valid_bbox(shape_a.get("bbox_pt")) or not _valid_bbox(shape_b.get("bbox_pt")):
            return None
        dx_norm, dy_norm = _overlap_delta(shape_a["bbox_pt"], shape_b["bbox_pt"])
        sx, sy = _slide_scale_xy(prs)
        current = _shape_geometry_pt(shape)
        geometry = {
            **current,
            "left": round(current["left"] + dx_norm * sx, 4),
            "top": round(current["top"] + dy_norm * sy, 4),
        }
        return _geometry_action(rule, finding, shape, geometry, ["mechanical_minimum_separation"])

    if rule == "spacing" and check == "object_gap_too_small":
        shape = _shape_from_finding(prs, finding, "shape_b")
        shape_a = detail.get("shape_a") or {}
        shape_b = detail.get("shape_b") or {}
        if shape is None or not _valid_bbox(shape_a.get("bbox_pt")) or not _valid_bbox(shape_b.get("bbox_pt")):
            return None
        gap = float(detail.get("gap_pt") or 0)
        threshold = float(detail.get("threshold_pt") or 8)
        move = max(0.0, threshold - gap + 4.0)
        ax, ay, aw, ah = [float(v) for v in shape_a["bbox_pt"]]
        bx, by, bw, bh = [float(v) for v in shape_b["bbox_pt"]]
        dx_norm = dy_norm = 0.0
        if detail.get("axis") == "horizontal":
            dx_norm = move if bx >= ax else -move
        else:
            dy_norm = move if by >= ay else -move
        sx, sy = _slide_scale_xy(prs)
        current = _shape_geometry_pt(shape)
        geometry = {
            **current,
            "left": round(current["left"] + dx_norm * sx, 4),
            "top": round(current["top"] + dy_norm * sy, 4),
        }
        return _geometry_action(rule, finding, shape, geometry, ["mechanical_minimum_gap"])

    if rule == "bbox_fit":
        if shape is None:
            return None
        sx, sy = _slide_scale_xy(prs)
        if check == "safe_text_area_text":
            norm_area = tuple(float(v) for v in detail.get("safe_text_area_pt", [81, 40, 1278, 690]))
        elif check == "safe_margins":
            margins = detail.get("safe_margin_pt") or {}
            norm_area = (
                float(margins.get("left", 81)),
                float(margins.get("top", 40)),
                SLIDE_W_PT - float(margins.get("left", 81)) - float(margins.get("right", 81)),
                SLIDE_H_PT - float(margins.get("top", 40)) - float(margins.get("bottom", 40)),
            )
        else:
            norm_area = (0.0, 0.0, SLIDE_W_PT, SLIDE_H_PT)
        actual_area = (norm_area[0] * sx, norm_area[1] * sy, norm_area[2] * sx, norm_area[3] * sy)
        geometry = _fit_bbox_into_area(_shape_geometry_pt(shape), actual_area)
        return _geometry_action(rule, finding, shape, geometry, ["mechanical_bbox_fit"])

    if rule == "image_crop":
        if shape is None or shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return None
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            before={"crop": {side: getattr(shape, f"crop_{side}") for side in ("left", "right", "top", "bottom")}},
            after={"crop": {side: 0.0 for side in ("left", "right", "top", "bottom")}},
        )

    if rule == "image_aspect":
        if shape is None or shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return None
        source = detail.get("source_px") or []
        if not (isinstance(source, Sequence) and len(source) == 2 and source[0] and source[1]):
            return None
        source_aspect = float(source[0]) / float(source[1])
        current = _shape_geometry_pt(shape)
        display_aspect = current["width"] / current["height"] if current["height"] else source_aspect
        geometry = dict(current)
        if display_aspect > source_aspect:
            new_width = current["height"] * source_aspect
            geometry["left"] = round(current["left"] + (current["width"] - new_width) / 2, 4)
            geometry["width"] = round(new_width, 4)
        else:
            new_height = current["width"] / source_aspect
            geometry["top"] = round(current["top"] + (current["height"] - new_height) / 2, 4)
            geometry["height"] = round(new_height, 4)
        return _geometry_action(rule, finding, shape, geometry, ["preserve_source_aspect"])

    if rule == "font_family":
        if shape is None:
            return None
        candidate = _finding_detail_value(finding, "candidate_values") or {}
        font_name = candidate.get("candidate_font_family") if isinstance(candidate, dict) else None
        if not font_name:
            font_name = "Noto Sans JP"
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            before={"font": detail.get("font"), "runs": detail.get("runs")},
            after={"font_name": font_name},
        )

    if rule == "text_color":
        if shape is None:
            return None
        candidate = _finding_detail_value(finding, "candidate_values") or {}
        after_hex = _candidate_hex(candidate, "color_hex", "foreground_hex")
        before_hex = _normalize_hex(detail.get("color_hex") or detail.get("text_hex"))
        if not after_hex:
            return None
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            before={"color_hex": before_hex},
            after={"color_hex": after_hex},
        )

    if rule == "fill_color":
        if shape is None:
            return None
        candidate = _finding_detail_value(finding, "candidate_values") or {}
        after_hex = _candidate_hex(candidate, "color_hex", "fill_hex", "background_hex")
        if not after_hex:
            return None
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            before={"color_hex": detail.get("color_hex")},
            after={"color_hex": after_hex},
        )

    if rule == "text_wrap":
        if shape is None or not getattr(shape, "has_text_frame", False):
            return None
        candidates = detail.get("candidate_values")
        widen: Optional[dict] = None
        candidate_entries: list = []
        if isinstance(candidates, list):
            candidate_entries = candidates
        elif isinstance(candidates, dict):
            candidate_entries = [candidates]
        for entry in candidate_entries:
            if (
                isinstance(entry, dict)
                and entry.get("strategy") == "widen_to_fit"
                and _valid_bbox(entry.get("bbox_pt"))
            ):
                widen = entry
                break
        after: dict = {"replace_breaks": True}
        before_geometry = _shape_geometry_pt(shape)
        if widen is not None:
            sx, _ = _slide_scale_xy(prs)
            after["widen_to_fit"] = {
                "width_pt_norm": float(widen["bbox_pt"][2]),
                "width_pt_actual": float(widen["bbox_pt"][2]) * sx,
                "bbox_pt": list(widen["bbox_pt"]),
            }
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            before={"text": shape.text_frame.text, "geometry": before_geometry},
            after=after,
        )

    if rule == "badge_alignment":
        if shape is None or not getattr(shape, "has_text_frame", False):
            return None
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            before={
                "vertical_anchor": str(shape.text_frame.vertical_anchor)
                if shape.text_frame.vertical_anchor is not None
                else None,
                "first_paragraph_alignment": (
                    str(shape.text_frame.paragraphs[0].alignment)
                    if shape.text_frame.paragraphs
                    and shape.text_frame.paragraphs[0].alignment is not None
                    else None
                ),
            },
            after={"alignment": "CENTER", "vertical_anchor": "MIDDLE"},
        )

    if rule == "text_canvas_reflow":
        # DS-OVERFLOW-001 段階3 (2026-05-23): text_canvas_overflow finding を
        # 受けて、SPA で promoted された 1 つの strategy を apply する。
        # candidate strategies: enable_word_wrap / shrink_box_width_to_canvas /
        # shrink_font_size。SPA が chosen_strategy を指定しなければ 先頭採用。
        if shape is None or not getattr(shape, "has_text_frame", False):
            return None
        candidates = detail.get("candidate_values")
        chosen = None
        if isinstance(candidates, list):
            chosen_name = detail.get("chosen_strategy")
            if isinstance(chosen_name, str):
                for entry in candidates:
                    if isinstance(entry, dict) and entry.get("strategy") == chosen_name:
                        chosen = entry
                        break
            if chosen is None:
                for entry in candidates:
                    if isinstance(entry, dict) and "strategy" in entry:
                        chosen = entry
                        break
        if chosen is None:
            return None
        strategy = chosen.get("strategy")
        before_state = {
            "shape_id": getattr(shape, "shape_id", None),
            "shape_name": getattr(shape, "name", None),
        }
        if strategy == "enable_word_wrap":
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before={**before_state, "word_wrap": shape.text_frame.word_wrap},
                after={"strategy": "enable_word_wrap"},
            )
        if strategy == "expand_box_to_canvas_and_wrap":
            # 3 段積み上げ combined apply ([[feedback-overflow-fix-priority]]):
            # box.width 拡張 + word_wrap=True + box.height 拡張 を 1 アクションで。
            target_w_norm = chosen.get("target_width_pt")
            target_h_norm = chosen.get("target_height_pt")
            target_word_wrap = chosen.get("word_wrap", True)
            if not (
                isinstance(target_w_norm, (int, float))
                and isinstance(target_h_norm, (int, float))
            ):
                return None
            sx, sy = _slide_scale_xy(prs)
            before_geometry = _shape_geometry_pt(shape)
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before={
                    **before_state,
                    "geometry": before_geometry,
                    "word_wrap": shape.text_frame.word_wrap,
                },
                after={
                    "strategy": "expand_box_to_canvas_and_wrap",
                    "geometry": {
                        "left": before_geometry["left"],
                        "top": before_geometry["top"],
                        "width": round(float(target_w_norm) * sx, 4),
                        "height": round(float(target_h_norm) * sy, 4),
                    },
                    "word_wrap": bool(target_word_wrap),
                    "applied_steps": chosen.get("applied_steps"),
                },
            )
        if strategy == "shrink_box_width_to_canvas":
            target_w_norm = chosen.get("target_width_pt")
            if not isinstance(target_w_norm, (int, float)):
                return None
            sx, sy = _slide_scale_xy(prs)
            before_geometry = _shape_geometry_pt(shape)
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before={**before_state, "geometry": before_geometry},
                after={
                    "strategy": "shrink_box_width_to_canvas",
                    "geometry": {
                        "left": before_geometry["left"],
                        "top": before_geometry["top"],
                        "width": round(float(target_w_norm) * sx, 4),
                        "height": before_geometry["height"],
                    },
                },
            )
        if strategy == "shrink_font_size":
            new_font = chosen.get("font_size_pt")
            if not isinstance(new_font, (int, float)):
                return None
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before=before_state,
                after={"strategy": "shrink_font_size", "font_size_pt": float(new_font)},
            )
        if strategy == "expand_box_width_to_canvas":
            new_w_norm = chosen.get("target_width_pt")
            if not isinstance(new_w_norm, (int, float)):
                return None
            sx, sy = _slide_scale_xy(prs)
            before_geometry = _shape_geometry_pt(shape)
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before={**before_state, "geometry": before_geometry},
                after={
                    "strategy": "expand_box_width_to_canvas",
                    "geometry": {
                        "left": before_geometry["left"],
                        "top": before_geometry["top"],
                        "width": round(float(new_w_norm) * sx, 4),
                        "height": before_geometry["height"],
                    },
                },
            )
        if strategy == "expand_box_height":
            new_h_norm = chosen.get("target_height_pt")
            if not isinstance(new_h_norm, (int, float)):
                return None
            sx, sy = _slide_scale_xy(prs)
            before_geometry = _shape_geometry_pt(shape)
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before={**before_state, "geometry": before_geometry},
                after={
                    "strategy": "expand_box_height",
                    "geometry": {
                        "left": before_geometry["left"],
                        "top": before_geometry["top"],
                        "width": before_geometry["width"],
                        "height": round(float(new_h_norm) * sy, 4),
                    },
                },
            )
        return None

    if rule == "text_box_resize":
        # DS-OVERFLOW-001 段階2 (2026-05-22): text_box_overflow finding を受け、
        # SPA で promoted (= judgement_reason=auto_fixable + chosen_strategy)
        # された 1 つの strategy を apply する。SPA が strategy を明示しなければ、
        # multi_step candidates の **先頭** (= shrink_font_size 優先) を採用する。
        if shape is None or not getattr(shape, "has_text_frame", False):
            return None
        candidates = detail.get("candidate_values")
        chosen = None
        if isinstance(candidates, list):
            chosen_name = detail.get("chosen_strategy")
            if isinstance(chosen_name, str):
                for entry in candidates:
                    if isinstance(entry, dict) and entry.get("strategy") == chosen_name:
                        chosen = entry
                        break
            if chosen is None:
                for entry in candidates:
                    if isinstance(entry, dict) and "strategy" in entry:
                        chosen = entry
                        break
        if chosen is None:
            return None
        strategy = chosen.get("strategy")
        before_state = {
            "shape_id": getattr(shape, "shape_id", None),
            "shape_name": getattr(shape, "name", None),
        }
        if strategy == "shrink_font_size":
            new_font = chosen.get("font_size_pt")
            if not isinstance(new_font, (int, float)):
                return None
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before=before_state,
                after={"strategy": "shrink_font_size", "font_size_pt": float(new_font)},
            )
        if strategy == "compress_line_height":
            new_lh = chosen.get("line_height_pt")
            if not isinstance(new_lh, (int, float)):
                return None
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before=before_state,
                after={"strategy": "compress_line_height", "line_height_pt": float(new_lh)},
            )
        if strategy == "expand_box_height":
            new_h_norm = chosen.get("target_height_pt")
            if not isinstance(new_h_norm, (int, float)):
                return None
            sx, sy = _slide_scale_xy(prs)
            before_geometry = _shape_geometry_pt(shape)
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before={**before_state, "geometry": before_geometry},
                after={
                    "strategy": "expand_box_height",
                    "geometry": {
                        "left": before_geometry["left"],
                        "top": before_geometry["top"],
                        "width": before_geometry["width"],
                        "height": round(float(new_h_norm) * sy, 4),
                    },
                },
            )
        if strategy == "expand_box_width_to_canvas":
            new_w_norm = chosen.get("target_width_pt")
            if not isinstance(new_w_norm, (int, float)):
                return None
            sx, sy = _slide_scale_xy(prs)
            before_geometry = _shape_geometry_pt(shape)
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before={**before_state, "geometry": before_geometry},
                after={
                    "strategy": "expand_box_width_to_canvas",
                    "geometry": {
                        "left": before_geometry["left"],
                        "top": before_geometry["top"],
                        "width": round(float(new_w_norm) * sx, 4),
                        "height": before_geometry["height"],
                    },
                },
            )
        if strategy == "expand_box_width_to_fit_text":
            new_w_norm = chosen.get("target_width_pt")
            if not isinstance(new_w_norm, (int, float)):
                return None
            sx, sy = _slide_scale_xy(prs)
            before_geometry = _shape_geometry_pt(shape)
            return FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=getattr(shape, "shape_id", None),
                shape_name=getattr(shape, "name", None),
                before={**before_state, "geometry": before_geometry},
                after={
                    "strategy": "expand_box_width_to_fit_text",
                    "geometry": {
                        "left": before_geometry["left"],
                        "top": before_geometry["top"],
                        "width": round(float(new_w_norm) * sx, 4),
                        "height": before_geometry["height"],
                    },
                },
            )
        return None

    if rule == "box_canvas_clip":
        # DS-OVERFLOW-001 ([[feedback-overflow-fix-priority]]):
        # lint は text 描画範囲 (margin 込み) で判定。fix も全方向の overflow
        # に対し box 全体を shift して text 描画範囲を canvas 内に戻す
        # (width/height 不変)。
        if shape is None:
            return None
        overflow = detail.get("overflow_sides_pt") or {}
        if not isinstance(overflow, dict):
            return None
        before_geometry = _shape_geometry_pt(shape)
        sx, sy = _slide_scale_xy(prs)
        left_norm = before_geometry["left"] / sx
        top_norm = before_geometry["top"] / sy
        width_norm = before_geometry["width"] / sx
        height_norm = before_geometry["height"] / sy
        right_over = float(overflow.get("right", 0) or 0)
        bottom_over = float(overflow.get("bottom", 0) or 0)
        left_over = float(overflow.get("left", 0) or 0)
        top_over = float(overflow.get("top", 0) or 0)
        new_left = left_norm + max(0.0, left_over) - max(0.0, right_over)
        new_top = top_norm + max(0.0, top_over) - max(0.0, bottom_over)
        new_w = width_norm
        new_h = height_norm
        if abs(new_left - left_norm) < 0.05 and abs(new_top - top_norm) < 0.05:
            return None
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            before={"geometry": before_geometry, "overflow_sides_pt": overflow},
            after={
                "geometry": {
                    "left": round(new_left * sx, 4),
                    "top": round(new_top * sy, 4),
                    "width": round(new_w * sx, 4),
                    "height": round(new_h * sy, 4),
                },
            },
        )

    if rule == "decorative_remove":
        # FIX-013: lint が出した remove_shape candidate を SPA judgement で
        # auto_fixable に promote した finding が来ると、対象 shape を
        # slide の spTree から物理削除する。promote されないと strict gate
        # が manual_required で skip するので、ここまで来た時点で削除可。
        if shape is None:
            return None
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            before={
                "shape_id": getattr(shape, "shape_id", None),
                "shape_name": getattr(shape, "name", None),
            },
            after={"action": "remove_shape"},
        )

    if rule == "heading_hierarchy":
        title = _shape_from_finding(prs, finding, "title_candidate")
        body = _shape_from_finding(prs, finding, "largest_body_candidate")
        if title is None and body is None:
            return None
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(title or body, "shape_id", None),
            shape_name=getattr(title or body, "name", None),
            before={
                "title_shape_id": getattr(title, "shape_id", None),
                "body_shape_id": getattr(body, "shape_id", None),
            },
            after={"title_size_pt": 40, "body_size_pt": 24},
        )

    if rule == "inner_padding":
        children = detail.get("children") or []
        padding = detail.get("padding_pt") or {}
        if not children or not isinstance(padding, dict):
            return None
        target = float(detail.get("target_padding_pt") or 0)
        tolerance = float(detail.get("target_tolerance_pt") or 0)
        if target <= 0:
            return None
        left_padding = float(padding.get("left", 0))
        right_padding = float(padding.get("right", 0))
        top_padding = float(padding.get("top", 0))
        bottom_padding = float(padding.get("bottom", 0))
        dx_norm = target - left_padding if abs(left_padding - target) > tolerance else 0.0
        dy_norm = target - top_padding if abs(top_padding - target) > tolerance else 0.0
        if right_padding - dx_norm < target - tolerance:
            return None
        if bottom_padding - dy_norm < target - tolerance:
            return None
        if abs(dx_norm) < 0.1 and abs(dy_norm) < 0.1:
            return None
        sx, sy = _slide_scale_xy(prs)
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=None,
            shape_name=None,
            before={"children": children, "padding_pt": padding},
            after={
                "move_children": [child.get("shape_id") for child in children],
                "dx_pt": dx_norm * sx,
                "dy_pt": dy_norm * sy,
                "target_padding_pt": target,
            },
        )

    if rule == "card_grid":
        candidate_values = detail.get("candidate_values") if isinstance(detail.get("candidate_values"), dict) else {}
        medians = detail.get("group_medians") or candidate_values.get("group_medians") or {}
        inconsistent = detail.get("inconsistent_containers") or []
        if not isinstance(medians, dict) or not inconsistent:
            return None

        row_containers = detail.get("row_containers") or []
        # FIX-012: compute role-aware per-child moves across the whole row.
        # Falls back to the v3 conservative behaviour (= container geometry
        # equalization, no child role inference) when row_containers is
        # missing (= older lint JSON).
        all_moves: dict[tuple[int, int], tuple[float, float]] = {}
        container_id_to_row_index: dict[Any, int] = {}
        if row_containers:
            all_moves = _card_grid_alignment_moves(row_containers)
            for i, rc in enumerate(row_containers):
                if isinstance(rc, dict):
                    cid = rc.get("shape_id")
                    if cid is not None:
                        container_id_to_row_index[cid] = i

        actions: list[FixAction] = []
        for item in inconsistent:
            if not isinstance(item, dict):
                continue
            container = item.get("container") or {}
            cid = container.get("shape_id")

            child_moves: list[dict] = []
            row_idx = container_id_to_row_index.get(cid)
            if row_idx is not None:
                row_children = row_containers[row_idx].get("children") or []
                for (ci, chi), (target_left, target_top) in sorted(all_moves.items()):
                    if ci != row_idx:
                        continue
                    if chi >= len(row_children):
                        continue
                    child = row_children[chi]
                    child_moves.append(
                        {
                            "shape_id": child.get("shape_id"),
                            "role": _card_role_map(row_containers[row_idx]).get(chi),
                            "before_bbox_pt": list(child.get("bbox_pt") or []),
                            "after_left_pt": float(target_left),
                            "after_top_pt": float(target_top),
                        }
                    )

            action = FixAction(
                rule=rule,
                slide_index=int(_finding_field(finding, "slide_index") or 1),
                slide_id=_finding_field(finding, "slide_id"),
                shape_id=cid,
                shape_name=None,
                before={"inconsistent_containers": [item]},
                after={
                    "group_medians": medians,
                    "child_moves": child_moves,
                    "child_policy": (
                        "role_aware_top_alignment"
                        if child_moves
                        else "container_equalize_only"
                    ),
                },
            )
            # The v3 "child_group_exceeds_inner_box" guard assumes the
            # container will be shrunk to row-median size. FIX-012 keeps
            # container sizes (medians equalization only adjusts mismatched
            # ones) and moves individual children based on role, so the
            # guard is only relevant when no role-aware moves are emitted.
            reasons = (
                _card_grid_item_manual_reasons(item, medians) if not child_moves else []
            )
            if reasons:
                action.status = "manual_required"
                action.reasons = reasons
            actions.append(action)
        return actions or None

    if rule == "text_vertical_balance":
        if shape is None or not getattr(shape, "has_text_frame", False):
            return None
        return FixAction(
            rule=rule,
            slide_index=int(_finding_field(finding, "slide_index") or 1),
            slide_id=_finding_field(finding, "slide_id"),
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            before={"vertical_anchor": str(shape.text_frame.vertical_anchor)},
            after={"vertical_anchor": "MIDDLE"},
        )

    return None


# ---- Apply -----------------------------------------------------------------


def _backup_once(path: Path) -> None:
    backup_path = Path(str(path) + ".bak")
    if backup_path.exists():
        return
    shutil.copy2(path, backup_path)


def _apply_action(shape, action: FixAction, slide=None) -> None:
    if action.status != "apply":
        return
    if action.rule == "autofit":
        shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    elif action.rule == "geometry":
        for c, pt in action.after.items():
            setattr(shape, c, int(round(pt * EMU_PER_PT)))
    elif action.rule == "font_size":
        for update in action.after.get("updates", []):
            if update["kind"] == "shape":
                tf = shape.text_frame
            else:
                tf = shape.table.cell(update["row"] - 1, update["col"] - 1).text_frame
            run = tf.paragraphs[update["paragraph"] - 1].runs[update["run"] - 1]
            run.font.size = Pt(update["after_size_pt"])
    elif action.rule == "line_height":
        for update in action.after.get("updates", []):
            para = shape.text_frame.paragraphs[update["paragraph"] - 1]
            para.line_spacing = Pt(update["after_line_height_pt"])
    elif action.rule == "alignment":
        for update in action.after.get("updates", []):
            if update["property"] == "vertical_anchor":
                shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            elif update["property"] == "alignment":
                para = shape.text_frame.paragraphs[update["paragraph"] - 1]
                para.alignment = PP_ALIGN.LEFT
    elif action.rule == "contrast":
        for update in action.after.get("updates", []):
            mode = update.get("mode") or (
                "foreground_run" if "paragraph" in update else None
            )
            if mode == "background_fill":
                target_id = update.get("shape_id")
                if target_id is not None and target_id == getattr(shape, "shape_id", None):
                    target = shape
                elif slide is not None:
                    target = _shape_by_id(slide, target_id)
                else:
                    target = None
                if target is None:
                    continue
                try:
                    target.fill.solid()
                    target.fill.fore_color.rgb = RGBColor.from_string(
                        update["after_hex"].lstrip("#")
                    )
                except (AttributeError, TypeError, ValueError):
                    continue
                continue
            # foreground_run (default)
            if update["kind"] == "shape":
                tf = shape.text_frame
            else:
                tf = shape.table.cell(update["row"] - 1, update["col"] - 1).text_frame
            run = tf.paragraphs[update["paragraph"] - 1].runs[update["run"] - 1]
            run.font.color.rgb = RGBColor.from_string(update["after_hex"].lstrip("#"))


def _remove_slide_animation(slide) -> int:
    removed = 0
    for tag in ("transition", "timing"):
        for element in list(slide.element.xpath(f"./p:{tag}")):
            parent = element.getparent()
            if parent is not None:
                parent.remove(element)
                removed += 1
    return removed


def _reorder_slide_shapes_top_left(slide) -> None:
    shape_elements = [shape._element for shape in slide.shapes]
    ordered_shapes = sorted(slide.shapes, key=lambda shape: (shape.top, shape.left))
    tree = slide.shapes._spTree
    for element in shape_elements:
        tree.remove(element)
    for shape in ordered_shapes:
        tree.append(shape._element)


def _set_shape_text_size(shape, size_pt: float) -> int:
    if shape is None:
        return 0
    count = 0
    for tf in _shape_text_frames(shape):
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text:
                    run.font.size = Pt(size_pt)
                    count += 1
    return count


def _replace_text_breaks(shape) -> int:
    replaced = 0
    if not getattr(shape, "has_text_frame", False):
        return replaced
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if "\n" in run.text or "\v" in run.text:
                text = re.sub(r"([A-Za-z])[\n\v]([a-z])", r"\1\2", run.text)
                run.text = re.sub(r"[\n\v]+", " ", text)
                replaced += 1
    return replaced


def _apply_finding_action(prs, action: FixAction) -> None:
    if action.status != "apply":
        return
    slide = _slide_by_index(prs, action.slide_index)
    if slide is None:
        return

    if action.rule == "animation":
        _remove_slide_animation(slide)
        return
    if action.rule == "reading_order":
        _reorder_slide_shapes_top_left(slide)
        return
    if action.rule == "inner_padding":
        for shape_id in action.after.get("move_children", []):
            shape = _shape_by_id(slide, shape_id)
            if shape is not None:
                shape.left = Pt(shape.left / EMU_PER_PT + float(action.after.get("dx_pt", 0)))
                shape.top = Pt(shape.top / EMU_PER_PT + float(action.after.get("dy_pt", 0)))
        return
    if action.rule == "card_grid":
        sx, sy = _slide_scale_xy(prs)
        medians = action.after.get("group_medians") or {}

        # 1. Container geometry equalization to row medians (v3 behavior).
        for item in action.before.get("inconsistent_containers", []):
            container = item.get("container") or {}
            shape = _shape_by_id(slide, container.get("shape_id"))
            if shape is None:
                continue
            container_bbox = container.get("bbox_pt") or []
            if not _valid_bbox(container_bbox):
                continue
            geometry = _shape_geometry_pt(shape)
            if "top" in medians:
                geometry["top"] = float(medians["top"]) * sy
            if "width" in medians:
                geometry["width"] = float(medians["width"]) * sx
            if "height" in medians:
                geometry["height"] = float(medians["height"]) * sy
            _apply_geometry(shape, geometry)

        # 2. FIX-012: role-aware child moves (new in v5).
        for move in action.after.get("child_moves") or []:
            child_shape = _shape_by_id(slide, move.get("shape_id"))
            if child_shape is None:
                continue
            before_bbox = move.get("before_bbox_pt") or []
            if len(before_bbox) != 4:
                continue
            target_left_norm = float(move.get("after_left_pt", before_bbox[0]))
            target_top_norm = float(move.get("after_top_pt", before_bbox[1]))
            _apply_geometry(
                child_shape,
                {
                    "left": target_left_norm * sx,
                    "top": target_top_norm * sy,
                },
            )
        return

    shape = _shape_by_id(slide, action.shape_id)
    if shape is None:
        return

    if "geometry" in action.after:
        _apply_geometry(shape, action.after["geometry"])
        strategy_after = action.after.get("strategy") if isinstance(action.after, dict) else None
        # text_box_resize の expand_box_width_to_canvas は geometry だけだと
        # text 量によっては inner_h を超える残存があり得る。box 拡張後の
        # 状態で再評価し、必要なら font を追加 shrink する (compound apply)。
        if action.rule == "text_box_resize" and strategy_after == "expand_box_width_to_canvas":
            _compound_resolve_text_box_overflow(shape)
        # text_canvas_reflow の combined strategy: geometry に加えて
        # word_wrap も同時に設定する (3 段積み上げの step 2 を実現)。
        if action.rule == "text_canvas_reflow" and strategy_after == "expand_box_to_canvas_and_wrap":
            if action.after.get("word_wrap") is True:
                shape.text_frame.word_wrap = True
    elif action.rule == "image_crop":
        for side, value in action.after.get("crop", {}).items():
            setattr(shape, f"crop_{side}", float(value))
    elif action.rule == "font_family":
        _set_text_shape_font_family(shape, action.after["font_name"])
    elif action.rule == "text_color":
        _set_text_shape_color(shape, action.before.get("color_hex"), action.after["color_hex"])
    elif action.rule == "fill_color":
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(action.after["color_hex"].lstrip("#"))
        except (AttributeError, TypeError, ValueError):
            return
    elif action.rule == "text_wrap":
        _replace_text_breaks(shape)
        widen = action.after.get("widen_to_fit") if isinstance(action.after, dict) else None
        if isinstance(widen, dict):
            try:
                shape.width = Pt(float(widen["width_pt_actual"]))
            except (KeyError, TypeError, ValueError):
                pass
    elif action.rule == "heading_hierarchy":
        title = _shape_by_id(slide, action.before.get("title_shape_id"))
        body = _shape_by_id(slide, action.before.get("body_shape_id"))
        _set_shape_text_size(title, float(action.after["title_size_pt"]))
        _set_shape_text_size(body, float(action.after["body_size_pt"]))
    elif action.rule == "text_vertical_balance":
        shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    elif action.rule == "badge_alignment":
        shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                para.alignment = PP_ALIGN.CENTER
    elif action.rule == "decorative_remove":
        # FIX-013: physically remove the shape element from its slide spTree.
        # python-pptx does not expose a high-level delete API, so we drop the
        # underlying lxml element.
        sp_element = shape._element  # noqa: SLF001 — intentional XML access
        parent = sp_element.getparent()
        if parent is not None:
            parent.remove(sp_element)
    elif action.rule == "text_box_resize":
        # DS-OVERFLOW-001 段階2: SPA で選ばれた 1 つの strategy を apply。
        # box_canvas_clip と違い "geometry" を action.after に置いた場合は
        # 上の `if "geometry" in action.after` 分岐で処理済 (= expand_box_height)。
        # ここでは font_size / line_height の strategy のみを扱う。
        strategy = action.after.get("strategy") if isinstance(action.after, dict) else None
        if strategy == "shrink_font_size":
            target_pt = float(action.after.get("font_size_pt") or 0)
            if target_pt > 0:
                _set_shape_text_size(shape, target_pt)
        elif strategy == "compress_line_height":
            target_pt = float(action.after.get("line_height_pt") or 0)
            if target_pt > 0:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        para.line_spacing = Pt(target_pt)
        elif strategy == "expand_box_width_to_canvas":
            # geometry は上の `if "geometry" in action.after` 分岐で apply 済 +
            # compound 処理もそちらで実行済。この elif は dead path だが
            # 防御的に同じ compound を呼んでおく。
            _compound_resolve_text_box_overflow(shape)
    elif action.rule == "text_canvas_reflow":
        # DS-OVERFLOW-001 段階3: enable_word_wrap / shrink_font_size を扱う。
        # shrink_box_width_to_canvas は action.after.geometry 経由で上の
        # `if "geometry" in action.after` 分岐で apply 済。
        strategy = action.after.get("strategy") if isinstance(action.after, dict) else None
        if strategy == "enable_word_wrap":
            shape.text_frame.word_wrap = True
        elif strategy == "shrink_font_size":
            target_pt = float(action.after.get("font_size_pt") or 0)
            if target_pt > 0:
                _set_shape_text_size(shape, target_pt)


# ---- Driver ----------------------------------------------------------------


def fix_pptx(
    path: Path,
    *,
    apply: bool = False,
    backup: bool = False,
    rules: Sequence[str] = DEFAULT_RULES,
    findings: Optional[Sequence[Any]] = None,
    judgement_gate: bool = True,
) -> List[FixAction]:
    """Apply registered fix rules to ``path``.

    ``judgement_gate`` (POLICY-001 段階 3, default ``True``): when ``True``
    the judgement_fix policy is enforced — findings whose lint output is
    still ``manual_required`` are skipped unless the SPA judgement layer
    has already promoted them to ``auto_fix_candidate`` (typically via
    ``apply_finding_judgements_overrides``). Pass ``False`` to recover the
    pre-段階 3 behavior that auto-applied any ``manual_required`` finding
    with concrete ``candidate_values`` (kept for legacy callers such as
    the REV-017/REV-019 pipelines).
    """
    prs = Presentation(str(path))
    actions: List[FixAction] = []
    active_rules = _enabled_rules(rules)
    shape_rules = tuple(rule for rule in active_rules if rule not in FINDING_DRIVEN_RULES)

    for shape, idx, sid, slide in _walk(prs):
        for rule in shape_rules:
            if rule == "contrast":
                action = _detect_contrast(shape, idx, sid, slide, findings=findings)
            else:
                det: Optional[Callable] = DETECTORS.get(rule)
                if det is None:
                    continue
                action = det(shape, idx, sid, slide)
            if action is None:
                continue
            if rule != "contrast":
                action = _apply_matching_finding_fixability(
                    action, findings, judgement_gate=judgement_gate
                )
            actions.append(action)
            if action.status == "apply":
                _apply_action(shape, action, slide=slide)

    if findings:
        for finding in findings:
            rule = _finding_rule(finding)
            if rule not in active_rules or rule not in FINDING_DRIVEN_RULES:
                continue
            detected = _detect_finding_action(prs, finding)
            if detected is None:
                continue
            detected_actions = detected if isinstance(detected, list) else [detected]
            for action in detected_actions:
                action = _apply_finding_fixability(
                    action, finding, judgement_gate=judgement_gate
                )
                actions.append(action)
                if action.status == "apply":
                    _apply_finding_action(prs, action)

    if apply and actions:
        if backup:
            _backup_once(path)
        prs.save(str(path))

    return actions


def verify_pptx(
    path: Path,
    *,
    rules: Sequence[str] = DEFAULT_RULES,
    findings: Optional[Sequence[Any]] = None,
) -> List[FixAction]:
    """Re-read the file from disk and report any actions still pending.

    A non-empty result after a successful --apply means the change was not
    durable on disk. Common causes: corrupted source PPTX with duplicate zip
    entries (python-pptx writes both copies, only one fixed), or inherited
    bodyPr that the slide-level setter cannot override.
    """
    prs = Presentation(str(path))
    residual: List[FixAction] = []
    active_rules = _enabled_rules(rules)
    shape_rules = tuple(rule for rule in active_rules if rule not in FINDING_DRIVEN_RULES)
    for shape, idx, sid, slide in _walk(prs):
        for rule in shape_rules:
            if rule == "contrast":
                action = _detect_contrast(shape, idx, sid, slide, findings=findings)
            else:
                det = DETECTORS.get(rule)
                if det is None:
                    continue
                action = det(shape, idx, sid, slide)
            if action is not None:
                if rule != "contrast":
                    action = _apply_matching_finding_fixability(action, findings)
            if action is not None and action.status == "apply":
                residual.append(action)
    return residual


# ---- Output ----------------------------------------------------------------


def _format_loc(a: FixAction) -> str:
    if a.shape_name and a.shape_id is not None:
        return f"{a.shape_name} (id={a.shape_id})"
    return a.shape_name or "-"


def format_actions(actions: List[FixAction], applied: bool) -> str:
    if not actions:
        return "OK: no fixable issues found.\n"
    by_rule: dict = {}
    for a in actions:
        by_rule.setdefault(a.rule, []).append(a)
    head = "Applied" if applied else "Would apply"
    applyable = [a for a in actions if a.status == "apply"]
    manual = [a for a in actions if a.status == "manual_required"]
    counts = ", ".join(f"{r}: {len(by_rule[r])}" for r in sorted(by_rule))
    lines = [
        f"{head} {len(applyable)} fixes; manual_required: {len(manual)} ({counts})",
        "",
    ]
    for rule in sorted(by_rule):
        lines.append(f"--- {rule} ---")
        for a in by_rule[rule]:
            loc = _format_loc(a)
            status = "" if a.status == "apply" else "  manual_required: " + ", ".join(a.reasons)
            if rule == "autofit":
                lines.append(
                    f"slide {a.slide_index}: {loc}  {a.before['auto_size']} -> NONE{status}"
                )
            elif rule == "geometry":
                diff = ", ".join(
                    f"{k}: {a.before[k]:g}->{a.after[k]:g}pt" for k in sorted(a.after)
                )
                lines.append(f"slide {a.slide_index}: {loc}  {diff}{status}")
            elif rule == "font_size":
                examples = a.after.get("updates", [])[:2]
                diff = ", ".join(
                    (
                        f"{r['before_size_pt']:g}->{r['after_size_pt']:g}pt "
                        f"({r['normalized_size_pt']:g}->{r['target_normalized_size_pt']:g}pt normalized)"
                    )
                    for r in examples
                )
                if len(a.after.get("updates", [])) > 2:
                    diff += f", ... {len(a.after.get('updates', []))} run(s)"
                lines.append(f"slide {a.slide_index}: {loc}  {diff}{status}")
            elif rule == "line_height":
                examples = a.after.get("updates", [])[:2]
                diff = ", ".join(
                    (
                        f"{r['before_line_height_pt']:g}->{r['after_line_height_pt']:g}pt "
                        f"({r['normalized_line_height_pt']:g}->{r['target_normalized_line_height_pt']:g}pt normalized)"
                    )
                    for r in examples
                )
                if len(a.after.get("updates", [])) > 2:
                    diff += f", ... {len(a.after.get('updates', []))} paragraph(s)"
                lines.append(f"slide {a.slide_index}: {loc}  {diff}{status}")
            elif rule == "alignment":
                examples = a.after.get("updates", [])[:2]
                diff = ", ".join(
                    f"{r['property']}: {r['before']}->{r['after']}" for r in examples
                )
                if len(a.after.get("updates", [])) > 2:
                    diff += f", ... {len(a.after.get('updates', []))} update(s)"
                lines.append(f"slide {a.slide_index}: {loc}  {diff}{status}")
            elif rule == "contrast":
                examples = a.after.get("updates", [])[:2]
                diff = ", ".join(
                    f"{r['before_hex']}->{r['after_hex']}" for r in examples
                )
                if len(a.after.get("updates", [])) > 2:
                    diff += f", ... {len(a.after.get('updates', []))} run(s)"
                lines.append(f"slide {a.slide_index}: {loc}  {diff}{status}")
            else:
                if "geometry" in a.after:
                    geometry = a.after["geometry"]
                    diff = ", ".join(
                        f"{k}: {a.before.get(k)}->{geometry[k]}pt"
                        for k in ("left", "top", "width", "height")
                        if k in geometry
                    )
                else:
                    diff = json.dumps(a.after, ensure_ascii=False, sort_keys=True)
                lines.append(f"slide {a.slide_index}: {loc}  {diff}{status}")
        lines.append("")
    return "\n".join(lines)


def format_residual(residual: List[FixAction]) -> str:
    by_rule: dict = {}
    for a in residual:
        by_rule.setdefault(a.rule, []).append(a)
    counts = ", ".join(f"{r}: {len(by_rule[r])}" for r in sorted(by_rule))
    lines = [
        f"WARNING: self-check found {len(residual)} residual actions after save ({counts}).",
        "         The change did not persist on disk (corrupted PPTX zip or inherited bodyPr).",
        "         First few:",
    ]
    for a in residual[:5]:
        lines.append(
            f"  slide {a.slide_index}: {_format_loc(a)}  rule={a.rule}  before={a.before}"
        )
    if len(residual) > 5:
        lines.append(f"  ... and {len(residual) - 5} more")
    lines.append("")
    return "\n".join(lines)


def load_findings_json(path: Path) -> list[Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if isinstance(payload, list):
        return payload
    if isinstance(payload, dict):
        findings = payload.get("findings")
        if isinstance(findings, list):
            return findings
    raise ValueError("findings JSON must be a lint JSON array or an object with findings[]")


def _judgement_finding_key(check: Any, slide_index: Any, shape_id: Any) -> Optional[str]:
    """Reproduce the SPA's fallbackKey(check, slideIndex, shapeId) format.

    The web UI uses `${check}:${slideIndex}:${shapeId ?? "unknown"}` whenever a
    finding lacks an explicit `group_key` field, and finding-judgements.json
    is keyed by that same string.
    """
    if not isinstance(check, str) or not check:
        return None
    if slide_index is None:
        return None
    sid = shape_id if shape_id is not None else "unknown"
    return f"{check}:{slide_index}:{sid}"


def apply_finding_judgements_overrides(
    findings: Sequence[Any],
    judgements_data: Any,
) -> int:
    """Promote `fixability=manual_required` to `auto_fix_candidate` when a
    reviewer judgement says `judgement_reason=auto_fixable`.

    Mutates the matched finding dicts in place. `judgements_data` is the parsed
    JSON object from `doc/reviews/<deck>/rev-<rev>-finding-judgements.json`
    (shape: `{"deck": ..., "rev": ..., "judgements": {key: {...}}}`).
    Returns the number of findings whose fixability was promoted.

    Only structural-XML lint findings carry stable group_key strings the UI
    can echo; rendered-image findings without `group_key`/`shape_id` fall back
    to `:unknown` and will not collide.
    """
    if not isinstance(judgements_data, dict):
        return 0
    raw_judgements = judgements_data.get("judgements")
    if not isinstance(raw_judgements, dict) or not raw_judgements:
        return 0
    promoted = 0
    for finding in findings:
        check = _finding_field(finding, "check")
        slide_index = _finding_field(finding, "slide_index")
        shape_id = _finding_field(finding, "shape_id")
        # Prefer explicit group_key if the JSON carries one.
        detail = _finding_detail(finding) or {}
        key = (
            (detail.get("group_key") if isinstance(detail, dict) else None)
            or _judgement_finding_key(check, slide_index, shape_id)
        )
        if not key:
            continue
        judgement = raw_judgements.get(key)
        if not isinstance(judgement, dict):
            continue
        if judgement.get("judgement_reason") != "auto_fixable":
            continue
        if detail.get("fixability") == "auto_fix_candidate":
            continue
        detail["fixability"] = "auto_fix_candidate"
        detail.setdefault(
            "fixability_reason",
            "promoted by reviewer judgement_reason=auto_fixable",
        )
        detail.setdefault("manual_required_reason", None)
        if isinstance(finding, dict):
            finding["detail"] = detail
        promoted += 1
    return promoted


def auto_rules_from_findings(findings: Sequence[Any]) -> tuple[str, ...]:
    """Return fixer rules that can be attempted from lint evidence.

    `manual_required` means a reviewer must decide whether to accept the
    repair, not that the mechanical mutation is impossible. Select every rule
    with an implementation; rules without enough evidence will return no action.
    """
    selected: set[str] = set()
    for finding in findings:
        detail = _finding_detail(finding)
        rule = detail.get("fixability_rule") or _finding_rule(finding)
        if rule in ALL_RULES and RULE_ENABLED.get(rule, True):
            selected.add(rule)
    return tuple(rule for rule in ALL_RULES if rule in selected)


# ---- CLI -------------------------------------------------------------------


def main(argv: List[str]) -> int:
    ap = argparse.ArgumentParser(description="PPTX auto-fixer (v1 guideline)")
    ap.add_argument("pptx", type=Path, help="path to .pptx")
    ap.add_argument(
        "--apply",
        action="store_true",
        help="write changes back to the file (default: dry-run)",
    )
    ap.add_argument(
        "--backup",
        action="store_true",
        help="with --apply, copy the original to <path>.bak before saving if absent",
    )
    ap.add_argument(
        "--rules",
        default=",".join(DEFAULT_RULES),
        help=(
            f"comma-separated subset of {{{','.join(ALL_RULES)}}} "
            f"(default: {','.join(DEFAULT_RULES)})"
        ),
    )
    ap.add_argument(
        "--auto",
        action="store_true",
        help=(
            "derive rules and candidate values from lint evidence and process every "
            "finding with an implemented mechanical fixer"
        ),
    )
    ap.add_argument(
        "--profile",
        choices=["default", "strict"],
        default="default",
        help="lint policy profile used with --auto (default: default)",
    )
    ap.add_argument(
        "--rendered-image-dir",
        type=Path,
        help="rendered slide PNG directory passed to lint when using --auto",
    )
    ap.add_argument("--json", action="store_true", help="emit actions as JSON")
    ap.add_argument(
        "--findings-json",
        type=Path,
        help=(
            "lint JSON to consume for evidence-schema fixability and candidate values; "
            "required for --rules contrast"
        ),
    )
    ap.add_argument(
        "--finding-judgements-json",
        type=Path,
        help=(
            "reviewer judgement JSON (doc/reviews/<deck>/rev-<rev>-finding-judgements.json); "
            "promotes manual_required → auto_fix_candidate for findings whose "
            "judgement_reason is auto_fixable"
        ),
    )
    ap.add_argument(
        "--no-judgement-gate",
        action="store_true",
        help=(
            "disable POLICY-001 段階3 strict gate: when set, judgement_fix-policy "
            "findings with manual_required + candidate_values are auto-applied "
            "without requiring an SPA judgement. Default is strict (manual_required "
            "findings under judgement_fix mode are skipped unless promoted by "
            "--finding-judgements-json)."
        ),
    )
    args = ap.parse_args(argv)

    if not args.pptx.exists():
        print(f"file not found: {args.pptx}", file=sys.stderr)
        return 1

    rules = [r.strip() for r in args.rules.split(",") if r.strip()]
    bad = [r for r in rules if r not in ALL_RULES]
    if bad:
        print(
            f"unknown rule(s): {bad}; valid: {list(ALL_RULES)}",
            file=sys.stderr,
        )
        return 1

    if args.backup and not args.apply:
        print("note: --backup has no effect without --apply", file=sys.stderr)

    findings: Optional[list[Any]] = None
    if args.findings_json:
        try:
            findings = load_findings_json(args.findings_json)
        except (OSError, ValueError, json.JSONDecodeError) as exc:
            print(f"failed to load --findings-json: {exc}", file=sys.stderr)
            return 1

    if args.auto:
        if findings is None:
            try:
                import pptx_lint  # type: ignore
            except ImportError as exc:
                print(f"failed to import pptx_lint for --auto: {exc}", file=sys.stderr)
                return 1
            raw_findings = pptx_lint.lint_pptx(
                args.pptx,
                profile=args.profile,
                rendered_image_dir=args.rendered_image_dir,
            )
            findings = [pptx_lint.finding_to_json_dict(finding) for finding in raw_findings]

    if args.finding_judgements_json:
        if findings is None:
            print(
                "--finding-judgements-json requires --findings-json or --auto",
                file=sys.stderr,
            )
            return 1
        try:
            judgements_data = json.loads(
                args.finding_judgements_json.read_text(encoding="utf-8")
            )
        except (OSError, ValueError, json.JSONDecodeError) as exc:
            print(
                f"failed to load --finding-judgements-json: {exc}",
                file=sys.stderr,
            )
            return 1
        promoted = apply_finding_judgements_overrides(findings, judgements_data)
        if promoted:
            print(
                f"promoted {promoted} finding(s) to auto_fix_candidate via judgement_reason",
                file=sys.stderr,
            )

    if args.auto:
        rules = list(auto_rules_from_findings(findings))

    if "contrast" in rules and findings is None:
        print("--rules contrast requires --findings-json from pptx_lint.py --json --no-consolidate, or use --auto", file=sys.stderr)
        return 1

    actions = fix_pptx(
        args.pptx,
        apply=args.apply,
        backup=args.backup,
        rules=rules,
        findings=findings,
        judgement_gate=not args.no_judgement_gate,
    )

    residual: List[FixAction] = []
    if args.apply and actions:
        residual = verify_pptx(args.pptx, rules=rules, findings=findings)

    if args.json:
        payload = {
            "applied": args.apply,
            "actions": [asdict(a) for a in actions],
            "residual": [asdict(a) for a in residual],
        }
        print(json.dumps(payload, ensure_ascii=False, indent=2))
    else:
        print(format_actions(actions, applied=args.apply), end="")
        if residual:
            print(format_residual(residual), end="", file=sys.stderr)

    return 2 if residual else 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
