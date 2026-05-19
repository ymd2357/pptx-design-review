#!/usr/bin/env python3
"""Smoke test for pptx_fix against generated fixtures.

Asserts:
- bad.pptx text_autofit_disabled findings drop to 0 after autofit fix
- non-targeted lint checks are not regressed (count unchanged)
- geometry rounding fixture: drifted (81.05pt) coord becomes exactly 81pt;
  intentional half-pt (81.5pt) is left untouched
- dry-run does not modify the file on disk
- existing .bak files are not overwritten by --backup
- self-check (verify_pptx) reports zero residual on a successfully fixed file

Exit code: 0 on success, 1 on any failed assertion.
"""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))

import make_examples  # noqa: E402
import pptx_fix  # noqa: E402
import pptx_lint  # noqa: E402


DRIFT_LEFT_EMU = round(81.05 * 12700)  # 1029335 -> should round to 81pt
HALF_LEFT_EMU = round(81.5 * 12700)    # 1034605 -> drift 0.5pt, must stay


def _build_geometry_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Emu(DRIFT_LEFT_EMU), Pt(40), Pt(200), Pt(50))
    slide.shapes.add_textbox(Emu(HALF_LEFT_EMU), Pt(120), Pt(200), Pt(50))
    prs.save(str(out))


def _build_safe_font_size_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40), Pt(40), Pt(300), Pt(80))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Short label"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(14.75)
    prs.save(str(out))


def _build_manual_font_size_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40), Pt(40), Pt(300), Pt(27))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Cover label"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(22.5)
    prs.save(str(out))


def _build_line_height_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40), Pt(40), Pt(420), Pt(90))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    para = box.text_frame.paragraphs[0]
    para.line_spacing = Pt(16.4)
    run = para.add_run()
    run.text = "Line height fixture"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(12)
    prs.save(str(out))


def _build_wrap_break_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(120), Pt(180), Pt(900), Pt(60))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Auto\nmation improves review speed"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    prs.save(str(out))


def _build_wrap_break_widen_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(120), Pt(180), Pt(200), Pt(60))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    box.text_frame.word_wrap = True
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Auto\nmation improves review speed"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    prs.save(str(out))


def _build_alignment_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(40), Pt(40), Pt(300), Pt(80))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    para = box.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "Alignment fixture"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(12)
    prs.save(str(out))


def _build_slide_size_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1000)
    prs.slide_height = Pt(500)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Pt(80), Pt(60), Pt(600), Pt(80))
    prs.save(str(out))


def _build_contrast_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(120), Pt(120), Pt(300), Pt(60))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Low contrast"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor.from_string("999999")
    box = slide.shapes.add_textbox(Pt(120), Pt(210), Pt(300), Pt(60))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Moderate contrast"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor.from_string("858585")
    prs.save(str(out))


def _build_bg_mode_contrast_fixture(out: Path) -> int:
    """White text inside a light-gray-filled rectangle.

    Foreground swap would change white→dark (huge luminance flip), so
    `_contrast_candidate` should prefer the background-mode strategy:
    keep the white run and darken the rectangle's fill to a passing color.
    """
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(120), Pt(120), Pt(400), Pt(80))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor.from_string("EBEBEB")
    rect.line.fill.background()
    tf = rect.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "White on light gray"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor.from_string("FFFFFF")
    prs.save(str(out))
    return rect.shape_id


def _build_behind_solid_fill_contrast_fixture(out: Path) -> tuple[int, int]:
    """White textbox sitting on top of a separate red-fill rectangle.

    Differs from `_build_bg_mode_contrast_fixture` (= white text *inside*
    the rect → `shape_solid_fill`) by putting the text in its own shape
    *above* the rect, with no own fill. Lint must report
    `background_source = behind_solid_fill:<rect-label>` and fix must
    repaint the *rect* fill while leaving the textbox run color white.
    Returns `(rect_id, textbox_id)`.
    """
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(120), Pt(120), Pt(400), Pt(120))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor.from_string("FF5757")
    rect.line.fill.background()

    textbox = slide.shapes.add_textbox(Pt(140), Pt(150), Pt(360), Pt(60))
    textbox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = textbox.text_frame.paragraphs[0].add_run()
    run.text = "White on red"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor.from_string("FFFFFF")

    prs.save(str(out))
    return rect.shape_id, textbox.shape_id


def _build_text_color_allowlist_fixture(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Pt(120), Pt(120), Pt(300), Pt(60))
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Brand color review"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor.from_string("FF0000")
    prs.save(str(out))


def _build_text_overlap_fixture(out: Path) -> tuple[int, int]:
    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    first = slide.shapes.add_textbox(Pt(120), Pt(100), Pt(260), Pt(60))
    first.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = first.text_frame.paragraphs[0].add_run()
    run.text = "Text overlap A"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    second = slide.shapes.add_textbox(Pt(140), Pt(120), Pt(260), Pt(60))
    second.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    run = second.text_frame.paragraphs[0].add_run()
    run.text = "Text overlap B"
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(24)
    ids = (first.shape_id, second.shape_id)
    prs.save(str(out))
    return ids


def _build_card_grid_picture_fixture(out: Path, tmp_dir: Path) -> tuple[dict, tuple[int, int, int]]:
    img = tmp_dir / "card-grid-picture.png"
    Image.new("RGB", (64, 64), (30, 120, 180)).save(img)

    prs = Presentation()
    prs.slide_width = Pt(1440)
    prs.slide_height = Pt(810)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(90), Pt(250), Pt(300), Pt(200))
    pic = slide.shapes.add_picture(str(img), Pt(120), Pt(290), width=Pt(31.5), height=Pt(31.5))
    title = slide.shapes.add_textbox(Pt(160), Pt(280), Pt(180), Pt(40))
    title.text_frame.text = "Card title"
    body = slide.shapes.add_textbox(Pt(120), Pt(340), Pt(240), Pt(60))
    body.text_frame.text = "Card body"
    prs.save(str(out))

    finding = {
        "check": "card_grid_consistency",
        "slide_index": 1,
        "slide_id": 256,
        "detail": {
            "fixability": "auto_fix_candidate",
            "fixability_rule": "card_grid",
            "candidate_values": {
                "group_medians": {
                    "top": 229.5,
                    "width": 298.0,
                    "height": 218.67,
                    "padding_left": 28.5,
                    "padding_right": 16.32,
                    "padding_top": 28.5,
                    "padding_bottom": 28.5,
                }
            },
            "group_medians": {
                "top": 229.5,
                "width": 298.0,
                "height": 218.67,
                "padding_left": 28.5,
                "padding_right": 16.32,
                "padding_top": 28.5,
                "padding_bottom": 28.5,
            },
            "inconsistent_containers": [
                {
                    "container": {
                        "shape_id": card.shape_id,
                        "kind": "shape",
                        "bbox_pt": [90.0, 250.0, 300.0, 200.0],
                    },
                    "children": [
                        {"shape_id": pic.shape_id, "kind": "image", "bbox_pt": [120.0, 290.0, 31.5, 31.5]},
                        {"shape_id": title.shape_id, "kind": "text", "bbox_pt": [160.0, 280.0, 180.0, 40.0]},
                        {"shape_id": body.shape_id, "kind": "text", "bbox_pt": [120.0, 340.0, 240.0, 60.0]},
                    ],
                    "padding_pt": {"left": 30.0, "right": 30.0, "top": 30.0, "bottom": 50.0},
                }
            ],
        },
    }
    return finding, (pic.shape_id, title.shape_id, body.shape_id)


def main() -> int:
    failures: list = []

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)

        # --- autofit fix on bad.pptx ---
        bad = tmp_dir / "bad.pptx"
        make_examples.make_bad(bad)

        pre = pptx_lint.lint_pptx(bad)
        pre_autofit = [f for f in pre if f.check == "text_autofit_disabled"]
        if not pre_autofit:
            failures.append("bad.pptx had no autofit findings to fix (fixture drift?)")

        actions = pptx_fix.fix_pptx(bad, apply=True, rules=("autofit",))
        if not any(a.rule == "autofit" for a in actions):
            failures.append("fixer reported no autofit actions on bad.pptx")

        post = pptx_lint.lint_pptx(bad)
        post_autofit = [f for f in post if f.check == "text_autofit_disabled"]
        if post_autofit:
            failures.append(
                f"autofit findings remain after fix: {len(post_autofit)}"
            )

        pre_other = sorted(f.check for f in pre if f.check != "text_autofit_disabled")
        post_other = sorted(f.check for f in post if f.check != "text_autofit_disabled")
        if pre_other != post_other:
            failures.append(
                f"non-targeted checks changed: pre={pre_other} post={post_other}"
            )

        # --- geometry fix on a drift+half fixture ---
        geo = tmp_dir / "geo.pptx"
        _build_geometry_fixture(geo)

        actions = pptx_fix.fix_pptx(geo, apply=True, rules=("geometry",))
        if not any(a.rule == "geometry" for a in actions):
            failures.append("geometry fixer reported no actions on drifted fixture")

        prs = Presentation(str(geo))
        shapes = list(prs.slides[0].shapes)
        drift_left_pt = shapes[0].left / 12700
        half_left_pt = shapes[1].left / 12700
        if abs(drift_left_pt - 81.0) > 1e-6:
            failures.append(
                f"drift_box.left expected 81pt; got {drift_left_pt}"
            )
        if abs(half_left_pt - 81.5) > 1e-6:
            failures.append(
                f"half_box.left should stay at 81.5pt; got {half_left_pt}"
            )

        # --- dry-run must not write ---
        bad2 = tmp_dir / "bad2.pptx"
        make_examples.make_bad(bad2)
        before_mtime = bad2.stat().st_mtime_ns
        actions = pptx_fix.fix_pptx(bad2, apply=False, rules=("autofit",))
        if not actions:
            failures.append("dry-run produced no actions on bad.pptx")
        if bad2.stat().st_mtime_ns != before_mtime:
            failures.append("dry-run modified the file on disk")

        # --- contrast fix consumes lint candidate_values explicitly ---
        contrast = tmp_dir / "contrast.pptx"
        _build_contrast_fixture(contrast)
        contrast_findings = [
            pptx_lint.finding_to_json_dict(f)
            for f in pptx_lint.lint_pptx(contrast)
        ]
        low_contrast = [
            f
            for f in contrast_findings
            if f["check"] == "low_contrast"
        ]
        if not low_contrast:
            failures.append("contrast fixture did not trigger low_contrast")
        elif low_contrast[0]["detail"].get("fixability_rule") != "contrast":
            failures.append(
                "low_contrast did not expose contrast fixability: "
                f"{low_contrast[0]['detail']}"
            )
        auto_rules = pptx_fix.auto_rules_from_findings(contrast_findings)
        if "contrast" not in auto_rules:
            failures.append(f"auto rules did not include contrast: {auto_rules}")
        actions = pptx_fix.fix_pptx(
            contrast,
            apply=True,
            rules=auto_rules,
            findings=contrast_findings,
        )
        contrast_actions = [a for a in actions if a.rule == "contrast" and a.status == "apply"]
        if len(contrast_actions) != 2:
            failures.append("contrast fixer did not apply lint candidate color")
        prs = Presentation(str(contrast))
        fixed_color = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.color.rgb
        if str(fixed_color).upper() != "707070":
            failures.append(f"contrast fixture expected #707070; got #{fixed_color}")
        fixed_color = prs.slides[0].shapes[1].text_frame.paragraphs[0].runs[0].font.color.rgb
        if str(fixed_color).upper() != "707070":
            failures.append(f"moderate contrast fixture expected #707070; got #{fixed_color}")
        red_candidate = pptx_lint._contrast_candidate("#FF0000", "#FFFFFF", 4.5)
        if red_candidate is None or red_candidate.get("foreground_hex") != "#E6033D":
            failures.append(f"red contrast repair expected #E6033D; got {red_candidate}")
        residual = pptx_fix.verify_pptx(
            contrast,
            rules=auto_rules,
            findings=contrast_findings,
        )
        if residual:
            failures.append(f"contrast verify reported residual: {len(residual)}")

        # --- bg-mode contrast fix: white run + light-fill rect → fill is darkened ---
        bg_mode_pptx = tmp_dir / "contrast-bg-mode.pptx"
        rect_id = _build_bg_mode_contrast_fixture(bg_mode_pptx)
        bg_findings = [
            pptx_lint.finding_to_json_dict(f)
            for f in pptx_lint.lint_pptx(bg_mode_pptx)
            if f.check in {"low_contrast", "contrast_ratio"}
        ]
        if not bg_findings:
            failures.append("bg-mode fixture did not trigger any contrast finding")
        else:
            cv = bg_findings[0]["detail"].get("candidate_values") or {}
            if cv.get("preferred_strategy") != "background":
                failures.append(
                    "bg-mode fixture expected preferred_strategy=background; "
                    f"got {cv.get('preferred_strategy')}"
                )
            if bg_findings[0]["detail"].get("fixability") != "auto_fix_candidate":
                failures.append(
                    "bg-mode contrast finding was not marked auto_fix_candidate: "
                    f"{bg_findings[0]['detail']}"
                )
            bg_rules = pptx_fix.auto_rules_from_findings(bg_findings)
            bg_actions = pptx_fix.fix_pptx(
                bg_mode_pptx,
                apply=True,
                rules=bg_rules,
                findings=bg_findings,
            )
            applied = [a for a in bg_actions if a.rule == "contrast" and a.status == "apply"]
            if not applied:
                failures.append("bg-mode contrast was not applied as a fill change")
            else:
                updates = applied[0].after.get("updates", []) or []
                if not any(u.get("mode") == "background_fill" for u in updates):
                    failures.append(
                        "bg-mode update did not record background_fill mode: "
                        f"{updates}"
                    )
            prs = Presentation(str(bg_mode_pptx))
            shapes_by_id = {s.shape_id: s for s in prs.slides[0].shapes}
            target_rect = shapes_by_id.get(rect_id)
            if target_rect is None:
                failures.append("bg-mode fixture rectangle disappeared after fix")
            else:
                new_fill = pptx_fix._shape_solid_fill_hex(target_rect)
                if new_fill == "#EBEBEB":
                    failures.append("bg-mode contrast did not change the rectangle fill")
                else:
                    run = target_rect.text_frame.paragraphs[0].runs[0]
                    if str(run.font.color.rgb).upper() != "FFFFFF":
                        failures.append(
                            "bg-mode contrast must keep the original white run color; "
                            f"got #{run.font.color.rgb}"
                        )

        # --- FIX-010: behind_solid_fill bg-mode contrast (textbox on top of rect) ---
        behind_pptx = tmp_dir / "contrast-behind-solid-fill.pptx"
        behind_rect_id, behind_textbox_id = _build_behind_solid_fill_contrast_fixture(behind_pptx)
        behind_findings = [
            pptx_lint.finding_to_json_dict(f)
            for f in pptx_lint.lint_pptx(behind_pptx)
            if f.check in {"low_contrast", "contrast_ratio"}
        ]
        if not behind_findings:
            failures.append("behind_solid_fill fixture did not trigger any contrast finding")
        else:
            bg_source = behind_findings[0]["detail"].get("background_source") or ""
            if not bg_source.startswith("behind_solid_fill"):
                failures.append(
                    "behind_solid_fill fixture expected background_source="
                    f"behind_solid_fill:*; got {bg_source}"
                )
            cv = behind_findings[0]["detail"].get("candidate_values") or {}
            if cv.get("preferred_strategy") != "background":
                failures.append(
                    "behind_solid_fill fixture expected preferred_strategy=background; "
                    f"got {cv.get('preferred_strategy')}"
                )
            if behind_findings[0]["detail"].get("fixability") != "auto_fix_candidate":
                failures.append(
                    "behind_solid_fill contrast finding was not marked auto_fix_candidate: "
                    f"{behind_findings[0]['detail']}"
                )
            behind_rules = pptx_fix.auto_rules_from_findings(behind_findings)
            behind_actions = pptx_fix.fix_pptx(
                behind_pptx,
                apply=True,
                rules=behind_rules,
                findings=behind_findings,
            )
            applied = [a for a in behind_actions if a.rule == "contrast" and a.status == "apply"]
            if not applied:
                failures.append("behind_solid_fill contrast was not applied as a fill change")
            else:
                updates = applied[0].after.get("updates", []) or []
                if not any(u.get("mode") == "background_fill" for u in updates):
                    failures.append(
                        "behind_solid_fill update did not record background_fill mode: "
                        f"{updates}"
                    )
            prs = Presentation(str(behind_pptx))
            shapes_by_id = {s.shape_id: s for s in prs.slides[0].shapes}
            rect_after = shapes_by_id.get(behind_rect_id)
            textbox_after = shapes_by_id.get(behind_textbox_id)
            if rect_after is None:
                failures.append("behind_solid_fill fixture rectangle disappeared after fix")
            else:
                new_fill = pptx_fix._shape_solid_fill_hex(rect_after)
                if not new_fill or new_fill.upper() == "#FF5757":
                    failures.append(
                        "behind_solid_fill bg-mode did not change the rectangle fill: "
                        f"{new_fill}"
                    )
            if textbox_after is None:
                failures.append("behind_solid_fill fixture textbox disappeared after fix")
            else:
                run = textbox_after.text_frame.paragraphs[0].runs[0]
                if str(run.font.color.rgb).upper() != "FFFFFF":
                    failures.append(
                        "behind_solid_fill must keep the original white run color; "
                        f"got #{run.font.color.rgb}"
                    )
                if pptx_fix._shape_solid_fill_hex(textbox_after):
                    failures.append(
                        "behind_solid_fill must not give the textbox a new solid fill"
                    )

        # --- FIX-007: judgement_reason=auto_fixable promotes manual to auto ---
        red_text = tmp_dir / "red-text.pptx"
        _build_text_color_allowlist_fixture(red_text)
        red_findings = [
            pptx_lint.finding_to_json_dict(f)
            for f in pptx_lint.lint_pptx(red_text)
            if f.check == "text_color_allowlist"
        ]
        if not red_findings:
            failures.append("text_color_allowlist fixture did not trigger a finding")
        else:
            base = red_findings[0]
            if base["detail"].get("fixability") != "manual_required":
                failures.append(
                    "text_color_allowlist should default to manual_required; got "
                    f"{base['detail'].get('fixability')}"
                )
            check = base["check"]
            slide_index = base["slide_index"]
            shape_id = base.get("shape_id")
            judgement_key = pptx_fix._judgement_finding_key(
                check, slide_index, shape_id
            )
            judgements_data = {
                "deck": "synthetic",
                "rev": "test",
                "judgements": {
                    judgement_key: {
                        "review_status": "fixed",
                        "judgement_reason": "auto_fixable",
                    }
                },
            }
            promoted = pptx_fix.apply_finding_judgements_overrides(
                red_findings, judgements_data
            )
            if promoted != 1:
                failures.append(f"expected 1 judgement promotion; got {promoted}")
            if red_findings[0]["detail"].get("fixability") != "auto_fix_candidate":
                failures.append(
                    "judgement promotion did not flip fixability to auto_fix_candidate"
                )
            promoted_rules = pptx_fix.auto_rules_from_findings(red_findings)
            if "text_color" not in promoted_rules:
                failures.append(
                    "auto_rules_from_findings did not pick up text_color after "
                    f"promotion; got {promoted_rules}"
                )
            actions = pptx_fix.fix_pptx(
                red_text,
                apply=True,
                rules=promoted_rules,
                findings=red_findings,
            )
            applied = [a for a in actions if a.rule == "text_color" and a.status == "apply"]
            if not applied:
                failures.append(
                    "FIX-007 promotion did not trigger any text_color apply action"
                )
            prs = Presentation(str(red_text))
            run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
            if str(run.font.color.rgb).upper() == "FF0000":
                failures.append(
                    "judgement-promoted text_color fix did not replace #FF0000"
                )

        # --- finding-driven text_overlap fixes the detected shape_b only ---
        overlap = tmp_dir / "text-overlap.pptx"
        first_id, second_id = _build_text_overlap_fixture(overlap)
        overlap_findings = [
            pptx_lint.finding_to_json_dict(f)
            for f in pptx_lint.lint_pptx(overlap)
            if f.check == "text_overlap"
        ]
        if not overlap_findings:
            failures.append("text_overlap fixture did not trigger text_overlap")
        overlap_rules = pptx_fix.auto_rules_from_findings(overlap_findings)
        if "overlap" not in overlap_rules:
            failures.append(f"auto rules did not include overlap: {overlap_rules}")
        actions = pptx_fix.fix_pptx(
            overlap,
            apply=True,
            rules=overlap_rules,
            findings=overlap_findings,
        )
        overlap_actions = [a for a in actions if a.rule == "overlap" and a.status == "apply"]
        if len(overlap_actions) != 1:
            failures.append(f"text_overlap fixer expected one action; got {len(overlap_actions)}")
        prs = Presentation(str(overlap))
        shapes = {shape.shape_id: shape for shape in prs.slides[0].shapes}
        if shapes[first_id].top != Pt(100):
            failures.append("text_overlap fixer moved shape_a; expected only shape_b to move")
        if shapes[second_id].top == Pt(120) and shapes[second_id].left == Pt(140):
            failures.append("text_overlap fixer did not move shape_b")
        remaining_overlap = [f for f in pptx_lint.lint_pptx(overlap) if f.check == "text_overlap"]
        if remaining_overlap:
            failures.append(f"text_overlap remained after fix: {len(remaining_overlap)}")

        # --- card_grid preserves child sizes, especially pictures ---
        card_grid = tmp_dir / "card-grid-picture.pptx"
        card_grid_finding, (pic_id, title_id, body_id) = _build_card_grid_picture_fixture(card_grid, tmp_dir)
        prs = Presentation(str(card_grid))
        shapes = {shape.shape_id: shape for shape in prs.slides[0].shapes}
        before = {sid: pptx_fix._shape_geometry_pt(shapes[sid]) for sid in (pic_id, title_id, body_id)}
        actions = pptx_fix.fix_pptx(
            card_grid,
            apply=True,
            rules=("card_grid",),
            findings=[card_grid_finding],
        )
        if not any(a.rule == "card_grid" and a.status == "apply" for a in actions):
            failures.append("card_grid fixer did not report an apply action")
        prs = Presentation(str(card_grid))
        shapes = {shape.shape_id: shape for shape in prs.slides[0].shapes}
        after = {sid: pptx_fix._shape_geometry_pt(shapes[sid]) for sid in (pic_id, title_id, body_id)}
        for sid, label in ((pic_id, "picture"), (title_id, "title text"), (body_id, "body text")):
            if after[sid]["width"] != before[sid]["width"] or after[sid]["height"] != before[sid]["height"]:
                failures.append(f"card_grid changed {label} size: before={before[sid]} after={after[sid]}")
        if abs((after[title_id]["left"] - after[pic_id]["left"]) - (before[title_id]["left"] - before[pic_id]["left"])) > 0.01:
            failures.append("card_grid did not preserve child horizontal offsets")

        unsafe_card_grid = tmp_dir / "card-grid-picture-unsafe.pptx"
        unsafe_finding, (unsafe_pic_id, unsafe_title_id, unsafe_body_id) = _build_card_grid_picture_fixture(
            unsafe_card_grid,
            tmp_dir,
        )
        unsafe_finding["detail"]["group_medians"]["height"] = 120.0
        unsafe_finding["detail"]["candidate_values"]["group_medians"]["height"] = 120.0
        prs = Presentation(str(unsafe_card_grid))
        shapes = {shape.shape_id: shape for shape in prs.slides[0].shapes}
        unsafe_before = {
            sid: pptx_fix._shape_geometry_pt(shapes[sid])
            for sid in (unsafe_pic_id, unsafe_title_id, unsafe_body_id)
        }
        unsafe_actions = pptx_fix.fix_pptx(
            unsafe_card_grid,
            apply=True,
            rules=("card_grid",),
            findings=[unsafe_finding],
        )
        if not any(a.rule == "card_grid" and a.status == "manual_required" for a in unsafe_actions):
            failures.append("unsafe card_grid fixer did not report manual_required")
        prs = Presentation(str(unsafe_card_grid))
        shapes = {shape.shape_id: shape for shape in prs.slides[0].shapes}
        unsafe_after = {
            sid: pptx_fix._shape_geometry_pt(shapes[sid])
            for sid in (unsafe_pic_id, unsafe_title_id, unsafe_body_id)
        }
        if unsafe_after != unsafe_before:
            failures.append("manual_required card_grid changed child geometry")

        # --- backup must preserve the oldest saved state ---
        backup_deck = tmp_dir / "backup.pptx"
        make_examples.make_bad(backup_deck)
        backup_path = Path(str(backup_deck) + ".bak")
        sentinel = b"existing original backup"
        backup_path.write_bytes(sentinel)
        pptx_fix.fix_pptx(backup_deck, apply=True, backup=True, rules=("autofit",))
        if backup_path.read_bytes() != sentinel:
            failures.append("--backup overwrote an existing .bak file")

        # --- font_size applies only when fit checks pass ---
        safe_font = tmp_dir / "safe-font.pptx"
        _build_safe_font_size_fixture(safe_font)
        actions = pptx_fix.fix_pptx(safe_font, apply=True, rules=("font_size",))
        if not any(a.rule == "font_size" and a.status == "apply" for a in actions):
            failures.append("font_size fixer did not report an apply action for safe fixture")
        prs = Presentation(str(safe_font))
        fixed_size = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.size.pt
        if abs(fixed_size - 14.0) > 1e-6:
            failures.append(f"font_size safe fixture expected 14pt; got {fixed_size}")

        manual_font = tmp_dir / "manual-font.pptx"
        _build_manual_font_size_fixture(manual_font)
        actions = pptx_fix.fix_pptx(manual_font, apply=True, rules=("font_size",))
        if not any(a.rule == "font_size" and a.status == "manual_required" for a in actions):
            failures.append("font_size fixer did not report manual_required for unsafe fixture")
        prs = Presentation(str(manual_font))
        unchanged_size = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.size.pt
        if abs(unchanged_size - 22.5) > 1e-6:
            failures.append(f"font_size manual fixture should remain 22.5pt; got {unchanged_size}")

        # --- line_height and alignment safe mechanical fixes ---
        line_height = tmp_dir / "line-height.pptx"
        _build_line_height_fixture(line_height)
        actions = pptx_fix.fix_pptx(line_height, apply=True, rules=("line_height",))
        if not any(a.rule == "line_height" and a.status == "apply" for a in actions):
            failures.append("line_height fixer did not report an apply action")
        prs = Presentation(str(line_height))
        fixed_line_height = prs.slides[0].shapes[0].text_frame.paragraphs[0].line_spacing.pt
        if abs(fixed_line_height - 15.0) > 1e-6:
            failures.append(f"line_height fixture expected 15pt; got {fixed_line_height}")

        alignment = tmp_dir / "alignment.pptx"
        _build_alignment_fixture(alignment)
        actions = pptx_fix.fix_pptx(alignment, apply=True, rules=("alignment",))
        if not any(a.rule == "alignment" and a.status == "apply" for a in actions):
            failures.append("alignment fixer did not report an apply action")
        prs = Presentation(str(alignment))
        box = prs.slides[0].shapes[0]
        if box.text_frame.vertical_anchor != MSO_VERTICAL_ANCHOR.TOP:
            failures.append(f"alignment fixture expected TOP anchor; got {box.text_frame.vertical_anchor}")
        if box.text_frame.paragraphs[0].alignment != PP_ALIGN.LEFT:
            failures.append(f"alignment fixture expected LEFT paragraph; got {box.text_frame.paragraphs[0].alignment}")

        # --- wrap-break latin word split joins the word, not a space ---
        wrap = tmp_dir / "wrap-break.pptx"
        _build_wrap_break_fixture(wrap)
        findings = [
            pptx_lint.finding_to_json_dict(f)
            for f in pptx_lint.lint_pptx(wrap)
            if f.check == "wrap_break_changes_meaning"
        ]
        if not findings:
            failures.append("wrap-break fixture did not trigger wrap_break_changes_meaning")
        actions = pptx_fix.fix_pptx(wrap, apply=True, rules=("text_wrap",), findings=findings)
        if not any(a.rule == "text_wrap" and a.status == "apply" for a in actions):
            failures.append("text_wrap fixer did not report an apply action")
        prs = Presentation(str(wrap))
        fixed_text = prs.slides[0].shapes[0].text_frame.text
        if fixed_text != "Automation improves review speed":
            failures.append(f"text_wrap fixture expected joined word; got {fixed_text!r}")
        residual_wrap = [f for f in pptx_lint.lint_pptx(wrap) if f.check == "wrap_break_changes_meaning"]
        if residual_wrap:
            failures.append(f"wrap_break_changes_meaning remained after fix: {len(residual_wrap)}")

        # --- wrap-break widen-to-fit candidate widens the shape [FIX-001] ---
        wrap_widen = tmp_dir / "wrap-break-widen.pptx"
        _build_wrap_break_widen_fixture(wrap_widen)
        widen_findings_raw = [
            f
            for f in pptx_lint.lint_pptx(wrap_widen)
            if f.check == "wrap_break_changes_meaning"
        ]
        if not widen_findings_raw:
            failures.append("widen fixture did not trigger wrap_break_changes_meaning")
        else:
            widen_detail = widen_findings_raw[0].detail
            cands = widen_detail.get("candidate_values") or []
            widen_candidate = next(
                (c for c in cands if isinstance(c, dict) and c.get("strategy") == "widen_to_fit"),
                None,
            )
            if widen_candidate is None:
                failures.append(
                    "widen fixture: expected candidate_values.widen_to_fit; got "
                    f"{cands!r}"
                )
            elif widen_detail.get("fixability") != "auto_fix_candidate":
                failures.append(
                    "widen fixture: expected fixability=auto_fix_candidate; got "
                    f"{widen_detail.get('fixability')!r}"
                )
            else:
                pre_prs = Presentation(str(wrap_widen))
                pre_width = pre_prs.slides[0].shapes[0].width
                widen_findings = [pptx_lint.finding_to_json_dict(f) for f in widen_findings_raw]
                actions = pptx_fix.fix_pptx(
                    wrap_widen, apply=True, rules=("text_wrap",), findings=widen_findings
                )
                if not any(a.rule == "text_wrap" and a.status == "apply" for a in actions):
                    failures.append("widen fixture: text_wrap fixer did not apply")
                post_prs = Presentation(str(wrap_widen))
                post_shape = post_prs.slides[0].shapes[0]
                post_width = post_shape.width
                expected_pt = float(widen_candidate["width_pt"])
                actual_pt = post_width / 12700.0
                if abs(actual_pt - expected_pt) > 1.0:
                    failures.append(
                        f"widen fixture: shape width expected ~{expected_pt}pt; "
                        f"got {actual_pt:.2f}pt (before {pre_width / 12700:.2f}pt)"
                    )
                if post_shape.text_frame.text != "Automation improves review speed":
                    failures.append(
                        "widen fixture: expected joined word; got "
                        f"{post_shape.text_frame.text!r}"
                    )

        # --- text color allowlist follows design-system candidate even if marked manual_required ---
        text_color = tmp_dir / "text-color.pptx"
        _build_text_color_allowlist_fixture(text_color)
        text_color_findings = [
            pptx_lint.finding_to_json_dict(f)
            for f in pptx_lint.lint_pptx(text_color)
            if f.check == "text_color_allowlist"
        ]
        if not text_color_findings:
            failures.append("text-color fixture did not trigger text_color_allowlist")
        actions = pptx_fix.fix_pptx(
            text_color,
            apply=True,
            rules=("text_color",),
            findings=text_color_findings,
        )
        if not any(a.rule == "text_color" and a.status == "apply" for a in actions):
            failures.append("text_color_allowlist did not apply design-system candidate")
        prs = Presentation(str(text_color))
        color = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.color.rgb
        if str(color).upper() != "E6033D":
            failures.append(f"text_color fixture expected #E6033D; got #{color}")

        # --- slide_size is intentionally not auto-fixed ---
        slide_size = tmp_dir / "slide-size.pptx"
        _build_slide_size_fixture(slide_size)
        slide_size_findings = [
            pptx_lint.finding_to_json_dict(f)
            for f in pptx_lint.lint_pptx(slide_size)
            if f.check == "slide_size"
        ]
        if not slide_size_findings:
            failures.append("slide-size fixture did not trigger slide_size")
        slide_size_rules = pptx_fix.auto_rules_from_findings(slide_size_findings)
        if "slide_size" in slide_size_rules:
            failures.append(f"slide_size must not be auto-selected; got {slide_size_rules}")
        actions = pptx_fix.fix_pptx(slide_size, apply=True, rules=pptx_fix.ALL_RULES)
        if any(a.rule == "slide_size" and a.status == "apply" for a in actions):
            failures.append("slide_size fixer applied an action even though size strategy is manual")
        prs = Presentation(str(slide_size))
        if prs.slide_width != Pt(1000) or prs.slide_height != Pt(500):
            failures.append("slide_size fixture dimensions changed unexpectedly")

        # --- self-check (verify_pptx) reports no residual on fixed file ---
        bad3 = tmp_dir / "bad3.pptx"
        make_examples.make_bad(bad3)
        pptx_fix.fix_pptx(bad3, apply=True, rules=pptx_fix.ALL_RULES)
        residual = pptx_fix.verify_pptx(bad3, rules=pptx_fix.ALL_RULES)
        if residual:
            failures.append(
                f"verify_pptx reported residual on a clean fix: {len(residual)}"
            )

    if failures:
        print("FAIL:")
        for line in failures:
            print(f"- {line}")
        return 1

    print(
        "OK: pptx_fix removes autofit issues, rounds drifted geometry, "
        "and respects dry-run"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
